import sys
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

sys.stdout.reconfigure(encoding='utf-8')

original_pptx = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2.pptx"
output_pptx = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - дополненный.pptx"

prs = Presentation(original_pptx)
blank_layout = prs.slide_layouts[12]

# Colors matching the webinar style
ORANGE = RGBColor(0xFF, 0x6F, 0x03)      # #FF6F03
ORANGE_DOT = RGBColor(0xFE, 0x60, 0x02)  # #FE6002
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)   # #222222
BODY_TEXT = RGBColor(0x44, 0x44, 0x44)   # #444444
CARD_BG = RGBColor(0xF9, 0xF9, 0xF9)     # #F9F9F9
CARD_BORDER = RGBColor(0xDC, 0xDF, 0xE4) # Border
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x18, 0x1A, 0x1F)     # Dark editor theme
CODE_TEXT = RGBColor(0xD4, 0xD4, 0xD8)   # Default text
CODE_KW = RGBColor(0x38, 0xBD, 0xF8)     # Cyan/Blue keywords
CODE_TAG = RGBColor(0xF4, 0x72, 0xB6)    # Pink tags
CODE_STR = RGBColor(0xFC, 0xD3, 0x4D)    # Amber strings
CODE_COMM = RGBColor(0x71, 0x71, 0x7A)   # Gray comments
CODE_PROP = RGBColor(0xA7, 0x8B, 0xFA)   # Violet properties
CODE_VAL = RGBColor(0x4A, 0xDE, 0x80)    # Green values

def add_header_and_footer(slide, category_text, question_text):
    tb_top = slide.shapes.add_textbox(Pt(260), Pt(20), Pt(430), Pt(28))
    tf_top = tb_top.text_frame
    tf_top.word_wrap = False  # Без переноса текста по словам
    tf_top.margin_left = tf_top.margin_top = tf_top.margin_right = tf_top.margin_bottom = 0
    p_top = tf_top.paragraphs[0]
    p_top.alignment = PP_ALIGN.RIGHT
    run_top = p_top.add_run()
    run_top.text = category_text.upper()
    run_top.font.name = "Inter"
    run_top.font.size = Pt(14)  # Размер 14 Inter
    run_top.font.bold = False
    run_top.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # Белый цвет текста

    tb_bot = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(500), Pt(20))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    tf_bot.margin_left = tf_bot.margin_top = tf_bot.margin_right = tf_bot.margin_bottom = 0
    p_bot = tf_bot.paragraphs[0]
    p_bot.alignment = PP_ALIGN.LEFT
    run_bot = p_bot.add_run()
    run_bot.text = question_text
    run_bot.font.name = "Montserrat"
    run_bot.font.size = Pt(10)
    run_bot.font.color.rgb = RGBColor(0xA6, 0xA1, 0xA1)  # #A6A1A1

def create_code_explanation_slide(slide, category, question, left_title, code_lines, right_title, bullet_items):
    add_header_and_footer(slide, category, question)
    
    # Left Card (Code snippet / editor card)
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(24), Pt(56), Pt(325), Pt(298))
    left_card.adjustments[0] = 0.04  # Уменьшенное скругление углов (было 0.16667)
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CODE_BG
    left_card.line.color.rgb = RGBColor(0x2D, 0x31, 0x39)
    left_card.line.width = Pt(1)

    # Title of Left Card
    tb_l_header = slide.shapes.add_textbox(Pt(38), Pt(64), Pt(295), Pt(18))
    tf_lh = tb_l_header.text_frame
    tf_lh.word_wrap = True
    tf_lh.margin_left = tf_lh.margin_top = tf_lh.margin_right = tf_lh.margin_bottom = 0
    p_lh = tf_lh.paragraphs[0]
    r_lh = p_lh.add_run()
    r_lh.text = "📄 " + left_title
    r_lh.font.name = "Consolas"
    r_lh.font.size = Pt(9)
    r_lh.font.bold = True
    r_lh.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)

    line_lh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(36), Pt(84), Pt(301), Pt(1))
    line_lh.fill.solid()
    line_lh.fill.fore_color.rgb = ORANGE
    line_lh.line.fill.background()

    # Code text box
    tb_code = slide.shapes.add_textbox(Pt(36), Pt(90), Pt(301), Pt(256))
    tf_code = tb_code.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = tf_code.margin_top = tf_code.margin_right = tf_code.margin_bottom = 0
    
    first_p = True
    for line in code_lines:
        if first_p:
            p = tf_code.paragraphs[0]
            first_p = False
        else:
            p = tf_code.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        
        if isinstance(line, str):
            r = p.add_run()
            r.text = line
            r.font.name = "Consolas"
            r.font.size = Pt(8)
            r.font.color.rgb = CODE_TEXT
        else:
            for text_chunk, chunk_color in line:
                r = p.add_run()
                r.text = text_chunk
                r.font.name = "Consolas"
                r.font.size = Pt(8)
                r.font.color.rgb = chunk_color

    # Правый фоновый серый блок (Rounded Rectangle 7) удален по запросу пользователя.
    # Текст и оранжевые акценты теперь располагаются прямо на чистом фоне слайда.

    tb_r_header = slide.shapes.add_textbox(Pt(380), Pt(66), Pt(300), Pt(20))
    tf_rh = tb_r_header.text_frame
    tf_rh.word_wrap = True
    tf_rh.margin_left = tf_rh.margin_top = tf_rh.margin_right = tf_rh.margin_bottom = 0
    p_rh = tf_rh.paragraphs[0]
    r_rh = p_rh.add_run()
    r_rh.text = right_title
    r_rh.font.name = "Montserrat"
    r_rh.font.size = Pt(10.5)
    r_rh.font.bold = True
    r_rh.font.color.rgb = ORANGE

    line_rh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(380), Pt(87), Pt(300), Pt(1))
    line_rh.fill.solid()
    line_rh.fill.fore_color.rgb = ORANGE
    line_rh.line.fill.background()

    y_offset = Pt(94)
    for b_idx, item in enumerate(bullet_items):
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(380), y_offset + Pt(3), Pt(6.4), Pt(6.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_DOT
        dot.line.fill.background()

        tb_b = slide.shapes.add_textbox(Pt(394), y_offset, Pt(285), Pt(46))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        p_b1 = tf_b.paragraphs[0]
        p_b1.space_after = Pt(2)
        r_title = p_b1.add_run()
        r_title.text = item["title"]
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = DARK_TEXT

        p_b2 = tf_b.add_paragraph()
        p_b2.line_spacing = 1.15
        p_b2.space_after = Pt(0)
        
        if isinstance(item["desc"], str):
            r_desc = p_b2.add_run()
            r_desc.text = item["desc"]
            r_desc.font.name = "Inter"
            r_desc.font.size = Pt(8.5)
            r_desc.font.color.rgb = BODY_TEXT
        else:
            for d_text, d_bold, d_color in item["desc"]:
                r_desc = p_b2.add_run()
                r_desc.text = d_text
                r_desc.font.name = "Inter"
                r_desc.font.size = Pt(8.5)
                r_desc.font.bold = d_bold
                r_desc.font.color.rgb = d_color

        y_offset += Pt(item.get("height", 50))

def create_grid_cards_slide(slide, category, question, title_text, cards_data):
    add_header_and_footer(slide, category, question)
    
    # 2x2 grid for errors
    coords = [
        (Pt(24), Pt(56), Pt(325), Pt(142)),   # Top Left
        (Pt(365), Pt(56), Pt(330), Pt(142)),  # Top Right
        (Pt(24), Pt(210), Pt(325), Pt(142)),  # Bottom Left
        (Pt(365), Pt(210), Pt(330), Pt(142)), # Bottom Right
    ]
    
    for idx, (x, y, w, h) in enumerate(coords):
        data = cards_data[idx]
        
        # Header text box
        tb_h = slide.shapes.add_textbox(x + Pt(14), y + Pt(10), w - Pt(28), Pt(20))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
        p_h = tf_h.paragraphs[0]
        r_num = p_h.add_run()
        r_num.text = f"⚠️ {data['title']}"
        r_num.font.name = "Montserrat"
        r_num.font.size = Pt(10)
        r_num.font.bold = True
        r_num.font.color.rgb = ORANGE
        
        # Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(14), y + Pt(32), w - Pt(28), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE
        line.line.fill.background()
        
        # Body text box
        tb_body = slide.shapes.add_textbox(x + Pt(14), y + Pt(38), w - Pt(28), Pt(94))
        tf_b = tb_body.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        # Problem
        p1 = tf_b.paragraphs[0]
        p1.space_after = Pt(3)
        r1_lbl = p1.add_run()
        r1_lbl.text = "Ошибка: "
        r1_lbl.font.name = "Inter"
        r1_lbl.font.size = Pt(8.5)
        r1_lbl.font.bold = True
        r1_lbl.font.color.rgb = DARK_TEXT
        r1_txt = p1.add_run()
        r1_txt.text = data['problem']
        r1_txt.font.name = "Inter"
        r1_txt.font.size = Pt(8.5)
        r1_txt.font.color.rgb = BODY_TEXT
        
        # Solution
        p2 = tf_b.add_paragraph()
        p2.space_after = Pt(0)
        r2_lbl = p2.add_run()
        r2_lbl.text = "Решение: "
        r2_lbl.font.name = "Inter"
        r2_lbl.font.size = Pt(8.5)
        r2_lbl.font.bold = True
        r2_lbl.font.color.rgb = RGBColor(0x16, 0xA3, 0x4A) # Green
        r2_txt = p2.add_run()
        r2_txt.text = data['solution']
        r2_txt.font.name = "Inter"
        r2_txt.font.size = Pt(8.5)
        r2_txt.font.color.rgb = BODY_TEXT

def create_checklist_slide(slide, category, question, title_text, items):
    add_header_and_footer(slide, category, question)
    # Фоновый серый Rounded Rectangle удален - чек-лист отображается на чистом фоне

    tb_h = slide.shapes.add_textbox(Pt(54), Pt(68), Pt(600), Pt(22))
    tf_h = tb_h.text_frame
    tf_h.word_wrap = True
    tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
    p_h = tf_h.paragraphs[0]
    r_h = p_h.add_run()
    r_h.text = title_text
    r_h.font.name = "Montserrat"
    r_h.font.size = Pt(11)
    r_h.font.bold = True
    r_h.font.color.rgb = ORANGE

    line_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(54), Pt(92), Pt(612), Pt(1))
    line_h.fill.solid()
    line_h.fill.fore_color.rgb = ORANGE
    line_h.line.fill.background()

    y_pos = Pt(102)
    for idx, item in enumerate(items):
        # Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(54), y_pos + Pt(2), Pt(22), Pt(20))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(0x10, 0xB9, 0x81) # Green
        badge.line.fill.background()
        
        # Checkmark text
        tb_chk = slide.shapes.add_textbox(Pt(54), y_pos + Pt(4), Pt(22), Pt(20))
        tf_chk = tb_chk.text_frame
        p_chk = tf_chk.paragraphs[0]
        p_chk.alignment = PP_ALIGN.CENTER
        r_chk = p_chk.add_run()
        r_chk.text = "✔"
        r_chk.font.name = "Segoe UI Symbol"
        r_chk.font.size = Pt(10)
        r_chk.font.bold = True
        r_chk.font.color.rgb = WHITE
        
        # Content box
        tb_item = slide.shapes.add_textbox(Pt(86), y_pos, Pt(570), Pt(38))
        tf_item = tb_item.text_frame
        tf_item.word_wrap = True
        tf_item.margin_left = tf_item.margin_top = tf_item.margin_right = tf_item.margin_bottom = 0
        
        p1 = tf_item.paragraphs[0]
        p1.space_after = Pt(2)
        r1 = p1.add_run()
        r1.text = item['title']
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT
        
        p2 = tf_item.add_paragraph()
        p2.space_after = Pt(0)
        r2 = p2.add_run()
        r2.text = item['desc']
        r2.font.name = "Inter"
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = BODY_TEXT
        
        y_pos += Pt(46)

# ==========================================
# 15 NEW SLIDES DATA
# ==========================================

slides_config = [
    # 1. Структура проекта
    {
        "type": "code_exp",
        "category": "СТРУКТУРА ПРОЕКТА",
        "question": "Как правильно организовать файлы в веб-разработке?",
        "left_title": "Файловая система: smart-office / basic",
        "code": [
            [("📁 smart-office / basic/", CODE_STR)],
            [("│", CODE_COMM)],
            [("├── 📄 ", CODE_COMM), ("index.html", CODE_KW), ("        # Главная разметка страницы", CODE_COMM)],
            [("│", CODE_COMM)],
            [("├── 📁 ", CODE_COMM), ("css/", CODE_PROP), ("               # Стили оформления", CODE_COMM)],
            [("│   └── 📄 ", CODE_COMM), ("style.css", CODE_VAL), ("        # Таблица каскадных стилей", CODE_COMM)],
            [("│", CODE_COMM)],
            [("├── 📁 ", CODE_COMM), ("js/", CODE_PROP), ("                # Клиентские скрипты", CODE_COMM)],
            [("│   └── 📄 ", CODE_COMM), ("script.js", CODE_VAL), ("       # Интерактивная логика", CODE_COMM)],
            [("│", CODE_COMM)],
            [("└── 📁 ", CODE_COMM), ("img/", CODE_PROP), ("               # Графика и медиафайлы", CODE_COMM)],
            [("    ├── 🖼️ ", CODE_COMM), ("logo.svg", CODE_TEXT), ("        # Векторный логотип бренда", CODE_COMM)],
            [("    └── 🖼️ ", CODE_COMM), ("no-image.svg", CODE_TEXT), ("    # Заглушка для фото комнат", CODE_COMM)],
        ],
        "right_title": "ПРИНЦИП ЧИСТОЙ АРХИТЕКТУРЫ",
        "bullets": [
            {
                "title": "Файл index.html в корне:",
                "desc": "Браузер по умолчанию ищет именно index.html как стартовую страницу. Он лежит в корне проекта.",
                "height": 45
            },
            {
                "title": "Изоляция стилей в папке css/:",
                "desc": "Стили хранятся в файле style.css. Это разделяет разметку (HTML) и визуальный дизайн (CSS).",
                "height": 45
            },
            {
                "title": "Папка js/ для скриптов:",
                "desc": "Здесь хранятся файлы JavaScript (script.js). В них пишется логика: клики, слайдеры, фильтры и интерактив.",
                "height": 52
            },
            {
                "title": "Графика в папке img/:",
                "desc": "Все логотипы, иконки и фото комнат лежат отдельно, чтобы не засорять корень проекта.",
                "height": 45
            },
            {
                "title": "Относительные пути к файлам:",
                "desc": "В коде мы пишем путь от текущей папки: css/style.css, js/script.js, img/logo.svg.",
                "height": 45
            }
        ]
    },

    # 2. HTML5 Семантика
    {
        "type": "code_exp",
        "category": "HTML5 СЕМАНТИКА",
        "question": "Из каких смысловых блоков строится современный сайт?",
        "left_title": "index.html (семантический каркас)",
        "code": [
            [("<!DOCTYPE html>", CODE_COMM)],
            [("<html ", CODE_TAG), ("lang=", CODE_PROP), ("\"ru\"", CODE_STR), (">", CODE_TAG)],
            [("<head>", CODE_TAG), (" ... ", CODE_COMM), ("</head>", CODE_TAG)],
            [("<body>", CODE_TAG)],
            [("  <!-- 1. Шапка сайта (лого + меню) -->", CODE_COMM)],
            [("  <header ", CODE_TAG), ("class=", CODE_PROP), ("\"header\"", CODE_STR), (">", CODE_TAG)],
            [("    <div ", CODE_TAG), ("class=", CODE_PROP), ("\"container header-container\"", CODE_STR), (">...</div>", CODE_TAG)],
            [("  </header>", CODE_TAG)],
            [("  ", CODE_TEXT)],
            [("  <!-- 2. Главный уникальный контент -->", CODE_COMM)],
            [("  <main ", CODE_TAG), ("class=", CODE_PROP), ("\"main\"", CODE_STR), (">", CODE_TAG)],
            [("    <div ", CODE_TAG), ("class=", CODE_PROP), ("\"container\"", CODE_STR), (">", CODE_TAG)],
            [("      <section ", CODE_TAG), ("class=", CODE_PROP), ("\"hero-section\"", CODE_STR), (">...</section>", CODE_TAG)],
            [("    </div>", CODE_TAG)],
            [("  </main>", CODE_TAG)],
            [("  ", CODE_TEXT)],
            [("  <!-- 3. Подвал сайта (копирайт и контакты) -->", CODE_COMM)],
            [("  <footer ", CODE_TAG), ("class=", CODE_PROP), ("\"footer\"", CODE_STR), (">...</footer>", CODE_TAG)],
            [("</body>", CODE_TAG)],
            [("</html>", CODE_TAG)]
        ],
        "right_title": "СЕМАНТИЧЕСКИЕ ТЕГИ HTML5",
        "bullets": [
            {
                "title": "Зачем нужны <header>, <main>, <footer>?",
                "desc": "Раньше сайты верстали сплошными тегами <div>. Семантические теги делают код понятным людям, поисковикам и скринридерам.",
                "height": 56
            },
            {
                "title": "Тег <header> (Шапка):",
                "desc": "Верхняя часть страницы: содержит бренд, логотип, навигационные ссылки по разделам и кнопки профиля.",
                "height": 52
            },
            {
                "title": "Тег <main> (Главный блок):",
                "desc": "Центральная область. Здесь располагается уникальное содержимое страницы, не повторяющееся на других вкладках.",
                "height": 54
            },
            {
                "title": "Тег <footer> (Подвал):",
                "desc": "Нижняя часть сайта: юридическая информация, авторские права (копирайт), контакты техподдержки и соцсети.",
                "height": 54
            }
        ]
    },

    # 3. Служебный блок <head>
    {
        "type": "code_exp",
        "category": "СЛУЖЕБНЫЙ БЛОК <HEAD>",
        "question": "Что подключается в невидимой части документа?",
        "left_title": "index.html (настройка <head>)",
        "code": [
            [("<head>", CODE_TAG)],
            [("  <meta ", CODE_TAG), ("charset=", CODE_PROP), ("\"UTF-8\"", CODE_STR), (">", CODE_TAG)],
            [("  <meta ", CODE_TAG), ("name=", CODE_PROP), ("\"viewport\"", CODE_STR)],
            [("        ", CODE_TEXT), ("content=", CODE_PROP), ("\"width=device-width, initial-scale=1.0\"", CODE_STR), (">", CODE_TAG)],
            [("  <title>", CODE_TAG), ("СмартОфис — Бронирование", CODE_TEXT), ("</title>", CODE_TAG)],
            [("  ", CODE_TEXT)],
            [("  <!-- Подключение шрифта Inter из Google Fonts -->", CODE_COMM)],
            [("  <link ", CODE_TAG), ("rel=", CODE_PROP), ("\"preconnect\"", CODE_STR), (" href=", CODE_PROP), ("\"https://fonts.googleapis.com\"", CODE_STR), (">", CODE_TAG)],
            [("  <link ", CODE_TAG), ("rel=", CODE_PROP), ("\"preconnect\"", CODE_STR), (" href=", CODE_PROP), ("\"https://fonts.gstatic.com\"", CODE_STR), (" crossorigin>", CODE_TAG)],
            [("  <link ", CODE_TAG), ("href=", CODE_PROP), ("\"https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap\"", CODE_STR)],
            [("        ", CODE_TEXT), ("rel=", CODE_PROP), ("\"stylesheet\"", CODE_STR), (">", CODE_TAG)],
            [("  ", CODE_TEXT)],
            [("  <!-- Подключение таблицы стилей и иконки -->", CODE_COMM)],
            [("  <link ", CODE_TAG), ("rel=", CODE_PROP), ("\"stylesheet\"", CODE_STR), (" href=", CODE_PROP), ("\"css/style.css\"", CODE_STR), (">", CODE_TAG)],
            [("  <link ", CODE_TAG), ("rel=", CODE_PROP), ("\"icon\"", CODE_STR), (" type=", CODE_PROP), ("\"image/svg+xml\"", CODE_STR), (" href=", CODE_PROP), ("\"img/logo.svg\"", CODE_STR), (">", CODE_TAG)],
            [("</head>", CODE_TAG)]
        ],
        "right_title": "СЛУЖЕБНЫЕ МЕТАТЕГИ",
        "bullets": [
            {
                "title": "Кодировка UTF-8:",
                "desc": "Указывает браузеру стандарт символов. Без <meta charset=\"UTF-8\"> кириллические буквы превратятся в иероглифы.",
                "height": 52
            },
            {
                "title": "Метатег Viewport:",
                "desc": "Говорит браузеру отображать страницу в масштабе 1:1 в соответствии с реальной шириной экрана экрана устройства.",
                "height": 52
            },
            {
                "title": "Google Fonts (Inter):",
                "desc": "Шрифт Inter — мировой стандарт интерфейсов (Figma, Linear). preconnect заранее открывает скоростное соединение с сервером.",
                "height": 56
            },
            {
                "title": "Фавикон сайта (Favicon):",
                "desc": "Векторная иконка img/logo.svg в теге <link rel=\"icon\"> украшает вкладку в браузере и закладки.",
                "height": 52
            }
        ]
    },

    # 4. CSS: Сброс стилей
    {
        "type": "code_exp",
        "category": "CSS: СБРОС СТИЛЕЙ (*)",
        "question": "Зачем сбрасывать стандартные отступы браузеров?",
        "left_title": "css/style.css (базовый сброс)",
        "code": [
            [("/* ========================================= */", CODE_COMM)],
            [("/* СБРОС СТИЛЕЙ И БАЗОВЫЕ НАСТРОЙКИ         */", CODE_COMM)],
            [("/* ========================================= */", CODE_COMM)],
            [(" ", CODE_TEXT)],
            [("/* Селектор * выбирает абсолютно все теги */", CODE_COMM)],
            [("* {", CODE_KW)],
            [("  /* Включаем padding и border в общую ширину */", CODE_COMM)],
            [("  box-sizing", CODE_PROP), (": ", CODE_TEXT), ("border-box", CODE_VAL), (";", CODE_TEXT)],
            [("  ", CODE_TEXT)],
            [("  /* Сбрасываем браузерные внешние отступы */", CODE_COMM)],
            [("  margin", CODE_PROP), (": ", CODE_TEXT), ("0", CODE_VAL), (";", CODE_TEXT)],
            [("  ", CODE_TEXT)],
            [("  /* Сбрасываем браузерные внутренние поля */", CODE_COMM)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("0", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)]
        ],
        "right_title": "МАГИЯ BOX-SIZING: BORDER-BOX",
        "bullets": [
            {
                "title": "Проблема браузерных стилей:",
                "desc": "У каждого браузера (Chrome, Safari, Firefox, Edge) свои отступы для <body>, <h1>, <p>, <ul>. Сброс делает их одинаковыми.",
                "height": 54
            },
            {
                "title": "Что делает margin: 0; padding: 0;?",
                "desc": "Убирает паразитные белые поля по краям экрана и у списков, давая нам полный контроль над каждым пикселем.",
                "height": 54
            },
            {
                "title": "Главное свойство — box-sizing: border-box:",
                "desc": "По умолчанию (content-box) если блоку шириной 100px задать padding: 20px, он станет 140px! Это ломает сетку верстки.",
                "height": 56
            },
            {
                "title": "Результат border-box:",
                "desc": "Ширина блока всегда остается ровно 100px, а внутренние отступы аккуратно сжимают контент внутри него.",
                "height": 50
            }
        ]
    },

    # 5. Базовые стили body
    {
        "type": "code_exp",
        "category": "БАЗОВЫЕ НАСТРОЙКИ BODY",
        "question": "Как настроить тело страницы и заложить прижатие подвала?",
        "left_title": "css/style.css (body)",
        "code": [
            [("body {", CODE_KW)],
            [("  /* Базовый шрифт и цвет текста */", CODE_COMM)],
            [("  font-family", CODE_PROP), (": ", CODE_TEXT), ("'Inter', sans-serif", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#222222", CODE_VAL), (";", CODE_TEXT)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#ffffff", CODE_VAL), (";", CODE_TEXT)],
            [("  line-height", CODE_PROP), (": ", CODE_TEXT), ("1.5", CODE_VAL), (";", CODE_TEXT)],
            [("  ", CODE_TEXT)],
            [("  /* Flexbox-раскладка для прижатия футера */", CODE_COMM)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  flex-direction", CODE_PROP), (": ", CODE_TEXT), ("column", CODE_VAL), (";", CODE_TEXT)],
            [("  min-height", CODE_PROP), (": ", CODE_TEXT), ("100vh", CODE_VAL), (";", CODE_TEXT)],
            [("  min-width", CODE_PROP), (": ", CODE_TEXT), ("1200px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)]
        ],
        "right_title": "ФУНДАМЕНТ СТРАНИЦЫ (BODY)",
        "bullets": [
            {
                "title": "font-family: 'Inter', sans-serif:",
                "desc": "Браузер использует загруженный шрифт Inter, а если интернета нет — системный шрифт без засечек sans-serif.",
                "height": 54
            },
            {
                "title": "line-height: 1.5 (Интерлиньяж):",
                "desc": "Полуторный межстрочный интервал делает даже длинный текст легким и приятным для чтения студентами.",
                "height": 52
            },
            {
                "title": "min-height: 100vh (Высота viewport):",
                "desc": "100vh означает 100% высоты экрана. Страница гарантированно не сожмется, даже если на ней всего одно слово.",
                "height": 54
            },
            {
                "title": "display: flex; flex-direction: column:",
                "desc": "Выстраивает шапку, контент и футер вертикально друг под другом, подготавливая почву для идеального прижатия футера.",
                "height": 56
            }
        ]
    },

    # 6. Центрирующий контейнер
    {
        "type": "code_exp",
        "category": "ЦЕНТРИРУЮЩИЙ КОНТЕЙНЕР",
        "question": "Как удержать контент в центре экрана на любом мониторе?",
        "left_title": "css/style.css (.container)",
        "code": [
            [("/* Ограничивающий и центрирующий контейнер */", CODE_COMM)],
            [(".container {", CODE_KW)],
            [("  /* Фиксированная ширина контентной зоны */", CODE_COMM)],
            [("  width", CODE_PROP), (": ", CODE_TEXT), ("1200px", CODE_VAL), (";", CODE_TEXT)],
            [("  ", CODE_TEXT)],
            [("  /* Автоматическое центрирование по горизонтали */", CODE_COMM)],
            [("  margin", CODE_PROP), (": ", CODE_TEXT), ("0 auto", CODE_VAL), (";", CODE_TEXT)],
            [("  ", CODE_TEXT)],
            [("  /* Внутренние поля безопасности слева и справа */", CODE_COMM)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("0 15px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)]
        ],
        "right_title": "АНАТОМИЯ КЛАССА .CONTAINER",
        "bullets": [
            {
                "title": "width: 1200px (Ширина контента):",
                "desc": "На мониторах 2K и 4K текст без контейнера растянулся бы на 2 метра. 1200px — золотой стандарт читаемости.",
                "height": 54
            },
            {
                "title": "Секрет margin: 0 auto:",
                "desc": "Первое значение (0) — отступы сверху и снизу. Второе (auto) — браузер делит оставшееся место поровну слева и справа!",
                "height": 56
            },
            {
                "title": "padding: 0 15px (Боковые поля):",
                "desc": "Если окно уменьшится, текст не упрется в край экрана, а сохранит эстетичный зазор в 15 пикселей.",
                "height": 52
            },
            {
                "title": "Где применяется:",
                "desc": "Мы оборачиваем в .container контент внутри <header>, <main> и <footer>, выстраивая всю страницу по единой линейке.",
                "height": 52
            }
        ]
    },

    # 7. Шапка сайта: разметка
    {
        "type": "code_exp",
        "category": "ШАПКА САЙТА (HEADER): HTML",
        "question": "Как структурировать элементы внутри шапки сайта?",
        "left_title": "index.html (<header>)",
        "code": [
            [("<header ", CODE_TAG), ("class=", CODE_PROP), ("\"header\"", CODE_STR), (">", CODE_TAG)],
            [("  <div ", CODE_TAG), ("class=", CODE_PROP), ("\"container header-container\"", CODE_STR), (">", CODE_TAG)],
            [("    ", CODE_TEXT)],
            [("    <!-- Логотип со ссылкой на Главную -->", CODE_COMM)],
            [("    <a ", CODE_TAG), ("href=", CODE_PROP), ("\"index.html\"", CODE_STR), (" class=", CODE_PROP), ("\"logo\"", CODE_STR), (">", CODE_TAG)],
            [("      <img ", CODE_TAG), ("src=", CODE_PROP), ("\"img/logo.svg\"", CODE_STR), (" alt=", CODE_PROP), ("\"Логотип\"", CODE_STR), (" class=", CODE_PROP), ("\"logo-icon\"", CODE_STR), (">", CODE_TAG)],
            [("      <span>", CODE_TAG), ("СмартОфис", CODE_TEXT), ("</span>", CODE_TAG)],
            [("    </a>", CODE_TAG)],
            [("    ", CODE_TEXT)],
            [("    <!-- Навигационное меню -->", CODE_COMM)],
            [("    <nav ", CODE_TAG), ("class=", CODE_PROP), ("\"nav\"", CODE_STR), (">", CODE_TAG)],
            [("      <ul ", CODE_TAG), ("class=", CODE_PROP), ("\"nav-list\"", CODE_STR), (">", CODE_TAG)],
            [("        <li><a ", CODE_TAG), ("href=", CODE_PROP), ("\"index.html\"", CODE_STR), (" class=", CODE_PROP), ("\"nav-link active\"", CODE_STR), (">Главная</a></li>", CODE_TAG)],
            [("        <li><a ", CODE_TAG), ("href=", CODE_PROP), ("\"pages/catalog.html\"", CODE_STR), (" class=", CODE_PROP), ("\"nav-link\"", CODE_STR), (">Каталог</a></li>", CODE_TAG)],
            [("        <li><a ", CODE_TAG), ("href=", CODE_PROP), ("\"pages/login.html\"", CODE_STR), (" class=", CODE_PROP), ("\"nav-link nav-btn\"", CODE_STR), (">Войти</a></li>", CODE_TAG)],
            [("      </ul>", CODE_TAG)],
            [("    </nav>", CODE_TAG)],
            [("  </div>", CODE_TAG)],
            [("</header>", CODE_TAG)]
        ],
        "right_title": "СТРУКТУРА И ВЛОЖЕННОСТЬ ШАПКИ",
        "bullets": [
            {
                "title": "Двухуровневая обертка:",
                "desc": "Внешний тег <header> отвечает за фон и нижнюю границу на всю ширину. Внутренний .header-container центрирует элементы.",
                "height": 56
            },
            {
                "title": "Логотип как ссылка (<a>):",
                "desc": "Логотип объединяет SVG-иконку и текст. Клик по любой его части перезагружает сайт или возвращает пользователя на главную.",
                "height": 56
            },
            {
                "title": "Семантический тег <nav>:",
                "desc": "Группирует основные ссылки сайта. Поисковые роботы сразу понимают, что здесь находится карта разделов.",
                "height": 50
            },
            {
                "title": "Класс .active и кнопка .nav-btn:",
                "desc": ".active показывает пользователю, на какой странице он находится. .nav-btn превращает третью ссылку в акцентную кнопку.",
                "height": 52
            }
        ]
    },

    # 8. Шапка: Стили и Flexbox
    {
        "type": "code_exp",
        "category": "ШАПКА САЙТА: CSS & FLEXBOX",
        "question": "Как разнести логотип и меню по противоположным краям?",
        "left_title": "css/style.css (header flexbox)",
        "code": [
            [(".header {", CODE_KW)],
            [("  border-bottom", CODE_PROP), (": ", CODE_TEXT), ("1px solid #dddddd", CODE_VAL), (";", CODE_TEXT)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("18px 0", CODE_VAL), (";", CODE_TEXT)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#ffffff", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(" ", CODE_TEXT)],
            [(".header-container {", CODE_KW)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  justify-content", CODE_PROP), (": ", CODE_TEXT), ("space-between", CODE_VAL), (";", CODE_TEXT)],
            [("  align-items", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(" ", CODE_TEXT)],
            [(".logo {", CODE_KW)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  align-items", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";", CODE_TEXT)],
            [("  gap", CODE_PROP), (": ", CODE_TEXT), ("10px", CODE_VAL), (";", CODE_TEXT)],
            [("  text-decoration", CODE_PROP), (": ", CODE_TEXT), ("none", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#222222", CODE_VAL), (";", CODE_TEXT)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("20px", CODE_VAL), (";", CODE_TEXT)],
            [("  font-weight", CODE_PROP), (": ", CODE_TEXT), ("800", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(".logo-icon", CODE_KW), (" { ", CODE_TEXT), ("width", CODE_PROP), (": ", CODE_TEXT), ("32px", CODE_VAL), ("; ", CODE_TEXT), ("height", CODE_PROP), (": ", CODE_TEXT), ("32px", CODE_VAL), ("; }", CODE_TEXT)]
        ],
        "right_title": "СИЛА FLEXBOX В ДЕЙСТВИИ",
        "bullets": [
            {
                "title": "display: flex (Включаем гибкую сетку):",
                "desc": "Превращает .header-container в флекс-контейнер. Все прямые потомки (логотип и меню) встают в одну ровную линию.",
                "height": 56
            },
            {
                "title": "justify-content: space-between:",
                "desc": "Главный трюк верстки! Первый элемент (лого) прижимается к левой стенке, последний (меню) — к правой стенке.",
                "height": 56
            },
            {
                "title": "align-items: center:",
                "desc": "Выравнивает элементы по вертикальной оси. Логотип и кнопки встают строго по центру без ручных отступов.",
                "height": 52
            },
            {
                "title": "Свойство gap: 10px в .logo:",
                "desc": "Современный способ задать отступ ровно в 10px между иконкой логотипа и текстом бренда без margin-right.",
                "height": 50
            }
        ]
    },

    # 9. Навигация, ссылки и кнопка
    {
        "type": "code_exp",
        "category": "НАВИГАЦИОННОЕ МЕНЮ И КНОПКА",
        "question": "Как оформить ссылки, состояния hover и кнопку «Войти»?",
        "left_title": "css/style.css (навигация)",
        "code": [
            [(".nav-list {", CODE_KW)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  list-style", CODE_PROP), (": ", CODE_TEXT), ("none", CODE_VAL), ("; /* Убираем маркеры */", CODE_COMM)],
            [("  gap", CODE_PROP), (": ", CODE_TEXT), ("12px", CODE_VAL), (";", CODE_TEXT)],
            [("  align-items", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(" ", CODE_TEXT)],
            [(".nav-link {", CODE_KW)],
            [("  text-decoration", CODE_PROP), (": ", CODE_TEXT), ("none", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#222222", CODE_VAL), (";", CODE_TEXT)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("8px 14px", CODE_VAL), (";", CODE_TEXT)],
            [("  border-radius", CODE_PROP), (": ", CODE_TEXT), ("4px", CODE_VAL), (";", CODE_TEXT)],
            [("  font-weight", CODE_PROP), (": ", CODE_TEXT), ("500", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(".nav-link:hover {", CODE_KW)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), (";", CODE_TEXT)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#eaf2ff", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(".nav-btn {", CODE_KW)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#ffffff", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(".nav-btn:hover", CODE_KW), (" { ", CODE_TEXT), ("background-color", CODE_PROP), (": ", CODE_TEXT), ("#0056b3", CODE_STR), ("; }", CODE_TEXT)]
        ],
        "right_title": "СТИЛИЗАЦИЯ И ИНТЕРАКТИВНОСТЬ",
        "bullets": [
            {
                "title": "list-style: none:",
                "desc": "Убирает браузерные черные точки у тегов <li>. Меню превращается в аккуратную группу кнопок.",
                "height": 50
            },
            {
                "title": "Псевдокласс :hover:",
                "desc": "Срабатывает в момент наведения курсора мыши. Ссылка меняет цвет на синий и получает легкий голубой фон #eaf2ff.",
                "height": 56
            },
            {
                "title": "Кнопка призыва к действию (.nav-btn):",
                "desc": "Класс .nav-btn окрашивает кнопку «Войти» в насыщенный цвет #007bff с белым текстом, делая её заметной.",
                "height": 56
            },
            {
                "title": "border-radius: 4px:",
                "desc": "Легкое скругление углов придает элементам интерфейса современный мягкий вид.",
                "height": 50
            }
        ]
    },

    # 10. Hero-секция: Заголовок
    {
        "type": "code_exp",
        "category": "ГЛАВНЫЙ ЭКРАН (HERO SECTION)",
        "question": "Как сверстать эффектный презентационный заголовок?",
        "left_title": "index.html & style.css (Hero)",
        "code": [
            [("<!-- Главный баннер (Hero-блок) -->", CODE_COMM)],
            [("<section ", CODE_TAG), ("class=", CODE_PROP), ("\"hero-section\"", CODE_STR), (">", CODE_TAG)],
            [("  <h1 ", CODE_TAG), ("class=", CODE_PROP), ("\"hero-title\"", CODE_STR), (">", CODE_TAG)],
            [("    Портал бронирования офисных комнат", CODE_TEXT)],
            [("    <span ", CODE_TAG), ("class=", CODE_PROP), ("\"brand-highlight\"", CODE_STR), (">«СмартОфис»</span>", CODE_TAG)],
            [("  </h1>", CODE_TAG)],
            [("  <div ", CODE_TAG), ("class=", CODE_PROP), ("\"hero-bottom\"", CODE_STR), (">...</div>", CODE_TAG)],
            [("</section>", CODE_TAG)],
            [(" ", CODE_TEXT)],
            [("/* Стили гигантского заголовка */", CODE_COMM)],
            [(".hero-title {", CODE_KW)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("108px", CODE_VAL), (";", CODE_TEXT)],
            [("  line-height", CODE_PROP), (": ", CODE_TEXT), ("0.95", CODE_VAL), (";", CODE_TEXT)],
            [("  letter-spacing", CODE_PROP), (": ", CODE_TEXT), ("-0.04em", CODE_VAL), (";", CODE_TEXT)],
            [("  font-weight", CODE_PROP), (": ", CODE_TEXT), ("800", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#222222", CODE_VAL), (";", CODE_TEXT)],
            [("  margin-bottom", CODE_PROP), (": ", CODE_TEXT), ("25px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(".brand-highlight {", CODE_KW)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)]
        ],
        "right_title": "ПЛАКАТНАЯ ТИПОГРАФИКА HERO",
        "bullets": [
            {
                "title": "Единственный заголовок <h1>:",
                "desc": "По стандартам веб-доступности и поисковых систем Google/Яндекс на странице должен быть только один <h1>.",
                "height": 54
            },
            {
                "title": "font-size: 108px:",
                "desc": "Смелый плакатный размер сразу приковывает внимание студента и формирует премиальное ощущение от сервиса.",
                "height": 54
            },
            {
                "title": "line-height: 0.95 и letter-spacing: -0.04em:",
                "desc": "При экстремально больших шрифтах буквы разъезжаются. Отрицательный кернинг и плотный интерлиньяж делают текст монолитным.",
                "height": 56
            },
            {
                "title": "Подсветка бренда через <span>:",
                "desc": "Тег <span> позволяет изменить цвет слова «СмартОфис» на фирменный синий, не разбивая заголовок на части.",
                "height": 50
            }
        ]
    },

    # 11. Блок преимуществ и метрик
    {
        "type": "code_exp",
        "category": "БЛОК ПРЕИМУЩЕСТВ И МЕТРИК",
        "question": "Как скомпоновать цифры и показатели сервиса в одну линию?",
        "left_title": "index.html & style.css (Metrics)",
        "code": [
            [("<div ", CODE_TAG), ("class=", CODE_PROP), ("\"hero-bottom\"", CODE_STR), (">", CODE_TAG)],
            [("  <p ", CODE_TAG), ("class=", CODE_PROP), ("\"hero-subtitle\"", CODE_STR), (">", CODE_TAG)],
            [("    Удобный выбор и быстрое бронирование...", CODE_TEXT)],
            [("  </p>", CODE_TAG)],
            [("  <div ", CODE_TAG), ("class=", CODE_PROP), ("\"hero-metrics\"", CODE_STR), (">", CODE_TAG)],
            [("    <div ", CODE_TAG), ("class=", CODE_PROP), ("\"metric-item\"", CODE_STR), (">", CODE_TAG)],
            [("      <span ", CODE_TAG), ("class=", CODE_PROP), ("\"metric-val\"", CODE_STR), (">", CODE_TAG), ("24/7", CODE_TEXT), ("</span>", CODE_TAG)],
            [("      <span ", CODE_TAG), ("class=", CODE_PROP), ("\"metric-lbl\"", CODE_STR), (">", CODE_TAG), ("Доступ", CODE_TEXT), ("</span>", CODE_TAG)],
            [("    </div>", CODE_TAG)],
            [("    <div ", CODE_TAG), ("class=", CODE_PROP), ("\"metric-item\"", CODE_STR), (">", CODE_TAG)],
            [("      <span ", CODE_TAG), ("class=", CODE_PROP), ("\"metric-val\"", CODE_STR), (">", CODE_TAG), ("от 250 ₽", CODE_TEXT), ("</span>", CODE_TAG)],
            [("      <span ", CODE_TAG), ("class=", CODE_PROP), ("\"metric-lbl\"", CODE_STR), (">", CODE_TAG), ("Почасовая аренда", CODE_TEXT), ("</span>", CODE_TAG)],
            [("    </div>", CODE_TAG)],
            [("  </div>", CODE_TAG)],
            [("</div>", CODE_TAG)],
            [(" ", CODE_TEXT)],
            [(".hero-bottom", CODE_KW), (" {", CODE_TEXT)],
            [("  ", CODE_TEXT), ("border-top", CODE_PROP), (": ", CODE_TEXT), ("2px", CODE_VAL), (" ", CODE_TEXT), ("solid", CODE_VAL), (" ", CODE_TEXT), ("#222222", CODE_STR), ("; ", CODE_TEXT), ("padding-top", CODE_PROP), (": ", CODE_TEXT), ("20px", CODE_VAL), (";", CODE_TEXT)],
            [("  ", CODE_TEXT), ("display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), ("; ", CODE_TEXT), ("justify-content", CODE_PROP), (": ", CODE_TEXT), ("space-between", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [(".metric-val", CODE_KW), (" { ", CODE_TEXT), ("font-size", CODE_PROP), (": ", CODE_TEXT), ("20px", CODE_VAL), ("; ", CODE_TEXT), ("font-weight", CODE_PROP), (": ", CODE_TEXT), ("800", CODE_VAL), ("; ", CODE_TEXT), ("color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_STR), ("; }", CODE_TEXT)]
        ],
        "right_title": "АРХИТЕКТУРА НИЖНЕЙ ЧАСТИ HERO",
        "bullets": [
            {
                "title": "Графическая черта border-top: 2px solid:",
                "desc": "Контрастная черная линия отделяет масштабный заголовок от подробностей, задавая строгий журнальный стиль.",
                "height": 56
            },
            {
                "title": "Раскладка justify-content: space-between:",
                "desc": "Подзаголовок аккуратно размещается слева (max-width: 600px), а метрики выстраиваются в правой части экрана.",
                "height": 54
            },
            {
                "title": "Вертикальная колонка .metric-item:",
                "desc": "Свойство display: flex; flex-direction: column ставит яркую цифру 24/7 строго над маленькой подписью «Доступ».",
                "height": 56
            },
            {
                "title": "Контраст типографики:",
                "desc": "Значение (.metric-val): 20px, жирный шрифт 800, синий цвет. Подпись (.metric-lbl): 12px, серый нейтральный цвет #777777.",
                "height": 52
            }
        ]
    },

    # 12. Заголовки внутренних страниц
    {
        "type": "code_exp",
        "category": "УНИВЕРСАЛЬНЫЕ ЗАГОЛОВКИ СТРАНИЦ",
        "question": "Зачем в стилях классы .page-title и .page-subtitle?",
        "left_title": "css/style.css (.page-title)",
        "code": [
            [("/* ========================================= */", CODE_COMM)],
            [("/* СТИЛИ ДЛЯ ВНУТРЕННИХ СТРАНИЦ ПРОЕКТА      */", CODE_COMM)],
            [("/* ========================================= */", CODE_COMM)],
            [(" ", CODE_TEXT)],
            [("/* Главный заголовок страниц Каталога и Входа */", CODE_COMM)],
            [(".page-title {", CODE_KW)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("28px", CODE_VAL), (";", CODE_TEXT)],
            [("  font-weight", CODE_PROP), (": ", CODE_TEXT), ("800", CODE_VAL), (";", CODE_TEXT)],
            [("  letter-spacing", CODE_PROP), (": ", CODE_TEXT), ("-0.02em", CODE_VAL), (";", CODE_TEXT)],
            [("  margin-bottom", CODE_PROP), (": ", CODE_TEXT), ("8px", CODE_VAL), (";", CODE_TEXT)],
            [("  text-align", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(" ", CODE_TEXT)],
            [("/* Серый центрированный подзаголовок */", CODE_COMM)],
            [(".page-subtitle {", CODE_KW)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#666666", CODE_VAL), (";", CODE_TEXT)],
            [("  text-align", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";", CODE_TEXT)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("15px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)]
        ],
        "right_title": "ПРОЕКТИРОВАНИЕ НА ОПЕРЕЖЕНИЕ",
        "bullets": [
            {
                "title": "Задел для следующих вебинаров:",
                "desc": "На 3 и 4 вебинарах мы создадим catalog.html и login.html. Классы .page-title и .page-subtitle уже готовы для них!",
                "height": 56
            },
            {
                "title": "text-align: center:",
                "desc": "Центрирует текст по горизонтали внутри контейнера, создавая симметричную шапку внутренних разделов.",
                "height": 52
            },
            {
                "title": "Принцип DRY (Don't Repeat Yourself):",
                "desc": "Мы не пишем заново оформление заголовков для каждой страницы. Мы создаем единый переиспользуемый класс.",
                "height": 54
            },
            {
                "title": "Визуальная иерархия:",
                "desc": "Крупный жирный заголовок 28px и деликатный серый подзаголовок 15px создают гармоничную композицию.",
                "height": 52
            }
        ]
    },

    # 13. Подвал сайта: HTML и контакты
    {
        "type": "code_exp",
        "category": "ПОДВАЛ САЙТА (FOOTER)",
        "question": "Как оформить кликабельные контакты и юридический блок?",
        "left_title": "index.html & style.css (footer)",
        "code": [
            [("<footer ", CODE_TAG), ("class=", CODE_PROP), ("\"footer\"", CODE_STR), (">", CODE_TAG)],
            [("  <div ", CODE_TAG), ("class=", CODE_PROP), ("\"container footer-container\"", CODE_STR), (">", CODE_TAG)],
            [("    ", CODE_TEXT)],
            [("    <!-- Левая колонка: описание и копирайт -->", CODE_COMM)],
            [("    <div ", CODE_TAG), ("class=", CODE_PROP), ("\"footer-info\"", CODE_STR), (">", CODE_TAG)],
            [("      <p><strong>СмартОфис</strong> — Сервис бронирования</p>", CODE_TEXT)],
            [("      <p>© 2026 СмартОфис. Все права защищены.</p>", CODE_TEXT)],
            [("    </div>", CODE_TAG)],
            [("    ", CODE_TEXT)],
            [("    <!-- Правая колонка: интерактивные контакты -->", CODE_COMM)],
            [("    <div ", CODE_TAG), ("class=", CODE_PROP), ("\"footer-contacts\"", CODE_STR), (">", CODE_TAG)],
            [("      <p>Email: <a ", CODE_TEXT), ("href=", CODE_PROP), ("\"mailto:info@smartoffice.ru\"", CODE_STR), (">info@smartoffice.ru</a></p>", CODE_TEXT)],
            [("      <p>Телефон: <a ", CODE_TEXT), ("href=", CODE_PROP), ("\"tel:+78005553535\"", CODE_STR), (">+7 (800) 555-35-35</a></p>", CODE_TEXT)],
            [("    </div>", CODE_TAG)],
            [("  </div>", CODE_TAG)],
            [("</footer>", CODE_TAG)],
            [(" ", CODE_TEXT)],
            [(".footer", CODE_KW), (" { ", CODE_TEXT), ("background-color", CODE_PROP), (": ", CODE_TEXT), ("#f8f9fa", CODE_STR), ("; ", CODE_TEXT), ("border-top", CODE_PROP), (": ", CODE_TEXT), ("1px solid #dddddd", CODE_VAL), ("; ", CODE_TEXT), ("padding", CODE_PROP), (": ", CODE_TEXT), ("25px 0", CODE_VAL), ("; }", CODE_TEXT)],
            [(".footer-contacts a", CODE_KW), (" { ", CODE_TEXT), ("color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_STR), ("; ", CODE_TEXT), ("text-decoration", CODE_PROP), (": ", CODE_TEXT), ("none", CODE_VAL), ("; }", CODE_TEXT)]
        ],
        "right_title": "СТАНДАРТЫ ОФОРМЛЕНИЯ ПОДВАЛА",
        "bullets": [
            {
                "title": "Протокол mailto: в ссылках:",
                "desc": "При клике на ссылку href=\"mailto:...\" устройство пользователя автоматически открывает почтовую программу с заполненным адресом.",
                "height": 58
            },
            {
                "title": "Протокол tel: для быстрых звонков:",
                "desc": "Ссылка href=\"tel:+78005553535\" на смартфонах сразу запускает набор телефонного номера. Это стандарт мобильного UX.",
                "height": 58
            },
            {
                "title": "Тег <strong> для акцента:",
                "desc": "Выделяет название сервиса жирным шрифтом в первом абзаце футера.",
                "height": 48
            },
            {
                "title": "Фон #f8f9fa и разделитель #dddddd:",
                "desc": "Мягкий светлый фон и тонкая верхняя рамка деликатно завершают страницу в самом низу.",
                "height": 50
            }
        ]
    },

    # 14. Техника Sticky Footer
    {
        "type": "code_exp",
        "category": "ТЕХНИКА FLEXBOX: ПРИЖАТИЕ ФУТЕРА",
        "question": "Почему футер не улетает вверх при малом количестве контента?",
        "left_title": "css/style.css (Sticky Footer)",
        "code": [
            [("/* ШАГ 1: Тело сайта тянется минимум на 100vh */", CODE_COMM)],
            [("body {", CODE_KW)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  flex-direction", CODE_PROP), (": ", CODE_TEXT), ("column", CODE_VAL), (";", CODE_TEXT)],
            [("  min-height", CODE_PROP), (": ", CODE_TEXT), ("100vh", CODE_VAL), ("; /* 100% высоты окна */", CODE_COMM)],
            [("}", CODE_KW)],
            [(" ", CODE_TEXT)],
            [("/* ШАГ 2: Блок main забирает всё свободное место */", CODE_COMM)],
            [(".main {", CODE_KW)],
            [("  flex", CODE_PROP), (": ", CODE_TEXT), ("1", CODE_VAL), ("; /* Жадно растягивается */", CODE_COMM)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("40px 0", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(" ", CODE_TEXT)],
            [("/* ШАГ 3: Подвал гарантированно прижат вниз */", CODE_COMM)],
            [(".footer {", CODE_KW)],
            [("  margin-top", CODE_PROP), (": ", CODE_TEXT), ("auto", CODE_VAL), ("; /* Выталкивается вниз */", CODE_COMM)],
            [("}", CODE_KW)]
        ],
        "right_title": "МЕХАНИКА «ЛИПКОГО» ПОДВАЛА",
        "bullets": [
            {
                "title": "Классическая беда новичков:",
                "desc": "Если на странице мало текста, футер поднимается к середине монитора, а под ним висит пустота. Это выглядит непрофессионально.",
                "height": 56
            },
            {
                "title": "Шаг 1: body на 100vh:",
                "desc": "Окно браузера делится на флекс-колонку, которая всегда занимает как минимум 100% высоты экрана.",
                "height": 52
            },
            {
                "title": "Шаг 2: flex: 1 для .main:",
                "desc": "Магическое свойство flex: 1 говорит браузеру: «Растяни блок main на всю оставшуюся пустоту страницы».",
                "height": 54
            },
            {
                "title": "Шаг 3: margin-top: auto:",
                "desc": "Финальная страховка. Футер выталкивается в самый нижний край экрана при любом объеме информации!",
                "height": 52
            }
        ]
    },

    # 15. Toast-уведомления и анимации
    {
        "type": "code_exp",
        "category": "TOAST-УВЕДОМЛЕНИЯ И CSS-АНИМАЦИИ",
        "question": "Как работают всплывающие системные уведомления?",
        "left_title": "css/style.css (Toast & Animation)",
        "code": [
            [("/* Контейнер прикреплен к правому нижнему углу */", CODE_COMM)],
            [(".toast-container {", CODE_KW)],
            [("  position", CODE_PROP), (": ", CODE_TEXT), ("fixed", CODE_VAL), (";", CODE_TEXT)],
            [("  bottom", CODE_PROP), (": ", CODE_TEXT), ("24px", CODE_VAL), ("; ", CODE_TEXT), ("right", CODE_PROP), (": ", CODE_TEXT), ("24px", CODE_VAL), (";", CODE_TEXT)],
            [("  z-index", CODE_PROP), (": ", CODE_TEXT), ("1000", CODE_VAL), (";", CODE_TEXT)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), ("; ", CODE_TEXT), ("flex-direction", CODE_PROP), (": ", CODE_TEXT), ("column", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(".toast {", CODE_KW)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#222222", CODE_VAL), ("; ", CODE_TEXT), ("color", CODE_PROP), (": ", CODE_TEXT), ("#ffffff", CODE_VAL), (";", CODE_TEXT)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("14px 18px", CODE_VAL), ("; ", CODE_TEXT), ("border-radius", CODE_PROP), (": ", CODE_TEXT), ("6px", CODE_VAL), (";", CODE_TEXT)],
            [("  animation", CODE_PROP), (": ", CODE_TEXT), ("toastFadeIn 0.3s ease", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_KW)],
            [(".toast-success", CODE_KW), (" { ", CODE_TEXT), ("border-left", CODE_PROP), (": ", CODE_TEXT), ("4px solid #28a745", CODE_VAL), ("; }", CODE_TEXT)],
            [(".toast-danger", CODE_KW), ("  { ", CODE_TEXT), ("border-left", CODE_PROP), (": ", CODE_TEXT), ("4px solid #dc3545", CODE_VAL), ("; }", CODE_TEXT)],
            [("/* Плавная анимация всплытия снизу */", CODE_COMM)],
            [("@keyframes toastFadeIn", CODE_KW), (" {", CODE_TEXT)],
            [("  ", CODE_TEXT), ("from", CODE_KW), (" { ", CODE_TEXT), ("opacity", CODE_PROP), (": ", CODE_TEXT), ("0", CODE_VAL), ("; ", CODE_TEXT), ("transform", CODE_PROP), (": ", CODE_TEXT), ("translateY(15px)", CODE_VAL), ("; }", CODE_TEXT)],
            [("  ", CODE_TEXT), ("to", CODE_KW), ("   { ", CODE_TEXT), ("opacity", CODE_PROP), (": ", CODE_TEXT), ("1", CODE_VAL), ("; ", CODE_TEXT), ("transform", CODE_PROP), (": ", CODE_TEXT), ("translateY(0)", CODE_VAL), ("; }", CODE_TEXT)],
            [("}", CODE_KW)]
        ],
        "right_title": "ПОЗИЦИОНИРОВАНИЕ И ЭФФЕКТЫ",
        "bullets": [
            {
                "title": "position: fixed (Фиксация на экране):",
                "desc": "Контейнер уведомлений вырван из общего потока и привязан к углу монитора. При прокрутке страницы он остается на месте.",
                "height": 56
            },
            {
                "title": "Свойство z-index: 1000:",
                "desc": "Определяет слой по высоте. Значение 1000 гарантирует, что всплывающее сообщение перекроет любой другой элемент сайта.",
                "height": 56
            },
            {
                "title": "Цветовой индикатор статуса:",
                "desc": "border-left: 4px solid #28a745 красит левую грань в зеленый при успехе и в красный #dc3545 при ошибке.",
                "height": 52
            },
            {
                "title": "Анимация @keyframes toastFadeIn:",
                "desc": "Сообщение появляется из прозрачности (opacity: 0 -> 1) и сдвигается на 15px снизу вверх за треть секунды (0.3s).",
                "height": 54
            }
        ]
    },

    # 16. Типичные ошибки новичков
    {
        "type": "grid_errors",
        "category": "РАЗБОР ТИПИЧНЫХ ОШИБОК",
        "question": "На чем чаще всего спотыкаются студенты при первой верстке?",
        "cards": [
            {
                "title": "Забыли rel=\"stylesheet\"",
                "problem": "Написали <link href=\"css/style.css\"> без атрибута rel=\"stylesheet\".",
                "solution": "Браузер скачает файл, но не поймет, что это таблица стилей. Обязательно пишите rel=\"stylesheet\"!"
            },
            {
                "title": "Опечатка в пути к стилям или иконке",
                "problem": "Указали href=\"style.css\" вместо правильного относительного пути href=\"css/style.css\".",
                "solution": "Стили лежат во вложенной папке css/. Проверяйте пути в проводнике VS Code или WebStorm."
            },
            {
                "title": "Незакрытые парные теги",
                "problem": "Забыли закрыть тег </div>, <header>, <section> или <ul>.",
                "solution": "Один незакрытый тег ломает всю сетку Flexbox дальше по коду. Пользуйтесь подсветкой парных тегов в IDE."
            },
            {
                "title": "Пропустили сброс box-sizing: border-box",
                "problem": "Не написали сброс * { box-sizing: border-box; } в самом начале CSS.",
                "solution": "Любой padding и border будут увеличивать ширину блоков, создавая неприятный горизонтальный скролл."
            }
        ]
    },

    # 17. Чек-лист самопроверки
    {
        "type": "checklist",
        "category": "ИТОГИ И ЧЕК-ЛИСТ САМОПРОВЕРКИ",
        "question": "Как убедиться, что код вебинара написан на 100% правильно?",
        "title": "ЧЕК-ЛИСТ ГОТОВНОСТИ ПРАКТИЧЕСКОЙ ЧАСТИ",
        "items": [
            {
                "title": "Файловая структура проекта организована корректно:",
                "desc": "Файл index.html находится в корне, таблица стилей в папке css/style.css, векторный логотип в img/logo.svg."
            },
            {
                "title": "Кодировка и шрифты подключены без сбоев:",
                "desc": "Указан <meta charset=\"UTF-8\">, шрифт Inter подгружается из Google Fonts, текст не содержит символов-кракозябр."
            },
            {
                "title": "Семантический каркас страницы валиден:",
                "desc": "Использованы семантические теги HTML5: <header>, <nav>, <main>, <section class=\"hero-section\">, <footer>."
            },
            {
                "title": "Контент центрирован внутри .container на 1200px:",
                "desc": "Шапка, hero-блок и футер аккуратно держатся по центру экрана с полями безопасности благодаря margin: 0 auto."
            },
            {
                "title": "Футер надежно прижат к низу страницы:",
                "desc": "Благодаря связке body { min-height: 100vh; flex-direction: column; } и main { flex: 1; } футер никогда не зависает в воздухе."
            }
        ]
    }
]

# Generate slides
for cfg in slides_config:
    slide = prs.slides.add_slide(blank_layout)
    if cfg["type"] == "code_exp":
        create_code_explanation_slide(
            slide,
            cfg["category"],
            cfg["question"],
            cfg["left_title"],
            cfg["code"],
            cfg["right_title"],
            cfg["bullets"]
        )
    elif cfg["type"] == "grid_errors":
        create_grid_cards_slide(
            slide,
            cfg["category"],
            cfg["question"],
            "ТИПИЧНЫЕ ОШИБКИ ПРИ ВЕРСТКЕ",
            cfg["cards"]
        )
    elif cfg["type"] == "checklist":
        create_checklist_slide(
            slide,
            cfg["category"],
            cfg["question"],
            cfg["title"],
            cfg["items"]
        )

# Now reorder slides:
# Original slides: [0: Title, 1: Plan, 2: Result, 3: Goodbye]
# We added 15 slides: indices 4 .. 18
# New order should be: Title (0), Plan (1), New Slides (4 .. 18), Result (2), Goodbye (3)
sldIdLst = prs.slides._sldIdLst
slide_ids = list(sldIdLst)
title_slide = slide_ids[0]
plan_slide = slide_ids[1]
result_slide = slide_ids[2]
goodbye_slide = slide_ids[3]
new_slide_ids = slide_ids[4:]

new_order = [title_slide, plan_slide] + new_slide_ids + [result_slide, goodbye_slide]
sldIdLst[:] = new_order

# Применяем форматирование для TextBox 1 в правом верхнем углу на ВСЕХ слайдах:
for slide in prs.slides:
    for shape in slide.shapes:
        # Проверяем TextBox 1 в верхней правой зоне
        if shape.name == "TextBox 1" and shape.has_text_frame and shape.top.pt < 60 and shape.left.pt > 200:
            tf = shape.text_frame
            tf.word_wrap = False  # Без переноса по словам
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.RIGHT
                for r in p.runs:
                    r.font.name = "Inter"
                    r.font.size = Pt(14)
                    r.font.bold = False
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # Белый цвет

fallback_pptx = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx"

# Сохраняем в оба файла (обновленный и дополненный), если они доступны
for target_path in [output_pptx, fallback_pptx]:
    try:
        prs.save(target_path)
        print(f"Презентация успешно сохранена в: {target_path}")
    except PermissionError:
        print(f"Файл '{os.path.basename(target_path)}' заблокирован для записи.")

print(f"Всего слайдов в обновленной презентации: {len(prs.slides)}")
