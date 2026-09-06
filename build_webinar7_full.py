import os
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

# -------------------------------------------------------------
# PATH CONFIGURATION
# -------------------------------------------------------------
base_template_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx"
output_pptx_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-07-slider-timers\вебинар 7.pptx"
result_img_path = r"d:\IEEU International East-European University\fullstack-web-development-course\webinar7_result_exact.png"

# -------------------------------------------------------------
# PALETTE & STYLES (Strict corporate college theme)
# -------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x18, 0x1A, 0x1F)
CODE_BORDER = RGBColor(0x2D, 0x31, 0x39)
CODE_TEXT = RGBColor(0xD4, 0xD4, 0xD8)
ORANGE_PILL = RGBColor(0xFE, 0x60, 0x02)
ORANGE_LINE = RGBColor(0xFE, 0x60, 0x02)
QUESTION_TEXT = RGBColor(0xA6, 0xA1, 0xA1) # Strictly #A6A1A1
TEXT_PRIMARY = RGBColor(0x18, 0x1A, 0x1F)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
KEYWORD_COLOR = RGBColor(0xC6, 0x78, 0xDD) # Purple
FUNC_COLOR = RGBColor(0x61, 0xAF, 0xEF)    # Blue
STR_COLOR = RGBColor(0x98, 0xC3, 0x79)     # Green
COMMENT_COLOR = RGBColor(0x5C, 0x63, 0x70) # Gray
NUM_COLOR = RGBColor(0xD1, 0x9A, 0x66)     # Orange / number
TAG_COLOR = RGBColor(0xE0, 0x6C, 0x75)     # Red/Coral
DANGER_COLOR = RGBColor(0xDC, 0x26, 0x26)  # Red for invalid state

def strip_shape_styles_and_shadows(shape):
    """Zero tolerance to shadows and PowerPoint theme styles."""
    spPr = shape._element.spPr
    for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst'):
        spPr.remove(shd)
    for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw'):
        spPr.remove(shd)
    style = shape._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
    if style is not None:
        shape._element.remove(style)

def add_header_and_footer(slide, category_text, question_text):
    """Standard header pill badge with dynamic width and bottom question prompt."""
    pill_w = Pt(max(140, len(category_text) * 7.8 + 24))
    pill_h = Pt(22.6)
    pill_left = Pt(700) - pill_w
    pill_top = Pt(20.4)
    
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_left, pill_top, pill_w, pill_h)
    badge.name = "Google Shape;69;p14"
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE_PILL
    badge.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    badge.line.width = Pt(1)
    strip_shape_styles_and_shadows(badge)
    
    tf_b = badge.text_frame
    tf_b.word_wrap = False
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = category_text
    r_b.font.name = "Inter"
    r_b.font.size = Pt(12)
    r_b.font.bold = False
    r_b.font.color.rgb = WHITE
    
    # TextBox 2 at bottom
    tb2 = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2.name = "TextBox 2"
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = question_text
    r2.font.name = "Montserrat"
    r2.font.size = Pt(10)
    r2.font.bold = False
    r2.font.color.rgb = QUESTION_TEXT

def create_code_explanation_slide(slide, category, question, left_title, code_lines, right_title, bullet_items):
    """Standard 2-column layout: Left code card, Right clean explanation with bullet items."""
    add_header_and_footer(slide, category, question)
    
    # Left Card
    card_x = Pt(24)
    card_y = Pt(56)
    card_w = Pt(325)
    card_h = Pt(296)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
    card.name = "Rounded Rectangle 1"
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG
    card.line.color.rgb = CODE_BORDER
    card.line.width = Pt(1)
    strip_shape_styles_and_shadows(card)
    
    # Left Card Title
    tb_lt = slide.shapes.add_textbox(card_x + Pt(14), card_y + Pt(10), card_w - Pt(28), Pt(20))
    tb_lt.name = "TextBox Left Title"
    tf_lt = tb_lt.text_frame
    tf_lt.margin_left = tf_lt.margin_top = tf_lt.margin_right = tf_lt.margin_bottom = 0
    p_lt = tf_lt.paragraphs[0]
    r_lt = p_lt.add_run()
    r_lt.text = left_title
    r_lt.font.name = "Montserrat"
    r_lt.font.size = Pt(11)
    r_lt.font.bold = True
    r_lt.font.color.rgb = WHITE
    
    # Inner Code Box
    box_x = card_x + Pt(14)
    box_y = card_y + Pt(32)
    box_w = card_w - Pt(28)
    box_h = card_h - Pt(42)
    
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, box_y, box_w, box_h)
    code_box.name = "Rounded Rectangle Code Box"
    code_box.adjustments[0] = 0.03
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x2A)
    code_box.line.color.rgb = RGBColor(0x36, 0x3B, 0x44)
    code_box.line.width = Pt(1)
    strip_shape_styles_and_shadows(code_box)
    
    # Code Text
    tb_code = slide.shapes.add_textbox(box_x + Pt(10), box_y + Pt(8), box_w - Pt(20), box_h - Pt(16))
    tb_code.name = "TextBox Code Text"
    tf_c = tb_code.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    
    for i, line_runs in enumerate(code_lines):
        p = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
        p.space_after = Pt(2)
        for token_text, token_color in line_runs:
            r = p.add_run()
            r.text = token_text
            r.font.name = "Consolas"
            r.font.size = Pt(8.5)
            r.font.bold = False
            r.font.color.rgb = token_color
            
    # Right Column
    rcol_x = Pt(365)
    rcol_y = Pt(56)
    rcol_w = Pt(335)
    
    tb_rt = slide.shapes.add_textbox(rcol_x, rcol_y, rcol_w, Pt(22))
    tb_rt.name = "TextBox Right Title"
    tf_rt = tb_rt.text_frame
    tf_rt.margin_left = tf_rt.margin_top = tf_rt.margin_right = tf_rt.margin_bottom = 0
    p_rt = tf_rt.paragraphs[0]
    r_rt = p_rt.add_run()
    r_rt.text = right_title
    r_rt.font.name = "Montserrat"
    r_rt.font.size = Pt(13)
    r_rt.font.bold = True
    r_rt.font.color.rgb = TEXT_PRIMARY
    
    # Orange divider
    rdiv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rcol_x, rcol_y + Pt(25), rcol_w, Pt(1))
    rdiv.name = "Rectangle 8"
    rdiv.fill.solid()
    rdiv.fill.fore_color.rgb = ORANGE_LINE
    rdiv.line.fill.background()
    strip_shape_styles_and_shadows(rdiv)
    
    cur_y = rcol_y + Pt(34)
    for idx, item in enumerate(bullet_items):
        item_h = Pt(item.get("height", 42))
        
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rcol_x, cur_y + Pt(3), Pt(7), Pt(7))
        dot.name = f"Rounded Rectangle Dot {idx}"
        dot.adjustments[0] = 0.2
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)
        
        tb_item = slide.shapes.add_textbox(rcol_x + Pt(14), cur_y, rcol_w - Pt(14), item_h)
        tb_item.name = f"TextBox Bullet {idx}"
        tf_i = tb_item.text_frame
        tf_i.word_wrap = True
        tf_i.margin_left = tf_i.margin_top = tf_i.margin_right = tf_i.margin_bottom = 0
        
        p = tf_i.paragraphs[0]
        p.space_after = Pt(0)
        
        r_title = p.add_run()
        r_title.text = item["title"] + ":\n"
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = TEXT_PRIMARY
        
        r_desc = p.add_run()
        r_desc.text = item["desc"]
        r_desc.font.name = "Inter"
        r_desc.font.size = Pt(8.5)
        r_desc.font.bold = False
        r_desc.font.color.rgb = TEXT_MUTED
        
        cur_y += item_h + Pt(6)

def create_function_arch_diagram_slide(slide, category, question):
    """Custom Diagram 1: Function Architecture (Inputs -> Body -> Output & DOM)."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. ВХОД (АРГУМЕНТЫ)", "subtitle": "Параметры и defaults", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. ОБРАБОТКА (ТЕЛО)", "subtitle": "Чистая логика и расчет", "accent": NUM_COLOR},
        {"x": Pt(476), "title": "3. ВЫХОД И ЭФФЕКТЫ", "subtitle": "Return и реакция DOM", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = c["title"]
        r_t1.font.name = "Montserrat"
        r_t1.font.size = Pt(10)
        r_t1.font.bold = True
        r_t1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = c["subtitle"]
        r_t2.font.name = "Inter"
        r_t2.font.size = Pt(8.5)
        r_t2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = c["accent"]
        div.line.fill.background()
        strip_shape_styles_and_shadows(div)
        
    box_y = top_y + Pt(52)
    box_h = Pt(54)
    
    # Box 1
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(36), box_y, Pt(186), box_h)
    b1.adjustments[0] = 0.08
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(0x1F, 0x24, 0x2E)
    b1.line.color.rgb = FUNC_COLOR
    b1.line.width = Pt(1)
    strip_shape_styles_and_shadows(b1)
    tf_b1 = b1.text_frame
    tf_b1.margin_top = Pt(6)
    p_b1 = tf_b1.paragraphs[0]
    p_b1.alignment = PP_ALIGN.CENTER
    r_b1 = p_b1.add_run()
    r_b1.text = "function showSlide(\n  index = 0,\n  step = 1\n) {"
    r_b1.font.name = "Consolas"
    r_b1.font.size = Pt(7.5)
    r_b1.font.color.rgb = CODE_TEXT
    
    # Box 2
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(262), box_y, Pt(186), box_h)
    b2.adjustments[0] = 0.08
    b2.fill.solid()
    b2.fill.fore_color.rgb = RGBColor(0x28, 0x24, 0x1E)
    b2.line.color.rgb = NUM_COLOR
    b2.line.width = Pt(1)
    strip_shape_styles_and_shadows(b2)
    tf_b2 = b2.text_frame
    tf_b2.margin_top = Pt(6)
    p_b2 = tf_b2.paragraphs[0]
    p_b2.alignment = PP_ALIGN.CENTER
    r_b2 = p_b2.add_run()
    r_b2.text = "const len = slides.length;\nlet next = (index + len) % len;\ncurrentSlide = next;"
    r_b2.font.name = "Consolas"
    r_b2.font.size = Pt(7.5)
    r_b2.font.color.rgb = CODE_TEXT
    
    # Box 3
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(488), box_y, Pt(186), box_h)
    b3.adjustments[0] = 0.08
    b3.fill.solid()
    b3.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x20)
    b3.line.color.rgb = STR_COLOR
    b3.line.width = Pt(1)
    strip_shape_styles_and_shadows(b3)
    tf_b3 = b3.text_frame
    tf_b3.margin_top = Pt(6)
    p_b3 = tf_b3.paragraphs[0]
    p_b3.alignment = PP_ALIGN.CENTER
    r_b3 = p_b3.add_run()
    r_b3.text = "slides.forEach((s, i) => {\n  s.classList.toggle('active', ...);\n});\nreturn next;"
    r_b3.font.name = "Consolas"
    r_b3.font.size = Pt(7.5)
    r_b3.font.color.rgb = CODE_TEXT
    
    bullets = [
        {"x": Pt(36), "header": "Параметры и изоляция:", "items": [
            "Входные данные передаются через аргументы",
            "Дефолтные значения (index = 0) для защиты",
            "Параметры изолированы внутри тела функции",
            "Исключение зависимости от глобальных данных",
            "Функция как универсальный черный ящик"
        ]},
        {"x": Pt(262), "header": "Вычислительная логика:", "items": [
            "Расчет нового индекса без мутации данных",
            "Кольцевой переход: после 3 идет 0, до 0 идет 3",
            "Исключение непредвиденных side-effects",
            "Детерминированный предсказуемый расчет",
            "Легкость проверки юнит-тестами"
        ]},
        {"x": Pt(488), "header": "Эффект в интерфейсе:", "items": [
            "Возврат вычисленного значения return",
            "Синхронизация классов .active у слайдов",
            "Подсветка круглой точки-индикатора (Dot)",
            "Плавный transition прозрачности 0.3s",
            "Интерфейс приведен в консистентное состояние"
        ]}
    ]
    
    bullets_y = box_y + box_h + Pt(10)
    for b in bullets:
        tb = slide.shapes.add_textbox(b["x"], bullets_y, Pt(186), Pt(165))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_head = tf.paragraphs[0]
        r_head = p_head.add_run()
        r_head.text = b["header"] + "\n"
        r_head.font.name = "Montserrat"
        r_head.font.size = Pt(8.5)
        r_head.font.bold = True
        r_head.font.color.rgb = WHITE
        
        for it in b["items"]:
            p_it = tf.add_paragraph()
            p_it.space_after = Pt(2)
            r_dot = p_it.add_run()
            r_dot.text = "• "
            r_dot.font.name = "Inter"
            r_dot.font.size = Pt(8.0)
            r_dot.font.color.rgb = ORANGE_PILL
            
            r_txt = p_it.add_run()
            r_txt.text = it
            r_txt.font.name = "Inter"
            r_txt.font.size = Pt(8.0)
            r_txt.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD8)

def create_timers_lifecycle_diagram_slide(slide, category, question):
    """Custom Diagram 2: Timers & Event Lifecycle (Auto-scroll -> User Click -> Reset & Restart)."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. АВТОПРОКРУТКА", "subtitle": "setInterval(next, 3000)", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. РУЧНОЙ КЛИК", "subtitle": "Стрелки ‹ › и Dots", "accent": DANGER_COLOR},
        {"x": Pt(476), "title": "3. СБРОС И РЕСТАРТ", "subtitle": "clearInterval(timerId)", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = c["title"]
        r_t1.font.name = "Montserrat"
        r_t1.font.size = Pt(10)
        r_t1.font.bold = True
        r_t1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = c["subtitle"]
        r_t2.font.name = "Inter"
        r_t2.font.size = Pt(8.5)
        r_t2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = c["accent"]
        div.line.fill.background()
        strip_shape_styles_and_shadows(div)
        
    box_y = top_y + Pt(52)
    box_h = Pt(54)
    
    # Box 1
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(36), box_y, Pt(186), box_h)
    b1.adjustments[0] = 0.08
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(0x1F, 0x24, 0x2E)
    b1.line.color.rgb = FUNC_COLOR
    b1.line.width = Pt(1)
    strip_shape_styles_and_shadows(b1)
    tf_b1 = b1.text_frame
    tf_b1.margin_top = Pt(6)
    p_b1 = tf_b1.paragraphs[0]
    p_b1.alignment = PP_ALIGN.CENTER
    r_b1 = p_b1.add_run()
    r_b1.text = "function startAuto() {\n  stopAuto();\n  timerId = setInterval(next, 3000);\n}"
    r_b1.font.name = "Consolas"
    r_b1.font.size = Pt(7.5)
    r_b1.font.color.rgb = CODE_TEXT
    
    # Box 2
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(262), box_y, Pt(186), box_h)
    b2.adjustments[0] = 0.08
    b2.fill.solid()
    b2.fill.fore_color.rgb = RGBColor(0x2D, 0x1E, 0x20)
    b2.line.color.rgb = DANGER_COLOR
    b2.line.width = Pt(1)
    strip_shape_styles_and_shadows(b2)
    tf_b2 = b2.text_frame
    tf_b2.margin_top = Pt(6)
    p_b2 = tf_b2.paragraphs[0]
    p_b2.alignment = PP_ALIGN.CENTER
    r_b2 = p_b2.add_run()
    r_b2.text = "nextBtn.addEventListener('click',\n  () => {\n    next();\n    startAuto();\n});"
    r_b2.font.name = "Consolas"
    r_b2.font.size = Pt(7.5)
    r_b2.font.color.rgb = CODE_TEXT
    
    # Box 3
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(488), box_y, Pt(186), box_h)
    b3.adjustments[0] = 0.08
    b3.fill.solid()
    b3.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x20)
    b3.line.color.rgb = STR_COLOR
    b3.line.width = Pt(1)
    strip_shape_styles_and_shadows(b3)
    tf_b3 = b3.text_frame
    tf_b3.margin_top = Pt(6)
    p_b3 = tf_b3.paragraphs[0]
    p_b3.alignment = PP_ALIGN.CENTER
    r_b3 = p_b3.add_run()
    r_b3.text = "function stopAuto() {\n  if (timerId) {\n    clearInterval(timerId);\n  }\n}"
    r_b3.font.name = "Consolas"
    r_b3.font.size = Pt(7.5)
    r_b3.font.color.rgb = CODE_TEXT
    
    bullets = [
        {"x": Pt(36), "header": "Регулярный цикл:", "items": [
            "Автоматический тик каждые 3000 миллисекунд",
            "Фоновая работа через Web API браузера",
            "Вызов функции next() и смена кадра",
            "Не блокирует основной поток выполнения UI",
            "Постоянное движение медиа-контента"
        ]},
        {"x": Pt(262), "header": "Событие пользователя:", "items": [
            "Клик по кнопке «‹» или «›», либо по Dot",
            "Мгновенный показ выбранного изображения",
            "Прерывание стандартного тайминга",
            "Запрос на рестарт интервала",
            "Пользователь получает мгновенный отклик"
        ]},
        {"x": Pt(488), "header": "Защита от сбоев UX:", "items": [
            "Остановка таймера перед созданием нового",
            "Исключение наложения нескольких setInterval",
            "Полные 3 секунды на просмотр после клика",
            "Слайдер не переключается сразу после клика",
            "Полный контроль над утечками ресурсов"
        ]}
    ]
    
    bullets_y = box_y + box_h + Pt(10)
    for b in bullets:
        tb = slide.shapes.add_textbox(b["x"], bullets_y, Pt(186), Pt(165))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_head = tf.paragraphs[0]
        r_head = p_head.add_run()
        r_head.text = b["header"] + "\n"
        r_head.font.name = "Montserrat"
        r_head.font.size = Pt(8.5)
        r_head.font.bold = True
        r_head.font.color.rgb = WHITE
        
        for it in b["items"]:
            p_it = tf.add_paragraph()
            p_it.space_after = Pt(2)
            r_dot = p_it.add_run()
            r_dot.text = "• "
            r_dot.font.name = "Inter"
            r_dot.font.size = Pt(8.0)
            r_dot.font.color.rgb = ORANGE_PILL
            
            r_txt = p_it.add_run()
            r_txt.text = it
            r_txt.font.name = "Inter"
            r_txt.font.size = Pt(8.0)
            r_txt.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD8)

# -------------------------------------------------------------
# SLIDE CONTENT DEFINITIONS (21 content slides)
# -------------------------------------------------------------
slides_data = [
    # Slide 3: Module 1
    {
        "type": "code",
        "category": "ФУНКЦИИ И АРГУМЕНТЫ",
        "question": "В чем принципиальная разница между параметрами и аргументами функции?",
        "left_title": "Объявление и вызов функции",
        "code_lines": [
            [("// Параметры функции при объявлении", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("setSlideActive", FUNC_COLOR), ("(slideIndex, isSmooth) {", CODE_TEXT)],
            [("  console.log(", CODE_TEXT), ("'Слайд:'", STR_COLOR), (", slideIndex);", CODE_TEXT)],
            [("  console.log(", CODE_TEXT), ("'Анимация:'", STR_COLOR), (", isSmooth);", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Передача аргументов при вызове", COMMENT_COLOR)],
            [("setSlideActive(", CODE_TEXT), ("2", NUM_COLOR), (", ", CODE_TEXT), ("true", KEYWORD_COLOR), (");", CODE_TEXT)],
            [("setSlideActive(", CODE_TEXT), ("0", NUM_COLOR), (", ", CODE_TEXT), ("false", KEYWORD_COLOR), (");", CODE_TEXT)],
        ],
        "right_title": "Параметры vs Аргументы",
        "bullet_items": [
            {"title": "Параметры функции", "desc": "Именованные переменные в круглых скобках объявления функции. Служат шаблоном для входящих данных.", "height": 38},
            {"title": "Аргументы вызова", "desc": "Фактические конкретные значения, которые передаются в функцию в момент её исполнения.", "height": 38},
            {"title": "Локальная видимость", "desc": "Параметры slideIndex и isSmooth существуют исключительно внутри тела функции.", "height": 38},
            {"title": "Повторное использование", "desc": "Единая функция переключения обслуживает любые слайды каталога «СмартОфис».", "height": 38},
        ]
    },

    # Slide 4: Module 1
    {
        "type": "code",
        "category": "ФУНКЦИИ И АРГУМЕНТЫ",
        "question": "Зачем указывать значения параметров по умолчанию при объявлении функции?",
        "left_title": "Параметры по умолчанию (ES6)",
        "code_lines": [
            [("// Синтаксис дефолтных значений", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("changeSlide", FUNC_COLOR), ("(step = ", CODE_TEXT), ("1", NUM_COLOR), (", autoLoop = ", CODE_TEXT), ("true", KEYWORD_COLOR), (") {", CODE_TEXT)],
            [("  let ", KEYWORD_COLOR), ("next = currentSlide + step;", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(autoLoop && next >= totalSlides) {", CODE_TEXT)],
            [("    next = ", CODE_TEXT), ("0", NUM_COLOR), (";", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("  return ", KEYWORD_COLOR), ("next;", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Вызовы с пропуском аргументов", COMMENT_COLOR)],
            [("changeSlide();      ", CODE_TEXT), ("// step=1, autoLoop=true", COMMENT_COLOR)],
            [("changeSlide(", CODE_TEXT), ("2", NUM_COLOR), (");     ", CODE_TEXT), ("// step=2, autoLoop=true", COMMENT_COLOR)],
        ],
        "right_title": "Значения по умолчанию (ES6)",
        "bullet_items": [
            {"title": "Синтаксис ES6", "desc": "Значение по умолчанию задается через оператор = прямо в списке параметров заголовка.", "height": 38},
            {"title": "Защита от undefined", "desc": "Если аргумент не передан или равен undefined, функция автоматически берет дефолт.", "height": 38},
            {"title": "Опциональные параметры", "desc": "Устраняет необходимость в громоздких проверках вида step = step || 1.", "height": 38},
            {"title": "Порядок параметров", "desc": "Обязательные параметры всегда располагают в начале, а параметры с дефолтами — в конце.", "height": 38},
        ]
    },

    # Slide 5: Diagram 1
    {
        "type": "diagram_arch",
        "category": "СХЕМА: АРХИТЕКТУРА ФУНКЦИИ",
        "question": "Как правильно разделять входные аргументы, тело функции и возвращаемый результат?"
    },

    # Slide 6: Module 1
    {
        "type": "code",
        "category": "ЧИСТЫЕ ФУНКЦИИ",
        "question": "Почему чистые функции проще тестировать и поддерживать в интерфейсах?",
        "left_title": "Чистая функция расчета индекса",
        "code_lines": [
            [("// Чистая функция без побочных эффектов", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("calcIndex", FUNC_COLOR), ("(current, total, step = ", CODE_TEXT), ("1", NUM_COLOR), (") {", CODE_TEXT)],
            [("  return ", KEYWORD_COLOR), ("(current + step + total) % total;", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Детерминированный результат:", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("next1 = calcIndex(", CODE_TEXT), ("0", NUM_COLOR), (", ", CODE_TEXT), ("4", NUM_COLOR), (", ", CODE_TEXT), ("1", NUM_COLOR), (");  ", CODE_TEXT), ("// 1", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("prev1 = calcIndex(", CODE_TEXT), ("0", NUM_COLOR), (", ", CODE_TEXT), ("4", NUM_COLOR), (", -", CODE_TEXT), ("1", NUM_COLOR), ("); ", CODE_TEXT), ("// 3 (конец)", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("next2 = calcIndex(", CODE_TEXT), ("3", NUM_COLOR), (", ", CODE_TEXT), ("4", NUM_COLOR), (", ", CODE_TEXT), ("1", NUM_COLOR), (");  ", CODE_TEXT), ("// 0 (начало)", COMMENT_COLOR)],
        ],
        "right_title": "Концепция Pure Functions",
        "bullet_items": [
            {"title": "Детерминированность", "desc": "При одних и тех же аргументах чистая функция всегда возвращает строго одинаковый ответ.", "height": 38},
            {"title": "Без Side-Effects", "desc": "Функция не изменяет внешние переменные, не делает запросы и не манипулирует DOM напрямую.", "height": 38},
            {"title": "Изоляция математики", "desc": "Математика кольцевого переключения отделена от графического рендеринга карточек.", "height": 38},
            {"title": "Надежность кода", "desc": "Такие функции легко тестировать, рефакторить и переносить между проектами.", "height": 38},
        ]
    },

    # Slide 7: Module 2
    {
        "type": "code",
        "category": "АСИНХРОННЫЕ ТАЙМЕРЫ",
        "question": "Почему setTimeout и setInterval не замораживают пользовательский интерфейс?",
        "left_title": "Event Loop и однопоточность",
        "code_lines": [
            [("console.log(", CODE_TEXT), ("'1. Инициализация слайдера'", STR_COLOR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Отправка колбэка в таймеры Web API", COMMENT_COLOR)],
            [("setTimeout(() => {", KEYWORD_COLOR)],
            [("  console.log(", CODE_TEXT), ("'2. Сработал таймер задержки'", STR_COLOR), (");", CODE_TEXT)],
            [("}, ", CODE_TEXT), ("1000", NUM_COLOR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("console.log(", CODE_TEXT), ("'3. Страница готова к кликам'", STR_COLOR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Порядок в консоли: 1 -> 3 -> 2", COMMENT_COLOR)],
        ],
        "right_title": "Асинхронность в браузере",
        "bullet_items": [
            {"title": "Однопоточный JS", "desc": "JavaScript выполняет один поток инструкций, предотвращая конфликты при изменении DOM.", "height": 38},
            {"title": "Web API таймеров", "desc": "Отсчет времени ведет браузер в фоновом режиме, разгружая поток выполнения скрипта.", "height": 38},
            {"title": "Очередь задач (Queue)", "desc": "Функция таймера помещается в очередь и ждет завершения синхронного кода.", "height": 38},
            {"title": "Отзывчивый UI", "desc": "Кнопки, скролл и поля ввода остаются активными во время работы автопрокрутки.", "height": 38},
        ]
    },

    # Slide 8: Module 2
    {
        "type": "code",
        "category": "МЕТОД SETINTERVAL",
        "question": "Что возвращает метод setInterval и как этот результат используется в скрипте?",
        "left_title": "Запуск интервала автопрокрутки",
        "code_lines": [
            [("// Циклический запуск каждые 3000 мс (3 сек)", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("INTERVAL_MS = ", CODE_TEXT), ("3000", NUM_COLOR), (";", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("const ", KEYWORD_COLOR), ("timerId = setInterval(() => {", FUNC_COLOR)],
            [("  console.log(", CODE_TEXT), ("'Автоматический тик слайдера'", STR_COLOR), (");", CODE_TEXT)],
            [("  nextSlide();", FUNC_COLOR)],
            [("}, INTERVAL_MS);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("console.log(", CODE_TEXT), ("'ID таймера:'", STR_COLOR), (", timerId); ", CODE_TEXT), ("// Число (напр. 1)", COMMENT_COLOR)],
        ],
        "right_title": "Циклический вызов setInterval",
        "bullet_items": [
            {"title": "Назначение метода", "desc": "setInterval вызывает переданную функцию-колбэк повторно через равные промежутки времени.", "height": 38},
            {"title": "Миллисекунды", "desc": "Время задается в миллисекундах: 1000 мс = 1 секунда, 3000 мс = 3 секунды.", "height": 38},
            {"title": "Числовой дескриптор", "desc": "Возвращает уникальный числовой ID, по которому этот таймер можно остановить.", "height": 38},
            {"title": "Оптимальная задержка", "desc": "Для витрины «СмартОфис» 3 секунды идеальны: фото успевают рассмотреть.", "height": 38},
        ]
    },

    # Slide 9: Module 2
    {
        "type": "code",
        "category": "ОСТАНОВКА ТАЙМЕРОВ",
        "question": "Что произойдет со страницей, если многократно вызывать setInterval без clearInterval?",
        "left_title": "Остановка через clearInterval",
        "code_lines": [
            [("let ", KEYWORD_COLOR), ("autoSlideTimer = ", CODE_TEXT), ("null", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("function ", KEYWORD_COLOR), ("stopAutoPlay", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(autoSlideTimer) {", CODE_TEXT)],
            [("    clearInterval(autoSlideTimer);", FUNC_COLOR)],
            [("    autoSlideTimer = ", CODE_TEXT), ("null", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("    console.log(", CODE_TEXT), ("'Таймер остановлен'", STR_COLOR), (");", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Очистка с clearInterval",
        "bullet_items": [
            {"title": "Остановка по ID", "desc": "Функция clearInterval(timerId) находит таймер в Web API и навсегда гасит его интервал.", "height": 38},
            {"title": "Сброс в null", "desc": "Присвоение null переменной таймера позволяет надежно проверять, запущен ли слайдер.", "height": 38},
            {"title": "Утечки памяти", "desc": "Без clearInterval неиспользуемые таймеры продолжают грузить память и процессор браузера.", "height": 38},
            {"title": "Защитное условие", "desc": "Проверка if (autoSlideTimer) исключает вызов очистки для уже выключенного таймера.", "height": 38},
        ]
    },

    # Slide 10: Module 2
    {
        "type": "code",
        "category": "МЕТОД SETTIMEOUT",
        "question": "В каких сценариях используется setTimeout вместо setInterval?",
        "left_title": "Отложенные вызовы setTimeout",
        "code_lines": [
            [("// Однократное выполнение с задержкой", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("showNotificationWithAutoClose", FUNC_COLOR), ("(msg) {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("toast = document.createElement(", CODE_TEXT), ("'div'", STR_COLOR), (");", CODE_TEXT)],
            [("  toast.textContent = msg;", CODE_TEXT)],
            [("  document.body.appendChild(toast);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // Удаление через 3.5 секунды", COMMENT_COLOR)],
            [("  setTimeout(() => {", KEYWORD_COLOR)],
            [("    toast.remove();", FUNC_COLOR)],
            [("  }, ", CODE_TEXT), ("3500", NUM_COLOR), (");", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Разовая задержка setTimeout",
        "bullet_items": [
            {"title": "Однократный запуск", "desc": "setTimeout(fn, delay) выполняет функцию ровно 1 раз по истечении указанного времени.", "height": 38},
            {"title": "Применение в UI", "desc": "Скрытие уведомлений Toast, закрытие модальных окон, пауза перед возобновлением слайдера.", "height": 38},
            {"title": "Очистка clearTimeout", "desc": "При необходимости запланированный вызов можно отменить с помощью clearTimeout(id).", "height": 38},
            {"title": "Debounce и задержки", "desc": "Служит основой для ожидания окончания ввода пользователя в строку поиска.", "height": 38},
        ]
    },

    # Slide 11: Module 2
    {
        "type": "code",
        "category": "УПРАВЛЕНИЕ ТАЙМЕРАМИ",
        "question": "Почему перед каждым запуском setInterval необходимо вызывать clearInterval?",
        "left_title": "Безопасный перезапуск таймера",
        "code_lines": [
            [("let ", KEYWORD_COLOR), ("timerId = ", CODE_TEXT), ("null", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("function ", KEYWORD_COLOR), ("startAuto", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  // КРИТИЧЕСКИ ВАЖНО: сбросить старый перед новым", COMMENT_COLOR)],
            [("  stopAuto();", FUNC_COLOR)],
            [("  timerId = setInterval(next, ", CODE_TEXT), ("3000", NUM_COLOR), (");", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("function ", KEYWORD_COLOR), ("stopAuto", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(timerId) {", CODE_TEXT)],
            [("    clearInterval(timerId);", FUNC_COLOR)],
            [("    timerId = ", CODE_TEXT), ("null", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Паттерн анти-двоения таймеров",
        "bullet_items": [
            {"title": "Проблема наложения", "desc": "Если вызвать setInterval 5 раз подряд, запустятся 5 таймеров, и кадры замелькают хаотично.", "height": 38},
            {"title": "Сброс перед стартом", "desc": "Вызов stopAuto() перед каждым запуском гарантирует, что активен ровно один таймер.", "height": 38},
            {"title": "Управление из кликов", "desc": "Когда пользователь щелкает стрелку, таймер перезапускается на чистые 3 секунды.", "height": 38},
            {"title": "Предсказуемость", "desc": "Поведение интерфейса остается стабильным даже при бешеном кликании по кнопкам.", "height": 38},
        ]
    },

    # Slide 12: Diagram 2
    {
        "type": "diagram_lifecycle",
        "category": "СХЕМА: ЦИКЛ ТАЙМЕРОВ",
        "question": "Как взаимодействуют автопрокрутка setInterval, клики пользователя и сброс таймера?"
    },

    # Slide 13: Module 3
    {
        "type": "code",
        "category": "ОБЛАСТЬ ВИДИМОСТИ",
        "question": "Чем блочная область видимости let и const безопаснее функциональной видимости var?",
        "left_title": "Блочная видимость let и const",
        "code_lines": [
            [("// Блочная область видимости (ES6)", COMMENT_COLOR)],
            [("if ", KEYWORD_COLOR), ("(true) {", CODE_TEXT)],
            [("  let ", KEYWORD_COLOR), ("sliderSpeed = ", CODE_TEXT), ("300", NUM_COLOR), (";", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("TOTAL_SLIDES = ", CODE_TEXT), ("4", NUM_COLOR), (";", CODE_TEXT)],
            [("  var ", KEYWORD_COLOR), ("legacyVar = ", CODE_TEXT), ("'Я протекаю наружу'", STR_COLOR), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("console.log(legacyVar);   ", CODE_TEXT), ("// Работает: var протек!", COMMENT_COLOR)],
            [("// console.log(sliderSpeed); // ReferenceError: not defined", COMMENT_COLOR)],
        ],
        "right_title": "Блочная vs Функциональная",
        "bullet_items": [
            {"title": "Блочный Scope", "desc": "Переменные let и const существуют строго внутри своего блока фигурных скобок {}.", "height": 38},
            {"title": "Утечка var", "desc": "Устаревший var игнорирует блоки if/for и проникает наружу, вызывая трудноуловимые баги.", "height": 38},
            {"title": "Защита от перезаписи", "desc": "Блочная видимость не дает случайно переписать счетчик цикла или настройки слайдера.", "height": 38},
            {"title": "Стандарт индустрии", "desc": "В современном JS используют только const для констант и let для переменных состояния.", "height": 38},
        ]
    },

    # Slide 14: Module 3
    {
        "type": "code",
        "category": "ЗАМЫКАНИЯ (CLOSURES)",
        "question": "Как замыкание позволяет функции сохранять доступ к переменным внешней функции?",
        "left_title": "Механизм замыкания в JS",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("createSlideManager", FUNC_COLOR), ("(total) {", CODE_TEXT)],
            [("  let ", KEYWORD_COLOR), ("current = ", CODE_TEXT), ("0", NUM_COLOR), ("; ", CODE_TEXT), ("// Закрытая переменная", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("  return ", KEYWORD_COLOR), ("function ", KEYWORD_COLOR), ("next", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("    current = (current + ", CODE_TEXT), ("1", NUM_COLOR), (") % total;", CODE_TEXT)],
            [("    return ", KEYWORD_COLOR), ("current; ", CODE_TEXT), ("// Помнит 'current' из родителя", COMMENT_COLOR)],
            [("  };", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("const ", KEYWORD_COLOR), ("sliderNext = createSlideManager(", CODE_TEXT), ("4", NUM_COLOR), (");", CODE_TEXT)],
            [("sliderNext(); ", CODE_TEXT), ("// 1", COMMENT_COLOR)],
            [("sliderNext(); ", CODE_TEXT), ("// 2", COMMENT_COLOR)],
        ],
        "right_title": "Замыкания (Closures)",
        "bullet_items": [
            {"title": "Определение замыкания", "desc": "Способность внутренней функции помнить и использовать переменные своего создателя.", "height": 38},
            {"title": "Жизнь после возврата", "desc": "Даже когда внешняя функция отработала, переменная current сохраняется в памяти.", "height": 38},
            {"title": "Приватные данные", "desc": "Никакой внешний скрипт не может изменить current напрямую — только через вызов next().", "height": 38},
            {"title": "Инженерная база", "desc": "Замыкания лежат в основе модульного JS, обработчиков событий и фабричных функций.", "height": 38},
        ]
    },

    # Slide 15: Module 3
    {
        "type": "code",
        "category": "ИНКАПСУЛЯЦИЯ",
        "question": "Почему переменные слайдера currentSlide и timerId нельзя делать глобальными?",
        "left_title": "Изоляция модуля initSlider()",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("initSlider", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  // Локальное изолированное состояние:", COMMENT_COLOR)],
            [("  const ", KEYWORD_COLOR), ("slides = document.querySelectorAll(", CODE_TEXT), ("'.slide'", STR_COLOR), (");", CODE_TEXT)],
            [("  let ", KEYWORD_COLOR), ("currentSlide = ", CODE_TEXT), ("0", NUM_COLOR), (";", CODE_TEXT)],
            [("  let ", KEYWORD_COLOR), ("timerId = ", CODE_TEXT), ("null", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  function ", KEYWORD_COLOR), ("showSlide", FUNC_COLOR), ("(index) { /* ... */ }", CODE_TEXT)],
            [("  function ", KEYWORD_COLOR), ("next", FUNC_COLOR), ("() { showSlide(currentSlide + ", CODE_TEXT), ("1", NUM_COLOR), ("); }", CODE_TEXT)],
            [("  startAuto();", FUNC_COLOR)],
            [("}", CODE_TEXT)],
            [("// currentSlide и timerId снаружи недоступны!", COMMENT_COLOR)],
        ],
        "right_title": "Инкапсуляция слайдера",
        "bullet_items": [
            {"title": "Чистота global scope", "desc": "Отказ от глобальных переменных window исключает случайные конфликты с другими скриптами.", "height": 38},
            {"title": "Автономия модуля", "desc": "Слайдер становится законченным независимым виджетом, который можно внедрить куда угодно.", "height": 38},
            {"title": "Безопасность состояния", "desc": "Внешние функции не смогут поломать счетчик currentSlide или обнулить таймер.", "height": 38},
            {"title": "Простота инициализации", "desc": "Для запуска на странице достаточно вызвать одну функцию initSlider().", "height": 38},
        ]
    },

    # Slide 16: Module 4
    {
        "type": "code",
        "category": "АРХИТЕКТУРА ПРОЕКТА",
        "question": "К каким проблемам приводит смешивание данных, стилей и логики в одном файле?",
        "left_title": "Разделение зон ответственности",
        "code_lines": [
            [("// Структура проекта «СмартОфис»:", COMMENT_COLOR)],
            [("webinar-07-slider-timers/", TAG_COLOR)],
            [("├── index.html        ", CODE_TEXT), ("# Каркас и разметка", COMMENT_COLOR)],
            [("├── css/style.css     ", CODE_TEXT), ("# Оформление и transitions", COMMENT_COLOR)],
            [("├── js/", TAG_COLOR)],
            [("│   ├── data.js       ", FUNC_COLOR), ("# Mock-данные каталога", COMMENT_COLOR)],
            [("│   └── main.js       ", STR_COLOR), ("# Поведение и таймеры", COMMENT_COLOR)],
            [("└── img/              ", NUM_COLOR), ("# Офисные фото и графика", COMMENT_COLOR)],
        ],
        "right_title": "Separation of Concerns",
        "bullet_items": [
            {"title": "HTML — Каркас", "desc": "Отвечает за семантику, DOM-дерево и наличие контейнеров слайдера и кнопок.", "height": 38},
            {"title": "CSS — Визуал", "desc": "Задает абсолютное позиционирование слайдов, размеры и плавность transition: opacity.", "height": 38},
            {"title": "JS — Поведение", "desc": "Обрабатывает клики пользователей, управляет переключением классов и таймерами.", "height": 38},
            {"title": "Масштабируемость", "desc": "Разработчик интерфейса может менять стили стрелок, не трогая JavaScript-код.", "height": 38},
        ]
    },

    # Slide 17: Module 4
    {
        "type": "code",
        "category": "ФАЙЛОВАЯ СТРУКТУРА",
        "question": "В чем преимущество вынесения mock-данных в отдельный файл data.js?",
        "left_title": "Разделение data.js и main.js",
        "code_lines": [
            [("// js/data.js — только данные комнат", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("roomsData = [", CODE_TEXT)],
            [("  { id: ", CODE_TEXT), ("'focus-1'", STR_COLOR), (", title: ", CODE_TEXT), ("'Мини-офис Focus'", STR_COLOR), (", price: ", CODE_TEXT), ("450", NUM_COLOR), (" },", CODE_TEXT)],
            [("  { id: ", CODE_TEXT), ("'smart-2'", STR_COLOR), (", title: ", CODE_TEXT), ("'Смарт-офис Grand'", STR_COLOR), (", price: ", CODE_TEXT), ("900", NUM_COLOR), (" }", CODE_TEXT)],
            [("];", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// js/main.js — только управляющая логика", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("renderCatalog", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  // Использует готовый массив roomsData", COMMENT_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Модульное разделение скриптов",
        "bullet_items": [
            {"title": "Чистота логики", "desc": "main.js не загроможден длинными массивами объектов и текстами описаний офисов.", "height": 38},
            {"title": "Простая замена", "desc": "В будущем файл data.js легко заменить сетевым запросом fetch('/api/rooms').", "height": 38},
            {"title": "Командная работа", "desc": "Контент-менеджер может редактировать цены комнат в data.js без риска поломать JS.", "height": 38},
            {"title": "Удобство отладки", "desc": "Ошибки в коде логики локализуются в main.js, опечатки в данных — в data.js.", "height": 38},
        ]
    },

    # Slide 18: Module 4
    {
        "type": "code",
        "category": "ПОДКЛЮЧЕНИЕ СКРИПТОВ",
        "question": "Почему при запуске слайдера в main.js критически важно дождаться DOMContentLoaded?",
        "left_title": "Порядок тегов script и DOMContentLoaded",
        "code_lines": [
            [("<!-- Порядок подключения в index.html -->", COMMENT_COLOR)],
            [("<script src=\"js/data.js\"></script>", TAG_COLOR)],
            [("<script src=\"js/main.js\"></script>", TAG_COLOR)],
            [("", CODE_TEXT)],
            [("// Точка входа в js/main.js:", COMMENT_COLOR)],
            [("document.addEventListener(", CODE_TEXT), ("'DOMContentLoaded'", STR_COLOR), (", () => {", CODE_TEXT)],
            [("  initNavigation();", FUNC_COLOR)],
            [("  initSlider();     ", FUNC_COLOR), ("// Слайдер запускается здесь", COMMENT_COLOR)],
            [("  renderCatalog();", FUNC_COLOR)],
            [("});", CODE_TEXT)],
        ],
        "right_title": "Безопасный запуск приложения",
        "bullet_items": [
            {"title": "Очередность тегов", "desc": "Скрипт data.js объявляется первым, чтобы глобальный массив комнат был готов до старта main.js.", "height": 38},
            {"title": "DOMContentLoaded", "desc": "Гарантирует, что браузер полностью построил HTML-дерево перед выполнением querySelector.", "height": 38},
            {"title": "Защита от null", "desc": "Предотвращает фатальную ошибку Cannot read properties of null (reading 'addEventListener').", "height": 38},
            {"title": "Централизация старта", "desc": "Единый обработчик DOMContentLoaded прозрачно управляет порядком запуска всех систем.", "height": 38},
        ]
    },

    # Slide 19: Module 5
    {
        "type": "code",
        "category": "ВЕРСТКА СЛАЙДЕРА",
        "question": "Почему для плавного появления слайдов используется opacity, а не display: none?",
        "left_title": "CSS-наложение и плавность",
        "code_lines": [
            [("/* Родительский контейнер слайдера */", COMMENT_COLOR)],
            [(".slider { position: relative; overflow: hidden; height: 420px; }", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Все слайды лежат друг под другом */", COMMENT_COLOR)],
            [(".slide {", CODE_TEXT)],
            [("  position: absolute; top: 0; left: 0; width: 100%; height: 100%;", CODE_TEXT)],
            [("  opacity: 0;", KEYWORD_COLOR)],
            [("  transition: opacity 0.3s ease-in-out;", FUNC_COLOR)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Активный видимый слайд */", COMMENT_COLOR)],
            [(".slide.active { opacity: 1; z-index: 2; }", STR_COLOR)],
        ],
        "right_title": "Принцип наложения слайдов",
        "bullet_items": [
            {"title": "Абсолютное наложение", "desc": "Все 4 слайда занимают одинаковые координаты top: 0, left: 0 внутри .slider.", "height": 38},
            {"title": "Скрытие через opacity", "desc": "В отличие от display: none, изменение opacity плавно анимируется браузером через transition.", "height": 38},
            {"title": "Стек слоев z-index", "desc": "Класс .active поднимает видимый кадр выше остальных слоев, исключая артефакты.", "height": 38},
            {"title": "Быстродействие GPU", "desc": "Анимация opacity ускоряется видеокартой и обеспечивает стабильные 60 кадров/сек.", "height": 38},
        ]
    },

    # Slide 20: Module 5
    {
        "type": "code",
        "category": "АЛГОРИТМ СЛАЙДЕРА",
        "question": "Как работает условие закольцовывания слайдов при выходе за пределы массива?",
        "left_title": "Кольцевая функция showSlide()",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("showSlide", FUNC_COLOR), ("(index) {", CODE_TEXT)],
            [("  // Проверка правого и левого краев:", COMMENT_COLOR)],
            [("  if ", KEYWORD_COLOR), ("(index >= slides.length) {", CODE_TEXT)],
            [("    currentSlide = ", CODE_TEXT), ("0", NUM_COLOR), (";             ", CODE_TEXT), ("// Зацикливание в начало", COMMENT_COLOR)],
            [("  } else if ", KEYWORD_COLOR), ("(index < ", CODE_TEXT), ("0", NUM_COLOR), (") {", CODE_TEXT)],
            [("    currentSlide = slides.length - ", CODE_TEXT), ("1", NUM_COLOR), ("; ", CODE_TEXT), ("// Зацикливание в конец", COMMENT_COLOR)],
            [("  } else {", KEYWORD_COLOR)],
            [("    currentSlide = index;", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("  updateVisuals();", FUNC_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Кольцевая навигация",
        "bullet_items": [
            {"title": "Бесконечный цикл", "desc": "После слайда 3 открывается слайд 0. При клике назад со слайда 0 открывается слайд 3.", "height": 38},
            {"title": "Единая точка правды", "desc": "Переменная currentSlide всегда хранит точный номер кадра, показанного пользователю.", "height": 38},
            {"title": "Универсальный вход", "desc": "Функция принимает целевой номер слайда и от стрелок, и от нижних индикаторов точек.", "height": 38},
            {"title": "Разделение логики", "desc": "Сначала рассчитывается корректный индекс, а затем обновляется графика в updateVisuals().", "height": 38},
        ]
    },

    # Slide 21: Module 5
    {
        "type": "code",
        "category": "ИНДИКАТОРЫ ТОЧЕК",
        "question": "В чем преимущество использования classList.toggle со вторым аргументом-условием?",
        "left_title": "Синхронизация слайдов и Dots",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("updateVisuals", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  // Элегантный переключатель классов одной строкой", COMMENT_COLOR)],
            [("  slides.forEach((slide, i) => {", KEYWORD_COLOR)],
            [("    slide.classList.toggle(", CODE_TEXT), ("'active'", STR_COLOR), (", i === currentSlide);", CODE_TEXT)],
            [("  });", KEYWORD_COLOR)],
            [("", CODE_TEXT)],
            [("  dots.forEach((dot, i) => {", KEYWORD_COLOR)],
            [("    dot.classList.toggle(", CODE_TEXT), ("'active'", STR_COLOR), (", i === currentSlide);", CODE_TEXT)],
            [("  });", KEYWORD_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Метод classList.toggle()",
        "bullet_items": [
            {"title": "Второй аргумент", "desc": "classList.toggle('active', condition) добавляет класс при true и удаляет при false.", "height": 38},
            {"title": "Никаких дубликатов", "desc": "Исключает ситуацию, когда из-за бага активными одновременно оказываются 2 кадра.", "height": 38},
            {"title": "Идеальная связка", "desc": "Слайды и круглые индикаторы снизу меняются одновременно без рассинхронизации.", "height": 38},
            {"title": "Лаконичность кода", "desc": "Заменяет длинные блоки if/else с add() и remove(), повышая читаемость.", "height": 38},
        ]
    },

    # Slide 22: Module 5
    {
        "type": "code",
        "category": "РУЧНОЕ УПРАВЛЕНИЕ",
        "question": "Зачем сбрасывать интервал таймера при ручном переключении слайда пользователем?",
        "left_title": "Слушатели кликов со сбросом таймера",
        "code_lines": [
            [("// Кнопки вперед и назад:", COMMENT_COLOR)],
            [("nextBtn.addEventListener(", CODE_TEXT), ("'click'", STR_COLOR), (", () => {", CODE_TEXT)],
            [("  next();", FUNC_COLOR)],
            [("  startAuto(); ", FUNC_COLOR), ("// Сброс таймера на полные 3 сек!", COMMENT_COLOR)],
            [("});", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Нижние круглые индикаторы Dots:", COMMENT_COLOR)],
            [("dots.forEach((dot, idx) => {", KEYWORD_COLOR)],
            [("  dot.addEventListener(", CODE_TEXT), ("'click'", STR_COLOR), (", () => {", CODE_TEXT)],
            [("    showSlide(idx);", FUNC_COLOR)],
            [("    startAuto();", FUNC_COLOR)],
            [("  });", CODE_TEXT)],
            [("});", CODE_TEXT)],
        ],
        "right_title": "Комфортный интерактивный UX",
        "bullet_items": [
            {"title": "Мгновенная реакция", "desc": "Клик по стрелке или точке сразу же переключает медиа-контент без задержки.", "height": 38},
            {"title": "Уважение к пользователю", "desc": "Вызов startAuto() сбрасывает старый таймер: слайд не уедет через долю секунды.", "height": 38},
            {"title": "Полные 3 секунды", "desc": "После ручного клика пользователь получает полные 3 секунды на изучение фотографии.", "height": 38},
            {"title": "Бесшовный гибрид", "desc": "Автопрокрутка и ручное управление работают в полной гармонии без сбоев.", "height": 38},
        ]
    },

    # Slide 23: Module 5
    {
        "type": "code",
        "category": "ИТОГИ ВЕБИНАРА",
        "question": "Какие критерии определяют надежность и стабильность интерактивного слайдера?",
        "left_title": "Финальная сборка медиа-компонента",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("initSlider", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("slides = document.querySelectorAll(", CODE_TEXT), ("'.slide'", STR_COLOR), (");", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(!slides.length) return; ", CODE_TEXT), ("// Защита от ошибок", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("  // 1. Инкапсулированные let currentSlide, timerId", COMMENT_COLOR)],
            [("  // 2. Кольцевой алгоритм showSlide(index)", COMMENT_COLOR)],
            [("  // 3. Безопасные startAuto() и stopAuto()", COMMENT_COLOR)],
            [("  // 4. Поддержка стрелок и точек-индикаторов", COMMENT_COLOR)],
            [("  startAuto();", FUNC_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Чек-лист качества компонента",
        "bullet_items": [
            {"title": "Защита от сбоев", "desc": "Проверка if (!slides.length) return защищает другие страницы сайта от ошибок консоли.", "height": 38},
            {"title": "Изоляция переменных", "desc": "Вся логика запечатана внутри initSlider() без риска конфликтов в global scope.", "height": 38},
            {"title": "Контроль таймеров", "desc": "Отсутствие утечек памяти и наложения параллельных интервалов автопрокрутки.", "height": 38},
            {"title": "Готовый медиа-блок", "desc": "Главная страница «СмартОфис» получила эффектный интерактивный фото-слайдер.", "height": 38},
        ]
    }
]

# -------------------------------------------------------------
# PRESENTATION GENERATOR EXECUTION
# -------------------------------------------------------------
print(f"Total content slides to generate: {len(slides_data)}")

prs = Presentation(base_template_path)
blank_layout = prs.slide_layouts[6]

orig_slide_ids = list(prs.slides._sldIdLst)
title_slide_id = orig_slide_ids[0]
plan_slide_id = orig_slide_ids[1]
result_slide_id = orig_slide_ids[25]
goodbye_slide_id = orig_slide_ids[26]

# 0. Update Layouts footer to Вебинар 7
for layout in prs.slide_layouts:
    for s in layout.shapes:
        if "Google Shape;59;p13" in s.name and s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip() == "2":
                        r.text = "7"

# 1. Update Slide 1 (Title slide)
slide1 = prs.slides[0]
for sh in slide1.shapes:
    if sh.name == "Google Shape;59;p13" and sh.has_text_frame:
        sh.text_frame.text = "Вебинар 7 "
    elif sh.shape_type == MSO_SHAPE_TYPE.GROUP and sh.name == "Группа 12":
        for sub in sh.shapes:
            if sub.name == "TextBox 7":
                sub.text_frame.text = "ФУНКЦИОНАЛЬНЫЙ ПОДХОД И ТАЙМЕРЫ"
                sub.text_frame.paragraphs[0].runs[0].font.name = "Montserrat"
                sub.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
                sub.text_frame.paragraphs[0].runs[0].font.bold = True
            elif sub.name == "TextBox 9":
                sub.text_frame.text = "Проектирование \nи разработка интерфейсов пользователя"
print("Slide 1 updated successfully.")

# 2. Update Slide 2 (Plan slide)
slide2 = prs.slides[1]
plan_titles = {
    "TextBox 12": "1. Аргументы и функции",
    "TextBox 15": "2. Асинхронные таймеры",
    "TextBox 18": "3. Области видимости",
    "TextBox 21": "4. Модульная структура",
    "TextBox 24": "5. Практика: Слайдер"
}
plan_subtitles = {
    "TextBox 13": "Параметры функций, значения по умолчанию и чистые вызовы",
    "TextBox 16": "Методы setInterval и clearInterval: циклы времени в браузере",
    "TextBox 19": "Глобальная, блочная видимость (let/const) и концепция замыкания",
    "TextBox 22": "Организация файлов проекта и разделение зон ответственности",
    "TextBox 25": "Разработка интерактивного авто-слайдера с ручным управлением"
}

for sh in slide2.shapes:
    strip_shape_styles_and_shadows(sh)
    if sh.name in plan_titles and sh.has_text_frame:
        p = sh.text_frame.paragraphs[0]
        p.text = plan_titles[sh.name]
        p.runs[0].font.name = "Montserrat"
        p.runs[0].font.size = Pt(11.25)
        p.runs[0].font.bold = True
    elif sh.name in plan_subtitles and sh.has_text_frame:
        sh.width = Pt(450)
        tf = sh.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = plan_subtitles[sh.name]
        p.runs[0].font.name = "Inter"
        p.runs[0].font.size = Pt(9.75)
        p.runs[0].font.bold = False

print("Slide 2 updated successfully.")

# 3. Create 21 Content Slides
content_slide_elements = []
for idx, data in enumerate(slides_data, 1):
    new_slide = prs.slides.add_slide(blank_layout)
    content_slide_elements.append(new_slide._element)
    
    stype = data.get("type")
    if stype == "code":
        create_code_explanation_slide(
            new_slide,
            data["category"],
            data["question"],
            data["left_title"],
            data["code_lines"],
            data["right_title"],
            data["bullet_items"]
        )
    elif stype == "diagram_arch":
        create_function_arch_diagram_slide(
            new_slide,
            data["category"],
            data["question"]
        )
    elif stype == "diagram_lifecycle":
        create_timers_lifecycle_diagram_slide(
            new_slide,
            data["category"],
            data["question"]
        )
    print(f"  Created Slide {idx + 2}: {data['category']}")

# 4. Update Result Slide (Slide 26 in original template, prs.slides[25])
slide_res = prs.slides[25]
pic_to_remove = None
tb2_res = None
for sh in slide_res.shapes:
    strip_shape_styles_and_shadows(sh)
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        pic_to_remove = sh
    elif sh.name == "TextBox 2":
        tb2_res = sh

if pic_to_remove is not None:
    spTree = slide_res.shapes._spTree
    spTree.remove(pic_to_remove._element)

# Add Webinar 7 result screenshot: aspect ratio 1210:670 (1.806), centered nicely
res_w = Pt(510)
res_h = Pt(282)
res_left = Pt(105)
res_top = Pt(68)
slide_res.shapes.add_picture(result_img_path, res_left, res_top, res_w, res_h)

# Add or update TextBox 2 on Result slide
if tb2_res is None:
    tb2_res = slide_res.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2_res.name = "TextBox 2"

tf_res = tb2_res.text_frame
tf_res.word_wrap = True
tf_res.margin_left = tf_res.margin_top = tf_res.margin_right = tf_res.margin_bottom = 0
p_res = tf_res.paragraphs[0]
p_res.alignment = PP_ALIGN.LEFT
p_res.text = ""
r_res = p_res.add_run()
r_res.text = "Результат: Завершенный медиа-компонент для главной страницы"
r_res.font.name = "Montserrat"
r_res.font.size = Pt(10)
r_res.font.color.rgb = QUESTION_TEXT
print("Result slide updated with webinar 7 slider screenshot and TextBox 2.")

# 5. Clean shadows on Goodbye slide (Slide 27 in original template, prs.slides[26])
slide_goodbye = prs.slides[26]
for sh in slide_goodbye.shapes:
    strip_shape_styles_and_shadows(sh)
print("Goodbye slide shadows stripped.")

# 6. Reorder slide list in presentation XML:
# [Title, Plan] + new_slide_ids + [Result, Goodbye]
all_slide_ids = list(prs.slides._sldIdLst)
new_slide_ids = all_slide_ids[len(orig_slide_ids):]
new_order = [title_slide_id, plan_slide_id] + new_slide_ids + [result_slide_id, goodbye_slide_id]
prs.slides._sldIdLst[:] = new_order

# 7. Save presentation
os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
prs.save(output_pptx_path)

print(f"\nFinal presentation slide count: {len(prs.slides)}")
print(f"\nУСПЕХ! Презентация Вебинара 7 сохранена: {output_pptx_path}")
