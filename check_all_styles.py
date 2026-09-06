import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation(r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx")

for i in range(1, 19): # slides 2 to 19
    slide = prs.slides[i]
    shapes_with_style = []
    for s in slide.shapes:
        style_elem = s._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
        if style_elem is not None:
            # check what effects
            eff = style_elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}effectRef')
            eff_idx = eff.attrib.get('idx', 'none') if eff is not None else 'no_eff'
            shapes_with_style.append(f"{s.name} (eff={eff_idx})")
    print(f"Slide {i+1}: {', '.join(shapes_with_style)}")
