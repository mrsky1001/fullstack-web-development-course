import sys
import pptx
from pptx import Presentation
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

pptx_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-09-final-assembly-my-bookings\вебинар 9.pptx"
prs = Presentation(pptx_path)

print(f"=== АУДИТ ПРЕЗЕНТАЦИИ ВЕБИНАРА 9: {len(prs.slides)} СЛАЙДОВ ===")

forbidden_words = ["школьник", "школьники", "школьника", "школьниками", "школьникам", "для школьников"]
found_forbidden = []
shadow_errors = []
textbox2_errors = []
pills_checked = 0

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    has_tb2 = False
    
    for s in slide.shapes:
        # Check forbidden words
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    txt_lower = r.text.lower()
                    for fw in forbidden_words:
                        if fw in txt_lower:
                            found_forbidden.append((slide_num, s.name, r.text))

        # Check shadows
        style_elem = s._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
        if style_elem is not None:
            eff = style_elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}effectRef')
            eff_idx = eff.attrib.get('idx', '0') if eff is not None else '0'
            if eff_idx != '0':
                shadow_errors.append((slide_num, s.name, f"style effectRef idx={eff_idx}"))
        
        for effect in s._element.xpath('.//a:outerShdw'):
            shadow_errors.append((slide_num, s.name, "direct outerShdw found"))

        # Check TextBox 2 color
        if s.name == "TextBox 2" and s.has_text_frame:
            has_tb2 = True
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    col = r.font.color
                    if col and col.type == 1 and col.rgb:
                        hex_col = str(col.rgb)
                        if hex_col.upper() != "A6A1A1":
                            textbox2_errors.append((slide_num, s.name, hex_col))
                    else:
                        textbox2_errors.append((slide_num, s.name, "NO_RGB"))

        if "Google Shape;69;p14" in s.name:
            pills_checked += 1

    if slide_num not in [1, 2, 25] and not has_tb2:
        textbox2_errors.append((slide_num, "MISSING", "TextBox 2 not found"))

print(f"1. Запрещенные слова ('школьники'): {len(found_forbidden)}")
if found_forbidden:
    for item in found_forbidden:
        print(f"   [ERR] Slide {item[0]} {item[1]}: '{item[2]}'")
else:
    print("   ✅ СТРОГО СОБЛЮДЕНО: ни одного упоминания слова 'школьники'!")

print(f"2. Проверка артефактов теней (<p:style>, <a:outerShdw>): {len(shadow_errors)}")
if shadow_errors:
    for item in shadow_errors:
        print(f"   [ERR] Slide {item[0]} {item[1]}: {item[2]}")
else:
    print("   ✅ ЧИСТО: ни одной паразитной тени в презентации нет!")

print(f"3. Проверка цвета вопроса TextBox 2 (#A6A1A1): {len(textbox2_errors)}")
if textbox2_errors:
    for item in textbox2_errors:
        print(f"   [ERR] Slide {item[0]} {item[1]}: {item[2]}")
else:
    print("   ✅ 100% ТОЧНОСТЬ: все вопросы окрашены строго в #A6A1A1!")

print(f"4. Проверено плашек категорий: {pills_checked}")

print("\n--- ПРОВЕРКА КЛЮЧЕВЫХ СЛАЙДОВ ---")
for s_idx in [0, 1, 2, 4, 11, 23, 24]:
    sl = prs.slides[s_idx]
    texts = [sh.text.strip().replace('\n', ' ') for sh in sl.shapes if sh.has_text_frame and sh.text.strip()]
    first_t = texts[0][:80] if texts else 'EMPTY'
    print(f"Слайд {s_idx+1}: {first_t}")

# Check layout footer
print("\n--- ПРОВЕРКА ФУТЕРА В МАКЕТАХ ---")
for li, layout in enumerate(prs.slide_layouts):
    for s in layout.shapes:
        if "Google Shape;59;p13" in s.name and s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    print(f"Layout {li} footer run: '{r.text}'")
