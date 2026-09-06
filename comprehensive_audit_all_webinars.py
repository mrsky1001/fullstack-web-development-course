import os
import sys
from pptx import Presentation
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

webinars = [
    ("Вебинар 2", r"05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx", 2),
    ("Вебинар 3", r"05. step by step 2026 & 2027\01. frontend\webinar-03-flexbox-grid-cards\вебинар 3.pptx", 3),
    ("Вебинар 4", r"05. step by step 2026 & 2027\01. frontend\webinar-04-js-dom-navigation\вебинар 4.pptx", 4),
    ("Вебинар 5", r"05. step by step 2026 & 2027\01. frontend\webinar-05-dynamic-catalog-data\вебинар 5.pptx", 5),
    ("Вебинар 6", r"05. step by step 2026 & 2027\01. frontend\webinar-06-forms-validation\вебинар 6.pptx", 6),
    ("Вебинар 7", r"05. step by step 2026 & 2027\01. frontend\webinar-07-slider-timers\вебинар 7.pptx", 7),
    ("Вебинар 8", r"05. step by step 2026 & 2027\01. frontend\webinar-08-filters-and-booking-calc\вебинар 8.pptx", 8),
    ("Вебинар 9", r"05. step by step 2026 & 2027\01. frontend\webinar-09-final-assembly-my-bookings\вебинар 9.pptx", 9),
]

forbidden_words = ["школьник", "школьники", "школьника", "школьниками", "школьникам", "для школьников"]

print("================================================================================")
print("     КОМПЛЕКСНЫЙ АУДИТ ПРЕЗЕНТАЦИЙ ПО ПРАВИЛАМ АГЕНТА И ДИЗАЙН-СИСТЕМЫ")
print("================================================================================\n")

overall_summary = []

for name, rel_path, web_num in webinars:
    full_path = os.path.abspath(rel_path)
    if not os.path.exists(full_path):
        print(f"❌ {name}: ФАЙЛ НЕ НАЙДЕН: {full_path}")
        overall_summary.append((name, "NOT_FOUND", []))
        continue
        
    prs = Presentation(full_path)
    slide_count = len(prs.slides)
    
    issues = []
    
    # 1. Slide count
    if slide_count != 25:
        issues.append(f"Количество слайдов: {slide_count} (ожидается ровно 25)")
        
    # 2. Forbidden words
    found_forbidden = []
    shadow_errors = []
    tb2_color_errors = []
    pills_count = 0
    overlap_warnings = []
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        has_tb2 = False
        
        # Check right-column bullet items vertical coordinates
        bullet_boxes = []
        
        for s in slide.shapes:
            # Check text for forbidden words
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    for r in p.runs:
                        txt_l = r.text.lower()
                        for fw in forbidden_words:
                            if fw in txt_l:
                                found_forbidden.append((slide_num, s.name, r.text))
                                
            # Check shadows
            spPr = getattr(s._element, 'spPr', None)
            if spPr is not None:
                if spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw'):
                    shadow_errors.append((slide_num, s.name, "spPr outerShdw"))
            
            style_elem = s._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
            if style_elem is not None:
                eff = style_elem.find('{http://schemas.openxmlformats.org/drawingml/2006/main}effectRef')
                eff_idx = eff.attrib.get('idx', '0') if eff is not None else '0'
                if eff_idx != '0':
                    shadow_errors.append((slide_num, s.name, f"style effectRef idx={eff_idx}"))
                    
            for effect in s._element.xpath('.//a:outerShdw'):
                shadow_errors.append((slide_num, s.name, "xpath outerShdw"))
                
            # Check TextBox 2
            if s.name == "TextBox 2" and s.has_text_frame:
                has_tb2 = True
                for p in s.text_frame.paragraphs:
                    for r in p.runs:
                        col = r.font.color
                        if col and col.type == 1 and col.rgb:
                            hex_c = str(col.rgb).upper()
                            if hex_c != "A6A1A1":
                                tb2_color_errors.append((slide_num, hex_c))
                        else:
                            tb2_color_errors.append((slide_num, "NO_RGB"))
                            
            if "Google Shape;69;p14" in s.name:
                pills_count += 1
                
            # Track bullet textboxes on standard content slides
            if slide_num >= 3 and slide_num <= 23:
                if "TextBox Bullet" in s.name:
                    if s.has_text_frame and s.text.strip():
                        bullet_boxes.append((s.left.pt, s.top.pt, s.height.pt, s.name, s.text.strip()[:40]))
                        
        if slide_num not in [1, 2, 25] and not has_tb2:
            tb2_color_errors.append((slide_num, "MISSING"))
            
        # Check vertical overlaps only for textboxes in the same vertical column (abs(x1 - x2) < 40)
        # Group by column X
        columns = {}
        for bx, by, bh, bname, btext in bullet_boxes:
            # find closest existing column key
            col_key = None
            for k in columns:
                if abs(k - bx) < 30:
                    col_key = k
                    break
            if col_key is None:
                col_key = bx
                columns[col_key] = []
            columns[col_key].append((by, bh, bname, btext))
            
        for col_x, boxes in columns.items():
            boxes.sort(key=lambda x: x[0])
            for b_i in range(len(boxes) - 1):
                top_curr, h_curr, name_curr, text_curr = boxes[b_i]
                top_next, h_next, name_next, text_next = boxes[b_i + 1]
                gap = top_next - (top_curr + h_curr)
                if gap < -1:
                    overlap_warnings.append((slide_num, f"Пересечение ({gap:.1f}pt): '{text_curr}' -> '{text_next}'"))
                elif gap < 5:
                    overlap_warnings.append((slide_num, f"Минимальный зазор ({gap:.1f}pt < 5pt): '{text_curr}'"))
                
    # 3. Footer check in layouts
    footer_ok = True
    found_footers = []
    for layout in prs.slide_layouts:
        for s in layout.shapes:
            if "Google Shape;59;p13" in s.name and s.has_text_frame:
                txt = "".join(r.text for p in s.text_frame.paragraphs for r in p.runs)
                found_footers.append(txt)
                if f"Вебинар {web_num}" not in txt and f"Вебинар  {web_num}" not in txt:
                    # Note: webinar 2 might have 'Вебинар 2'
                    footer_ok = False
                    
    # Compile issues
    if found_forbidden:
        issues.append(f"Запрещенные слова ('школьники'): найдено {len(found_forbidden)}")
    if shadow_errors:
        issues.append(f"Тени (<p:style>, <a:outerShdw>): {len(shadow_errors)}")
    if tb2_color_errors:
        issues.append(f"Цвет TextBox 2 (#A6A1A1): ошибок {len(tb2_color_errors)}")
    if pills_count < 20 and slide_count == 25:
        issues.append(f"Плашки категорий: {pills_count} (ожидается ~21)")
    if overlap_warnings:
        # Group by slide
        bad_slides = set(w[0] for w in overlap_warnings)
        issues.append(f"Наложения/сжатия текста в списках: слайды {sorted(list(bad_slides))}")
        
    status = "OK" if not issues else "WARNINGS/ERRORS"
    overall_summary.append((name, status, issues, overlap_warnings))
    
    print(f"--- {name} (Слайдов: {slide_count}) ---")
    if not issues:
        print("  ✅ ВСЕ ПРАВИЛА СОБЛЮДЕНЫ НА 100%!")
        print(f"     - Запрещенные слова: 0")
        print(f"     - Тени: 0")
        print(f"     - Цвет вопросов #A6A1A1: ОК")
        print(f"     - Плашек категорий: {pills_count}")
        print(f"     - Наложений текста в списках: 0")
    else:
        print(f"  ⚠️ ОБНАРУЖЕНЫ ЗАМЕЧАНИЯ ({len(issues)}):")
        for iss in issues:
            print(f"     • {iss}")
        if overlap_warnings:
            for w in overlap_warnings[:5]:
                print(f"       [Слайд {w[0]}] {w[1]}")
    print()

print("================================================================================")
print("ИТОГОВАЯ СВОДКА ПО ВСЕМ ПРЕЗЕНТАЦИЯМ:")
for name, status, issues, *rest in overall_summary:
    icon = "✅" if status == "OK" else "⚠️"
    print(f"{icon} {name:12}: {status} {('— ' + '; '.join(issues)) if issues else ''}")
print("================================================================================")
