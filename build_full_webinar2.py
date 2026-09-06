import sys
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Set stdout encoding to utf-8
sys.stdout.reconfigure(encoding='utf-8')

original_pptx = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2.pptx"
output_pptx = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - дополненный.pptx"

prs = Presentation(original_pptx)
blank_layout = prs.slide_layouts[12] # BLANK layout

# Color definitions matching webinar 1 & 2
ORANGE = RGBColor(0xFF, 0x6F, 0x03)      # #FF6F03
ORANGE_DOT = RGBColor(0xFE, 0x60, 0x02)  # #FE6002
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)   # #222222
BODY_TEXT = RGBColor(0x44, 0x44, 0x44)   # #444444
MUTED_TEXT = RGBColor(0x66, 0x66, 0x66)  # #666666
CARD_BG = RGBColor(0xF9, 0xF9, 0xF9)     # #F9F9F9
CARD_BORDER = RGBColor(0xE5, 0xE7, 0xEB) # #E5E7EB
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)     # Dark editor theme
CODE_TEXT = RGBColor(0xDD, 0xDD, 0xDD)
CODE_KW = RGBColor(0x56, 0x9C, 0xD6)     # Blue
CODE_TAG = RGBColor(0x4E, 0xC9, 0xB0)    # Teal
CODE_STR = RGBColor(0xCE, 0x91, 0x78)    # Orange/Salmon
CODE_COMM = RGBColor(0x6A, 0x99, 0x55)   # Green
CODE_PROP = RGBColor(0x9C, 0xDC, 0xFE)   # Light blue

def add_header_and_footer(slide, category_text, question_text):
    # Top right category header
    tb_top = slide.shapes.add_textbox(Pt(340), Pt(22), Pt(350), Pt(24))
    tf_top = tb_top.text_frame
    tf_top.word_wrap = True
    tf_top.margin_left = tf_top.margin_top = tf_top.margin_right = tf_top.margin_bottom = 0
    p_top = tf_top.paragraphs[0]
    p_top.alignment = PP_ALIGN.RIGHT
    run_top = p_top.add_run()
    run_top.text = category_text.upper()
    run_top.font.name = "Inter"
    run_top.font.size = Pt(11.5)
    run_top.font.color.rgb = DARK_TEXT

    # Bottom left question
    tb_bot = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(500), Pt(20))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    tf_bot.margin_left = tf_bot.margin_top = tf_bot.margin_right = tf_bot.margin_bottom = 0
    p_bot = tf_bot.paragraphs[0]
    p_bot.alignment = PP_ALIGN.LEFT
    run_bot = p_bot.add_run()
    run_bot.text = question_text
    run_bot.font.name = "Montserrat"
    run_bot.font.size = Pt(10)
    run_bot.font.color.rgb = RGBColor(0xA6, 0xA1, 0xA1)  # #A6A1A1

def create_two_column_slide(slide, category, question, left_title, code_runs, right_title, bullet_items):
    add_header_and_footer(slide, category, question)
    
    # Left Card (Code snippet / editor card)
    # Background card shape
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(24), Pt(56), Pt(325), Pt(298))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CODE_BG
    left_card.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    left_card.line.width = Pt(1)

    # Title of Left Card (Editor Header bar)
    tb_l_header = slide.shapes.add_textbox(Pt(36), Pt(64), Pt(300), Pt(20))
    tf_lh = tb_l_header.text_frame
    tf_lh.word_wrap = True
    tf_lh.margin_left = tf_lh.margin_top = tf_lh.margin_right = tf_lh.margin_bottom = 0
    p_lh = tf_lh.paragraphs[0]
    r_lh = p_lh.add_run()
    r_lh.text = left_title
    r_lh.font.name = "Consolas"
    r_lh.font.size = Pt(9.5)
    r_lh.font.bold = True
    r_lh.font.color.rgb = RGBColor(0x9C, 0xDC, 0xFE)

    # Orange accent line under editor header
    line_lh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(36), Pt(85), Pt(301), Pt(1))
    line_lh.fill.solid()
    line_lh.fill.fore_color.rgb = ORANGE
    line_lh.line.fill.background()

    # Code text box
    tb_code = slide.shapes.add_textbox(Pt(34), Pt(92), Pt(305), Pt(254))
    tf_code = tb_code.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = tf_code.margin_top = tf_code.margin_right = tf_code.margin_bottom = 0
    
    # Render code lines with runs
    first_p = True
    for line_data in code_runs:
        if first_p:
            p = tf_code.paragraphs[0]
            first_p = False
        else:
            p = tf_code.add_paragraph()
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        
        # line_data can be a list of tuples (text, color_rgb)
        if isinstance(line_data, str):
            r = p.add_run()
            r.text = line_data
            r.font.name = "Consolas"
            r.font.size = Pt(8.5)
            r.font.color.rgb = CODE_TEXT
        else:
            for text_chunk, chunk_color in line_data:
                r = p.add_run()
                r.text = text_chunk
                r.font.name = "Consolas"
                r.font.size = Pt(8.5)
                r.font.color.rgb = chunk_color

    # Right Card (Explanation & Theory card)
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(365), Pt(56), Pt(330), Pt(298))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG
    right_card.line.color.rgb = CARD_BORDER
    right_card.line.width = Pt(1)

    # Right card title
    tb_r_header = slide.shapes.add_textbox(Pt(380), Pt(66), Pt(300), Pt(22))
    tf_rh = tb_r_header.text_frame
    tf_rh.word_wrap = True
    tf_rh.margin_left = tf_rh.margin_top = tf_rh.margin_right = tf_rh.margin_bottom = 0
    p_rh = tf_rh.paragraphs[0]
    r_rh = p_rh.add_run()
    r_rh.text = right_title
    r_rh.font.name = "Montserrat"
    r_rh.font.size = Pt(11)
    r_rh.font.bold = True
    r_rh.font.color.rgb = ORANGE

    # Orange line under right header
    line_rh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(380), Pt(88), Pt(300), Pt(1))
    line_rh.fill.solid()
    line_rh.fill.fore_color.rgb = ORANGE
    line_rh.line.fill.background()

    # Right card bullet items
    # Each bullet item: {"badge": "Название", "text": "Текст объяснения", "extra": "Дополнительно"}
    y_offset = Pt(98)
    for b_idx, item in enumerate(bullet_items):
        # Orange dot badge
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(380), y_offset + Pt(3), Pt(6.4), Pt(6.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_DOT
        dot.line.fill.background()

        # Text box for bullet item
        tb_b = slide.shapes.add_textbox(Pt(394), y_offset, Pt(285), Pt(42))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        # Paragraph 1: Badge title
        p_b1 = tf_b.paragraphs[0]
        p_b1.space_after = Pt(2)
        r_title = p_b1.add_run()
        r_title.text = item["title"]
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = DARK_TEXT

        # Paragraph 2: Description
        p_b2 = tf_b.add_paragraph()
        p_b2.line_spacing = 1.15
        p_b2.space_after = Pt(0)
        
        # If desc is string or chunks
        if isinstance(item["desc"], str):
            r_desc = p_b2.add_run()
            r_desc.text = item["desc"]
            r_desc.font.name = "Inter"
            r_desc.font.size = Pt(8.5)
            r_desc.font.color.rgb = BODY_TEXT
        else:
            for d_text, d_bold, d_color in item["desc"]:
                r_desc = p_b2.add_run()
                r_desc.text = d_text
                r_desc.font.name = "Inter"
                r_desc.font.size = Pt(8.5)
                r_desc.font.bold = d_bold
                r_desc.font.color.rgb = d_color

        y_offset += Pt(item.get("height", 46))

print("Helper definitions complete!")
