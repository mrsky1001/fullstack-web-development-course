import sys
import os
import shutil
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

sys.stdout.reconfigure(encoding='utf-8')

# Source base template from perfected Webinar 2
base_pptx = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx"
output_dir = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-03-flexbox-grid-cards"
output_pptx = os.path.join(output_dir, "вебинар 3.pptx")
result_img_path = r"d:\IEEU International East-European University\fullstack-web-development-course\webinar3_result_exact.png"

prs = Presentation(base_pptx)
blank_layout = prs.slide_layouts[12]

# Brand Design Tokens & Colors
ORANGE_PILL   = RGBColor(0xFE, 0x60, 0x02)  # #FE6002
ORANGE_LINE   = RGBColor(0xFF, 0x6F, 0x03)  # #FF6F03
BORDER_GRAY   = RGBColor(0xEE, 0xEE, 0xEE)  # #EEEEEE
QUESTION_TEXT = RGBColor(0xA6, 0xA1, 0xA1)  # #A6A1A1 - TextBox 2
DARK_TEXT     = RGBColor(0x22, 0x22, 0x22)  # #222222
BODY_TEXT     = RGBColor(0x44, 0x44, 0x44)  # #444444
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)  # #FFFFFF
GREEN_BADGE   = RGBColor(0x10, 0xB9, 0x81)  # #10B981
GREEN_TEXT    = RGBColor(0x16, 0xA3, 0x4A)  # #16A34A

# Code Syntax Colors
CODE_BG       = RGBColor(0x18, 0x1A, 0x1F)  # #181A1F
CODE_BORDER   = RGBColor(0x2D, 0x31, 0x39)  # #2D3139
CODE_TEXT     = RGBColor(0xD4, 0xD4, 0xD8)  # #D4D4D8
CODE_KW       = RGBColor(0x38, 0xBD, 0xF8)  # Cyan
CODE_TAG      = RGBColor(0xF4, 0x72, 0xB6)  # Pink
CODE_STR      = RGBColor(0xFC, 0xD3, 0x4D)  # Amber
CODE_COMM     = RGBColor(0x71, 0x71, 0x7A)  # Gray
CODE_PROP     = RGBColor(0xA7, 0x8B, 0xFA)  # Violet
CODE_VAL      = RGBColor(0x4A, 0xDE, 0x80)  # Green

def strip_shape_styles_and_shadows(shape):
    """Ensure zero default PPT theme shadows or effects on created shapes."""
    elem = shape._element
    style = elem.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
    if style is not None:
        elem.remove(style)
    for effect in elem.xpath('.//a:outerShdw'):
        effect.getparent().remove(effect)

def add_header_and_footer(slide, category_text, question_text):
    """Creates top-right orange category pill badge and bottom-left question."""
    pill_width = Pt(max(150, min(330, len(category_text) * 9.8 + 36)))
    pill_left = Pt(700) - pill_width
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_left, Pt(20.4), pill_width, Pt(22.6))
    pill.name = "Google Shape;69;p14"
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = ORANGE_PILL
    pill.line.color.rgb = BORDER_GRAY
    pill.line.width = Pt(0.75)
    strip_shape_styles_and_shadows(pill)

    tf_p = pill.text_frame
    tf_p.word_wrap = False
    tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = 0
    p_p = tf_p.paragraphs[0]
    p_p.alignment = PP_ALIGN.CENTER
    r_p = p_p.add_run()
    r_p.text = category_text.upper()
    r_p.font.name = "Inter"
    r_p.font.size = Pt(12)
    r_p.font.bold = False
    r_p.font.color.rgb = WHITE

    tb_bot = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(500), Pt(20))
    tb_bot.name = "TextBox 2"
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    tf_bot.margin_left = tf_bot.margin_top = tf_bot.margin_right = tf_bot.margin_bottom = 0
    p_bot = tf_bot.paragraphs[0]
    p_bot.alignment = PP_ALIGN.LEFT
    r_bot = p_bot.add_run()
    r_bot.text = question_text
    r_bot.font.name = "Montserrat"
    r_bot.font.size = Pt(10)
    r_bot.font.color.rgb = QUESTION_TEXT

def create_code_explanation_slide(slide, category, question, left_title, code_lines, right_title, bullet_items):
    """Standard 2-column layout: Left code card, Right clean explanation with bullet items."""
    add_header_and_footer(slide, category, question)
    
    # Left Card
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(24), Pt(56), Pt(325), Pt(298))
    left_card.name = "Rounded Rectangle 3"
    left_card.adjustments[0] = 0.04
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CODE_BG
    left_card.line.color.rgb = CODE_BORDER
    left_card.line.width = Pt(1)
    strip_shape_styles_and_shadows(left_card)

    tb_lh = slide.shapes.add_textbox(Pt(38), Pt(64), Pt(295), Pt(18))
    tb_lh.name = "TextBox 4"
    tf_lh = tb_lh.text_frame
    tf_lh.word_wrap = True
    tf_lh.margin_left = tf_lh.margin_top = tf_lh.margin_right = tf_lh.margin_bottom = 0
    p_lh = tf_lh.paragraphs[0]
    r_lh = p_lh.add_run()
    r_lh.text = "📄 " + left_title
    r_lh.font.name = "Consolas"
    r_lh.font.size = Pt(9)
    r_lh.font.bold = True
    r_lh.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)

    line_lh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(36), Pt(84), Pt(301), Pt(1))
    line_lh.name = "Rectangle 5"
    line_lh.fill.solid()
    line_lh.fill.fore_color.rgb = ORANGE_LINE
    line_lh.line.fill.background()
    strip_shape_styles_and_shadows(line_lh)

    tb_code = slide.shapes.add_textbox(Pt(36), Pt(90), Pt(301), Pt(256))
    tb_code.name = "TextBox 6"
    tf_code = tb_code.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = tf_code.margin_top = tf_code.margin_right = tf_code.margin_bottom = 0
    
    first_p = True
    for line in code_lines:
        if first_p:
            p = tf_code.paragraphs[0]
            first_p = False
        else:
            p = tf_code.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        
        if isinstance(line, str):
            r = p.add_run()
            r.text = line
            r.font.name = "Consolas"
            r.font.size = Pt(8)
            r.font.color.rgb = CODE_TEXT
        else:
            for text_chunk, chunk_color in line:
                r = p.add_run()
                r.text = text_chunk
                r.font.name = "Consolas"
                r.font.size = Pt(8)
                r.font.color.rgb = chunk_color

    # Right side
    tb_rh = slide.shapes.add_textbox(Pt(380), Pt(66), Pt(300), Pt(20))
    tb_rh.name = "TextBox 7"
    tf_rh = tb_rh.text_frame
    tf_rh.word_wrap = True
    tf_rh.margin_left = tf_rh.margin_top = tf_rh.margin_right = tf_rh.margin_bottom = 0
    p_rh = tf_rh.paragraphs[0]
    r_rh = p_rh.add_run()
    r_rh.text = right_title
    r_rh.font.name = "Montserrat"
    r_rh.font.size = Pt(10.5)
    r_rh.font.bold = True
    r_rh.font.color.rgb = ORANGE_LINE

    line_rh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(380), Pt(87), Pt(300), Pt(1))
    line_rh.name = "Rectangle 8"
    line_rh.fill.solid()
    line_rh.fill.fore_color.rgb = ORANGE_LINE
    line_rh.line.fill.background()
    strip_shape_styles_and_shadows(line_rh)

    y_offset = Pt(94)
    for b_idx, item in enumerate(bullet_items):
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(380), y_offset + Pt(3), Pt(6.4), Pt(6.4))
        dot.name = f"Rounded Rectangle Dot {b_idx}"
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)

        tb_b = slide.shapes.add_textbox(Pt(394), y_offset, Pt(285), Pt(item.get("height", 46)))
        tb_b.name = f"TextBox Bullet {b_idx}"
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        p_b1 = tf_b.paragraphs[0]
        p_b1.space_after = Pt(2)
        r_title = p_b1.add_run()
        r_title.text = item["title"]
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = DARK_TEXT

        p_b2 = tf_b.add_paragraph()
        p_b2.line_spacing = 1.15
        p_b2.space_after = Pt(0)
        
        if isinstance(item["desc"], str):
            r_desc = p_b2.add_run()
            r_desc.text = item["desc"]
            r_desc.font.name = "Inter"
            r_desc.font.size = Pt(8.5)
            r_desc.font.color.rgb = BODY_TEXT
        else:
            for d_text, d_bold, d_color in item["desc"]:
                r_desc = p_b2.add_run()
                r_desc.text = d_text
                r_desc.font.name = "Inter"
                r_desc.font.size = Pt(8.5)
                r_desc.font.bold = d_bold
                r_desc.font.color.rgb = d_color

        y_offset += Pt(item.get("height", 46))

def create_diagram_flex_grid_slide(slide, category, question):
    """Full-width 3-column diagram comparing 1D Flexbox, 2D Grid, and Interface Synergy."""
    add_header_and_footer(slide, category, question)
    
    col_w = Pt(210)
    col_h = Pt(298)
    top_y = Pt(56)
    
    cols_data = [
        {
            "left": Pt(24),
            "accent": RGBColor(0x4A, 0xDE, 0x80), # Green
            "title": "1. FLEXBOX (1D)",
            "sub": "Линейная ось (ряд / колонка)",
            "code_badge": "display: flex; gap: 12px;",
            "visual_type": "flex",
            "summary_title": "Одномерная цепочка (1D):",
            "summary_text": "Идеально выстраивает элементы вдоль одной линии (шапка, меню, ряд кнопок, футер) и делит пустоту."
        },
        {
            "left": Pt(255),
            "accent": RGBColor(0x38, 0xBD, 0xF8), # Cyan
            "title": "2. CSS GRID (2D)",
            "sub": "Строки + Столбцы (Матрица)",
            "code_badge": "grid-template-columns: repeat(3, 1fr);",
            "visual_type": "grid",
            "summary_title": "Двумерная координатная сетка (2D):",
            "summary_text": "Управляет строками и колонками одновременно. Незаменима для витрин каталога, дашбордов и фотогалерей."
        },
        {
            "left": Pt(486),
            "accent": RGBColor(0xA7, 0x8B, 0xFA), # Violet
            "title": "3. СИНЕРГИЯ СЕТОК",
            "sub": "Совместная работа в UI",
            "code_badge": "Grid (каталог) + Flex (карточка)",
            "visual_type": "synergy",
            "summary_title": "Командная работа в реальном коде:",
            "summary_text": "CSS Grid держит общую 3-колоночную сетку каталога, а Flexbox организует кнопки и контент внутри каждой карточки!"
        }
    ]
    
    for c_idx, c in enumerate(cols_data):
        x = c["left"]
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.name = f"Card Col {c_idx+1}"
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = CODE_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)

        tb_th = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_th = tb_th.text_frame
        tf_th.word_wrap = True
        tf_th.margin_left = tf_th.margin_top = tf_th.margin_right = tf_th.margin_bottom = 0
        p_t1 = tf_th.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = c["title"]
        r_t1.font.name = "Montserrat"
        r_t1.font.size = Pt(11)
        r_t1.font.bold = True
        r_t1.font.color.rgb = c["accent"]

        p_t2 = tf_th.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = c["sub"]
        r_t2.font.name = "Inter"
        r_t2.font.size = Pt(8)
        r_t2.font.color.rgb = CODE_COMM

        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(46), col_w - Pt(24), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = c["accent"]
        div.line.fill.background()
        strip_shape_styles_and_shadows(div)

        vbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(12), top_y + Pt(54), col_w - Pt(24), Pt(176))
        vbox.adjustments[0] = 0.04
        vbox.fill.solid()
        vbox.fill.fore_color.rgb = RGBColor(0x13, 0x14, 0x18)
        vbox.line.color.rgb = CODE_BORDER
        vbox.line.width = Pt(0.75)
        strip_shape_styles_and_shadows(vbox)

        tb_vb = slide.shapes.add_textbox(x + Pt(16), top_y + Pt(56), col_w - Pt(32), Pt(14))
        tf_vb = tb_vb.text_frame
        tf_vb.word_wrap = True
        tf_vb.margin_left = tf_vb.margin_top = tf_vb.margin_right = tf_vb.margin_bottom = 0
        p_vb = tf_vb.paragraphs[0]
        r_vb = p_vb.add_run()
        r_vb.text = c["code_badge"]
        r_vb.font.name = "Consolas"
        r_vb.font.size = Pt(7.5)
        r_vb.font.color.rgb = CODE_COMM

        if c["visual_type"] == "flex":
            arrow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(20), top_y + Pt(74), col_w - Pt(40), Pt(18))
            arrow.adjustments[0] = 0.5
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0x1F, 0x24, 0x22)
            arrow.line.color.rgb = c["accent"]
            arrow.line.width = Pt(0.75)
            strip_shape_styles_and_shadows(arrow)
            p_ar = arrow.text_frame.paragraphs[0]
            p_ar.alignment = PP_ALIGN.CENTER
            r_ar = p_ar.add_run()
            r_ar.text = "ГЛАВНАЯ ОСЬ X (row) ──►"
            r_ar.font.name = "Consolas"
            r_ar.font.size = Pt(7.5)
            r_ar.font.color.rgb = c["accent"]

            item_w = Pt(51)
            item_h = Pt(76)
            for it_i, (it_title, it_prop) in enumerate([("Item 1", "auto"), ("Item 2", "flex: 1"), ("Item 3", "auto")]):
                it_x = x + Pt(20) + Pt(it_i * 58)
                it_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, it_x, top_y + Pt(98), item_w, item_h)
                it_shp.adjustments[0] = 0.08
                it_shp.fill.solid()
                it_shp.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
                it_shp.line.color.rgb = c["accent"]
                it_shp.line.width = Pt(1)
                strip_shape_styles_and_shadows(it_shp)
                
                tf_it = it_shp.text_frame
                tf_it.margin_left = tf_it.margin_top = tf_it.margin_right = tf_it.margin_bottom = 0
                p_it1 = tf_it.paragraphs[0]
                p_it1.alignment = PP_ALIGN.CENTER
                r_it1 = p_it1.add_run()
                r_it1.text = it_title
                r_it1.font.name = "Consolas"
                r_it1.font.size = Pt(8)
                r_it1.font.bold = True
                r_it1.font.color.rgb = WHITE
                
                p_it2 = tf_it.add_paragraph()
                p_it2.alignment = PP_ALIGN.CENTER
                r_it2 = p_it2.add_run()
                r_it2.text = it_prop
                r_it2.font.name = "Consolas"
                r_it2.font.size = Pt(7)
                r_it2.font.color.rgb = c["accent"]

            inf = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(20), top_y + Pt(182), col_w - Pt(40), Pt(38))
            inf.fill.solid()
            inf.fill.fore_color.rgb = RGBColor(0x16, 0x1E, 0x24)
            inf.line.fill.background()
            strip_shape_styles_and_shadows(inf)
            p_inf = inf.text_frame.paragraphs[0]
            p_inf.alignment = PP_ALIGN.CENTER
            r_inf = p_inf.add_run()
            r_inf.text = "justify-content: space-between;\nalign-items: center;"
            r_inf.font.name = "Consolas"
            r_inf.font.size = Pt(7)
            r_inf.font.color.rgb = CODE_TEXT

        elif c["visual_type"] == "grid":
            cell_w = Pt(79)
            cell_h = Pt(66)
            grid_cells = [
                (x + Pt(20), top_y + Pt(76), "Ячейка 1", "col 1 / row 1"),
                (x + Pt(108), top_y + Pt(76), "Ячейка 2", "col 2 / row 1"),
                (x + Pt(20), top_y + Pt(148), "Ячейка 3", "col 1 / row 2"),
                (x + Pt(108), top_y + Pt(148), "Ячейка 4", "col 2 / row 2"),
            ]
            for gx, gy, gtitle, gcoord in grid_cells:
                g_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gx, gy, cell_w, cell_h)
                g_shp.adjustments[0] = 0.08
                g_shp.fill.solid()
                g_shp.fill.fore_color.rgb = RGBColor(0x1B, 0x24, 0x33)
                g_shp.line.color.rgb = c["accent"]
                g_shp.line.width = Pt(1)
                strip_shape_styles_and_shadows(g_shp)
                
                tf_g = g_shp.text_frame
                tf_g.margin_left = tf_g.margin_top = tf_g.margin_right = tf_g.margin_bottom = 0
                p_g1 = tf_g.paragraphs[0]
                p_g1.alignment = PP_ALIGN.CENTER
                r_g1 = p_g1.add_run()
                r_g1.text = gtitle
                r_g1.font.name = "Consolas"
                r_g1.font.size = Pt(8)
                r_g1.font.bold = True
                r_g1.font.color.rgb = WHITE
                
                p_g2 = tf_g.add_paragraph()
                p_g2.alignment = PP_ALIGN.CENTER
                r_g2 = p_g2.add_run()
                r_g2.text = gcoord
                r_g2.font.name = "Inter"
                r_g2.font.size = Pt(6.5)
                r_g2.font.color.rgb = c["accent"]

        else: # synergy
            fr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(20), top_y + Pt(74), col_w - Pt(40), Pt(146))
            fr.adjustments[0] = 0.06
            fr.fill.solid()
            fr.fill.fore_color.rgb = RGBColor(0x20, 0x1A, 0x2E)
            fr.line.color.rgb = c["accent"]
            fr.line.width = Pt(1)
            strip_shape_styles_and_shadows(fr)
            
            p_fr = fr.text_frame.paragraphs[0]
            p_fr.alignment = PP_ALIGN.LEFT
            r_fr = p_fr.add_run()
            r_fr.text = " ВНЕШНИЙ GRID (3 колонки)"
            r_fr.font.name = "Consolas"
            r_fr.font.size = Pt(7)
            r_fr.font.color.rgb = c["accent"]

            in1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(28), top_y + Pt(94), col_w - Pt(56), Pt(56))
            in1.adjustments[0] = 0.08
            in1.fill.solid()
            in1.fill.fore_color.rgb = RGBColor(0x2D, 0x25, 0x3F)
            in1.line.color.rgb = RGBColor(0x6D, 0x28, 0xD9)
            in1.line.width = Pt(1)
            strip_shape_styles_and_shadows(in1)
            
            tf_in1 = in1.text_frame
            p_in1 = tf_in1.paragraphs[0]
            r_in1 = p_in1.add_run()
            r_in1.text = "🖼️ .card-img (200px)\n"
            r_in1.font.name = "Consolas"
            r_in1.font.size = Pt(7)
            r_in1.font.color.rgb = WHITE
            r_in2 = p_in1.add_run()
            r_in2.text = "📄 .card-content (flex: 1)"
            r_in2.font.name = "Consolas"
            r_in2.font.size = Pt(7)
            r_in2.font.color.rgb = RGBColor(0x38, 0xBD, 0xF8)

            in2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(28), top_y + Pt(156), col_w - Pt(56), Pt(54))
            in2.adjustments[0] = 0.08
            in2.fill.solid()
            in2.fill.fore_color.rgb = RGBColor(0x1B, 0x24, 0x33)
            in2.line.color.rgb = RGBColor(0x10, 0xB9, 0x81)
            in2.line.width = Pt(1)
            strip_shape_styles_and_shadows(in2)
            
            tf_in2 = in2.text_frame
            p_in2 = tf_in2.paragraphs[0]
            r_in2a = p_in2.add_run()
            r_in2a.text = "💰 450 ₽/час  |  [👁️] [Бронь]\n"
            r_in2a.font.name = "Consolas"
            r_in2a.font.size = Pt(7)
            r_in2a.font.color.rgb = WHITE
            r_in2b = p_in2.add_run()
            r_in2b.text = "⚙️ display: flex; space-between;"
            r_in2b.font.name = "Consolas"
            r_in2b.font.size = Pt(6.5)
            r_in2b.font.color.rgb = RGBColor(0x4A, 0xDE, 0x80)

        tb_sm = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(236), col_w - Pt(24), Pt(54))
        tf_sm = tb_sm.text_frame
        tf_sm.word_wrap = True
        tf_sm.margin_left = tf_sm.margin_top = tf_sm.margin_right = tf_sm.margin_bottom = 0
        p_s1 = tf_sm.paragraphs[0]
        r_s1 = p_s1.add_run()
        r_s1.text = c["summary_title"] + "\n"
        r_s1.font.name = "Inter"
        r_s1.font.size = Pt(8)
        r_s1.font.bold = True
        r_s1.font.color.rgb = DARK_TEXT

        p_s2 = tf_sm.add_paragraph()
        r_s2 = p_s2.add_run()
        r_s2.text = c["summary_text"]
        r_s2.font.name = "Inter"
        r_s2.font.size = Pt(7.5)
        r_s2.font.color.rgb = CODE_TEXT

def create_card_anatomy_slide(slide, category, question):
    """Full-width visual breakdown of the .room-card component architecture."""
    add_header_and_footer(slide, category, question)
    
    # Left: Mockup of the card
    card_mock = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(36), Pt(56), Pt(280), Pt(298))
    card_mock.adjustments[0] = 0.04
    card_mock.fill.solid()
    card_mock.fill.fore_color.rgb = WHITE
    card_mock.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card_mock.line.width = Pt(1.5)
    strip_shape_styles_and_shadows(card_mock)

    img_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(37), Pt(57), Pt(278), Pt(120))
    img_box.adjustments[0] = 0.04
    img_box.fill.solid()
    img_box.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xF6)
    img_box.line.color.rgb = RGBColor(0x00, 0x7B, 0xFF)
    img_box.line.width = Pt(1)
    strip_shape_styles_and_shadows(img_box)
    p_ib = img_box.text_frame.paragraphs[0]
    p_ib.alignment = PP_ALIGN.CENTER
    r_ib = p_ib.add_run()
    r_ib.text = "🖼️ .card-img-wrap (height: 200px)\nobject-fit: cover; overflow: hidden;"
    r_ib.font.name = "Consolas"
    r_ib.font.size = Pt(8)
    r_ib.font.color.rgb = RGBColor(0x00, 0x7B, 0xFF)

    tb_mt = slide.shapes.add_textbox(Pt(50), Pt(184), Pt(250), Pt(20))
    tf_mt = tb_mt.text_frame
    p_mt = tf_mt.paragraphs[0]
    r_mt = p_mt.add_run()
    r_mt.text = "Мини-офис Focus (.card-title)"
    r_mt.font.name = "Montserrat"
    r_mt.font.size = Pt(10)
    r_mt.font.bold = True
    r_mt.font.color.rgb = DARK_TEXT

    tb_me = slide.shapes.add_textbox(Pt(50), Pt(206), Pt(250), Pt(50))
    tf_me = tb_me.text_frame
    for eq_t in ["• Wi-Fi 500 Мбит/с", "• 4K Монитор", "• Эргономичное кресло"]:
        p_e = tf_me.add_paragraph() if tf_me.paragraphs[0].text else tf_me.paragraphs[0]
        r_e = p_e.add_run()
        r_e.text = eq_t
        r_e.font.name = "Inter"
        r_e.font.size = Pt(8)
        r_e.font.color.rgb = BODY_TEXT

    f_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(37), Pt(265), Pt(278), Pt(1))
    f_line.fill.solid()
    f_line.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    f_line.line.fill.background()
    strip_shape_styles_and_shadows(f_line)

    f_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(37), Pt(268), Pt(278), Pt(46))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = WHITE
    f_box.line.fill.background()
    strip_shape_styles_and_shadows(f_box)
    
    tb_pr = slide.shapes.add_textbox(Pt(48), Pt(276), Pt(100), Pt(24))
    p_pr = tb_pr.text_frame.paragraphs[0]
    r_pr1 = p_pr.add_run()
    r_pr1.text = "450 ₽ "
    r_pr1.font.name = "Montserrat"
    r_pr1.font.size = Pt(10)
    r_pr1.font.bold = True
    r_pr1.font.color.rgb = DARK_TEXT
    r_pr2 = p_pr.add_run()
    r_pr2.text = "/ час"
    r_pr2.font.name = "Inter"
    r_pr2.font.size = Pt(8)
    r_pr2.font.color.rgb = BODY_TEXT

    b_ic = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(170), Pt(274), Pt(22), Pt(22))
    b_ic.adjustments[0] = 0.2
    b_ic.fill.solid()
    b_ic.fill.fore_color.rgb = WHITE
    b_ic.line.color.rgb = RGBColor(0x00, 0x7B, 0xFF)
    b_ic.line.width = Pt(1)
    strip_shape_styles_and_shadows(b_ic)
    p_bic = b_ic.text_frame.paragraphs[0]
    p_bic.alignment = PP_ALIGN.CENTER
    r_bic = p_bic.add_run()
    r_bic.text = "👁️"
    r_bic.font.size = Pt(8)

    b_act = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(200), Pt(274), Pt(100), Pt(22))
    b_act.adjustments[0] = 0.2
    b_act.fill.solid()
    b_act.fill.fore_color.rgb = RGBColor(0x00, 0x7B, 0xFF)
    b_act.line.fill.background()
    strip_shape_styles_and_shadows(b_act)
    p_bact = b_act.text_frame.paragraphs[0]
    p_bact.alignment = PP_ALIGN.CENTER
    r_bact = p_bact.add_run()
    r_bact.text = "Забронировать"
    r_bact.font.name = "Inter"
    r_bact.font.size = Pt(8)
    r_bact.font.bold = True
    r_bact.font.color.rgb = WHITE

    # Right: Architectural Callouts
    tb_rh = slide.shapes.add_textbox(Pt(345), Pt(66), Pt(340), Pt(20))
    tb_rh.name = "TextBox 7"
    tf_rh = tb_rh.text_frame
    p_rh = tf_rh.paragraphs[0]
    r_rh = p_rh.add_run()
    r_rh.text = "4 СЛОЯ КОМПОНЕНТА .ROOM-CARD"
    r_rh.font.name = "Montserrat"
    r_rh.font.size = Pt(10.5)
    r_rh.font.bold = True
    r_rh.font.color.rgb = ORANGE_LINE

    line_rh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(345), Pt(87), Pt(340), Pt(1))
    line_rh.name = "Rectangle 8"
    line_rh.fill.solid()
    line_rh.fill.fore_color.rgb = ORANGE_LINE
    line_rh.line.fill.background()
    strip_shape_styles_and_shadows(line_rh)

    layers = [
        ("1. Медиа-контейнер (.card-img-wrap):", "Фиксированная высота 200px с overflow: hidden. Защищает карточку от вылезания картинок разного разрешения и формы."),
        ("2. Контентное тело (.card-content):", "Внутренний padding 15px и колоночный Flexbox (display: flex; flex-direction: column; flex: 1;). Гарантирует равную высоту всех карточек в каталоге."),
        ("3. Распорка списка (.card-equipment):", "Список тегов или оборудования с flex: 1 забирает излишки пустого места и выталкивает блок стоимости и кнопок строго в нижний край."),
        ("4. Подвал действия (.card-footer):", "Flexbox с justify-content: space-between разносит цену налево, а блок кнопок .card-btns (иконка детального просмотра + primary кнопка) — направо.")
    ]

    y_pos = Pt(94)
    for l_idx, (ltitle, ldesc) in enumerate(layers):
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(345), y_pos + Pt(3), Pt(6.4), Pt(6.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)

        tb_l = slide.shapes.add_textbox(Pt(360), y_pos, Pt(325), Pt(48))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
        
        p1 = tf_l.paragraphs[0]
        p1.space_after = Pt(2)
        r1 = p1.add_run()
        r1.text = ltitle
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT

        p2 = tf_l.add_paragraph()
        p2.line_spacing = 1.15
        r2 = p2.add_run()
        r2.text = ldesc
        r2.font.name = "Inter"
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = BODY_TEXT

        y_pos += Pt(58)

def create_grid_cards_slide(slide, category, question, title_text, cards_data):
    """2x2 error analysis cards."""
    add_header_and_footer(slide, category, question)
    
    coords = [
        (Pt(24), Pt(56), Pt(325), Pt(142)),
        (Pt(365), Pt(56), Pt(330), Pt(142)),
        (Pt(24), Pt(210), Pt(325), Pt(142)),
        (Pt(365), Pt(210), Pt(330), Pt(142)),
    ]
    
    for idx, (x, y, w, h) in enumerate(coords):
        data = cards_data[idx]
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.adjustments[0] = 0.04
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        box.line.width = Pt(1)
        strip_shape_styles_and_shadows(box)

        tb_h = slide.shapes.add_textbox(x + Pt(14), y + Pt(10), w - Pt(28), Pt(20))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
        p_h = tf_h.paragraphs[0]
        r_num = p_h.add_run()
        r_num.text = f"⚠️ {data['title']}"
        r_num.font.name = "Montserrat"
        r_num.font.size = Pt(10)
        r_num.font.bold = True
        r_num.font.color.rgb = ORANGE_LINE
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(14), y + Pt(32), w - Pt(28), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE_LINE
        line.line.fill.background()
        strip_shape_styles_and_shadows(line)
        
        tb_body = slide.shapes.add_textbox(x + Pt(14), y + Pt(38), w - Pt(28), Pt(94))
        tf_b = tb_body.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        p1 = tf_b.paragraphs[0]
        p1.space_after = Pt(3)
        r1_lbl = p1.add_run()
        r1_lbl.text = "Ошибка: "
        r1_lbl.font.name = "Inter"
        r1_lbl.font.size = Pt(8.5)
        r1_lbl.font.bold = True
        r1_lbl.font.color.rgb = DARK_TEXT
        r1_txt = p1.add_run()
        r1_txt.text = data['problem']
        r1_txt.font.name = "Inter"
        r1_txt.font.size = Pt(8.5)
        r1_txt.font.color.rgb = BODY_TEXT
        
        p2 = tf_b.add_paragraph()
        p2.space_after = Pt(0)
        r2_lbl = p2.add_run()
        r2_lbl.text = "Решение: "
        r2_lbl.font.name = "Inter"
        r2_lbl.font.size = Pt(8.5)
        r2_lbl.font.bold = True
        r2_lbl.font.color.rgb = GREEN_TEXT
        r2_txt = p2.add_run()
        r2_txt.text = data['solution']
        r2_txt.font.name = "Inter"
        r2_txt.font.size = Pt(8.5)
        r2_txt.font.color.rgb = BODY_TEXT

def create_checklist_slide(slide, category, question, title_text, items):
    """Clean 6-item checklist with checkmarks."""
    add_header_and_footer(slide, category, question)

    tb_h = slide.shapes.add_textbox(Pt(54), Pt(68), Pt(600), Pt(22))
    tf_h = tb_h.text_frame
    tf_h.word_wrap = True
    tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
    p_h = tf_h.paragraphs[0]
    r_h = p_h.add_run()
    r_h.text = title_text
    r_h.font.name = "Montserrat"
    r_h.font.size = Pt(11)
    r_h.font.bold = True
    r_h.font.color.rgb = ORANGE_LINE

    line_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(54), Pt(92), Pt(612), Pt(1))
    line_h.fill.solid()
    line_h.fill.fore_color.rgb = ORANGE_LINE
    line_h.line.fill.background()
    strip_shape_styles_and_shadows(line_h)

    y_pos = Pt(102)
    for idx, item in enumerate(items):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(54), y_pos + Pt(2), Pt(22), Pt(20))
        badge.adjustments[0] = 0.2
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN_BADGE
        badge.line.fill.background()
        strip_shape_styles_and_shadows(badge)
        
        tb_chk = slide.shapes.add_textbox(Pt(54), y_pos + Pt(4), Pt(22), Pt(20))
        tf_chk = tb_chk.text_frame
        p_chk = tf_chk.paragraphs[0]
        p_chk.alignment = PP_ALIGN.CENTER
        r_chk = p_chk.add_run()
        r_chk.text = "✔"
        r_chk.font.name = "Segoe UI Symbol"
        r_chk.font.size = Pt(10)
        r_chk.font.bold = True
        r_chk.font.color.rgb = WHITE
        
        tb_item = slide.shapes.add_textbox(Pt(86), y_pos, Pt(570), Pt(38))
        tf_item = tb_item.text_frame
        tf_item.word_wrap = True
        tf_item.margin_left = tf_item.margin_top = tf_item.margin_right = tf_item.margin_bottom = 0
        
        p1 = tf_item.paragraphs[0]
        p1.space_after = Pt(2)
        r1 = p1.add_run()
        r1.text = item["title"]
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT

        p2 = tf_item.add_paragraph()
        r2 = p2.add_run()
        r2.text = item["desc"]
        r2.font.name = "Inter"
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = BODY_TEXT

        y_pos += Pt(42)

print("Ready to assemble Webinar 3 presentation.")
