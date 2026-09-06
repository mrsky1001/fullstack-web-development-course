import os
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

# -------------------------------------------------------------
# PATH CONFIGURATION
# -------------------------------------------------------------
base_template_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx"
output_pptx_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-06-forms-validation\вебинар 6.pptx"
result_img_path = r"d:\IEEU International East-European University\fullstack-web-development-course\webinar6_result_exact.png"

# -------------------------------------------------------------
# PALETTE & STYLES (Strict corporate college theme)
# -------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x18, 0x1A, 0x1F)
CODE_BORDER = RGBColor(0x2D, 0x31, 0x39)
CODE_TEXT = RGBColor(0xD4, 0xD4, 0xD8)
ORANGE_PILL = RGBColor(0xFE, 0x60, 0x02)
ORANGE_LINE = RGBColor(0xFE, 0x60, 0x02)
QUESTION_TEXT = RGBColor(0xA6, 0xA1, 0xA1) # Strictly #A6A1A1
TEXT_PRIMARY = RGBColor(0x18, 0x1A, 0x1F)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
KEYWORD_COLOR = RGBColor(0xC6, 0x78, 0xDD) # Purple
FUNC_COLOR = RGBColor(0x61, 0xAF, 0xEF)    # Blue
STR_COLOR = RGBColor(0x98, 0xC3, 0x79)     # Green
COMMENT_COLOR = RGBColor(0x5C, 0x63, 0x70) # Gray
NUM_COLOR = RGBColor(0xD1, 0x9A, 0x66)     # Orange / number
TAG_COLOR = RGBColor(0xE0, 0x6C, 0x75)     # Red/Coral
DANGER_COLOR = RGBColor(0xDC, 0x26, 0x26)  # Red for invalid state

def strip_shape_styles_and_shadows(shape):
    """Zero tolerance to shadows and PowerPoint theme styles."""
    spPr = shape._element.spPr
    for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst'):
        spPr.remove(shd)
    for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw'):
        spPr.remove(shd)
    style = shape._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
    if style is not None:
        shape._element.remove(style)

def add_header_and_footer(slide, category_text, question_text):
    """Standard header pill badge with dynamic width and bottom question prompt."""
    pill_w = Pt(max(140, len(category_text) * 7.8 + 24))
    pill_h = Pt(22.6)
    pill_left = Pt(700) - pill_w
    pill_top = Pt(20.4)
    
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_left, pill_top, pill_w, pill_h)
    badge.name = "Google Shape;69;p14"
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE_PILL
    badge.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    badge.line.width = Pt(1)
    strip_shape_styles_and_shadows(badge)
    
    tf_b = badge.text_frame
    tf_b.word_wrap = False
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = category_text
    r_b.font.name = "Inter"
    r_b.font.size = Pt(12)
    r_b.font.bold = False
    r_b.font.color.rgb = WHITE
    
    # TextBox 2 at bottom
    tb2 = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2.name = "TextBox 2"
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = question_text
    r2.font.name = "Montserrat"
    r2.font.size = Pt(10)
    r2.font.bold = False
    r2.font.color.rgb = QUESTION_TEXT

def create_code_explanation_slide(slide, category, question, left_title, code_lines, right_title, bullet_items):
    """Standard 2-column layout: Left code card, Right clean explanation with bullet items."""
    add_header_and_footer(slide, category, question)
    
    # Left Card
    card_x = Pt(24)
    card_y = Pt(56)
    card_w = Pt(325)
    card_h = Pt(296)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
    card.name = "Rounded Rectangle 3"
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG
    card.line.color.rgb = CODE_BORDER
    card.line.width = Pt(1)
    strip_shape_styles_and_shadows(card)
    
    # Header inside card
    tb_hdr = slide.shapes.add_textbox(card_x + Pt(14), card_y + Pt(10), card_w - Pt(28), Pt(18))
    tb_hdr.name = "TextBox 4"
    tf_h = tb_hdr.text_frame
    tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
    p_h = tf_h.paragraphs[0]
    r_h = p_h.add_run()
    r_h.text = left_title
    r_h.font.name = "Consolas"
    r_h.font.size = Pt(9.5)
    r_h.font.bold = True
    r_h.font.color.rgb = FUNC_COLOR
    
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_x + Pt(14), card_y + Pt(30), card_w - Pt(28), Pt(1))
    div.name = "Rectangle 5"
    div.fill.solid()
    div.fill.fore_color.rgb = ORANGE_LINE
    div.line.fill.background()
    strip_shape_styles_and_shadows(div)
    
    tb_code = slide.shapes.add_textbox(card_x + Pt(14), card_y + Pt(36), card_w - Pt(28), card_h - Pt(44))
    tb_code.name = "TextBox 6"
    tf_c = tb_code.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    
    for i, line_data in enumerate(code_lines):
        p = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
        p.space_after = Pt(2)
        if isinstance(line_data, str):
            r = p.add_run()
            r.text = line_data
            r.font.name = "Consolas"
            r.font.size = Pt(8.5)
            r.font.color.rgb = CODE_TEXT
        elif isinstance(line_data, list):
            for token_text, token_color in line_data:
                r = p.add_run()
                r.text = token_text
                r.font.name = "Consolas"
                r.font.size = Pt(8.5)
                r.font.color.rgb = token_color
                
    # Right Column
    rcol_x = Pt(365)
    rcol_y = Pt(66)
    rcol_w = Pt(330)
    
    tb_rtitle = slide.shapes.add_textbox(rcol_x, rcol_y, rcol_w, Pt(22))
    tb_rtitle.name = "TextBox 7"
    tf_rt = tb_rtitle.text_frame
    tf_rt.word_wrap = True
    tf_rt.margin_left = tf_rt.margin_top = tf_rt.margin_right = tf_rt.margin_bottom = 0
    p_rt = tf_rt.paragraphs[0]
    r_rt = p_rt.add_run()
    r_rt.text = right_title
    r_rt.font.name = "Montserrat"
    r_rt.font.size = Pt(12)
    r_rt.font.bold = True
    r_rt.font.color.rgb = ORANGE_LINE
    
    rdiv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rcol_x, rcol_y + Pt(25), rcol_w, Pt(1))
    rdiv.name = "Rectangle 8"
    rdiv.fill.solid()
    rdiv.fill.fore_color.rgb = ORANGE_LINE
    rdiv.line.fill.background()
    strip_shape_styles_and_shadows(rdiv)
    
    cur_y = rcol_y + Pt(34)
    for idx, item in enumerate(bullet_items):
        item_h = Pt(item.get("height", 42))
        
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rcol_x, cur_y + Pt(3), Pt(7), Pt(7))
        dot.name = f"Rounded Rectangle Dot {idx}"
        dot.adjustments[0] = 0.2
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)
        
        tb_item = slide.shapes.add_textbox(rcol_x + Pt(14), cur_y, rcol_w - Pt(14), item_h)
        tb_item.name = f"TextBox Bullet {idx}"
        tf_i = tb_item.text_frame
        tf_i.word_wrap = True
        tf_i.margin_left = tf_i.margin_top = tf_i.margin_right = tf_i.margin_bottom = 0
        
        p = tf_i.paragraphs[0]
        p.space_after = Pt(0)
        
        r_title = p.add_run()
        r_title.text = item["title"] + ":\n"
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = TEXT_PRIMARY
        
        r_desc = p.add_run()
        r_desc.text = item["desc"]
        r_desc.font.name = "Inter"
        r_desc.font.size = Pt(8.5)
        r_desc.font.bold = False
        r_desc.font.color.rgb = TEXT_MUTED
        
        cur_y += item_h + Pt(6)

def create_pipeline_diagram_slide(slide, category, question):
    """Custom Diagram 1: Form submission & validation pipeline."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. ПЕРЕХВАТ SUBMIT", "subtitle": "Отмена перезагрузки", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. ПРОВЕРКА ДАННЫХ", "subtitle": "Алгоритм валидации", "accent": NUM_COLOR},
        {"x": Pt(476), "title": "3. РЕАКЦИЯ ИНТЕРФЕЙСА", "subtitle": "UX ошибок и успеха", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r1 = p_t1.add_run()
        r1.text = c["title"]
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r2 = p_t2.add_run()
        r2.text = c["subtitle"]
        r2.font.name = "Inter"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CODE_TEXT
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = c["accent"]
        line.line.fill.background()
        strip_shape_styles_and_shadows(line)
        
        if c["title"].startswith("1."):
            b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b1.adjustments[0] = 0.1
            b1.fill.solid()
            b1.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b1.line.color.rgb = c["accent"]
            b1.line.width = Pt(1)
            strip_shape_styles_and_shadows(b1)
            p_b1 = b1.text_frame.paragraphs[0]
            p_b1.alignment = PP_ALIGN.CENTER
            r_b1 = p_b1.add_run()
            r_b1.text = "form.addEventListener(\n  'submit', (e) => {\n    e.preventDefault();\n});"
            r_b1.font.name = "Consolas"
            r_b1.font.size = Pt(7)
            r_b1.font.color.rgb = WHITE
            
            tb_desc = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d = tb_desc.text_frame
            tf_d.word_wrap = True
            tf_d.margin_left = tf_d.margin_top = tf_d.margin_right = tf_d.margin_bottom = 0
            p_d = tf_d.paragraphs[0]
            r_d = p_d.add_run()
            r_d.text = "Полный контроль браузера:\n\n• Пользователь жмет кнопку отправки\n• Браузер генерирует событие submit\n• e.preventDefault() отменяет reload\n• Страница сохраняет введенный ввод\n• Управление переходит к JavaScript\n• Исключена потеря данных формы"
            r_d.font.name = "Inter"
            r_d.font.size = Pt(8)
            r_d.font.color.rgb = CODE_TEXT
            
        elif c["title"].startswith("2."):
            b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b2.adjustments[0] = 0.1
            b2.fill.solid()
            b2.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b2.line.color.rgb = c["accent"]
            b2.line.width = Pt(1)
            strip_shape_styles_and_shadows(b2)
            p_b2 = b2.text_frame.paragraphs[0]
            p_b2.alignment = PP_ALIGN.CENTER
            r_b2 = p_b2.add_run()
            r_b2.text = "val = input.value.trim();\nisValid = val.length > 0;\npass === confirmPass;\nregex.test(email);"
            r_b2.font.name = "Consolas"
            r_b2.font.size = Pt(6.5)
            r_b2.font.color.rgb = NUM_COLOR
            
            tb_desc2 = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d2 = tb_desc2.text_frame
            tf_d2.word_wrap = True
            tf_d2.margin_left = tf_d2.margin_top = tf_d2.margin_right = tf_d2.margin_bottom = 0
            p_d2 = tf_d2.paragraphs[0]
            r_d2 = p_d2.add_run()
            r_d2.text = "Многоуровневый фильтр:\n\n• Очистка пробелов через .trim()\n• Проверка на пустые обязательные поля\n• Сравнение паролей на идентичность\n• Проверка минимальной длины (6 симв.)\n• Регулярные выражения для Email\n• Аккумуляция флага let isValid"
            r_d2.font.name = "Inter"
            r_d2.font.size = Pt(8)
            r_d2.font.color.rgb = CODE_TEXT
            
        elif c["title"].startswith("3."):
            b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b3.adjustments[0] = 0.1
            b3.fill.solid()
            b3.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b3.line.color.rgb = c["accent"]
            b3.line.width = Pt(1)
            strip_shape_styles_and_shadows(b3)
            p_b3 = b3.text_frame.paragraphs[0]
            p_b3.alignment = PP_ALIGN.CENTER
            r_b3 = p_b3.add_run()
            r_b3.text = "if (!isValid) {\n  input.classList.add('is-invalid');\n} else {\n  showNotification('Успех!');\n}"
            r_b3.font.name = "Consolas"
            r_b3.font.size = Pt(6.5)
            r_b3.font.color.rgb = STR_COLOR
            
            tb_desc3 = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d3 = tb_desc3.text_frame
            tf_d3.word_wrap = True
            tf_d3.margin_left = tf_d3.margin_top = tf_d3.margin_right = tf_d3.margin_bottom = 0
            p_d3 = tf_d3.paragraphs[0]
            r_d3 = p_d3.add_run()
            r_d3.text = "Визуальный отклик интерфейса:\n\n• ОШИБКА: красная рамка у полей\n• Сообщение под полем или в алерте\n• УСПЕХ: всплывающий Toast 'Успешно'\n• Сохранение сессии в localStorage\n• Очистка формы form.reset()\n• Плавный редирект в личный кабинет"
            r_d3.font.name = "Inter"
            r_d3.font.size = Pt(8)
            r_d3.font.color.rgb = CODE_TEXT

def create_input_states_diagram_slide(slide, category, question):
    """Custom Diagram 2: Input visual states comparison."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. DEFAULT / FOCUS", "subtitle": "Нейтральное состояние", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. INVALID (ОШИБКА)", "subtitle": "Класс .is-invalid", "accent": DANGER_COLOR},
        {"x": Pt(476), "title": "3. VALID (УСПЕХ)", "subtitle": "Данные корректны", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r1 = p_t1.add_run()
        r1.text = c["title"]
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r2 = p_t2.add_run()
        r2.text = c["subtitle"]
        r2.font.name = "Inter"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CODE_TEXT
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = c["accent"]
        line.line.fill.background()
        strip_shape_styles_and_shadows(line)
        
        if c["title"].startswith("1."):
            b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b1.adjustments[0] = 0.1
            b1.fill.solid()
            b1.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b1.line.color.rgb = RGBColor(0x64, 0x74, 0x8B)
            b1.line.width = Pt(1)
            strip_shape_styles_and_shadows(b1)
            p_b1 = b1.text_frame.paragraphs[0]
            p_b1.alignment = PP_ALIGN.CENTER
            r_b1 = p_b1.add_run()
            r_b1.text = "[ Логин: Введите логин... ]\nborder: 1px solid #E2E8F0;\noutline: none;"
            r_b1.font.name = "Consolas"
            r_b1.font.size = Pt(7)
            r_b1.font.color.rgb = WHITE
            
            tb_desc = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d = tb_desc.text_frame
            tf_d.word_wrap = True
            tf_d.margin_left = tf_d.margin_top = tf_d.margin_right = tf_d.margin_bottom = 0
            p_d = tf_d.paragraphs[0]
            r_d = p_d.add_run()
            r_d.text = "Исходный вид поля ввода:\n\n• Нейтральная тонкая серая граница\n• Серый плейсхолдер с подсказкой\n• Синее кольцо фокуса при клике\n• Отсутствие тревожных цветов\n• Подготовка к комфортному вводу\n• Поддержка автозаполнения браузера"
            r_d.font.name = "Inter"
            r_d.font.size = Pt(8)
            r_d.font.color.rgb = CODE_TEXT
            
        elif c["title"].startswith("2."):
            b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b2.adjustments[0] = 0.1
            b2.fill.solid()
            b2.fill.fore_color.rgb = RGBColor(0x3B, 0x14, 0x1A)
            b2.line.color.rgb = DANGER_COLOR
            b2.line.width = Pt(1.5)
            strip_shape_styles_and_shadows(b2)
            p_b2 = b2.text_frame.paragraphs[0]
            p_b2.alignment = PP_ALIGN.CENTER
            r_b2 = p_b2.add_run()
            r_b2.text = "[ Пароль: ••••••       ! ]\nborder-color: #DC2626;\nbackground: #FEF2F2;"
            r_b2.font.name = "Consolas"
            r_b2.font.size = Pt(7)
            r_b2.font.color.rgb = RGBColor(0xF8, 0x71, 0x71)
            
            tb_desc2 = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d2 = tb_desc2.text_frame
            tf_d2.word_wrap = True
            tf_d2.margin_left = tf_d2.margin_top = tf_d2.margin_right = tf_d2.margin_bottom = 0
            p_d2 = tf_d2.paragraphs[0]
            r_d2 = p_d2.add_run()
            r_d2.text = "Индикация ошибки (.is-invalid):\n\n• Яркая красная граница 1.5px\n• Светло-розовый оттенок фона\n• Мгновенно привлекает внимание\n• Снятие ошибки по событию input\n• Подсказка с конкретной причиной\n• Блокировка отправки формы"
            r_d2.font.name = "Inter"
            r_d2.font.size = Pt(8)
            r_d2.font.color.rgb = CODE_TEXT
            
        elif c["title"].startswith("3."):
            b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b3.adjustments[0] = 0.1
            b3.fill.solid()
            b3.fill.fore_color.rgb = RGBColor(0x14, 0x2E, 0x22)
            b3.line.color.rgb = STR_COLOR
            b3.line.width = Pt(1)
            strip_shape_styles_and_shadows(b3)
            p_b3 = b3.text_frame.paragraphs[0]
            p_b3.alignment = PP_ALIGN.CENTER
            r_b3 = p_b3.add_run()
            r_b3.text = "[ Логин: admin         ✔ ]\nborder-color: #0E9F6E;\nstatus: VALID;"
            r_b3.font.name = "Consolas"
            r_b3.font.size = Pt(7)
            r_b3.font.color.rgb = STR_COLOR
            
            tb_desc3 = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d3 = tb_desc3.text_frame
            tf_d3.word_wrap = True
            tf_d3.margin_left = tf_d3.margin_top = tf_d3.margin_right = tf_d3.margin_bottom = 0
            p_d3 = tf_d3.paragraphs[0]
            r_d3 = p_d3.add_run()
            r_d3.text = "Успешная валидация:\n\n• Зеленая граница подтверждения\n• Символ галочки завершения ввода\n• Разрешение отправки на сервер\n• Показ всплывающего Toast\n• Запись сессии в localStorage\n• Редирект в защищенный раздел"
            r_d3.font.name = "Inter"
            r_d3.font.size = Pt(8)
            r_d3.font.color.rgb = CODE_TEXT

def create_grid_cards_slide(slide, category, question, title_text, cards_data):
    """2x2 error analysis cards."""
    add_header_and_footer(slide, category, question)
    
    coords = [
        (Pt(24), Pt(56), Pt(325), Pt(142)),
        (Pt(365), Pt(56), Pt(330), Pt(142)),
        (Pt(24), Pt(210), Pt(325), Pt(142)),
        (Pt(365), Pt(210), Pt(330), Pt(142)),
    ]
    
    for idx, (x, y, w, h) in enumerate(coords):
        data = cards_data[idx]
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.adjustments[0] = 0.04
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        box.line.width = Pt(1)
        strip_shape_styles_and_shadows(box)

        tb_h = slide.shapes.add_textbox(x + Pt(14), y + Pt(10), w - Pt(28), Pt(20))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
        p_h = tf_h.paragraphs[0]
        r_num = p_h.add_run()
        r_num.text = f"{idx + 1}. {data['title']}"
        r_num.font.name = "Montserrat"
        r_num.font.size = Pt(10)
        r_num.font.bold = True
        r_num.font.color.rgb = ORANGE_LINE
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(14), y + Pt(32), w - Pt(28), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE_LINE
        line.line.fill.background()
        strip_shape_styles_and_shadows(line)
        
        tb_body = slide.shapes.add_textbox(x + Pt(14), y + Pt(38), w - Pt(28), Pt(94))
        tf_b = tb_body.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        p1 = tf_b.paragraphs[0]
        p1.space_after = Pt(3)
        r1_lbl = p1.add_run()
        r1_lbl.text = "Ошибка: "
        r1_lbl.font.name = "Inter"
        r1_lbl.font.size = Pt(8.5)
        r1_lbl.font.bold = True
        r1_lbl.font.color.rgb = TEXT_PRIMARY
        
        r1_txt = p1.add_run()
        r1_txt.text = data["problem"]
        r1_txt.font.name = "Inter"
        r1_txt.font.size = Pt(8)
        r1_txt.font.color.rgb = TEXT_MUTED
        
        p2 = tf_b.add_paragraph()
        r2_lbl = p2.add_run()
        r2_lbl.text = "Решение: "
        r2_lbl.font.name = "Inter"
        r2_lbl.font.size = Pt(8.5)
        r2_lbl.font.bold = True
        r2_lbl.font.color.rgb = RGBColor(0x0E, 0x9F, 0x6E) # Green
        
        r2_txt = p2.add_run()
        r2_txt.text = data["solution"]
        r2_txt.font.name = "Inter"
        r2_txt.font.size = Pt(8)
        r2_txt.font.color.rgb = TEXT_MUTED

def create_checklist_slide(slide, category, question, title_text, items):
    """Quality checklist with clean checkmark cards."""
    add_header_and_footer(slide, category, question)
    
    tb_t = slide.shapes.add_textbox(Pt(24), Pt(56), Pt(670), Pt(22))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    r_t = p_t.add_run()
    r_t.text = title_text
    r_t.font.name = "Montserrat"
    r_t.font.size = Pt(11)
    r_t.font.bold = True
    r_t.font.color.rgb = ORANGE_LINE
    
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(24), Pt(80), Pt(670), Pt(1))
    div.fill.solid()
    div.fill.fore_color.rgb = ORANGE_LINE
    div.line.fill.background()
    strip_shape_styles_and_shadows(div)
    
    y_start = Pt(92)
    row_h = Pt(40)
    for idx, item in enumerate(items):
        cur_y = y_start + idx * (row_h + Pt(4))
        
        # Check badge
        chk = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(24), cur_y + Pt(2), Pt(24), Pt(24))
        chk.adjustments[0] = 0.2
        chk.fill.solid()
        chk.fill.fore_color.rgb = RGBColor(0x0E, 0x9F, 0x6E) # Green
        chk.line.fill.background()
        strip_shape_styles_and_shadows(chk)
        
        p_c = chk.text_frame.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        r_c = p_c.add_run()
        r_c.text = "✔"
        r_c.font.name = "Segoe UI"
        r_c.font.size = Pt(10)
        r_c.font.color.rgb = WHITE
        
        # Text
        tb_item = slide.shapes.add_textbox(Pt(58), cur_y, Pt(636), row_h)
        tf_i = tb_item.text_frame
        tf_i.word_wrap = True
        tf_i.margin_left = tf_i.margin_top = tf_i.margin_right = tf_i.margin_bottom = 0
        
        p = tf_i.paragraphs[0]
        r_lbl = p.add_run()
        r_lbl.text = item["title"] + ": "
        r_lbl.font.name = "Montserrat"
        r_lbl.font.size = Pt(9.5)
        r_lbl.font.bold = True
        r_lbl.font.color.rgb = TEXT_PRIMARY
        
        p2 = tf_i.add_paragraph()
        r_desc = p2.add_run()
        r_desc.text = item["desc"]
        r_desc.font.name = "Inter"
        r_desc.font.size = Pt(8.5)
        r_desc.font.color.rgb = TEXT_MUTED

# -------------------------------------------------------------
# PRESENTATION DATA (21 content slides)
# -------------------------------------------------------------
slides_data = [
    # 1. Slide 3: Submit Event & Prevent Default
    {
        "type": "code",
        "category": "КОНТРОЛЬ ОТПРАВКИ",
        "question": "Почему браузер перезагружает страницу при отправке формы?",
        "left_title": "📄 js/main.js (Перехват submit)",
        "code_lines": [
            [("// Перехват стандартного поведения формы", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("form = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('registerForm');", STR_COLOR)],
            [("", CODE_TEXT)],
            [("form.", CODE_TEXT), ("addEventListener", FUNC_COLOR), ("('submit', (event) => {", STR_COLOR)],
            [("  // 1. Отменяем перезагрузку страницы!", COMMENT_COLOR)],
            [("  event.", CODE_TEXT), ("preventDefault", FUNC_COLOR), ("();", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // 2. Теперь мы можем безопасно проверить данные", COMMENT_COLOR)],
            [("  console.", FUNC_COLOR), ("log", FUNC_COLOR), ("('Форма перехвачена скриптом!');", STR_COLOR)],
            [("});", CODE_TEXT)]
        ],
        "right_title": "КОНТРОЛЬ ОТПРАВКИ: СОБЫТИЕ SUBMIT",
        "bullet_items": [
            {
                "title": "Поведение формы по умолчанию",
                "desc": "Исторически HTML-форма отправляет GET/POST запрос на сервер с полной перезагрузкой документа.",
                "height": 40
            },
            {
                "title": "Метод event.preventDefault()",
                "desc": "Отменяет действие браузера по умолчанию, позволяя выполнить клиентскую валидацию перед отправкой.",
                "height": 40
            },
            {
                "title": "Слушатель на форме, а не на кнопке",
                "desc": "Событие submit срабатывает как при клике на кнопку «Отправить», так и при нажатии клавиши Enter в поле.",
                "height": 40
            },
            {
                "title": "Сохранение введенных данных",
                "desc": "При ошибке страница не перезагружается, сохраняя введенные пользователем значения в полях.",
                "height": 40
            }
        ]
    },

    # 2. Slide 4: Input sanitization via .trim()
    {
        "type": "code",
        "category": "ОЧИСТКА ВВОДА",
        "question": "Как предотвратить отправку формы, если пользователь ввел одни пробелы?",
        "left_title": "📄 js/main.js (Метод String.trim)",
        "code_lines": [
            [("// Проверка поля с очисткой от пробелов", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("input = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('login');", STR_COLOR)],
            [("const ", KEYWORD_COLOR), ("rawVal = input.value;        ", CODE_TEXT), ("// '   admin   '", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("cleanVal = rawVal.", CODE_TEXT), ("trim", FUNC_COLOR), ("();     ", CODE_TEXT), ("// 'admin'", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("if ", KEYWORD_COLOR), ("(!cleanVal) {", CODE_TEXT)],
            [("  // Поле пустое или содержало одни пробелы!", COMMENT_COLOR)],
            [("  input.classList.", CODE_TEXT), ("add", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("} ", CODE_TEXT), ("else ", KEYWORD_COLOR), ("{", CODE_TEXT)],
            [("  input.classList.", CODE_TEXT), ("remove", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "САНИТИЗАЦИЯ СТРОК: МЕТОД .TRIM()",
        "bullet_items": [
            {
                "title": "Опасность ложной заполненности",
                "desc": "Строка из пробелов ('   ') имеет длину больше нуля, и простая проверка if (val) сочтет поле валидным.",
                "height": 40
            },
            {
                "title": "Очистка граничных пробелов",
                "desc": "Метод .trim() удаляет пробельные символы в начале и в конце строки, не затрагивая пробелы внутри текста.",
                "height": 40
            },
            {
                "title": "Проверка истинной пустоты",
                "desc": "Конструкция !input.value.trim() безошибочно выявляет пустые поля и попытки обмана формы.",
                "height": 40
            },
            {
                "title": "Чистота передаваемых данных",
                "desc": "Очищенное значение гарантирует, что логин и пароль запишутся в базу без случайных хвостовых пробелов.",
                "height": 40
            }
        ]
    },

    # 3. Slide 5: Diagram: Validation Pipeline
    {
        "type": "pipeline",
        "category": "СХЕМА: ВАЛИДАЦИЯ",
        "question": "Каков полный пошаговый цикл клиентской валидации данных формы?"
    },

    # 4. Slide 6: Regular Expressions Basics
    {
        "type": "code",
        "category": "РЕГУЛЯРНЫЕ ВЫРАЖЕНИЯ",
        "question": "Как работают регулярные выражения и метод .test() в JavaScript?",
        "left_title": "📄 js/main.js (Основы RegExp)",
        "code_lines": [
            [("// Создание регулярного выражения (литерал)", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("emailRegex = /^[^\s@]+@[^\s@]+\.[^\s@]+$/;", NUM_COLOR)],
            [("", CODE_TEXT)],
            [("// Метод .test() возвращает true или false", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("isValidEmail = emailRegex.", CODE_TEXT), ("test", FUNC_COLOR), ("('user@mail.ru');", STR_COLOR)],
            [("console.", FUNC_COLOR), ("log", FUNC_COLOR), ("(isValidEmail); ", CODE_TEXT), ("// true", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("const ", KEYWORD_COLOR), ("isBad = emailRegex.", CODE_TEXT), ("test", FUNC_COLOR), ("('wrong-email');", STR_COLOR)],
            [("console.", FUNC_COLOR), ("log", FUNC_COLOR), ("(isBad); ", CODE_TEXT), ("// false", COMMENT_COLOR)]
        ],
        "right_title": "РЕГУЛЯРНЫЕ ВЫРАЖЕНИЯ: МЕТОД REGEXP.TEST()",
        "bullet_items": [
            {
                "title": "Что такое RegExp (Регулярные выражения)",
                "desc": "Специальный шаблон для поиска и проверки соответствия строк определенному синтаксическому формату.",
                "height": 40
            },
            {
                "title": "Быстрый метод regex.test()",
                "desc": "Принимает строку и возвращает булево true, если строка соответствует паттерну, либо false при ошибке.",
                "height": 40
            },
            {
                "title": "Символы привязки ^ и $",
                "desc": "Символ ^ обозначает начало строки, а $ — конец. Они гарантируют проверку всей строки целиком.",
                "height": 40
            },
            {
                "title": "Эффективность в веб-формах",
                "desc": "Позволяет за одну строчку кода проверить корректность адреса почты, номера телефона или сложности пароля.",
                "height": 40
            }
        ]
    },

    # 5. Slide 7: Email & Phone Patterns
    {
        "type": "code",
        "category": "МАСКИ И ШАБЛОНЫ",
        "question": "Какие регулярные выражения используются для стандартных контактных данных?",
        "left_title": "📄 js/main.js (Паттерны контактов)",
        "code_lines": [
            [("// Паттерн адреса электронной почты", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("EMAIL_REGEX = /^[^\s@]+@[^\s@]+\.[^\s@]+$/;", NUM_COLOR)],
            [("", CODE_TEXT)],
            [("// Паттерн номера телефона (от 10 до 15 цифр)", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("PHONE_REGEX = /^\+?[0-9\s\-\(\)]{10,18}$/;", NUM_COLOR)],
            [("", CODE_TEXT)],
            [("function ", KEYWORD_COLOR), ("validateContact", FUNC_COLOR), ("(email, phone) {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("validMail = EMAIL_REGEX.", CODE_TEXT), ("test", FUNC_COLOR), ("(email.trim());", STR_COLOR)],
            [("  const ", KEYWORD_COLOR), ("validPhone = PHONE_REGEX.", CODE_TEXT), ("test", FUNC_COLOR), ("(phone.trim());", STR_COLOR)],
            [("  return ", KEYWORD_COLOR), ("validMail && validPhone;", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ВАЛИДАЦИЯ EMAIL И НОМЕРА ТЕЛЕФОНА",
        "bullet_items": [
            {
                "title": "Структура почтового паттерна",
                "desc": "Проверяет наличие имени пользователя, знака @, доменного имени и точки перед зоной (например, .ru / .com).",
                "height": 40
            },
            {
                "title": "Гибкость телефонного паттерна",
                "desc": "Разрешает международный знак плюс (+), пробелы, скобки и дефисы, проверяя минимальную длину цифр.",
                "height": 40
            },
            {
                "title": "Предотвращение некорректных лидов",
                "desc": "Исключает попадание случайных текстовых опечаток в базу данных клиентов коворкинга.",
                "height": 40
            },
            {
                "title": "Обратная совместимость",
                "desc": "Клиентская регулярка дает быстрый отклик, но обязательно должна дублироваться проверкой на сервере.",
                "height": 40
            }
        ]
    },

    # 6. Slide 8: Password matching algorithm
    {
        "type": "code",
        "category": "СВЕРКА ПАРОЛЕЙ",
        "question": "Как алгоритмически проверить совпадение пароля и его подтверждения?",
        "left_title": "📄 js/main.js (Сравнение паролей)",
        "code_lines": [
            [("// Сверка основного пароля и подтверждения", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("pass = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('password');", STR_COLOR)],
            [("const ", KEYWORD_COLOR), ("confirm = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('confirmPassword');", STR_COLOR)],
            [("", CODE_TEXT)],
            [("if ", KEYWORD_COLOR), ("(pass && confirm) {", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(pass.value !== confirm.value) {", CODE_TEXT)],
            [("    // Пароли не совпадают!", COMMENT_COLOR)],
            [("    confirm.classList.", CODE_TEXT), ("add", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("    isValid = false;", CODE_TEXT)],
            [("  } ", CODE_TEXT), ("else ", KEYWORD_COLOR), ("{", CODE_TEXT)],
            [("    confirm.classList.", CODE_TEXT), ("remove", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("  }", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "АЛГОРИТМ ПОДТВЕРЖДЕНИЯ ПАРОЛЯ",
        "bullet_items": [
            {
                "title": "Защита от опечаток пользователя",
                "desc": "Скрытые символы пароля легко набрать с ошибкой. Поле подтверждения гарантирует осознанный ввод.",
                "height": 40
            },
            {
                "title": "Строгое неравенство (!==)",
                "desc": "Оператор !== проверяет точное посимвольное совпадение строк с учетом регистра букв.",
                "height": 40
            },
            {
                "title": "Таргетированная подсветка ошибки",
                "desc": "Класс ошибки .is-invalid вешается именно на поле confirmPassword, подсказывая, где произошла нестыковка.",
                "height": 40
            },
            {
                "title": "Проверка обоих полей на заполненность",
                "desc": "Сравнение имеет смысл только если оба поля уже заполнены, исключая ложные срабатывания.",
                "height": 40
            }
        ]
    },

    # 7. Slide 9: Password length & security constraints
    {
        "type": "code",
        "category": "КРИТЕРИИ ПАРОЛЯ",
        "question": "Как установить ограничение на длину пароля (не менее 6 символов)?",
        "left_title": "📄 js/main.js (Проверка длины пароля)",
        "code_lines": [
            [("// Проверка минимальной длины пароля", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("passwordInput = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('password');", STR_COLOR)],
            [("const ", KEYWORD_COLOR), ("password = passwordInput.value;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("if ", KEYWORD_COLOR), ("(password.length < 6) {", CODE_TEXT)],
            [("  passwordInput.classList.", CODE_TEXT), ("add", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("  showNotification(", FUNC_COLOR), ("'Пароль должен содержать от 6 символов'", STR_COLOR), (", 'danger');", STR_COLOR)],
            [("  isValid = false;", CODE_TEXT)],
            [("} ", CODE_TEXT), ("else ", KEYWORD_COLOR), ("{", CODE_TEXT)],
            [("  passwordInput.classList.", CODE_TEXT), ("remove", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ТРЕБОВАНИЯ БЕЗОПАСНОСТИ К ПАРОЛЮ",
        "bullet_items": [
            {
                "title": "Свойство string.length",
                "desc": "Свойство .length возвращает количество введенных символов, позволяя легко ограничить минимальную длину.",
                "height": 40
            },
            {
                "title": "Стандарт 6+ символов для обучения",
                "desc": "Для студенческих проектов порог в 6 символов оптимален: он учит валидации, не перегружая студента.",
                "height": 40
            },
            {
                "title": "Точечный Toast-алерт",
                "desc": "Помимо рамки, функция showNotification выводит конкретную причину ошибки в углу экрана.",
                "height": 40
            },
            {
                "title": "Подсказка в placeholder",
                "desc": "В HTML атрибут placeholder=\"Минимум 6 символов\" заранее предупреждает пользователя о правиле.",
                "height": 40
            }
        ]
    },

    # 8. Slide 10: Visual error indicator (.is-invalid)
    {
        "type": "code",
        "category": "ИНДИКАЦИЯ ОШИБОК",
        "question": "Как визуально показать пользователю, в каком именно поле допущена ошибка?",
        "left_title": "📄 js/main.js (Класс .is-invalid)",
        "code_lines": [
            [("// Динамическое переключение статуса поля", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("setFieldError", FUNC_COLOR), ("(input, hasError) {", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(hasError) {", CODE_TEXT)],
            [("    input.classList.", CODE_TEXT), ("add", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("    input.", CODE_TEXT), ("focus", FUNC_COLOR), ("(); ", CODE_TEXT), ("// Переводим фокус на проблемное поле", COMMENT_COLOR)],
            [("  } ", CODE_TEXT), ("else ", KEYWORD_COLOR), ("{", CODE_TEXT)],
            [("    input.classList.", CODE_TEXT), ("remove", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("  }", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ИНДИКАЦИЯ ОШИБОК ЧЕРЕЗ CLASSLIST",
        "bullet_items": [
            {
                "title": "Класс состояния .is-invalid",
                "desc": "Общепринятый в вебе класс модификатора, сообщающий стилям CSS о наличии ошибки ввода.",
                "height": 40
            },
            {
                "title": "Фокус на первом ошибочном поле",
                "desc": "Вызов input.focus() автоматически ставит курсор в поле с ошибкой, экономя время пользователя.",
                "height": 40
            },
            {
                "title": "Своевременная очистка статуса",
                "desc": "Если ошибка исправлена, класс .is-invalid должен быть немедленно удален через .classList.remove().",
                "height": 40
            },
            {
                "title": "Доступность для скринридеров",
                "desc": "В паре с классом можно добавлять атрибут aria-invalid=\"true\" для людей с ограниченными возможностями.",
                "height": 40
            }
        ]
    },

    # 9. Slide 11: CSS styling for invalid inputs
    {
        "type": "code",
        "category": "CSS-ОФОРМЛЕНИЕ",
        "question": "Как стилизовать ошибочное состояние поля в css/style.css?",
        "left_title": "📄 css/style.css (Стили валидации)",
        "code_lines": [
            [("/* Красная граница и фон при ошибке */", COMMENT_COLOR)],
            [(".form-control.is-invalid {", FUNC_COLOR)],
            [("  border-color: ", KEYWORD_COLOR), ("#DC2626 !important;", DANGER_COLOR)],
            [("  background-color: ", KEYWORD_COLOR), ("#FEF2F2;", DANGER_COLOR)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Текст ошибки под полем */", COMMENT_COLOR)],
            [(".invalid-feedback {", FUNC_COLOR)],
            [("  display: ", KEYWORD_COLOR), ("none;", CODE_TEXT)],
            [("  color: ", KEYWORD_COLOR), ("#DC2626;", DANGER_COLOR)],
            [("  font-size: ", KEYWORD_COLOR), ("12px;", NUM_COLOR)],
            [("  margin-top: ", KEYWORD_COLOR), ("4px;", NUM_COLOR)],
            [("}", CODE_TEXT)],
            [(".is-invalid ~ .invalid-feedback {", FUNC_COLOR)],
            [("  display: ", KEYWORD_COLOR), ("block;", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ДИЗАЙН ОШИБОК: CSS-СЕЛЕКТОРЫ СОСТОЯНИЙ",
        "bullet_items": [
            {
                "title": "Контрастный красный акцент #DC2626",
                "desc": "Стандартный сигнальный цвет опасности сразу выделяет ошибочные поля среди остального контента.",
                "height": 40
            },
            {
                "title": "Мягкий фоновый оттенок #FEF2F2",
                "desc": "Тонирование внутренней области поля делает индикацию заметной даже боковым зрением.",
                "height": 40
            },
            {
                "title": "Соседский селектор (~)",
                "desc": "Селектор .is-invalid ~ .invalid-feedback автоматически отображает подсказку только при наличии ошибки.",
                "height": 40
            },
            {
                "title": "Плавные переходы transition",
                "desc": "Свойство transition: border-color 0.2s делает смену состояний аккуратной и приятной для глаз.",
                "height": 40
            }
        ]
    },

    # 10. Slide 12: Diagram: Input States Comparison
    {
        "type": "input_states",
        "category": "СХЕМА: СОСТОЯНИЯ",
        "question": "Как трансформируется визуальный статус поля ввода в процессе валидации?"
    },

    # 11. Slide 13: Loop validation across fields
    {
        "type": "code",
        "category": "ПЕРЕБОР ПОЛЕЙ",
        "question": "Как проверить 6 полей формы регистрации без дублирования кода?",
        "left_title": "📄 js/main.js (Массовая проверка)",
        "code_lines": [
            [("// Массив идентификаторов всех обязательных полей", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("fields = [", CODE_TEXT)],
            [("  'login', 'password', 'confirmPassword',", STR_COLOR)],
            [("  'fullName', 'email', 'phone'", STR_COLOR)],
            [("];", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("fields.", CODE_TEXT), ("forEach", FUNC_COLOR), ("(id => {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("input = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("(id);", STR_COLOR)],
            [("  if ", KEYWORD_COLOR), ("(!input) return;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(!input.value.", CODE_TEXT), ("trim", FUNC_COLOR), ("()) {", CODE_TEXT)],
            [("    input.classList.", CODE_TEXT), ("add", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("    isValid = false;", CODE_TEXT)],
            [("  } ", CODE_TEXT), ("else ", KEYWORD_COLOR), ("{", CODE_TEXT)],
            [("    input.classList.", CODE_TEXT), ("remove", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("  }", CODE_TEXT)],
            [("});", CODE_TEXT)]
        ],
        "right_title": "АВТОМАТИЗАЦИЯ ПРОВЕРКИ ЧЕРЕЗ FOREACH",
        "bullet_items": [
            {
                "title": "Принцип DRY (Don't Repeat Yourself)",
                "desc": "Вместо 6 копий однотипного кода мы описываем массив ID и обходим его в компактном цикле .forEach().",
                "height": 40
            },
            {
                "title": "Масштабируемость формы",
                "desc": "Чтобы добавить новое поле в форму (например, промокод), достаточно просто дописать его ID в массив.",
                "height": 40
            },
            {
                "title": "Защитная проверка if (!input)",
                "desc": "Предотвращает фатальные ошибки консоли, если разметка какого-то поля временно отсутствует на странице.",
                "height": 40
            },
            {
                "title": "Синхронная подсветка всех ошибок",
                "desc": "Пользователь за один клик видит сразу все пропущенные поля, а не исправляет их по одному.",
                "height": 40
            }
        ]
    },

    # 12. Slide 14: Accumulator Validity Flag (isValid)
    {
        "type": "code",
        "category": "ФЛАГ ВАЛИДНОСТИ",
        "question": "Как определить, что хотя бы одно поле заполнено неверно?",
        "left_title": "📄 js/main.js (Флаг isValid)",
        "code_lines": [
            [("// Паттерн флага-аккумулятора", COMMENT_COLOR)],
            [("let ", KEYWORD_COLOR), ("isValid = true; ", CODE_TEXT), ("// Предполагаем, что всё верно", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("// Если хотя бы одна проверка упала:", COMMENT_COLOR)],
            [("if ", KEYWORD_COLOR), ("(!loginVal) {", CODE_TEXT)],
            [("  isValid = false;", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("if ", KEYWORD_COLOR), ("(pass !== confirm) {", CODE_TEXT)],
            [("  isValid = false;", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Финальное решение об отправке:", COMMENT_COLOR)],
            [("if ", KEYWORD_COLOR), ("(isValid) {", CODE_TEXT)],
            [("  // Успех! Переходим к отправке или редиректу", COMMENT_COLOR)],
            [("  submitRegistration();", FUNC_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ПАТТЕРН ФЛАГА: НАКОПЛЕНИЕ СТАТУСА",
        "bullet_items": [
            {
                "title": "Оптимистичный старт (isValid = true)",
                "desc": "Изначально переменная устанавливается в true. Любая найденная ошибка переключает ее в false.",
                "height": 40
            },
            {
                "title": "Необратимость сброса",
                "desc": "Если флаг переключился в false, последующие корректные поля не должны возвращать его в true.",
                "height": 40
            },
            {
                "title": "Единая точка принятия решения",
                "desc": "Проверка if (isValid) в самом конце функции разделяет поток выполнения на успех и блокировку.",
                "height": 40
            },
            {
                "title": "Чистота архитектуры",
                "desc": "Избавляет от запутанных вложенных конструкций if-else и упрощает отладку валидации.",
                "height": 40
            }
        ]
    },

    # 13. Slide 15: Form Reset Method (form.reset())
    {
        "type": "code",
        "category": "СБРОС ФОРМЫ",
        "question": "Как очистить все поля формы после успешной регистрации пользователя?",
        "left_title": "📄 js/main.js (Метод form.reset)",
        "code_lines": [
            [("if ", KEYWORD_COLOR), ("(isValid) {", CODE_TEXT)],
            [("  // 1. Показываем сообщение об успехе", COMMENT_COLOR)],
            [("  showNotification('Пользователь зарегистрирован!', 'success');", STR_COLOR)],
            [("", CODE_TEXT)],
            [("  // 2. Очищаем форму встроенным методом браузера", COMMENT_COLOR)],
            [("  form.", CODE_TEXT), ("reset", FUNC_COLOR), ("();", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // 3. Снимаем возможные оставшиеся классы", COMMENT_COLOR)],
            [("  fields.forEach(id => {", CODE_TEXT)],
            [("    document.getElementById(id)?.classList.remove('is-invalid');", CODE_TEXT)],
            [("  });", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // 4. Редирект на страницу входа через 1.2 сек", COMMENT_COLOR)],
            [("  setTimeout(() => location.href = 'login.html', 1200);", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ОЧИСТКА СОСТОЯНИЯ: МЕТОД FORM.RESET()",
        "bullet_items": [
            {
                "title": "Нативный метод form.reset()",
                "desc": "Сбрасывает значения всех полей формы к исходным значениям по умолчанию за одну команду.",
                "height": 40
            },
            {
                "title": "Очистка классов ошибок",
                "desc": "Метод reset() не удаляет CSS-классы, поэтому вызов classList.remove('is-invalid') обязателен.",
                "height": 40
            },
            {
                "title": "Предотвращение повторной отправки",
                "desc": "Очистка формы защищает от случайного дублирования данных при случайном двойном клике.",
                "height": 40
            },
            {
                "title": "Плавный переход на шаг входа",
                "desc": "Задержка 1200 мс дает пользователю прочитать уведомление об успехе перед переходом на login.html.",
                "height": 40
            }
        ]
    },

    # 14. Slide 16: Login form verification
    {
        "type": "code",
        "category": "АВТОРИЗАЦИЯ",
        "question": "Как реализовать клиентскую проверку учетных данных (admin / 12345)?",
        "left_title": "📄 js/main.js (Функция initLoginForm)",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("initLoginForm", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("form = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('loginForm');", STR_COLOR)],
            [("  if ", KEYWORD_COLOR), ("(!form) return;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  form.", CODE_TEXT), ("addEventListener", FUNC_COLOR), ("('submit', (e) => {", STR_COLOR)],
            [("    e.", CODE_TEXT), ("preventDefault", FUNC_COLOR), ("();", CODE_TEXT)],
            [("    const ", KEYWORD_COLOR), ("login = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('login').value.trim();", STR_COLOR)],
            [("    const ", KEYWORD_COLOR), ("pass = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('password').value.trim();", STR_COLOR)],
            [("", CODE_TEXT)],
            [("    if ", KEYWORD_COLOR), ("(login === 'admin' && pass === '12345') {", STR_COLOR)],
            [("      // Успешная авторизация тестового аккаунта", COMMENT_COLOR)],
            [("      localStorage.setItem('currentUser', login);", CODE_TEXT)],
            [("      window.location.href = '../index.html';", CODE_TEXT)],
            [("    }", CODE_TEXT)],
            [("  });", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ПРОВЕРКА ВХОДА: ТЕСТОВЫЙ ПОЛЬЗОВАТЕЛЬ",
        "bullet_items": [
            {
                "title": "Эталонная учетная запись admin / 12345",
                "desc": "Учебный mock-пользователь позволяет отладить весь цикл входа и выхода без подключения сложного сервера.",
                "height": 40
            },
            {
                "title": "Обязательный .trim() логина и пароля",
                "desc": "Предотвращает случайные ошибки авторизации из-за пробелов, скопированных из буфера обмена.",
                "height": 40
            },
            {
                "title": "Сохранение сессии в браузере",
                "desc": "Команда localStorage.setItem запоминает пользователя, активируя кнопку «Выйти» и личный кабинет.",
                "height": 40
            },
            {
                "title": "Редирект на главную страницу",
                "desc": "После успешного входа пользователь мгновенно перенаправляется на главную страницу или в бронирования.",
                "height": 40
            }
        ]
    },

    # 15. Slide 17: Login Alert Box
    {
        "type": "code",
        "category": "СИСТЕМНЫЙ АЛЕРТ",
        "question": "Как вывести общее предупреждение «Неверный логин или пароль»?",
        "left_title": "📄 js/main.js (Блок loginAlert)",
        "code_lines": [
            [("// Вывод сообщения об ошибке авторизации", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("alertBox = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('loginAlert');", STR_COLOR)],
            [("", CODE_TEXT)],
            [("if ", KEYWORD_COLOR), ("(login === 'admin' && pass === '12345') {", STR_COLOR)],
            [("  if ", KEYWORD_COLOR), ("(alertBox) alertBox.style.display = 'none';", CODE_TEXT)],
            [("  showNotification('Успешный вход!', 'success');", STR_COLOR)],
            [("} ", CODE_TEXT), ("else ", KEYWORD_COLOR), ("{", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(alertBox) {", CODE_TEXT)],
            [("    alertBox.", CODE_TEXT), ("textContent", FUNC_COLOR), (" = 'Неверный логин или пароль';", STR_COLOR)],
            [("    alertBox.", CODE_TEXT), ("className", FUNC_COLOR), (" = 'form-alert alert-danger';", STR_COLOR)],
            [("    alertBox.style.", CODE_TEXT), ("display", FUNC_COLOR), (" = 'block';", STR_COLOR)],
            [("  }", CODE_TEXT)],
            [("  showNotification('Неверный логин или пароль', 'danger');", STR_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ОБЩЕЕ ОПОВЕЩЕНИЕ: БЛОК LOGINALERT",
        "bullet_items": [
            {
                "title": "Вывод над полями формы",
                "desc": "Блок #loginAlert расположен прямо над полями ввода, сразу объясняя пользователю причину отказа.",
                "height": 40
            },
            {
                "title": "Управление свойством style.display",
                "desc": "Изначально скрытый через display: none блок проявляется только при неверной комбинации логина и пароля.",
                "height": 40
            },
            {
                "title": "Безопасное сообщение об ошибке",
                "desc": "Формулировка «Неверный логин или пароль» не подсказывает злоумышленнику, что именно было неверно.",
                "height": 40
            },
            {
                "title": "Дублирование через системный Toast",
                "desc": "Совместная работа локального алерта и всплывающего Toast гарантирует, что пользователь заметит статус.",
                "height": 40
            }
        ]
    },

    # 16. Slide 18: Session Persistence in LocalStorage
    {
        "type": "code",
        "category": "СЕССИЯ ПОЛЬЗОВАТЕЛЯ",
        "question": "Как зафиксировать факт успешного входа для остальных страниц сайта?",
        "left_title": "📄 js/main.js (Запись в localStorage)",
        "code_lines": [
            [("// Запись имени авторизованного пользователя", COMMENT_COLOR)],
            [("localStorage.", CODE_TEXT), ("setItem", FUNC_COLOR), ("('currentUser', login);", STR_COLOR)],
            [("", CODE_TEXT)],
            [("// Проверка авторизации на любой другой странице:", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("checkAuth", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("user = localStorage.", CODE_TEXT), ("getItem", FUNC_COLOR), ("('currentUser');", STR_COLOR)],
            [("  if ", KEYWORD_COLOR), ("(user) {", CODE_TEXT)],
            [("    console.", FUNC_COLOR), ("log", FUNC_COLOR), ("('В системе пользователь:', user);", STR_COLOR)],
            [("  } ", CODE_TEXT), ("else ", KEYWORD_COLOR), ("{", CODE_TEXT)],
            [("    console.", FUNC_COLOR), ("log", FUNC_COLOR), ("('Гостевой режим (не авторизован)');", STR_COLOR)],
            [("  }", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ХРАНЕНИЕ СОСТОЯНИЯ: LOCALSTORAGE",
        "bullet_items": [
            {
                "title": "Связь между страницами",
                "desc": "Хранилище localStorage сохраняет ключ currentUser даже после закрытия вкладки и перехода по ссылкам.",
                "height": 40
            },
            {
                "title": "Реакция навигационной шапки",
                "desc": "Функция updateAuthNav() считывает этот ключ, показывая вкладку «Мои бронирования» и кнопку «Выйти».",
                "height": 40
            },
            {
                "title": "Простой строковый формат",
                "desc": "Для авторизации достаточно сохранить имя пользователя; пароли в открытом виде в localStorage никогда не хранят.",
                "height": 40
            },
            {
                "title": "Очистка при выходе (Logout)",
                "desc": "По нажатию на кнопку «Выйти» команда localStorage.removeItem('currentUser') сбрасывает сессию.",
                "height": 40
            }
        ]
    },

    # 17. Slide 19: Live feedback on input
    {
        "type": "code",
        "category": "ЖИВАЯ РЕАКЦИЯ",
        "question": "Как убрать красную рамку сразу же, когда пользователь начал печатать?",
        "left_title": "📄 js/main.js (Событие 'input')",
        "code_lines": [
            [("// Снятие ошибки сразу при начале исправления", COMMENT_COLOR)],
            [("fields.", CODE_TEXT), ("forEach", FUNC_COLOR), ("(id => {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("input = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("(id);", STR_COLOR)],
            [("  if ", KEYWORD_COLOR), ("(!input) return;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // Слушаем каждое нажатие клавиши в поле", COMMENT_COLOR)],
            [("  input.", CODE_TEXT), ("addEventListener", FUNC_COLOR), ("('input', () => {", STR_COLOR)],
            [("    if ", KEYWORD_COLOR), ("(input.classList.", CODE_TEXT), ("contains", FUNC_COLOR), ("('is-invalid')) {", STR_COLOR)],
            [("      input.classList.", CODE_TEXT), ("remove", FUNC_COLOR), ("('is-invalid');", STR_COLOR)],
            [("    }", CODE_TEXT)],
            [("  });", CODE_TEXT)],
            [("});", CODE_TEXT)]
        ],
        "right_title": "МГНОВЕННЫЙ ОТКЛИК: СОБЫТИЕ INPUT",
        "bullet_items": [
            {
                "title": "Разница между change и input",
                "desc": "Событие change ждет потери фокуса, а событие input срабатывает мгновенно при каждом нажатии клавиши.",
                "height": 40
            },
            {
                "title": "Снятие психологического барьера",
                "desc": "Пользователь видит, что интерфейс реагирует на исправление, и красная рамка не висит понапрасну.",
                "height": 40
            },
            {
                "title": "Проверка методом .contains()",
                "desc": "Снятие класса выполняется только если класс действительно был добавлен, экономя вызовы браузера.",
                "height": 40
            },
            {
                "title": "Современный стандарт UX",
                "desc": "Подход «ошибка при отправке, очистка при вводе» признан лучшей практикой в мировом веб-дизайне.",
                "height": 40
            }
        ]
    },

    # 18. Slide 20: Success notifications & redirection
    {
        "type": "code",
        "category": "УВЕДОМЛЕНИЯ",
        "question": "Как плавно показать сообщение об успехе перед переходом на другую страницу?",
        "left_title": "📄 js/main.js (Toast и задержка)",
        "code_lines": [
            [("// Комбинация уведомления и отложенного перехода", COMMENT_COLOR)],
            [("if ", KEYWORD_COLOR), ("(isValid) {", CODE_TEXT)],
            [("  // 1. Показываем зеленое всплывающее окно", COMMENT_COLOR)],
            [("  showNotification('Регистрация прошла успешно!', 'success');", STR_COLOR)],
            [("", CODE_TEXT)],
            [("  // 2. Блокируем кнопку, чтобы не было повторных кликов", COMMENT_COLOR)],
            [("  const ", KEYWORD_COLOR), ("btn = form.querySelector('button[type=\"submit\"]');", CODE_TEXT)],
            [("  if (btn) btn.disabled = true;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // 3. Через 1 секунду направляем на форму входа", COMMENT_COLOR)],
            [("  setTimeout(() => {", CODE_TEXT)],
            [("    window.location.href = 'login.html';", CODE_TEXT)],
            [("  }, 1000);", NUM_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ПОДТВЕРЖДЕНИЕ УСПЕХА И ПЕРЕНАПРАВЛЕНИЕ",
        "bullet_items": [
            {
                "title": "Визуальное подтверждение действия",
                "desc": "Зеленый Toast информирует студента об успешном создании профиля, создавая чувство завершенности.",
                "height": 40
            },
            {
                "title": "Блокировка кнопки (btn.disabled)",
                "desc": "Защищает сервер и локальное хранилище от повторных отправок при нервных множественных кликах.",
                "height": 40
            },
            {
                "title": "Задержка через setTimeout()",
                "desc": "Пауза в 1000–1200 миллисекунд позволяет прочитать текст сообщения до смены страницы.",
                "height": 40
            },
            {
                "title": "Естественный переход по воронке",
                "desc": "Пользователь логично попадает на страницу входа, где может сразу протестировать новый логин.",
                "height": 40
            }
        ]
    },

    # 19. Slide 21: HTML attributes & accessibility
    {
        "type": "code",
        "category": "ДОСТУПНОСТЬ И HTML",
        "question": "Какие атрибуты тегов <input> и <label> делают форму профессиональной?",
        "left_title": "📄 pages/login.html (Атрибуты доступности)",
        "code_lines": [
            [("<!-- Правильная семантическая связка -->", COMMENT_COLOR)],
            [("<div class=\"form-group\">", TAG_COLOR)],
            [("  <label for=\"userLogin\">Логин</label>", TAG_COLOR)],
            [("  <input", TAG_COLOR)],
            [("    type=\"text\"", TAG_COLOR)],
            [("    id=\"userLogin\"", TAG_COLOR)],
            [("    name=\"login\"", TAG_COLOR)],
            [("    autocomplete=\"username\"", TAG_COLOR)],
            [("    placeholder=\"admin\"", TAG_COLOR)],
            [("    required", KEYWORD_COLOR)],
            [("    class=\"form-control\"", TAG_COLOR)],
            [("  >", TAG_COLOR)],
            [("</div>", TAG_COLOR)]
        ],
        "right_title": "ПРОФЕССИОНАЛЬНАЯ ВЕРСТКА ПОЛЕЙ ВВОДА",
        "bullet_items": [
            {
                "title": "Связка label for и input id",
                "desc": "Клик по текстовой подписи <label> автоматически фокусирует курсор в соответствующем поле ввода.",
                "height": 40
            },
            {
                "title": "Атрибут autocomplete",
                "desc": "Значения 'username' и 'current-password' позволяют браузеру и менеджерам паролей безопасно подставлять данные.",
                "height": 40
            },
            {
                "title": "Атрибуты type=\"email\" и type=\"password\"",
                "desc": "На смартфонах активируют специализированные клавиатуры (с символом @ или скрытыми точками ввода).",
                "height": 40
            },
            {
                "title": "Атрибут required как базовый барьер",
                "desc": "Встроенная браузерная валидация служит дополнительной подстраховкой для клиентского JavaScript.",
                "height": 40
            }
        ]
    },

    # 20. Slide 22: Student Mistakes (Grid Cards 2x2)
    {
        "type": "grid_cards",
        "category": "РАЗБОР ОШИБОК",
        "question": "На чем чаще всего спотыкаются студенты при написании валидации?",
        "title_text": "ТОП-4 ТИПИЧНЫХ ОШИБОК СТУДЕНТОВ ПРИ ВАЛИДАЦИИ ФОРМ",
        "cards_data": [
            {
                "title": "Забытый вызов event.preventDefault()",
                "problem": "При клике на кнопку форма моментально перезагружает страницу, стирая всё, что студент написал в полях.",
                "solution": "Всегда принимать объект события (e) в слушателе submit и первой строчкой вызывать e.preventDefault()."
            },
            {
                "title": "Проверка значения без метода .trim()",
                "problem": "Если в поле ввести одни пробелы, обычная проверка if (input.value) посчитает поле корректно заполненным.",
                "solution": "Всегда проверять очищенную строку: if (!input.value.trim()) с добавлением класса ошибки .is-invalid."
            },
            {
                "title": "Слушатель 'click' на кнопке вместо 'submit' на форме",
                "problem": "Если пользователь отправляет форму клавишей Enter, валидация не запускается и страница ломается.",
                "solution": "Всегда вешать слушатель события 'submit' на сам тег <form>, перехватывая все способы отправки."
            },
            {
                "title": "Отсутствие снятия класса .is-invalid при вводе",
                "problem": "Студент подсветил ошибку красным, но забыл снять класс, и поле остается красным даже после исправления.",
                "solution": "Добавлять слушатель 'input' на каждое поле, вызывая input.classList.remove('is-invalid') при печати."
            }
        ]
    },

    # 21. Slide 23: Checklist
    {
        "type": "checklist",
        "category": "ЧЕК-ЛИСТ КАЧЕСТВА",
        "question": "Как убедиться, что валидация форм авторизации работает на 100% надежно?",
        "title_text": "КРИТЕРИИ КАЧЕСТВА ВАЛИДАЦИИ ИНТЕРФЕЙСОВ АВТОРИЗАЦИИ",
        "items": [
            {
                "title": "Отправка формы перехвачена через e.preventDefault()",
                "desc": "При клике на кнопку или нажатии Enter страница не перезагружается, сохраняя введенный текст."
            },
            {
                "title": "Все пустые обязательные поля подсвечиваются красным",
                "desc": "Метод .trim() выявляет пробелы, добавляя класс .is-invalid и активируя красную границу."
            },
            {
                "title": "Несовпадающие пароли вызывают ошибку подтверждения",
                "desc": "Поле confirmPassword помечается ошибкой, если пароли расходятся хотя бы на один символ."
            },
            {
                "title": "Красная подсветка исчезает сразу в процессе ввода",
                "desc": "Слушатель события input автоматически снимает класс .is-invalid, улучшая восприятие формы."
            },
            {
                "title": "Вход с данными admin / 12345 перенаправляет в систему",
                "desc": "Успешная авторизация сохраняет currentUser в localStorage и переводит пользователя на главную."
            },
            {
                "title": "Неверные данные выводят понятный алерт над формой",
                "desc": "Блок loginAlert корректно отображает текст 'Неверный логин или пароль' без раскрытия деталей."
            }
        ]
    }
]

print(f"Total content slides to generate: {len(slides_data)}")

# -------------------------------------------------------------
# PRESENTATION BUILDER EXECUTION
# -------------------------------------------------------------
prs = Presentation(base_template_path)
blank_layout = prs.slide_layouts[6]

orig_slide_ids = list(prs.slides._sldIdLst)
title_slide_id = orig_slide_ids[0]
plan_slide_id = orig_slide_ids[1]
result_slide_id = orig_slide_ids[25]
goodbye_slide_id = orig_slide_ids[26]

# 0. Update Layouts footer
for layout in prs.slide_layouts:
    for s in layout.shapes:
        if "Google Shape;59;p13" in s.name and s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip() == "2":
                        r.text = "6"

# 1. Update Slide 1 (Title slide)
slide1 = prs.slides[0]
for sh in slide1.shapes:
    if sh.name == "Google Shape;59;p13" and sh.has_text_frame:
        sh.text_frame.text = "Вебинар 6 "
    elif sh.shape_type == MSO_SHAPE_TYPE.GROUP and sh.name == "Группа 12":
        for sub in sh.shapes:
            if sub.name == "TextBox 7":
                sub.text_frame.text = "КЛИЕНТСКАЯ ВАЛИДАЦИЯ ФОРМ"
                sub.text_frame.paragraphs[0].runs[0].font.name = "Montserrat"
                sub.text_frame.paragraphs[0].runs[0].font.size = Pt(15)
                sub.text_frame.paragraphs[0].runs[0].font.bold = True
            elif sub.name == "TextBox 9":
                sub.text_frame.text = "Проектирование \nи разработка интерфейсов пользователя"
print("Slide 1 updated successfully.")

# 2. Update Slide 2 (Plan slide)
slide2 = prs.slides[1]
plan_titles = {
    "TextBox 12": "1. Контроль отправки",
    "TextBox 15": "2. Алгоритмы валидации",
    "TextBox 18": "3. Регулярные выражения",
    "TextBox 21": "4. Практика: Формы авторизации",
    "TextBox 24": "5. Результат и итоги",
}
plan_subtitles = {
    "TextBox 13": "Перехват события submit и отмена перезагрузки e.preventDefault()",
    "TextBox 16": "Очистка метода .trim(), обязательные поля и совпадение паролей",
    "TextBox 19": "RegExp: проверка формата Email, телефона и метод .test()",
    "TextBox 22": "Настройка форм регистрации и входа в систему СмартОфис",
    "TextBox 25": "Визуальная индикация ошибок .is-invalid, уведомления и чек-лист",
}

for sh in slide2.shapes:
    if sh.name in plan_titles:
        sh.text_frame.clear()
        p = sh.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = plan_titles[sh.name]
        r.font.name = "Montserrat SemiBold"
        r.font.size = Pt(11.25)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x18, 0x1A, 0x1F)
    elif sh.name in plan_subtitles:
        sh.width = Pt(450)
        sh.text_frame.clear()
        sh.text_frame.word_wrap = False
        p = sh.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = plan_subtitles[sh.name]
        r.font.name = "Inter"
        r.font.size = Pt(9.75)
        r.font.bold = False
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
print("Slide 2 updated successfully.")

# 3. Create all 21 new content slides
for idx, sdata in enumerate(slides_data):
    s = prs.slides.add_slide(blank_layout)
    stype = sdata["type"]
    
    if stype == "code":
        create_code_explanation_slide(
            s, sdata["category"], sdata["question"],
            sdata["left_title"], sdata["code_lines"],
            sdata["right_title"], sdata["bullet_items"]
        )
    elif stype == "pipeline":
        create_pipeline_diagram_slide(s, sdata["category"], sdata["question"])
    elif stype == "input_states":
        create_input_states_diagram_slide(s, sdata["category"], sdata["question"])
    elif stype == "grid_cards":
        create_grid_cards_slide(s, sdata["category"], sdata["question"], sdata["title_text"], sdata["cards_data"])
    elif stype == "checklist":
        create_checklist_slide(s, sdata["category"], sdata["question"], sdata["title_text"], sdata["items"])
    
    print(f"  Created Slide {idx+3}: {sdata['category']}")

# 4. Update Result Slide (Slide 24)
slide_res = prs.slides[25]
pic_to_remove = None
tb2_res = None
for sh in slide_res.shapes:
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        pic_to_remove = sh
    elif sh.name == "TextBox 2":
        tb2_res = sh

if pic_to_remove is not None:
    spTree = slide_res.shapes._spTree
    spTree.remove(pic_to_remove._element)

# Add Webinar 6 result screenshot: aspect ratio 1210:670 (1.806), centered nicely
res_w = Pt(510)
res_h = Pt(282)
res_left = Pt(105)
res_top = Pt(68)
slide_res.shapes.add_picture(result_img_path, res_left, res_top, res_w, res_h)

# Add or update TextBox 2 on Result slide
if tb2_res is None:
    tb2_res = slide_res.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2_res.name = "TextBox 2"

tf_res = tb2_res.text_frame
tf_res.word_wrap = True
tf_res.margin_left = tf_res.margin_top = tf_res.margin_right = tf_res.margin_bottom = 0
p_res = tf_res.paragraphs[0]
p_res.alignment = PP_ALIGN.LEFT
p_res.text = ""
r_res = p_res.add_run()
r_res.text = "Результат: Полностью валидируемые интерфейсы авторизации"
r_res.font.name = "Montserrat"
r_res.font.size = Pt(10)
r_res.font.color.rgb = QUESTION_TEXT
print("Result slide updated with webinar 6 screenshot and TextBox 2.")

# Clean shadows on Goodbye slide (Slide 25)
slide_goodbye = prs.slides[26]
for sh in slide_goodbye.shapes:
    strip_shape_styles_and_shadows(sh)
print("Goodbye slide shadows stripped.")

# 5. Reorder slide list in presentation XML:
# [Title, Plan] + new_slide_ids + [Result, Goodbye]
all_slide_ids = list(prs.slides._sldIdLst)
new_slide_ids = all_slide_ids[len(orig_slide_ids):]
new_order = [title_slide_id, plan_slide_id] + new_slide_ids + [result_slide_id, goodbye_slide_id]
prs.slides._sldIdLst[:] = new_order

# 6. Save presentation
os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
prs.save(output_pptx_path)

print(f"\nFinal presentation slide count: {len(prs.slides)}")
print(f"\nУСПЕХ! Презентация Вебинара 6 сохранена: {output_pptx_path}")
