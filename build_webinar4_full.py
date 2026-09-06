import sys
import os
import shutil
import pptx
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
import pygments
from pygments.lexers import JavascriptLexer, HtmlLexer, CssLexer
from pygments.token import Token

sys.stdout.reconfigure(encoding='utf-8')

base_pptx = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx"
output_dir = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-04-js-dom-navigation"
output_pptx = os.path.join(output_dir, "вебинар 4.pptx")
result_img_path = r"d:\IEEU International East-European University\fullstack-web-development-course\webinar4_result_exact.png"

prs = Presentation(base_pptx)
blank_layout = prs.slide_layouts[12]

# Brand Design Tokens & Colors
ORANGE_PILL   = RGBColor(0xFE, 0x60, 0x02)  # #FE6002
ORANGE_LINE   = RGBColor(0xFF, 0x6F, 0x03)  # #FF6F03
BORDER_GRAY   = RGBColor(0xEE, 0xEE, 0xEE)  # #EEEEEE
QUESTION_TEXT = RGBColor(0xA6, 0xA1, 0xA1)  # #A6A1A1 - TextBox 2
DARK_TEXT     = RGBColor(0x22, 0x22, 0x22)  # #222222
BODY_TEXT     = RGBColor(0x44, 0x44, 0x44)  # #444444
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)  # #FFFFFF
GREEN_BADGE   = RGBColor(0x10, 0xB9, 0x81)  # #10B981
GREEN_TEXT    = RGBColor(0x16, 0xA3, 0x4A)  # #16A34A

# Code Syntax Colors
CODE_BG       = RGBColor(0x18, 0x1A, 0x1F)  # #181A1F
CODE_BORDER   = RGBColor(0x2D, 0x31, 0x39)  # #2D3139
CODE_TEXT     = RGBColor(0xD4, 0xD4, 0xD8)  # #D4D4D8
CODE_KW       = RGBColor(0x38, 0xBD, 0xF8)  # Cyan (const, function, let)
CODE_TAG      = RGBColor(0xF4, 0x72, 0xB6)  # Pink (HTML tags, keywords)
CODE_STR      = RGBColor(0xFC, 0xD3, 0x4D)  # Amber strings
CODE_COMM     = RGBColor(0x71, 0x71, 0x7A)  # Gray comments
CODE_PROP     = RGBColor(0xA7, 0x8B, 0xFA)  # Violet (methods, properties)
CODE_VAL      = RGBColor(0x4A, 0xDE, 0x80)  # Green (values, variables)

def strip_shape_styles_and_shadows(shape):
    """Ensure zero default PPT theme shadows or effects on created shapes."""
    elem = shape._element
    style = elem.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
    if style is not None:
        elem.remove(style)
    for effect in elem.xpath('.//a:outerShdw'):
        effect.getparent().remove(effect)

def add_header_and_footer(slide, category_text, question_text):
    """Creates top-right orange category pill badge and bottom-left question."""
    pill_width = Pt(max(150, min(330, len(category_text) * 9.8 + 36)))
    pill_left = Pt(700) - pill_width
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_left, Pt(20.4), pill_width, Pt(22.6))
    pill.name = "Google Shape;69;p14"
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = ORANGE_PILL
    pill.line.color.rgb = BORDER_GRAY
    pill.line.width = Pt(0.75)
    strip_shape_styles_and_shadows(pill)

    tf_p = pill.text_frame
    tf_p.word_wrap = False
    tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = 0
    p_p = tf_p.paragraphs[0]
    p_p.alignment = PP_ALIGN.CENTER
    r_p = p_p.add_run()
    r_p.text = category_text.upper()
    r_p.font.name = "Inter"
    r_p.font.size = Pt(12)
    r_p.font.bold = False
    r_p.font.color.rgb = WHITE

    tb_bot = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(500), Pt(20))
    tb_bot.name = "TextBox 2"
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    tf_bot.margin_left = tf_bot.margin_top = tf_bot.margin_right = tf_bot.margin_bottom = 0
    p_bot = tf_bot.paragraphs[0]
    p_bot.alignment = PP_ALIGN.LEFT
    r_bot = p_bot.add_run()
    r_bot.text = question_text
    r_bot.font.name = "Montserrat"
    r_bot.font.size = Pt(10)
    r_bot.font.color.rgb = QUESTION_TEXT

KNOWN_PROPS = {
    'document', 'window', 'localStorage', 'console', 'classList', 'querySelectorAll',
    'querySelector', 'getElementById', 'addEventListener', 'getAttribute', 'setAttribute',
    'createElement', 'appendChild', 'remove', 'setTimeout', 'includes', 'endsWith',
    'preventDefault', 'pathname', 'location', 'add', 'item', 'setItem', 'getItem',
    'removeItem', 'target', 'event', 'href', 'dataset', 'length', 'push', 'join', 'split',
    'textContent', 'innerText', 'innerHTML', 'style', 'className', 'value'
}

KNOWN_VALS = {'true', 'false', 'null', 'undefined', 'NaN', 'Infinity'}

KNOWN_KWS = {'const', 'let', 'var', 'function', 'return', 'if', 'else', 'for', 'while', 'new', 'import', 'export', 'of', 'in'}

def highlight_code_line(line_str, left_title=""):
    """Tokenize a code line using Pygments into branded RGB runs."""
    if not line_str.strip():
        return [("", CODE_TEXT)]
    
    lt_lower = left_title.lower()
    if ".html" in lt_lower or line_str.strip().startswith("<"):
        lexer = HtmlLexer()
    elif ".css" in lt_lower:
        lexer = CssLexer()
    else:
        lexer = JavascriptLexer()
        
    tokens = list(pygments.lex(line_str, lexer))
    runs = []
    
    for ttype, val in tokens:
        if not val or val == '\n':
            continue
            
        color = CODE_TEXT
        if ttype in Token.Keyword or ttype in Token.Keyword.Constant or ttype in Token.Keyword.Declaration or ttype in Token.Keyword.Reserved or val in KNOWN_KWS:
            color = CODE_KW
        elif ttype in Token.Name.Tag or ttype in Token.Tag:
            color = CODE_TAG
        elif ttype in Token.String or ttype in Token.String.Single or ttype in Token.String.Double or ttype in Token.String.Backtick:
            color = CODE_STR
        elif ttype in Token.Comment or ttype in Token.Comment.Single or ttype in Token.Comment.Multiline:
            color = CODE_COMM
        elif ttype in Token.Number or ttype in Token.Number.Integer or ttype in Token.Number.Float or val in KNOWN_VALS:
            color = CODE_VAL
        elif ttype in Token.Name.Builtin or ttype in Token.Name.Function or ttype in Token.Name.Class or val in KNOWN_PROPS:
            color = CODE_PROP
        elif ttype in Token.Name.Attribute:
            color = CODE_PROP
        elif ttype in Token.Operator.Word:
            color = CODE_KW
        else:
            color = CODE_TEXT
            
        runs.append((val, color))
    return runs

def process_code_line(line, left_title=""):
    """Auto-detect and tokenize unhighlighted lines with Pygments."""
    if isinstance(line, str):
        return highlight_code_line(line, left_title)
    elif isinstance(line, list):
        all_default = all(chunk_col == CODE_TEXT for _, chunk_col in line)
        if all_default or (len(line) == 1 and line[0][1] == CODE_TEXT):
            combined = "".join(chunk_txt for chunk_txt, _ in line)
            return highlight_code_line(combined, left_title)
        return line
    return [("", CODE_TEXT)]

def create_code_explanation_slide(slide, category, question, left_title, code_lines, right_title, bullet_items):
    """Standard 2-column layout: Left code card, Right clean explanation with bullet items."""
    add_header_and_footer(slide, category, question)
    
    # Left Card
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(24), Pt(56), Pt(325), Pt(298))
    left_card.name = "Rounded Rectangle 3"
    left_card.adjustments[0] = 0.04
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CODE_BG
    left_card.line.color.rgb = CODE_BORDER
    left_card.line.width = Pt(1)
    strip_shape_styles_and_shadows(left_card)

    tb_lh = slide.shapes.add_textbox(Pt(38), Pt(64), Pt(295), Pt(18))
    tb_lh.name = "TextBox 4"
    tf_lh = tb_lh.text_frame
    tf_lh.word_wrap = True
    tf_lh.margin_left = tf_lh.margin_top = tf_lh.margin_right = tf_lh.margin_bottom = 0
    p_lh = tf_lh.paragraphs[0]
    r_lh = p_lh.add_run()
    r_lh.text = "📄 " + left_title
    r_lh.font.name = "Consolas"
    r_lh.font.size = Pt(9)
    r_lh.font.bold = True
    r_lh.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)

    line_lh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(36), Pt(84), Pt(301), Pt(1))
    line_lh.name = "Rectangle 5"
    line_lh.fill.solid()
    line_lh.fill.fore_color.rgb = ORANGE_LINE
    line_lh.line.fill.background()
    strip_shape_styles_and_shadows(line_lh)

    tb_code = slide.shapes.add_textbox(Pt(36), Pt(90), Pt(301), Pt(256))
    tb_code.name = "TextBox 6"
    tf_code = tb_code.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = tf_code.margin_top = tf_code.margin_right = tf_code.margin_bottom = 0
    
    first_p = True
    for line in code_lines:
        if first_p:
            p = tf_code.paragraphs[0]
            first_p = False
        else:
            p = tf_code.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        
        runs = process_code_line(line, left_title)
        for text_chunk, chunk_color in runs:
            r = p.add_run()
            r.text = text_chunk
            r.font.name = "Consolas"
            r.font.size = Pt(8)
            r.font.color.rgb = chunk_color

    # Right side: Dynamic Header & Line Placement to prevent collision on multi-line titles
    is_multi_line = len(right_title) > 34
    if is_multi_line:
        tb_rh = slide.shapes.add_textbox(Pt(374), Pt(54), Pt(307), Pt(28))
        font_sz = Pt(10)
        line_top = Pt(85)
        start_y = Pt(94)
    else:
        tb_rh = slide.shapes.add_textbox(Pt(374), Pt(60), Pt(307), Pt(18))
        font_sz = Pt(10.5)
        line_top = Pt(82)
        start_y = Pt(91)
        
    tb_rh.name = "TextBox 7"
    tf_rh = tb_rh.text_frame
    tf_rh.word_wrap = True
    tf_rh.margin_left = tf_rh.margin_top = tf_rh.margin_right = tf_rh.margin_bottom = 0
    p_rh = tf_rh.paragraphs[0]
    r_rh = p_rh.add_run()
    r_rh.text = right_title
    r_rh.font.name = "Montserrat"
    r_rh.font.size = font_sz
    r_rh.font.bold = True
    r_rh.font.color.rgb = ORANGE_LINE

    line_rh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(374), line_top, Pt(307), Pt(1))
    line_rh.name = "Rectangle 8"
    line_rh.fill.solid()
    line_rh.fill.fore_color.rgb = ORANGE_LINE
    line_rh.line.fill.background()
    strip_shape_styles_and_shadows(line_rh)

    y_offset = start_y
    for b_idx, item in enumerate(bullet_items):
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(374), y_offset + Pt(3.5), Pt(6.4), Pt(6.4))
        dot.name = f"Rounded Rectangle Dot {b_idx}"
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)

        tb_b = slide.shapes.add_textbox(Pt(388), y_offset, Pt(293), Pt(item.get("height", 52)))
        tb_b.name = f"TextBox Bullet {b_idx}"
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        p_b1 = tf_b.paragraphs[0]
        p_b1.space_after = Pt(2.5)
        r_title = p_b1.add_run()
        r_title.text = item["title"]
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = DARK_TEXT

        p_b2 = tf_b.add_paragraph()
        p_b2.line_spacing = 1.15
        p_b2.space_after = Pt(0)
        
        if isinstance(item["desc"], str):
            r_desc = p_b2.add_run()
            r_desc.text = item["desc"]
            r_desc.font.name = "Inter"
            r_desc.font.size = Pt(8.5)
            r_desc.font.color.rgb = BODY_TEXT
            desc_len = len(item["desc"])
        else:
            desc_len = 0
            for d_text, d_bold, d_color in item["desc"]:
                r_desc = p_b2.add_run()
                r_desc.text = d_text
                r_desc.font.name = "Inter"
                r_desc.font.size = Pt(8.5)
                r_desc.font.bold = d_bold
                r_desc.font.color.rgb = d_color
                desc_len += len(d_text)

        # Dynamic height calculation and guaranteed 10pt gap to prevent text overlapping
        computed_h = 54 if desc_len > 95 else (44 if desc_len > 45 else 34)
        item_h = Pt(max(item.get("height", 44), computed_h))
        gap = Pt(10)
        y_offset += item_h + gap

def create_dom_tree_diagram_slide(slide, category, question):
    """Visual hierarchical diagram of the DOM Tree structure."""
    add_header_and_footer(slide, category, question)
    
    col_w = Pt(210)
    col_h = Pt(298)
    top_y = Pt(56)
    
    cols_data = [
        {
            "left": Pt(24),
            "accent": RGBColor(0x38, 0xBD, 0xF8), # Cyan
            "title": "1. КОРЕНЬ И СТРУКТУРА",
            "sub": "Объекты window и document",
            "code_badge": "window.document",
            "summary_title": "Корень и окружение:",
            "summary_text": "window представляет окно браузера, а document — загруженную веб-страницу и точку входа в дерево."
        },
        {
            "left": Pt(255),
            "accent": RGBColor(0x4A, 0xDE, 0x80), # Green
            "title": "2. ВЕТВИ И ЭЛЕМЕНТЫ",
            "sub": "Древовидная иерархия тегов",
            "code_badge": "<html> ──► <head> / <body>",
            "summary_title": "Древовидная связь узлов:",
            "summary_text": "Каждый тег превращается в узел DOM (Node). Родители содержат детей, формируя строгую иерархию."
        },
        {
            "left": Pt(486),
            "accent": RGBColor(0xA7, 0x8B, 0xFA), # Violet
            "title": "3. УПРАВЛЕНИЕ ИЗ JS",
            "sub": "Свойства, классы и события",
            "code_badge": "node.classList.add('active')",
            "summary_title": "Живой интерактивный интерфейс:",
            "summary_text": "Через методы DOM JavaScript находит кнопки, меняет классы, выводит уведомления и управляет навигацией."
        }
    ]
    
    for c_idx, c in enumerate(cols_data):
        x = c["left"]
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.name = f"DOM Card Col {c_idx+1}"
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = CODE_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)

        tb_th = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_th = tb_th.text_frame
        tf_th.word_wrap = True
        tf_th.margin_left = tf_th.margin_top = tf_th.margin_right = tf_th.margin_bottom = 0
        p_t1 = tf_th.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = c["title"]
        r_t1.font.name = "Montserrat"
        r_t1.font.size = Pt(11)
        r_t1.font.bold = True
        r_t1.font.color.rgb = c["accent"]

        p_t2 = tf_th.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = c["sub"]
        r_t2.font.name = "Inter"
        r_t2.font.size = Pt(8)
        r_t2.font.color.rgb = CODE_COMM

        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(46), col_w - Pt(24), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = c["accent"]
        div.line.fill.background()
        strip_shape_styles_and_shadows(div)

        vbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(12), top_y + Pt(54), col_w - Pt(24), Pt(176))
        vbox.adjustments[0] = 0.04
        vbox.fill.solid()
        vbox.fill.fore_color.rgb = RGBColor(0x13, 0x14, 0x18)
        vbox.line.color.rgb = CODE_BORDER
        vbox.line.width = Pt(0.75)
        strip_shape_styles_and_shadows(vbox)

        # Draw specific visual diagram
        if c_idx == 0:
            # Window & Document tree root
            b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(20), top_y + Pt(68), col_w - Pt(40), Pt(28))
            b1.adjustments[0] = 0.1
            b1.fill.solid()
            b1.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b1.line.color.rgb = c["accent"]
            b1.line.width = Pt(1)
            strip_shape_styles_and_shadows(b1)
            p_b1 = b1.text_frame.paragraphs[0]
            p_b1.alignment = PP_ALIGN.CENTER
            r_b1 = p_b1.add_run()
            r_b1.text = "window (Глобальный объект)"
            r_b1.font.name = "Consolas"
            r_b1.font.size = Pt(7.5)
            r_b1.font.bold = True
            r_b1.font.color.rgb = WHITE

            # Down arrow
            tb_ar = slide.shapes.add_textbox(x + Pt(20), top_y + Pt(100), col_w - Pt(40), Pt(16))
            p_ar = tb_ar.text_frame.paragraphs[0]
            p_ar.alignment = PP_ALIGN.CENTER
            r_ar = p_ar.add_run()
            r_ar.text = "│\n▼"
            r_ar.font.name = "Consolas"
            r_ar.font.size = Pt(8)
            r_ar.font.color.rgb = c["accent"]

            b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(20), top_y + Pt(130), col_w - Pt(40), Pt(28))
            b2.adjustments[0] = 0.1
            b2.fill.solid()
            b2.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b2.line.color.rgb = c["accent"]
            b2.line.width = Pt(1)
            strip_shape_styles_and_shadows(b2)
            p_b2 = b2.text_frame.paragraphs[0]
            p_b2.alignment = PP_ALIGN.CENTER
            r_b2 = p_b2.add_run()
            r_b2.text = "document (DOM-дерево)"
            r_b2.font.name = "Consolas"
            r_b2.font.size = Pt(7.5)
            r_b2.font.bold = True
            r_b2.font.color.rgb = c["accent"]

            tb_prop = slide.shapes.add_textbox(x + Pt(20), top_y + Pt(168), col_w - Pt(40), Pt(50))
            p_pr = tb_prop.text_frame.paragraphs[0]
            p_pr.alignment = PP_ALIGN.LEFT
            r_pr = p_pr.add_run()
            r_pr.text = "• location.pathname\n• localStorage\n• document.body"
            r_pr.font.name = "Consolas"
            r_pr.font.size = Pt(7.5)
            r_pr.font.color.rgb = CODE_TEXT

        elif c_idx == 1:
            # HTML -> Head / Body hierarchy
            h_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(20), top_y + Pt(64), col_w - Pt(40), Pt(24))
            h_box.adjustments[0] = 0.1
            h_box.fill.solid()
            h_box.fill.fore_color.rgb = RGBColor(0x13, 0x2E, 0x22)
            h_box.line.color.rgb = c["accent"]
            h_box.line.width = Pt(1)
            strip_shape_styles_and_shadows(h_box)
            p_hb = h_box.text_frame.paragraphs[0]
            p_hb.alignment = PP_ALIGN.CENTER
            r_hb = p_hb.add_run()
            r_hb.text = "<html> (Корневой тег)"
            r_hb.font.name = "Consolas"
            r_hb.font.size = Pt(7.5)
            r_hb.font.color.rgb = WHITE

            # 2 branches: <head> and <body>
            sub_w = Pt(78)
            sub1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(20), top_y + Pt(100), sub_w, Pt(54))
            sub1.adjustments[0] = 0.1
            sub1.fill.solid()
            sub1.fill.fore_color.rgb = RGBColor(0x1B, 0x24, 0x33)
            sub1.line.color.rgb = CODE_BORDER
            sub1.line.width = Pt(1)
            strip_shape_styles_and_shadows(sub1)
            p_s1 = sub1.text_frame.paragraphs[0]
            p_s1.alignment = PP_ALIGN.CENTER
            r_s1 = p_s1.add_run()
            r_s1.text = "<head>\n"
            r_s1.font.name = "Consolas"
            r_s1.font.size = Pt(7.5)
            r_s1.font.bold = True
            r_s1.font.color.rgb = CODE_KW
            r_s1b = p_s1.add_run()
            r_s1b.text = "• <script defer>\n• <link css>"
            r_s1b.font.name = "Consolas"
            r_s1b.font.size = Pt(6.5)
            r_s1b.font.color.rgb = CODE_COMM

            sub2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(108), top_y + Pt(100), sub_w, Pt(54))
            sub2.adjustments[0] = 0.1
            sub2.fill.solid()
            sub2.fill.fore_color.rgb = RGBColor(0x1B, 0x24, 0x33)
            sub2.line.color.rgb = c["accent"]
            sub2.line.width = Pt(1)
            strip_shape_styles_and_shadows(sub2)
            p_s2 = sub2.text_frame.paragraphs[0]
            p_s2.alignment = PP_ALIGN.CENTER
            r_s2 = p_s2.add_run()
            r_s2.text = "<body>\n"
            r_s2.font.name = "Consolas"
            r_s2.font.size = Pt(7.5)
            r_s2.font.bold = True
            r_s2.font.color.rgb = c["accent"]
            r_s2b = p_s2.add_run()
            r_s2b.text = "• <header>\n• <main>\n• <footer>"
            r_s2b.font.name = "Consolas"
            r_s2b.font.size = Pt(6.5)
            r_s2b.font.color.rgb = CODE_TEXT

            ch_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(20), top_y + Pt(166), col_w - Pt(40), Pt(50))
            ch_box.adjustments[0] = 0.1
            ch_box.fill.solid()
            ch_box.fill.fore_color.rgb = RGBColor(0x13, 0x14, 0x18)
            ch_box.line.color.rgb = CODE_BORDER
            ch_box.line.width = Pt(0.75)
            strip_shape_styles_and_shadows(ch_box)
            p_ch = ch_box.text_frame.paragraphs[0]
            p_ch.alignment = PP_ALIGN.LEFT
            r_ch = p_ch.add_run()
            r_ch.text = " └─ <nav class=\"nav\">\n      └─ <a class=\"nav-link\">"
            r_ch.font.name = "Consolas"
            r_ch.font.size = Pt(7.5)
            r_ch.font.color.rgb = c["accent"]

        else:
            # JavaScript DOM manipulation methods
            methods = [
                ("1. ПОИСК:", "document.querySelectorAll('.nav-link')", c["accent"]),
                ("2. ЧТЕНИЕ:", "link.getAttribute('href')", WHITE),
                ("3. КЛАССЫ:", "link.classList.add('active')", RGBColor(0x4A, 0xDE, 0x80)),
                ("4. УВЕДОМЛЕНИЕ:", "document.createElement('div')", RGBColor(0xFC, 0xD3, 0x4D))
            ]
            my = top_y + Pt(64)
            for m_lbl, m_code, m_col in methods:
                mb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(18), my, col_w - Pt(36), Pt(34))
                mb.adjustments[0] = 0.08
                mb.fill.solid()
                mb.fill.fore_color.rgb = RGBColor(0x1E, 0x1B, 0x2A)
                mb.line.color.rgb = CODE_BORDER
                mb.line.width = Pt(0.75)
                strip_shape_styles_and_shadows(mb)
                
                tf_mb = mb.text_frame
                tf_mb.margin_left = tf_mb.margin_top = tf_mb.margin_right = tf_mb.margin_bottom = 0
                p_mb1 = tf_mb.paragraphs[0]
                r_mb1 = p_mb1.add_run()
                r_mb1.text = " " + m_lbl
                r_mb1.font.name = "Consolas"
                r_mb1.font.size = Pt(7)
                r_mb1.font.bold = True
                r_mb1.font.color.rgb = m_col

                p_mb2 = tf_mb.add_paragraph()
                r_mb2 = p_mb2.add_run()
                r_mb2.text = " " + m_code
                r_mb2.font.name = "Consolas"
                r_mb2.font.size = Pt(6.5)
                r_mb2.font.color.rgb = CODE_TEXT
                my += Pt(40)

        # Summary text
        tb_sm = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(236), col_w - Pt(24), Pt(54))
        tf_sm = tb_sm.text_frame
        tf_sm.word_wrap = True
        tf_sm.margin_left = tf_sm.margin_top = tf_sm.margin_right = tf_sm.margin_bottom = 0
        p_s1 = tf_sm.paragraphs[0]
        r_s1 = p_s1.add_run()
        r_s1.text = c["summary_title"] + "\n"
        r_s1.font.name = "Inter"
        r_s1.font.size = Pt(8)
        r_s1.font.bold = True
        r_s1.font.color.rgb = c["accent"]

        p_s2 = tf_sm.add_paragraph()
        r_s2 = p_s2.add_run()
        r_s2.text = c["summary_text"]
        r_s2.font.name = "Inter"
        r_s2.font.size = Pt(7.5)
        r_s2.font.color.rgb = CODE_TEXT

def create_lifecycle_diagram_slide(slide, category, question):
    """Visual timeline of script execution and DOMContentLoaded event."""
    add_header_and_footer(slide, category, question)
    
    # Timeline card
    col_w = Pt(672)
    col_h = Pt(298)
    top_y = Pt(56)
    
    steps = [
        ("ШАГ 1: ЗАГРУЗКА HTML", "Браузер считывает HTML сверху вниз и строит DOM-дерево в памяти.", RGBColor(0x38, 0xBD, 0xF8)),
        ("ШАГ 2: СКРИПТ DEFER", "Скрипт скачивается в фоне, не блокируя отрисовку стилей и картинок.", RGBColor(0xFC, 0xD3, 0x4D)),
        ("ШАГ 3: DOMCONTENTLOADED", "HTML полностью готов! Браузер мгновенно стреляет системным событием.", RGBColor(0x4A, 0xDE, 0x80)),
        ("ШАГ 4: ЗАПУСК ЛОГИКИ", "Скрипт main.js безопасно находит элементы .nav-link и подсвечивает меню.", RGBColor(0xA7, 0x8B, 0xFA))
    ]

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(24), top_y, col_w, col_h)
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = CODE_BG
    card.line.color.rgb = CODE_BORDER
    card.line.width = Pt(1)
    strip_shape_styles_and_shadows(card)

    tb_h = slide.shapes.add_textbox(Pt(40), top_y + Pt(14), col_w - Pt(32), Pt(24))
    p_h = tb_h.text_frame.paragraphs[0]
    r_h = p_h.add_run()
    r_h.text = "ЖИЗНЕННЫЙ ЦИКЛ СТРАНИЦЫ И БЕЗОПАСНЫЙ СТАРТ JAVASCRIPT"
    r_h.font.name = "Montserrat"
    r_h.font.size = Pt(11)
    r_h.font.bold = True
    r_h.font.color.rgb = ORANGE_LINE

    line_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(40), top_y + Pt(38), col_w - Pt(32), Pt(1))
    line_h.fill.solid()
    line_h.fill.fore_color.rgb = ORANGE_LINE
    line_h.line.fill.background()
    strip_shape_styles_and_shadows(line_h)

    # 4 horizontal cards along timeline
    c_w = Pt(150)
    c_h = Pt(180)
    start_x = Pt(40)
    gap_x = Pt(14)

    for s_idx, (stitle, sdesc, scol) in enumerate(steps):
        sx = start_x + Pt(s_idx * (150 + 14))
        sy = top_y + Pt(54)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, c_w, c_h)
        box.adjustments[0] = 0.06
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x13, 0x16, 0x1C)
        box.line.color.rgb = scol
        box.line.width = Pt(1)
        strip_shape_styles_and_shadows(box)

        # Number circle badge
        num_badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, sx + Pt(14), sy + Pt(14), Pt(24), Pt(24))
        num_badge.fill.solid()
        num_badge.fill.fore_color.rgb = scol
        num_badge.line.fill.background()
        strip_shape_styles_and_shadows(num_badge)
        p_nb = num_badge.text_frame.paragraphs[0]
        p_nb.alignment = PP_ALIGN.CENTER
        r_nb = p_nb.add_run()
        r_nb.text = str(s_idx + 1)
        r_nb.font.name = "Montserrat"
        r_nb.font.size = Pt(10)
        r_nb.font.bold = True
        r_nb.font.color.rgb = RGBColor(0x18, 0x1A, 0x1F)

        # Header
        tb_sh = slide.shapes.add_textbox(sx + Pt(14), sy + Pt(46), c_w - Pt(28), Pt(32))
        tf_sh = tb_sh.text_frame
        tf_sh.word_wrap = True
        tf_sh.margin_left = tf_sh.margin_top = tf_sh.margin_right = tf_sh.margin_bottom = 0
        p_sh = tf_sh.paragraphs[0]
        r_sh = p_sh.add_run()
        r_sh.text = stitle
        r_sh.font.name = "Montserrat"
        r_sh.font.size = Pt(8.5)
        r_sh.font.bold = True
        r_sh.font.color.rgb = scol

        # Line
        div_s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx + Pt(14), sy + Pt(82), c_w - Pt(28), Pt(1))
        div_s.fill.solid()
        div_s.fill.fore_color.rgb = CODE_BORDER
        div_s.line.fill.background()
        strip_shape_styles_and_shadows(div_s)

        # Desc
        tb_sd = slide.shapes.add_textbox(sx + Pt(14), sy + Pt(88), c_w - Pt(28), Pt(84))
        tf_sd = tb_sd.text_frame
        tf_sd.word_wrap = True
        tf_sd.margin_left = tf_sd.margin_top = tf_sd.margin_right = tf_sd.margin_bottom = 0
        p_sd = tf_sd.paragraphs[0]
        r_sd = p_sd.add_run()
        r_sd.text = sdesc
        r_sd.font.name = "Inter"
        r_sd.font.size = Pt(8)
        r_sd.font.color.rgb = CODE_TEXT

    # Bottom takeaway
    tb_bot_info = slide.shapes.add_textbox(Pt(40), top_y + Pt(246), col_w - Pt(32), Pt(40))
    tf_bi = tb_bot_info.text_frame
    tf_bi.word_wrap = True
    tf_bi.margin_left = tf_bi.margin_top = tf_bi.margin_right = tf_bi.margin_bottom = 0
    p_bi = tf_bi.paragraphs[0]
    r_bi1 = p_bi.add_run()
    r_bi1.text = "Ключевое правило: "
    r_bi1.font.name = "Montserrat"
    r_bi1.font.size = Pt(9)
    r_bi1.font.bold = True
    r_bi1.font.color.rgb = ORANGE_PILL
    r_bi2 = p_bi.add_run()
    r_bi2.text = "Никогда не обращайтесь к DOM до события DOMContentLoaded — всегда используйте атрибут defer."
    r_bi2.font.name = "Inter"
    r_bi2.font.size = Pt(8.5)
    r_bi2.font.color.rgb = CODE_TEXT

def create_grid_cards_slide(slide, category, question, title_text, cards_data):
    """2x2 error analysis cards without gray/transparent background plates."""
    add_header_and_footer(slide, category, question)
    
    coords = [
        (Pt(28), Pt(56), Pt(320), Pt(142)),
        (Pt(372), Pt(56), Pt(320), Pt(142)),
        (Pt(28), Pt(208), Pt(320), Pt(142)),
        (Pt(372), Pt(208), Pt(320), Pt(142)),
    ]
    
    for idx, (x, y, w, h) in enumerate(coords):
        data = cards_data[idx]
        
        tb_h = slide.shapes.add_textbox(x, y + Pt(6), w, Pt(20))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
        p_h = tf_h.paragraphs[0]
        r_num = p_h.add_run()
        r_num.text = f"{idx + 1}. {data['title']}"
        r_num.font.name = "Montserrat"
        r_num.font.size = Pt(10.5)
        r_num.font.bold = True
        r_num.font.color.rgb = ORANGE_LINE
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Pt(28), w, Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE_LINE
        line.line.fill.background()
        strip_shape_styles_and_shadows(line)
        
        tb_body = slide.shapes.add_textbox(x, y + Pt(36), w, Pt(96))
        tf_b = tb_body.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        p1 = tf_b.paragraphs[0]
        p1.space_after = Pt(4)
        r1_lbl = p1.add_run()
        r1_lbl.text = "Ошибка: "
        r1_lbl.font.name = "Inter"
        r1_lbl.font.size = Pt(8.5)
        r1_lbl.font.bold = True
        r1_lbl.font.color.rgb = DARK_TEXT
        r1_txt = p1.add_run()
        r1_txt.text = data['problem']
        r1_txt.font.name = "Inter"
        r1_txt.font.size = Pt(8.5)
        r1_txt.font.color.rgb = BODY_TEXT
        
        p2 = tf_b.add_paragraph()
        p2.space_after = Pt(0)
        p2.line_spacing = 1.15
        r2_lbl = p2.add_run()
        r2_lbl.text = "Решение: "
        r2_lbl.font.name = "Inter"
        r2_lbl.font.size = Pt(8.5)
        r2_lbl.font.bold = True
        r2_lbl.font.color.rgb = GREEN_TEXT
        r2_txt = p2.add_run()
        r2_txt.text = data['solution']
        r2_txt.font.name = "Inter"
        r2_txt.font.size = Pt(8.5)
        r2_txt.font.color.rgb = BODY_TEXT

def create_checklist_slide(slide, category, question, title_text, items):
    """Clean 6-item checklist with checkmarks."""
    add_header_and_footer(slide, category, question)

    tb_h = slide.shapes.add_textbox(Pt(54), Pt(68), Pt(600), Pt(22))
    tf_h = tb_h.text_frame
    tf_h.word_wrap = True
    tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
    p_h = tf_h.paragraphs[0]
    r_h = p_h.add_run()
    r_h.text = title_text
    r_h.font.name = "Montserrat"
    r_h.font.size = Pt(11)
    r_h.font.bold = True
    r_h.font.color.rgb = ORANGE_LINE

    line_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(54), Pt(92), Pt(612), Pt(1))
    line_h.fill.solid()
    line_h.fill.fore_color.rgb = ORANGE_LINE
    line_h.line.fill.background()
    strip_shape_styles_and_shadows(line_h)

    y_pos = Pt(102)
    for idx, item in enumerate(items):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(54), y_pos + Pt(2), Pt(22), Pt(20))
        badge.adjustments[0] = 0.2
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN_BADGE
        badge.line.fill.background()
        strip_shape_styles_and_shadows(badge)
        
        tb_chk = slide.shapes.add_textbox(Pt(54), y_pos + Pt(4), Pt(22), Pt(20))
        tf_chk = tb_chk.text_frame
        p_chk = tf_chk.paragraphs[0]
        p_chk.alignment = PP_ALIGN.CENTER
        r_chk = p_chk.add_run()
        r_chk.text = "✔"
        r_chk.font.name = "Segoe UI Symbol"
        r_chk.font.size = Pt(10)
        r_chk.font.bold = True
        r_chk.font.color.rgb = WHITE
        
        tb_item = slide.shapes.add_textbox(Pt(86), y_pos, Pt(570), Pt(38))
        tf_item = tb_item.text_frame
        tf_item.word_wrap = True
        tf_item.margin_left = tf_item.margin_top = tf_item.margin_right = tf_item.margin_bottom = 0
        
        p1 = tf_item.paragraphs[0]
        p1.space_after = Pt(2)
        r1 = p1.add_run()
        r1.text = item["title"]
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT

        p2 = tf_item.add_paragraph()
        r2 = p2.add_run()
        r2.text = item["desc"]
        r2.font.name = "Inter"
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = BODY_TEXT

        y_pos += Pt(42)

slides_data = [
    # 1. Slide 3: What is JavaScript in Browser
    {
        "type": "code",
        "category": "ЯЗЫК JAVASCRIPT",
        "question": "Какую роль выполняет JavaScript в современном веб-приложении?",
        "left_title": "Архитектура Web: HTML + CSS + JS",
        "code_lines": [
            [("// 1. HTML: Структура и семантика каркаса", CODE_COMM)],
            [("<a href=\"catalog.html\" class=\"nav-link\">Каталог</a>", CODE_COMM)],
            [("", CODE_TEXT)],
            [("/* 2. CSS: Визуальные стили и оформление */", CODE_COMM)],
            [(".nav-link.active { color: #007bff; }", CODE_COMM)],
            [("", CODE_TEXT)],
            [("// 3. JavaScript: Живая интерактивная логика", CODE_COMM)],
            [("const ", CODE_KW), ("navLinks = document.querySelectorAll(", CODE_TEXT), ("'.nav-link'", CODE_STR), (");", CODE_TEXT)],
            [("const ", CODE_KW), ("currentPath = window.location.pathname;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("navLinks.forEach", CODE_PROP), ("(", CODE_TEXT), ("link", CODE_VAL), (" => ", CODE_KW), ("{", CODE_TEXT)],
            [("  if (link.getAttribute(", CODE_TEXT), ("'href'", CODE_STR), (") === currentPath) {", CODE_TEXT)],
            [("    link.classList.add(", CODE_TEXT), ("'active'", CODE_STR), (");", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("});", CODE_TEXT)]
        ],
        "right_title": "ТРИЕДИНСТВО WEB: СТРУКТУРА, СТИЛЬ, ПОВЕДЕНИЕ",
        "bullet_items": [
            {
                "title": "Интерактивность в реальном времени:",
                "desc": "HTML создает скелет, CSS задает красоту, а JavaScript вдыхает жизнь: реагирует на клики, переключает страницы и проверяет формы.",
                "height": 46
            },
            {
                "title": "Клиентский язык (Client-Side):",
                "desc": "Код выполняется прямо на компьютере пользователя в браузере, не требуя перезагрузки страницы при каждом действии.",
                "height": 46
            },
            {
                "title": "Динамическое изменение стилей:",
                "desc": "JS может на лету добавлять или снимать CSS-классы (.active, .open, .loading), мгновенно меняя оформление элементов.",
                "height": 48
            },
            {
                "title": "Современный стандарт ES6+:",
                "desc": "Мы используем стрелочные функции (=>), строгие константы const и встроенные методы перебора коллекций .forEach().",
                "height": 44
            }
        ]
    },

    # 2. Slide 4: Theory: DOM Concept
    {
        "type": "code",
        "category": "ОСНОВЫ DOM",
        "question": "Что такое Document Object Model и как браузер видит страницу?",
        "left_title": "DOM: Представление тегов в виде объектов",
        "code_lines": [
            [("// HTML в исходном коде:", CODE_COMM)],
            [("<a class=\"nav-link\" href=\"catalog.html\">Каталог</a>", CODE_COMM)],
            [("", CODE_TEXT)],
            [("// В памяти браузера это JS-объект HTMLElement:", CODE_COMM)],
            [("const ", CODE_KW), ("linkObj = {", CODE_TEXT)],
            [("  tagName: ", CODE_PROP), ("'A'", CODE_STR), (",", CODE_TEXT)],
            [("  className: ", CODE_PROP), ("'nav-link'", CODE_STR), (",", CODE_TEXT)],
            [("  href: ", CODE_PROP), ("'http://site.ru/catalog.html'", CODE_STR), (",", CODE_TEXT)],
            [("  textContent: ", CODE_PROP), ("'Каталог'", CODE_STR), (",", CODE_TEXT)],
            [("  classList: { add: ƒ, remove: ƒ, toggle: ƒ },", CODE_COMM)],
            [("  style: { color: '', backgroundColor: '' }", CODE_COMM)],
            [("};", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Мы можем управлять объектом через JS:", CODE_COMM)],
            [("linkObj.classList.add(", CODE_TEXT), ("'active'", CODE_STR), (");", CODE_TEXT)]
        ],
        "right_title": "DOCUMENT OBJECT MODEL (DOM)",
        "bullet_items": [
            {
                "title": "Превращение HTML в дерево объектов:",
                "desc": "При чтении HTML браузер парсит каждый тег и создает в оперативной памяти древовидную модель объектов документа (DOM).",
                "height": 46
            },
            {
                "title": "Программный интерфейс (API):",
                "desc": "DOM предоставляет сотни встроенных функций (document.getElementById, addEventListener), через которые JS общается с HTML.",
                "height": 48
            },
            {
                "title": "Двусторонняя реакция:",
                "desc": "Когда JS меняет свойство объекта в DOM (например, textContent или classList), браузер тут же перерисовывает экран пользователя.",
                "height": 46
            },
            {
                "title": "Доступ через глобальный document:",
                "desc": "Служебный объект document служит входными воротами во все дерево элементов текущей загруженной веб-страницы.",
                "height": 44
            }
        ]
    },

    # 3. Slide 5: Diagram: DOM Tree (Diagram slide)
    {
        "type": "dom_tree",
        "category": "СХЕМА: ДЕРЕВО DOM",
        "question": "Как устроена иерархия объектов внутри окна браузера?"
    },

    # 4. Slide 6: Script Inclusion: Defer vs Async
    {
        "type": "code",
        "category": "ПОДКЛЮЧЕНИЕ JS",
        "question": "Почему атрибут defer является золотым стандартом подключения скриптов?",
        "left_title": "index.html (<script src=\"...\" defer>)",
        "code_lines": [
            [("<!-- 1. Без defer (ОПАСНО: блокирует парсинг) -->", CODE_COMM)],
            [("<script ", CODE_TAG), ("src=", CODE_PROP), ("\"js/main.js\"", CODE_STR), ("></script>", CODE_TAG)],
            [("", CODE_TEXT)],
            [("<!-- 2. С async (Хаотично: выполнится когда скачается) -->", CODE_COMM)],
            [("<script ", CODE_TAG), ("src=", CODE_PROP), ("\"js/main.js\"", CODE_STR), (" async", CODE_KW), ("></script>", CODE_TAG)],
            [("", CODE_TEXT)],
            [("<!-- 3. С defer (ИДЕАЛЬНО: фон + порядок + после DOM) -->", CODE_COMM)],
            [("<script ", CODE_TAG), ("src=", CODE_PROP), ("\"js/main.js\"", CODE_STR), (" defer", CODE_KW), ("></script>", CODE_TAG)],
            [("", CODE_TEXT)],
            [("<!-- Место подключения: внутри <head> -->", CODE_COMM)],
            [("<head>", CODE_TAG)],
            [("  <meta charset=\"UTF-8\">", CODE_TEXT)],
            [("  <link rel=\"stylesheet\" href=\"css/style.css\">", CODE_TEXT)],
            [("  <script src=\"js/main.js\" defer></script>", CODE_KW)],
            [("</head>", CODE_TAG)]
        ],
        "right_title": "АТРИБУТ DEFER: БЫСТРОДЕЙСТВИЕ И НАДЕЖНОСТЬ",
        "bullet_items": [
            {
                "title": "Проблема классического <script>:",
                "desc": "Без defer браузер останавливает построение страницы, скачивает JS-файл и только потом продолжает. Пользователь видит белый экран.",
                "height": 48
            },
            {
                "title": "Фоновая загрузка (Non-blocking):",
                "desc": "Атрибут defer указывает браузеру загружать скрипт параллельно в фоне, вообще не замедляя отображение стилей и верстки.",
                "height": 46
            },
            {
                "title": "Гарантия готовности DOM:",
                "desc": "Скрипт с defer гарантированно запускается только тогда, когда весь HTML-документ уже полностью прочитан и построен.",
                "height": 46
            },
            {
                "title": "Сохранение порядка вызовов:",
                "desc": "В отличие от async, несколько скриптов с defer всегда выполняются строго в том порядке, в котором они указаны в HTML.",
                "height": 44
            }
        ]
    },

    # 5. Slide 7: DOMContentLoaded Event
    {
        "type": "code",
        "category": "СТАРТ СКРИПТА",
        "question": "Как гарантировать, что JavaScript запустится вовремя?",
        "left_title": "js/main.js (DOMContentLoaded)",
        "code_lines": [
            [("// СмартОфис — Скрипт веб-приложения (Вебинар 4)", CODE_COMM)],
            [("", CODE_TEXT)],
            [("// Подписка на событие готовности DOM-дерева", CODE_COMM)],
            [("document.", CODE_TEXT), ("addEventListener", CODE_PROP), ("(", CODE_TEXT), ("'DOMContentLoaded'", CODE_STR), (", () => {", CODE_KW)],
            [("  // Вся разметка готова к поиску и манипуляциям!", CODE_COMM)],
            [("  console.log(", CODE_PROP), ("'DOM полностью загружен и готов!'", CODE_STR), (");", CODE_TEXT)],
            [("  initNavigation();", CODE_TEXT)],
            [("});", CODE_KW)],
            [("", CODE_TEXT)],
            [("// Функция инициализации навигации", CODE_COMM)],
            [("function ", CODE_KW), ("initNavigation() {", CODE_TEXT)],
            [("  // Безопасный поиск ссылок внутри DOM", CODE_COMM)],
            [("  const links = document.querySelectorAll('.nav-link');", CODE_TEXT)],
            [("  console.log(`Найдено ссылок меню: ${links.length}`);", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "СОБЫТИЕ ГОТОВНОСТИ DOMCONTENTLOADED",
        "bullet_items": [
            {
                "title": "Системное событие DOMContentLoaded:",
                "desc": "Срабатывает в тот миллисекундный момент, когда браузер завершил построение дерева тегов, не дожидаясь тяжелых фото и шрифтов.",
                "height": 48
            },
            {
                "title": "Метод addEventListener:",
                "desc": "Универсальный механизм подписки на события. Принимает имя события ('DOMContentLoaded') и функцию-обработчик (callback).",
                "height": 46
            },
            {
                "title": "Защита от фатального Cannot read null:",
                "desc": "Если вызвать querySelector до DOMContentLoaded, браузер не найдет элемент, вернет null, и вызов .classList вызовет сбой.",
                "height": 46
            },
            {
                "title": "Быстрый отклик интерфейса:",
                "desc": "Интерактивная логика меню подключается мгновенно, не заставляя пользователя ждать догрузки всех баннеров каталога.",
                "height": 44
            }
        ]
    },

    # 6. Slide 8: Diagram: Lifecycle Timeline (Diagram slide)
    {
        "type": "lifecycle_timeline",
        "category": "СХЕМА: ЖИЗНЕННЫЙ ЦИКЛ",
        "question": "Какова точная последовательность загрузки страницы и выполнения кода?"
    },

    # 7. Slide 9: Querying DOM Elements
    {
        "type": "code",
        "category": "ПОИСК В DOM",
        "question": "Как находить элементы страницы с помощью селекторов CSS?",
        "left_title": "js/main.js (Методы выборки элементов)",
        "code_lines": [
            [("// 1. Поиск ОДНОГО элемента по ID (быстрый доступ)", CODE_COMM)],
            [("const ", CODE_KW), ("authBtn = document.", CODE_TEXT), ("getElementById", CODE_PROP), ("(", CODE_TEXT), ("'authNavBtn'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 2. Поиск ПЕРВОГО совпадения по CSS-селектору", CODE_COMM)],
            [("const ", CODE_KW), ("mainTitle = document.", CODE_TEXT), ("querySelector", CODE_PROP), ("(", CODE_TEXT), ("'.hero-title'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 3. Поиск ВСЕХ элементов по CSS-селектору", CODE_COMM)],
            [("const ", CODE_KW), ("links = document.", CODE_TEXT), ("querySelectorAll", CODE_PROP), ("(", CODE_TEXT), ("'.nav-link'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Результат querySelectorAll — псевдомассив NodeList", CODE_COMM)],
            [("console.log(links);        // NodeList(3) [a, a, a]", CODE_COMM)],
            [("console.log(links.length); // 3 (количество найденных)", CODE_COMM)]
        ],
        "right_title": "ИНСТРУМЕНТЫ ВЫБОРКИ ЭЛЕМЕНТОВ В JAVASCRIPT",
        "bullet_items": [
            {
                "title": "document.getElementById('id'):",
                "desc": "Самый быстрый метод в движке браузера. Ищет единственный уникальный элемент по его идентификатору без решетки '#'.",
                "height": 46
            },
            {
                "title": "document.querySelector('селектор'):",
                "desc": "Принимает любой валидный CSS-селектор (.card, #btn, nav > ul a). Возвращает первый найденный элемент или null.",
                "height": 48
            },
            {
                "title": "document.querySelectorAll('селектор'):",
                "desc": "Находит абсолютно все элементы на странице, соответствующие CSS-правилу, и возвращает их в виде коллекции NodeList.",
                "height": 46
            },
            {
                "title": "Что такое NodeList:",
                "desc": "Это упорядоченный список элементов, похожий на массив. Он имеет свойство .length и поддерживает удобный перебор через .forEach().",
                "height": 44
            }
        ]
    },

    # 8. Slide 10: Iterating Elements with forEach
    {
        "type": "code",
        "category": "ПЕРЕБОР ЭЛЕМЕНТОВ",
        "question": "Как применить логику к каждому найденному пункту меню?",
        "left_title": "js/main.js (Цикл forEach)",
        "code_lines": [
            [("function ", CODE_KW), ("initNavigation() {", CODE_TEXT)],
            [("  const links = document.querySelectorAll('.nav-link');", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // Перебираем каждую ссылку по очереди:", CODE_COMM)],
            [("  links.", CODE_TEXT), ("forEach", CODE_PROP), ("((", CODE_TEXT), ("link, index", CODE_VAL), (") => {", CODE_KW)],
            [("    console.log(`Ссылка #${index}:`, link.textContent);", CODE_TEXT)],
            [("    ", CODE_TEXT)],
            [("    // Получаем адрес, куда ведет ссылка", CODE_COMM)],
            [("    const href = link.getAttribute('href');", CODE_TEXT)],
            [("    if (!href) return; // Пропуск пустых", CODE_TEXT)],
            [("    ", CODE_TEXT)],
            [("    // Подготовка элемента", CODE_COMM)],
            [("    link.classList.remove('active');", CODE_TEXT)],
            [("  });", CODE_KW)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "МЕТОД FOREACH: ПРОСТОТА И НАДЕЖНОСТЬ",
        "bullet_items": [
            {
                "title": "Перебор без счетчиков и i++:",
                "desc": "В отличие от старого for(let i = 0; ...), метод .forEach() автоматически передает каждый элемент в функцию без риска выйти за границы.",
                "height": 48
            },
            {
                "title": "Аргументы колбэк-функции:",
                "desc": "Первым параметром передается сам DOM-узел (link), вторым — порядковый индекс элемента (0, 1, 2...), если он необходим.",
                "height": 46
            },
            {
                "title": "Оператор return в forEach:",
                "desc": "Команда return внутри forEach работает как continue в обычном цикле: завершает текущую итерацию и переходит к следующей ссылке.",
                "height": 46
            },
            {
                "title": "Чистота функционального стиля:",
                "desc": "Стрелочные функции делают код компактным, легко читаемым для коллег и снижают количество синтаксических ошибок.",
                "height": 44
            }
        ]
    },

    # 9. Slide 11: Reading Attributes (getAttribute)
    {
        "type": "code",
        "category": "РАБОТА С АТРИБУТАМИ",
        "question": "Как извлечь значение атрибута href у тега ссылки?",
        "left_title": "js/main.js (getAttribute vs .href)",
        "code_lines": [
            [("// Ссылка в разметке: <a href=\"pages/catalog.html\">", CODE_COMM)],
            [("const ", CODE_KW), ("link = document.querySelector(", CODE_TEXT), ("'.nav-link'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Вариант А: Свойство объекта .href (Абсолютный URL)", CODE_COMM)],
            [("console.log(link.href);", CODE_TEXT)],
            [("// Выведет: 'http://127.0.0.1:5500/pages/catalog.html'", CODE_COMM)],
            [("", CODE_TEXT)],
            [("// Вариант Б: Метод getAttribute (Исходная строка)", CODE_COMM)],
            [("const ", CODE_KW), ("href = link.", CODE_TEXT), ("getAttribute", CODE_PROP), ("(", CODE_TEXT), ("'href'", CODE_STR), (");", CODE_TEXT)],
            [("console.log(href);", CODE_TEXT)],
            [("// Выведет строго то, что написано в HTML: 'pages/catalog.html'", CODE_COMM)],
            [("", CODE_TEXT)],
            [("// Защитная проверка:", CODE_COMM)],
            [("if (!href || href === '#') {", CODE_TEXT)],
            [("  return; // Ссылка без адреса, пропускаем", CODE_COMM)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "GETATTRIBUTE: ТОЧНОСТЬ И БЕЗОПАСНОСТЬ",
        "bullet_items": [
            {
                "title": "Разница между .href и getAttribute('href'):",
                "desc": "Свойство .href возвращает полный адрес с доменом и протоколом. getAttribute('href') возвращает ровно ту строку, что вбита в HTML.",
                "height": 48
            },
            {
                "title": "Независимость от сервера и домена:",
                "desc": "getAttribute возвращает относительный путь ('catalog.html'), благодаря чему скрипт навигации работает и локально, и на боевом сервере.",
                "height": 46
            },
            {
                "title": "Универсальность для любых атрибутов:",
                "desc": "Метод getAttribute считывает любые атрибуты: 'src', 'alt', 'data-id', 'title', 'aria-label', давая полный контроль над разметкой.",
                "height": 46
            },
            {
                "title": "Защитное программирование:",
                "desc": "Проверка if (!href) защищает код от падений, если верстальщик оставит в меню пустую ссылку без атрибута href.",
                "height": 44
            }
        ]
    },

    # 10. Slide 12: Theory: window.location.pathname
    {
        "type": "code",
        "category": "МАРШРУТИЗАЦИЯ В JS",
        "question": "Как скрипт определяет текущий адрес открытой веб-страницы?",
        "left_title": "Объект window.location в консоли",
        "code_lines": [
            [("// Допустим, адрес: https://smartoffice.ru/pages/catalog.html?sort=price", CODE_COMM)],
            [("", CODE_TEXT)],
            [("// 1. Полный URL страницы:", CODE_COMM)],
            [("console.log(window.location.href);", CODE_TEXT)],
            [("// 'https://smartoffice.ru/pages/catalog.html?sort=price'", CODE_COMM)],
            [("", CODE_TEXT)],
            [("// 2. Имя домена:", CODE_COMM)],
            [("console.log(window.location.hostname);", CODE_TEXT)],
            [("// 'smartoffice.ru'", CODE_COMM)],
            [("", CODE_TEXT)],
            [("// 3. Путь к текущему файлу (Pathname):", CODE_COMM)],
            [("const ", CODE_KW), ("current = window.location.pathname;", CODE_TEXT)],
            [("console.log(current);", CODE_TEXT)],
            [("// '/pages/catalog.html' — ИМЕННО ЭТО НАМ И НУЖНО!", CODE_COMM)]
        ],
        "right_title": "АНАТОМИЯ АДРЕСНОЙ СТРОКИ: WINDOW.LOCATION",
        "bullet_items": [
            {
                "title": "Служебный объект window.location:",
                "desc": "Встроенный объект браузера, содержащий полную информацию об URL текущей страницы и инструменты для перехода по ссылкам.",
                "height": 46
            },
            {
                "title": "Свойство .pathname:",
                "desc": "Возвращает путь к ресурсу, отсекая протокол (https://), домен и query-параметры. Например: '/index.html' или '/pages/catalog.html'.",
                "height": 48
            },
            {
                "title": "Устойчивость к портам локальных серверов:",
                "desc": "Запуск через Live Server (порт 5500) или Node.js не ломает сравнение путей, так как .pathname не зависит от номера порта.",
                "height": 46
            },
            {
                "title": "Основа SPA-маршрутизации:",
                "desc": "Именно на разборе pathname строятся современные роутеры во фреймворках React, Vue, Angular и нативных веб-приложениях.",
                "height": 44
            }
        ]
    },

    # 11. Slide 13: Practice JS: initNavigation Logic
    {
        "type": "code",
        "category": "УМНАЯ НАВИГАЦИЯ",
        "question": "Как устроена логика функции сопоставления ссылок с текущим адресом?",
        "left_title": "js/main.js (Функция initNavigation)",
        "code_lines": [
            [("function ", CODE_KW), ("initNavigation() {", CODE_TEXT)],
            [("  const links = document.querySelectorAll('.nav-link');", CODE_TEXT)],
            [("  const current = window.location.pathname;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  links.forEach(link => {", CODE_TEXT)],
            [("    const href = link.getAttribute('href');", CODE_TEXT)],
            [("    if (!href) return;", CODE_TEXT)],
            [("    link.classList.remove('active'); // Сброс", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("    // Главная страница (корень / или index.html)", CODE_COMM)],
            [("    const isHome = current.endsWith('index.html') || current.endsWith('/') || current === '';", CODE_TEXT)],
            [("    if (isHome && (href === 'index.html' || href === '../index.html')) {", CODE_TEXT)],
            [("      link.classList.add('active');", CODE_TEXT)],
            [("    } else if (href.includes('catalog.html') && current.includes('catalog.html')) {", CODE_TEXT)],
            [("      link.classList.add('active'); // Каталог", CODE_TEXT)],
            [("    }", CODE_TEXT)],
            [("  });", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "АЛГОРИТМ СОПОСТАВЛЕНИЯ ПУТЕЙ",
        "bullet_items": [
            {
                "title": "Шаг 1: Очистка старых классов:",
                "desc": "Перед проверкой с каждой ссылки снимается класс active (link.classList.remove). Это исключает подсветку двух пунктов одновременно.",
                "height": 46
            },
            {
                "title": "Метод endsWith('index.html'):",
                "desc": "Проверяет окончание строки адреса, корректно распознавая и прямые ссылки на файл, и заходы в корень сайта ('/').",
                "height": 46
            },
            {
                "title": "Метод includes('catalog.html'):",
                "desc": "Ищет вхождение подстроки в URL. Если в адресе есть catalog.html, пункт каталога моментально получает подсветку.",
                "height": 48
            },
            {
                "title": "Поддержка вложенных папок (/pages/):",
                "desc": "Логика учитывает ссылки с выходом наверх ('../index.html'), одинаково четко работая на главной и внутренних страницах.",
                "height": 44
            }
        ]
    },

    # 12. Slide 14: Practice JS: classList API
    {
        "type": "code",
        "category": "УПРАВЛЕНИЕ КЛАССАМИ",
        "question": "Как безопасно манипулировать CSS-классами без повреждения строки className?",
        "left_title": "js/main.js (Свойства объекта classList)",
        "code_lines": [
            [("const ", CODE_KW), ("link = document.querySelector(", CODE_TEXT), ("'.nav-link'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 1. Добавить класс (если его еще нет)", CODE_COMM)],
            [("link.", CODE_TEXT), ("classList.add", CODE_PROP), ("(", CODE_TEXT), ("'active'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 2. Удалить класс (если он присутствует)", CODE_COMM)],
            [("link.", CODE_TEXT), ("classList.remove", CODE_PROP), ("(", CODE_TEXT), ("'active'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 3. Переключить класс (добавить если нет, убрать если есть)", CODE_COMM)],
            [("link.", CODE_TEXT), ("classList.toggle", CODE_PROP), ("(", CODE_TEXT), ("'open'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 4. Проверить наличие класса (вернет true или false)", CODE_COMM)],
            [("if (link.", CODE_TEXT), ("classList.contains", CODE_PROP), ("(", CODE_TEXT), ("'active'", CODE_STR), (")) {", CODE_TEXT)],
            [("  console.log('Ссылка сейчас активна!');", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "CLASSLIST: СОВРЕМЕННЫЙ СТАНДАРТ РАБОТЫ С КЛАССАМИ",
        "bullet_items": [
            {
                "title": "Отказ от устаревшего className:",
                "desc": "Раньше писали elem.className += ' active'. Это часто приводило к слипанию слов 'nav-linkactive' и поломке стилей.",
                "height": 46
            },
            {
                "title": "Умная защита от дубликатов:",
                "desc": "Метод classList.add('active') проверяет список: если класс уже есть у тега, повторно он добавлен не будет.",
                "height": 46
            },
            {
                "title": "Метод classList.toggle():",
                "desc": "Идеален для бургер-меню и аккордеонов: по одному клику включает класс, по повторному — автоматически выключает.",
                "height": 48
            },
            {
                "title": "Предикат classList.contains():",
                "desc": "Позволяет писать понятные логические ветвления if (elem.classList.contains('active')) для проверок состояния интерфейса.",
                "height": 44
            }
        ]
    },

    # 13. Slide 15: Practice CSS: Styling Active Nav Link
    {
        "type": "code",
        "category": "СТИЛИЗАЦИЯ НАВИГАЦИИ",
        "question": "Как оформить активный пункт меню, чтобы он гармонировал с дизайн-системой?",
        "left_title": "css/style.css (.nav-link & .nav-link.active)",
        "code_lines": [
            [("/* Базовый вид пункта меню */", CODE_COMM)],
            [(".nav-link ", CODE_TAG), ("{", CODE_TEXT)],
            [("  text-decoration", CODE_PROP), (": ", CODE_TEXT), ("none", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#222222", CODE_VAL), (";", CODE_TEXT)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("8px 14px", CODE_VAL), (";", CODE_TEXT)],
            [("  border-radius", CODE_PROP), (": ", CODE_TEXT), ("4px", CODE_VAL), (";", CODE_TEXT)],
            [("  font-weight", CODE_PROP), (": ", CODE_TEXT), ("500", CODE_VAL), (";", CODE_TEXT)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("14px", CODE_VAL), (";", CODE_TEXT)],
            [("  transition", CODE_PROP), (": ", CODE_TEXT), ("all 0.2s ease", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Состояние наведения (:hover) и активная страница (.active) */", CODE_COMM)],
            [(".nav-link:hover,", CODE_TAG)],
            [(".nav-link.active ", CODE_TAG), ("{", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), (";           /* Акцентный синий */", CODE_COMM)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#eaf2ff", CODE_VAL), (";  /* Мягкая голубая подложка */", CODE_COMM)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ВИЗУАЛЬНЫЙ СТАТУС ТЕКУЩЕЙ СТРАНИЦЫ",
        "bullet_items": [
            {
                "title": "Сцепленный селектор .nav-link.active:",
                "desc": "Выбирает элемент, у которого одновременно присутствуют оба класса. Имеет более высокий вес специфичности, чем одиночный .nav-link.",
                "height": 48
            },
            {
                "title": "Согласованность с hover-состоянием:",
                "desc": "Объединение :hover и .active через запятую гарантирует, что активный пункт выглядит так же аккуратно, как при наведении курсора.",
                "height": 46
            },
            {
                "title": "Мягкая подложка #eaf2ff:",
                "desc": "Деликатный светлый фон выделяет активную вкладку, подсказывая пользователю его положение на сайте без кричащих контрастов.",
                "height": 46
            },
            {
                "title": "Плавный переход transition:",
                "desc": "Свойство all 0.2s обеспечивает мягкую смену цвета и фона при перемещении курсора между пунктами навигации.",
                "height": 44
            }
        ]
    },

    # 14. Slide 16: Dynamic Element Creation (createElement & appendChild)
    {
        "type": "code",
        "category": "СОЗДАНИЕ ЭЛЕМЕНТОВ",
        "question": "Как создавать новые HTML-теги и вставлять их на страницу из JavaScript?",
        "left_title": "js/main.js (Создание DOM-узлов)",
        "code_lines": [
            [("// 1. Создаем новый тег <div> в памяти", CODE_COMM)],
            [("const ", CODE_KW), ("toast = document.", CODE_TEXT), ("createElement", CODE_PROP), ("(", CODE_TEXT), ("'div'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 2. Настраиваем CSS-классы и параметры", CODE_COMM)],
            [("toast.className = ", CODE_TEXT), ("'toast toast-success'", CODE_STR), (";", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 3. Задаем безопасный текстовый контент", CODE_COMM)],
            [("toast.textContent = ", CODE_TEXT), ("'Бронирование подтверждено!'", CODE_STR), (";", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 4. Добавляем созданный элемент в дерево DOM", CODE_COMM)],
            [("const container = document.getElementById('toastContainer');", CODE_TEXT)],
            [("container.", CODE_TEXT), ("appendChild", CODE_PROP), ("(", CODE_TEXT), ("toast", CODE_VAL), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Элемент мгновенно появляется на экране пользователя!", CODE_COMM)]
        ],
        "right_title": "КОНВЕЙЕР ГЕНЕРАЦИИ ДИНАМИЧЕСКОГО КОНТЕНТА",
        "bullet_items": [
            {
                "title": "document.createElement('tag'):",
                "desc": "Создает полноценный узел DOM в оперативной памяти. До вызова appendChild элемент невидим и не нагружает страницу.",
                "height": 46
            },
            {
                "title": "textContent против innerHTML (Безопасность):",
                "desc": "Свойство textContent вставляет текст как есть, экранируя опасные теги и надежно защищая сайт от атак XSS-инъекций.",
                "height": 48
            },
            {
                "title": "Метод appendChild():",
                "desc": "Вставляет узел в конец списка дочерних элементов родителя. Элемент становится частью DOM и отрисовывается браузером.",
                "height": 46
            },
            {
                "title": "Основа динамических интерфейсов:",
                "desc": "Именно так создаются всплывающие подсказки, модальные окна, корзины заказов и карточки товаров при подгрузке с сервера.",
                "height": 44
            }
        ]
    },

    # 15. Slide 17: Toast Notifications (showNotification)
    {
        "type": "code",
        "category": "СИСТЕМНЫЕ TOAST",
        "question": "Как спроектировать компонент всплывающих уведомлений вместо устаревшего alert?",
        "left_title": "js/main.js (Функция showNotification)",
        "code_lines": [
            [("function ", CODE_KW), ("showNotification(message, type = ", CODE_TEXT), ("'success'", CODE_STR), (") {", CODE_KW)],
            [("  let container = document.getElementById('toastContainer');", CODE_TEXT)],
            [("  if (!container) {", CODE_TEXT)],
            [("    container = document.createElement('div');", CODE_TEXT)],
            [("    container.id = 'toastContainer';", CODE_TEXT)],
            [("    container.className = 'toast-container';", CODE_TEXT)],
            [("    document.body.appendChild(container);", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  const toast = document.createElement('div');", CODE_TEXT)],
            [("  toast.className = 'toast toast-' + type;", CODE_TEXT)],
            [("  toast.textContent = message;", CODE_TEXT)],
            [("  container.appendChild(toast);", CODE_TEXT)],
            [("  // Таймер авто-скрытия через 3.5 секунды...", CODE_COMM)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "АРХИТЕКТУРА СИСТЕМЫ УВЕДОМЛЕНИЙ",
        "bullet_items": [
            {
                "title": "Отказ от блокирующего window.alert():",
                "desc": "Стандартный alert замораживает всю страницу и выглядит архаично. Toast-уведомления всплывают плавно и не мешают пользователю.",
                "height": 48
            },
            {
                "title": "Ленивая инициализация контейнера:",
                "desc": "Контейнер #toastContainer создается автоматически при первом вызове функции. Если уведомлений нет, разметка не засоряется.",
                "height": 46
            },
            {
                "title": "Параметр по умолчанию (type = 'success'):",
                "desc": "Если тип не передан, сообщение автоматически оформляется зеленым цветом успеха. Можно передать 'danger' или 'info'.",
                "height": 46
            },
            {
                "title": "Стек сообщений:",
                "desc": "Если вызвать функцию несколько раз подряд, уведомления аккуратно встанут в вертикальную стопку в правом нижнем углу.",
                "height": 44
            }
        ]
    },

    # 16. Slide 18: Timers and Element Removal (setTimeout & .remove)
    {
        "type": "code",
        "category": "ТАЙМЕРЫ В JAVASCRIPT",
        "question": "Как настроить плавное автоматическое исчезновение уведомления через таймер?",
        "left_title": "js/main.js (setTimeout и toast.remove)",
        "code_lines": [
            [("  container.appendChild(toast);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // 1. Ждем 3.5 секунды (3500 мс)", CODE_COMM)],
            [("  setTimeout(() => {", CODE_KW)],
            [("    // Включаем плавное угасание в CSS", CODE_COMM)],
            [("    toast.style.transition = 'opacity 0.3s ease, transform 0.3s ease';", CODE_TEXT)],
            [("    toast.style.opacity = '0';", CODE_TEXT)],
            [("    toast.style.transform = 'translateY(10px)';", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("    // 2. После угасания (300 мс) удаляем узел из DOM", CODE_COMM)],
            [("    setTimeout(() => {", CODE_KW)],
            [("      toast.", CODE_TEXT), ("remove", CODE_PROP), ("(); // Очистка оперативной памяти", CODE_TEXT)],
            [("    }, 300);", CODE_KW)],
            [("  }, 3500);", CODE_KW)]
        ],
        "right_title": "ДВУХЭТАПНЫЙ ЖИЗНЕННЫЙ ЦИКЛ АНИМАЦИИ",
        "bullet_items": [
            {
                "title": "Функция setTimeout(callback, delay):",
                "desc": "Запускает переданную функцию ровно через указанное количество миллисекунд (3500 мс = 3.5 секунды), не блокируя браузер.",
                "height": 46
            },
            {
                "title": "Этап 1: Плавное угасание (Fade Out):",
                "desc": "Снижение opacity до 0 и сдвиг вниз делают исчезновение мягким и приятным для глаз пользователя.",
                "height": 46
            },
            {
                "title": "Этап 2: Физическое удаление elem.remove():",
                "desc": "Если элемент просто скрыть (display: none), он останется в памяти. Вызов .remove() полностью удаляет его из DOM-дерева.",
                "height": 48
            },
            {
                "title": "Защита от утечек памяти:",
                "desc": "При длительной работе приложения тысячи исчезнувших уведомлений не накапливаются в памяти браузера.",
                "height": 44
            }
        ]
    },

    # 17. Slide 19: Introduction to localStorage
    {
        "type": "code",
        "category": "ХРАНИЛИЩЕ LOCALSTORAGE",
        "question": "Как сохранять состояние пользователя между страницами без базы данных?",
        "left_title": "Консоль DevTools: localStorage API",
        "code_lines": [
            [("// 1. Запись данных в хранилище (Ключ - Значение)", CODE_COMM)],
            [("localStorage.", CODE_TEXT), ("setItem", CODE_PROP), ("(", CODE_TEXT), ("'currentUser'", CODE_STR), (", ", CODE_TEXT), ("'Иван Иванов'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 2. Чтение сохраненных данных", CODE_COMM)],
            [("const ", CODE_KW), ("user = localStorage.", CODE_TEXT), ("getItem", CODE_PROP), ("(", CODE_TEXT), ("'currentUser'", CODE_STR), (");", CODE_TEXT)],
            [("console.log(user); // 'Иван Иванов'", CODE_COMM)],
            [("", CODE_TEXT)],
            [("// 3. Удаление при выходе из системы (Logout)", CODE_COMM)],
            [("localStorage.", CODE_TEXT), ("removeItem", CODE_PROP), ("(", CODE_TEXT), ("'currentUser'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 4. Если данных нет, getItem вернет null", CODE_COMM)],
            [("console.log(localStorage.getItem('currentUser')); // null", CODE_COMM)]
        ],
        "right_title": "КЛИЕНТСКАЯ БАЗА ДАННЫХ В БРАУЗЕРЕ",
        "bullet_items": [
            {
                "title": "Энергонезависимое хранилище (Web Storage):",
                "desc": "Данные в localStorage сохраняются даже при закрытии вкладки или перезагрузке браузера. Объем хранилища — до 5-10 МБ.",
                "height": 48
            },
            {
                "title": "Формат хранения 'Ключ — Значение':",
                "desc": "Все данные сохраняются в виде строк. Для сложных объектов и массивов используется сериализация через JSON.stringify().",
                "height": 46
            },
            {
                "title": "Доступность для всех страниц домена:",
                "desc": "Пользователь залогинился на index.html — и страница catalog.html мгновенно видит его статус авторизации.",
                "height": 46
            },
            {
                "title": "Простой метод очистки removeItem():",
                "desc": "Позволяет одной строчкой реализовать выход из системы (Logout) и стереть временную сессию арендатора.",
                "height": 44
            }
        ]
    },

    # 18. Slide 20: Practice JS: Auth Navigation Updates (updateAuthNav)
    {
        "type": "code",
        "category": "АВТОРИЗАЦИЯ В МЕНЮ",
        "question": "Как динамически переключать кнопку «Войти» на «Выйти» в шапке сайта?",
        "left_title": "js/main.js (Функция updateAuthNav)",
        "code_lines": [
            [("function ", CODE_KW), ("updateAuthNav() {", CODE_TEXT)],
            [("  const currentUser = localStorage.getItem('currentUser');", CODE_TEXT)],
            [("  const myBookings = document.getElementById('myBookingsNavItem');", CODE_TEXT)],
            [("  const authBtn = document.getElementById('authNavBtn');", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  if (currentUser) {", CODE_KW)],
            [("    // Пользователь вошел: показываем брони и кнопку Выйти", CODE_COMM)],
            [("    if (myBookings) myBookings.style.display = 'block';", CODE_TEXT)],
            [("    if (authBtn) {", CODE_TEXT)],
            [("      authBtn.textContent = 'Выйти';", CODE_TEXT)],
            [("      authBtn.onclick = handleLogout; // Навешиваем клик", CODE_TEXT)],
            [("    }", CODE_TEXT)],
            [("  } else {", CODE_KW)],
            [("    // Гость: скрываем брони и предлагаем Войти", CODE_COMM)],
            [("    if (myBookings) myBookings.style.display = 'none';", CODE_TEXT)],
            [("    if (authBtn) authBtn.textContent = 'Войти';", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "УПРАВЛЕНИЕ ИНТЕРФЕЙСОМ НА ОСНОВЕ СОСТОЯНИЯ",
        "bullet_items": [
            {
                "title": "Проверка авторизации if (currentUser):",
                "desc": "Скрипт проверяет наличие ключа в localStorage. Если значение есть — режим авторизованного пользователя, если null — гость.",
                "height": 48
            },
            {
                "title": "Динамический показ пунктов меню:",
                "desc": "Свойство style.display = 'block' открывает вкладку 'Мои бронирования', которая не должна мозолить глаза анонимным гостям.",
                "height": 46
            },
            {
                "title": "Смена надписи textContent:",
                "desc": "Текст на синей кнопке в шапке мгновенно меняется с 'Войти' на 'Выйти' без изменения разметки HTML.",
                "height": 46
            },
            {
                "title": "Повторный вызов при смене страниц:",
                "desc": "Функция вызывается внутри initNavigation(), гарантируя актуальный вид шапки на любой странице сайта.",
                "height": 44
            }
        ]
    },

    # 19. Slide 21: Event Handling: onclick & preventDefault
    {
        "type": "code",
        "category": "ОБРАБОТКА КЛИКОВ",
        "question": "Как перехватить нажатие на ссылку и выполнить свой JS-сценарий выхода?",
        "left_title": "js/main.js (Обработчик выхода Logout)",
        "code_lines": [
            [("authBtn.onclick = (event) => {", CODE_KW)],
            [("  // 1. Отменяем стандартный переход по ссылке href=\"#\"", CODE_COMM)],
            [("  event.", CODE_TEXT), ("preventDefault", CODE_PROP), ("();", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // 2. Стираем сохраненную сессию", CODE_COMM)],
            [("  localStorage.", CODE_TEXT), ("removeItem", CODE_PROP), ("(", CODE_TEXT), ("'currentUser'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // 3. Показываем красивое системное уведомление", CODE_COMM)],
            [("  showNotification(", CODE_TEXT), ("'Вы успешно вышли из системы'", CODE_STR), (", ", CODE_TEXT), ("'info'", CODE_STR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // 4. Через 1 секунду перенаправляем на главную", CODE_COMM)],
            [("  const isPages = window.location.pathname.includes('/pages/');", CODE_TEXT)],
            [("  setTimeout(() => {", CODE_KW)],
            [("    window.location.href = isPages ? '../index.html' : 'index.html';", CODE_TEXT)],
            [("  }, 1000);", CODE_KW)],
            [("};", CODE_KW)]
        ],
        "right_title": "КОНТРОЛЬ ПОВЕДЕНИЯ: PREVENTDEFAULT И РЕДИРЕКТ",
        "bullet_items": [
            {
                "title": "Метод event.preventDefault():",
                "desc": "Критически важная команда: отменяет встроенное действие браузера (переход по ссылке href или отправку формы), передавая власть JS.",
                "height": 48
            },
            {
                "title": "Защита от прыжка страницы (#):",
                "desc": "Без preventDefault клик по ссылке href='#' мгновенно прокрутит страницу в самый верх, дергая экран пользователя.",
                "height": 46
            },
            {
                "title": "Обратная связь для человека:",
                "desc": "Вызов showNotification информирует арендатора об успешном выходе до того, как страница перезагрузится.",
                "height": 46
            },
            {
                "title": "Программный редирект (location.href):",
                "desc": "Присвоение значения window.location.href перенаправляет браузер на главную страницу, учитывая относительную вложенность папок.",
                "height": 44
            }
        ]
    },

    # 20. Slide 22: Student Mistakes (Grid Cards 2x2)
    {
        "type": "grid_cards",
        "category": "РАЗБОР ОШИБОК",
        "question": "На чем чаще всего спотыкаются студенты при первом написании JavaScript?",
        "title_text": "ТОП-4 ТИПИЧНЫХ ОШИБОК СТУДЕНТОВ В ПЕРВЫХ JS-СКРИПТАХ",
        "cards_data": [
            {
                "title": "Забытый атрибут defer у скрипта",
                "problem": "Скрипт в <head> запускается до того, как браузер построил теги <body>. querySelector возвращает null, и код ломается.",
                "solution": "Всегда добавлять атрибут defer в тег <script> и оборачивать старт в событие 'DOMContentLoaded'."
            },
            {
                "title": "Ошибка Cannot read properties of null",
                "problem": "Попытка вызвать .classList.add() у переменной, в которую querySelector ничего не нашел из-за опечатки в селекторе.",
                "solution": "Всегда проверять правильность точки перед классом ('.nav-link') и делать проверку if (elem) перед вызовом методов."
            },
            {
                "title": "Прыжок экрана при клике на ссылку",
                "problem": "При клике на кнопку 'Выйти' с href='#' страница скачет наверх, потому что студент забыл вызвать event.preventDefault().",
                "solution": "Всегда принимать объект event в обработчике клика и первой строчкой вызывать event.preventDefault()."
            },
            {
                "title": "Путаница с путями во вложенных папках",
                "problem": "На странице /pages/catalog.html скрипт пытается перейти на index.html вместо ../index.html, вызывая ошибку 404.",
                "solution": "Проверять window.location.pathname.includes('/pages/') и динамически подставлять '../' при выходе из папки."
            }
        ]
    },

    # 21. Slide 23: Checklist
    {
        "type": "checklist",
        "category": "ЧЕК-ЛИСТ КАЧЕСТВА",
        "question": "Как убедиться, что скрипт умной навигации работает на 100% стабильно?",
        "title_text": "КРИТЕРИИ КАЧЕСТВА ИНТЕРАКТИВНОЙ НАВИГАЦИИ САЙТА",
        "items": [
            {
                "title": "Идеально чистая консоль разработчика (F12):",
                "desc": "При переходе между страницами в консоли браузера нет красных ошибок (null, undefined, 404)."
            },
            {
                "title": "Автоматическая подсветка активного пункта:",
                "desc": "На главной подсвечена 'Главная', в каталоге — 'Каталог', без жесткого прописывания active в HTML."
            },
            {
                "title": "Скрипт подключен через defer:",
                "desc": "Тег <script src='js/main.js' defer> расположен в <head>, страница загружается плавно без задержек."
            },
            {
                "title": "Всплывающие уведомления Toast работают штатно:",
                "desc": "Сообщения появляются в правом нижнем углу и автоматически удаляются из DOM через 3.5 секунды."
            },
            {
                "title": "Корректная реакция на статус авторизации:",
                "desc": "При наличии currentUser в localStorage отображается вкладка броней и активна кнопка 'Выйти'."
            },
            {
                "title": "Безопасная обработка ссылок (preventDefault):",
                "desc": "Клик по сервисным кнопкам не скроллит страницу и корректно перенаправляет пользователя."
            }
        ]
    }
]

print(f"Total content slides to generate: {len(slides_data)}")

orig_slide_ids = list(prs.slides._sldIdLst)
title_slide_id = orig_slide_ids[0]
plan_slide_id = orig_slide_ids[1]
result_slide_id = orig_slide_ids[25]
goodbye_slide_id = orig_slide_ids[26]

# 1. Update Slide 1 (Title slide)
slide1 = prs.slides[0]
for sh in slide1.shapes:
    if "Google Shape;59;p13" in sh.name:
        sh.text_frame.text = "Вебинар 4 "
        p = sh.text_frame.paragraphs[0]
        p.runs[0].font.name = "PT Sans Caption"
        p.runs[0].font.size = Pt(10)
    elif sh.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in sh.shapes:
            if sub.name == "TextBox 7":
                sub.text_frame.text = "ВВЕДЕНИЕ В ПРОГРАММИРОВАНИЕ ИНТЕРФЕЙСОВ"
                sub.text_frame.paragraphs[0].runs[0].font.name = "Montserrat"
                sub.text_frame.paragraphs[0].runs[0].font.size = Pt(15)
                sub.text_frame.paragraphs[0].runs[0].font.bold = True
            elif sub.name == "TextBox 9":
                sub.text_frame.text = "Проектирование \nи разработка интерфейсов пользователя"
print("Slide 1 updated successfully.")

# 2. Update Slide 2 (Plan slide)
# Fix layout footers so "Вебинар 2" becomes "Вебинар 4" without resetting fonts
for lyt in prs.slide_layouts:
    for sh in lyt.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                if "Вебинар" in p.text:
                    for r in p.runs:
                        if r.text.strip() == "2":
                            r.text = "4"

slide2 = prs.slides[1]
plan_titles = {
    "TextBox 12": "1. Базовый синтаксис JS",
    "TextBox 15": "2. Манипуляция деревом DOM",
    "TextBox 18": "3. Обработка событий",
    "TextBox 21": "4. Практика: Переходы",
    "TextBox 24": "5. Практика: Стили элементов",
}
plan_subtitles = {
    "TextBox 13": "Базовый синтаксис JavaScript: переменные, функции и логика",
    "TextBox 16": "Методы манипуляции деревом DOM (селекторы querySelector)",
    "TextBox 19": "Обработка событий: клики, жизненный цикл и DOMContentLoaded",
    "TextBox 22": "Реализация переходов между страницами и сопоставление путей",
    "TextBox 25": "Динамическое управление стилями элементов и активное меню",
}

for sh in slide2.shapes:
    if sh.name in plan_titles:
        sh.text_frame.clear()
        p = sh.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = plan_titles[sh.name]
        r.font.name = "Montserrat SemiBold"
        r.font.size = Pt(11.25)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x18, 0x1A, 0x1F)
    elif sh.name in plan_subtitles:
        sh.width = Pt(450)
        sh.text_frame.clear()
        sh.text_frame.word_wrap = False
        p = sh.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = plan_subtitles[sh.name]
        r.font.name = "Inter"
        r.font.size = Pt(9.75)
        r.font.bold = False
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
print("Slide 2 updated successfully.")

# 3. Create all 21 new content slides
for idx, sdata in enumerate(slides_data):
    s = prs.slides.add_slide(blank_layout)
    stype = sdata["type"]
    
    if stype == "code":
        create_code_explanation_slide(
            s, sdata["category"], sdata["question"],
            sdata["left_title"], sdata["code_lines"],
            sdata["right_title"], sdata["bullet_items"]
        )
    elif stype == "dom_tree":
        create_dom_tree_diagram_slide(s, sdata["category"], sdata["question"])
    elif stype == "lifecycle_timeline":
        create_lifecycle_diagram_slide(s, sdata["category"], sdata["question"])
    elif stype == "grid_cards":
        create_grid_cards_slide(s, sdata["category"], sdata["question"], sdata["title_text"], sdata["cards_data"])
    elif stype == "checklist":
        create_checklist_slide(s, sdata["category"], sdata["question"], sdata["title_text"], sdata["items"])
    
    print(f"  Created Slide {idx+3}: {sdata['category']}")

# 4. Update Result Slide (Slide 24)
slide_res = prs.slides[25]
pic_to_remove = None
tb2_res = None
for sh in slide_res.shapes:
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        pic_to_remove = sh
    elif sh.name == "TextBox 2":
        tb2_res = sh

if pic_to_remove is not None:
    spTree = slide_res.shapes._spTree
    spTree.remove(pic_to_remove._element)

# Add Webinar 4 result screenshot: aspect ratio 1210:670 (1.806), centered nicely
res_w = Pt(510)
res_h = Pt(282)
res_left = Pt(105)
res_top = Pt(68)
slide_res.shapes.add_picture(result_img_path, res_left, res_top, res_w, res_h)

# Add or update TextBox 2 on Result slide
if tb2_res is None:
    tb2_res = slide_res.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2_res.name = "TextBox 2"

tf_res = tb2_res.text_frame
tf_res.word_wrap = True
tf_res.margin_left = tf_res.margin_top = tf_res.margin_right = tf_res.margin_bottom = 0
p_res = tf_res.paragraphs[0]
p_res.alignment = PP_ALIGN.LEFT
p_res.text = ""
r_res = p_res.add_run()
r_res.text = "Результат: Интерактивные элементы управления и логика навигации"
r_res.font.name = "Montserrat"
r_res.font.size = Pt(10)
r_res.font.color.rgb = QUESTION_TEXT
print("Result slide updated with webinar 4 screenshot and TextBox 2.")

# Clean shadows on Goodbye slide (Slide 25)
slide_goodbye = prs.slides[26]
for sh in slide_goodbye.shapes:
    strip_shape_styles_and_shadows(sh)
print("Goodbye slide shadows stripped.")

# 5. Reorder slide list in presentation XML:
# [Title, Plan] + new_slide_ids + [Result, Goodbye]
all_slide_ids = list(prs.slides._sldIdLst)
new_slide_ids = all_slide_ids[len(orig_slide_ids):]
new_order = [title_slide_id, plan_slide_id] + new_slide_ids + [result_slide_id, goodbye_slide_id]
prs.slides._sldIdLst[:] = new_order

print(f"\nFinal presentation slide count: {len(prs.slides)}")

# 6. Save presentation atomically
tmp_output = output_pptx + ".tmp"
prs.save(tmp_output)
if os.path.exists(output_pptx):
    os.remove(output_pptx)
os.replace(tmp_output, output_pptx)

print(f"\nУСПЕХ! Презентация Вебинара 4 сохранена: {output_pptx}")
