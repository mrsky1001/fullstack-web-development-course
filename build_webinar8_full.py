import os
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

# -------------------------------------------------------------
# PATH CONFIGURATION
# -------------------------------------------------------------
base_template_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx"
output_pptx_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-08-filters-and-booking-calc\вебинар 8.pptx"
result_img_path = r"d:\IEEU International East-European University\fullstack-web-development-course\webinar8_result_exact.png"

# -------------------------------------------------------------
# PALETTE & STYLES (Strict corporate college theme)
# -------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x18, 0x1A, 0x1F)
CODE_BORDER = RGBColor(0x2D, 0x31, 0x39)
CODE_TEXT = RGBColor(0xD4, 0xD4, 0xD8)
ORANGE_PILL = RGBColor(0xFE, 0x60, 0x02)
ORANGE_LINE = RGBColor(0xFE, 0x60, 0x02)
QUESTION_TEXT = RGBColor(0xA6, 0xA1, 0xA1) # Strictly #A6A1A1
TEXT_PRIMARY = RGBColor(0x18, 0x1A, 0x1F)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
KEYWORD_COLOR = RGBColor(0xC6, 0x78, 0xDD) # Purple
FUNC_COLOR = RGBColor(0x61, 0xAF, 0xEF)    # Blue
STR_COLOR = RGBColor(0x98, 0xC3, 0x79)     # Green
COMMENT_COLOR = RGBColor(0x5C, 0x63, 0x70) # Gray
NUM_COLOR = RGBColor(0xD1, 0x9A, 0x66)     # Orange / number
TAG_COLOR = RGBColor(0xE0, 0x6C, 0x75)     # Red/Coral
DANGER_COLOR = RGBColor(0xDC, 0x26, 0x26)  # Red for invalid state

def strip_shape_styles_and_shadows(shape):
    """Zero tolerance to shadows and PowerPoint theme styles."""
    spPr = getattr(shape._element, 'spPr', None)
    if spPr is not None:
        for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst'):
            spPr.remove(shd)
        for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw'):
            spPr.remove(shd)
    style = shape._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
    if style is not None:
        shape._element.remove(style)

def add_header_and_footer(slide, category_text, question_text):
    """Standard header pill badge with dynamic width and bottom question prompt."""
    pill_w = Pt(max(140, len(category_text) * 7.8 + 24))
    pill_h = Pt(22.6)
    pill_left = Pt(700) - pill_w
    pill_top = Pt(20.4)
    
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_left, pill_top, pill_w, pill_h)
    badge.name = "Google Shape;69;p14"
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE_PILL
    badge.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    badge.line.width = Pt(1)
    strip_shape_styles_and_shadows(badge)
    
    tf_b = badge.text_frame
    tf_b.word_wrap = False
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = category_text
    r_b.font.name = "Inter"
    r_b.font.size = Pt(12)
    r_b.font.bold = False
    r_b.font.color.rgb = WHITE
    
    # TextBox 2 at bottom
    tb2 = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2.name = "TextBox 2"
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = question_text
    r2.font.name = "Montserrat"
    r2.font.size = Pt(10)
    r2.font.bold = False
    r2.font.color.rgb = QUESTION_TEXT

def create_code_explanation_slide(slide, category, question, left_title, code_lines, right_title, bullet_items):
    """Standard 2-column layout: Left code card, Right clean explanation with bullet items."""
    add_header_and_footer(slide, category, question)
    
    # Left Card
    card_x = Pt(24)
    card_y = Pt(56)
    card_w = Pt(325)
    card_h = Pt(296)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
    card.name = "Rounded Rectangle 1"
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG
    card.line.color.rgb = CODE_BORDER
    card.line.width = Pt(1)
    strip_shape_styles_and_shadows(card)
    
    # Left Card Title
    tb_lt = slide.shapes.add_textbox(card_x + Pt(14), card_y + Pt(10), card_w - Pt(28), Pt(20))
    tb_lt.name = "TextBox Left Title"
    tf_lt = tb_lt.text_frame
    tf_lt.margin_left = tf_lt.margin_top = tf_lt.margin_right = tf_lt.margin_bottom = 0
    p_lt = tf_lt.paragraphs[0]
    r_lt = p_lt.add_run()
    r_lt.text = left_title
    r_lt.font.name = "Montserrat"
    r_lt.font.size = Pt(11)
    r_lt.font.bold = True
    r_lt.font.color.rgb = WHITE
    
    # Inner Code Box
    box_x = card_x + Pt(14)
    box_y = card_y + Pt(32)
    box_w = card_w - Pt(28)
    box_h = card_h - Pt(42)
    
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, box_y, box_w, box_h)
    code_box.name = "Rounded Rectangle Code Box"
    code_box.adjustments[0] = 0.03
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x2A)
    code_box.line.color.rgb = RGBColor(0x36, 0x3B, 0x44)
    code_box.line.width = Pt(1)
    strip_shape_styles_and_shadows(code_box)
    
    # Code Text
    tb_code = slide.shapes.add_textbox(box_x + Pt(10), box_y + Pt(8), box_w - Pt(20), box_h - Pt(16))
    tb_code.name = "TextBox Code Text"
    tf_c = tb_code.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    
    for i, line_runs in enumerate(code_lines):
        p = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
        p.space_after = Pt(2)
        for token_text, token_color in line_runs:
            r = p.add_run()
            r.text = token_text
            r.font.name = "Consolas"
            r.font.size = Pt(8.5)
            r.font.bold = False
            r.font.color.rgb = token_color
            
    # Right Column
    rcol_x = Pt(365)
    rcol_y = Pt(56)
    rcol_w = Pt(335)
    
    tb_rt = slide.shapes.add_textbox(rcol_x, rcol_y, rcol_w, Pt(22))
    tb_rt.name = "TextBox Right Title"
    tf_rt = tb_rt.text_frame
    tf_rt.margin_left = tf_rt.margin_top = tf_rt.margin_right = tf_rt.margin_bottom = 0
    p_rt = tf_rt.paragraphs[0]
    r_rt = p_rt.add_run()
    r_rt.text = right_title
    r_rt.font.name = "Montserrat"
    r_rt.font.size = Pt(13)
    r_rt.font.bold = True
    r_rt.font.color.rgb = TEXT_PRIMARY
    
    # Orange divider
    rdiv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rcol_x, rcol_y + Pt(25), rcol_w, Pt(1))
    rdiv.name = "Rectangle 8"
    rdiv.fill.solid()
    rdiv.fill.fore_color.rgb = ORANGE_LINE
    rdiv.line.fill.background()
    strip_shape_styles_and_shadows(rdiv)
    
    cur_y = rcol_y + Pt(34)
    for idx, item in enumerate(bullet_items):
        item_h = Pt(item.get("height", 42))
        
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rcol_x, cur_y + Pt(3), Pt(7), Pt(7))
        dot.name = f"Rounded Rectangle Dot {idx}"
        dot.adjustments[0] = 0.2
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)
        
        tb_item = slide.shapes.add_textbox(rcol_x + Pt(14), cur_y, rcol_w - Pt(14), item_h)
        tb_item.name = f"TextBox Bullet {idx}"
        tf_i = tb_item.text_frame
        tf_i.word_wrap = True
        tf_i.margin_left = tf_i.margin_top = tf_i.margin_right = tf_i.margin_bottom = 0
        
        p = tf_i.paragraphs[0]
        p.space_after = Pt(0)
        
        r_title = p.add_run()
        r_title.text = item["title"] + ":\n"
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = TEXT_PRIMARY
        
        r_desc = p.add_run()
        r_desc.text = item["desc"]
        r_desc.font.name = "Inter"
        r_desc.font.size = Pt(8.5)
        r_desc.font.bold = False
        r_desc.font.color.rgb = TEXT_MUTED
        
        cur_y += item_h + Pt(6)

def create_filter_pipeline_diagram_slide(slide, category, question):
    """Custom Diagram 1: Filtering & Sorting Pipeline (Raw Data -> Filter Search -> Sort & DOM)."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. ИСХОДНЫЕ ДАННЫЕ", "subtitle": "Массив OFFICE_ROOMS", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. ФИЛЬТРАЦИЯ ПОИСКА", "subtitle": ".filter() + .includes()", "accent": NUM_COLOR},
        {"x": Pt(476), "title": "3. СОРТИРОВКА И DOM", "subtitle": ".sort() + render()", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = c["title"]
        r_t1.font.name = "Montserrat"
        r_t1.font.size = Pt(10)
        r_t1.font.bold = True
        r_t1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = c["subtitle"]
        r_t2.font.name = "Inter"
        r_t2.font.size = Pt(8.5)
        r_t2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = c["accent"]
        div.line.fill.background()
        strip_shape_styles_and_shadows(div)
        
    box_y = top_y + Pt(52)
    box_h = Pt(54)
    
    # Box 1
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(36), box_y, Pt(186), box_h)
    b1.adjustments[0] = 0.08
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(0x1F, 0x24, 0x2E)
    b1.line.color.rgb = FUNC_COLOR
    b1.line.width = Pt(1)
    strip_shape_styles_and_shadows(b1)
    tf_b1 = b1.text_frame
    tf_b1.margin_top = Pt(6)
    p_b1 = tf_b1.paragraphs[0]
    p_b1.alignment = PP_ALIGN.CENTER
    r_b1 = p_b1.add_run()
    r_b1.text = "const rawRooms = [\n  ...OFFICE_ROOMS\n]; // Клон эталона"
    r_b1.font.name = "Consolas"
    r_b1.font.size = Pt(7.5)
    r_b1.font.color.rgb = CODE_TEXT
    
    # Box 2
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(262), box_y, Pt(186), box_h)
    b2.adjustments[0] = 0.08
    b2.fill.solid()
    b2.fill.fore_color.rgb = RGBColor(0x28, 0x24, 0x1E)
    b2.line.color.rgb = NUM_COLOR
    b2.line.width = Pt(1)
    strip_shape_styles_and_shadows(b2)
    tf_b2 = b2.text_frame
    tf_b2.margin_top = Pt(6)
    p_b2 = tf_b2.paragraphs[0]
    p_b2.alignment = PP_ALIGN.CENTER
    r_b2 = p_b2.add_run()
    r_b2.text = "const filtered = rawRooms.filter(\n  r => r.title.toLowerCase()\n       .includes(q)\n);"
    r_b2.font.name = "Consolas"
    r_b2.font.size = Pt(7.5)
    r_b2.font.color.rgb = CODE_TEXT
    
    # Box 3
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(488), box_y, Pt(186), box_h)
    b3.adjustments[0] = 0.08
    b3.fill.solid()
    b3.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x20)
    b3.line.color.rgb = STR_COLOR
    b3.line.width = Pt(1)
    strip_shape_styles_and_shadows(b3)
    tf_b3 = b3.text_frame
    tf_b3.margin_top = Pt(6)
    p_b3 = tf_b3.paragraphs[0]
    p_b3.alignment = PP_ALIGN.CENTER
    r_b3 = p_b3.add_run()
    r_b3.text = "filtered.sort((a, b) => {\n  return a.pricePerHour - b.pricePerHour;\n});\nrender(filtered);"
    r_b3.font.name = "Consolas"
    r_b3.font.size = Pt(7.5)
    r_b3.font.color.rgb = CODE_TEXT
    
    bullets = [
        {"x": Pt(36), "header": "Иммутабельность базы:", "items": [
            "Полный массив комнат импортируется из data.js",
            "Копирование через спред [...] защищает оригинал",
            "База не перезаписывается при поиске и сбросе",
            "Все карточки содержат ID, цены и параметры",
            "Единый источник правды для фильтрации"
        ]},
        {"x": Pt(262), "header": "Алгоритм поиска:", "items": [
            "Чтение строки поискового запроса из input",
            "Приведение к нижнему регистру toLowerCase()",
            "Очистка концевых пробелов методом trim()",
            "Проверка подстроки методом includes()",
            "Мгновенный отсев нерелевантных объектов"
        ]},
        {"x": Pt(488), "header": "Сортировка и рендер:", "items": [
            "Числовой компаратор цен (asc / desc)",
            "Сортировка только отобранных поиском комнат",
            "Генерация HTML-разметки карточек .map()",
            "Обновление контейнера #catalogContainer",
            "Показ Empty State при отсутствии совпадений"
        ]}
    ]
    
    bullets_y = box_y + box_h + Pt(10)
    for b in bullets:
        tb = slide.shapes.add_textbox(b["x"], bullets_y, Pt(186), Pt(165))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_head = tf.paragraphs[0]
        r_head = p_head.add_run()
        r_head.text = b["header"] + "\n"
        r_head.font.name = "Montserrat"
        r_head.font.size = Pt(8.5)
        r_head.font.bold = True
        r_head.font.color.rgb = WHITE
        
        for it in b["items"]:
            p_it = tf.add_paragraph()
            p_it.space_after = Pt(2)
            r_dot = p_it.add_run()
            r_dot.text = "• "
            r_dot.font.name = "Inter"
            r_dot.font.size = Pt(8.0)
            r_dot.font.color.rgb = ORANGE_PILL
            
            r_txt = p_it.add_run()
            r_txt.text = it
            r_txt.font.name = "Inter"
            r_txt.font.size = Pt(8.0)
            r_txt.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD8)

def create_calc_realtime_diagram_slide(slide, category, question):
    """Custom Diagram 2: Real-time Booking Calculator (Select Room -> Input Hours -> Calculate & DOM)."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. ВЫБОР КОМНАТЫ", "subtitle": "Событие change у <select>", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. ВВОД ЧАСОВ", "subtitle": "Событие input у <input>", "accent": NUM_COLOR},
        {"x": Pt(476), "title": "3. ИТОГОВЫЙ РАСЧЕТ", "subtitle": "price * hours -> #totalPrice", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = c["title"]
        r_t1.font.name = "Montserrat"
        r_t1.font.size = Pt(10)
        r_t1.font.bold = True
        r_t1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = c["subtitle"]
        r_t2.font.name = "Inter"
        r_t2.font.size = Pt(8.5)
        r_t2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = c["accent"]
        div.line.fill.background()
        strip_shape_styles_and_shadows(div)
        
    box_y = top_y + Pt(52)
    box_h = Pt(54)
    
    # Box 1
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(36), box_y, Pt(186), box_h)
    b1.adjustments[0] = 0.08
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(0x1F, 0x24, 0x2E)
    b1.line.color.rgb = FUNC_COLOR
    b1.line.width = Pt(1)
    strip_shape_styles_and_shadows(b1)
    tf_b1 = b1.text_frame
    tf_b1.margin_top = Pt(6)
    p_b1 = tf_b1.paragraphs[0]
    p_b1.alignment = PP_ALIGN.CENTER
    r_b1 = p_b1.add_run()
    r_b1.text = "const opt = select.options[select.selectedIndex];\nconst price = Number(\n  opt.dataset.price || 0\n);"
    r_b1.font.name = "Consolas"
    r_b1.font.size = Pt(7.5)
    r_b1.font.color.rgb = CODE_TEXT
    
    # Box 2
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(262), box_y, Pt(186), box_h)
    b2.adjustments[0] = 0.08
    b2.fill.solid()
    b2.fill.fore_color.rgb = RGBColor(0x28, 0x24, 0x1E)
    b2.line.color.rgb = NUM_COLOR
    b2.line.width = Pt(1)
    strip_shape_styles_and_shadows(b2)
    tf_b2 = b2.text_frame
    tf_b2.margin_top = Pt(6)
    p_b2 = tf_b2.paragraphs[0]
    p_b2.alignment = PP_ALIGN.CENTER
    r_b2 = p_b2.add_run()
    r_b2.text = "const raw = hoursInput.value;\nconst hours = Math.max(\n  1, Number(raw || 1)\n);"
    r_b2.font.name = "Consolas"
    r_b2.font.size = Pt(7.5)
    r_b2.font.color.rgb = CODE_TEXT
    
    # Box 3
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(488), box_y, Pt(186), box_h)
    b3.adjustments[0] = 0.08
    b3.fill.solid()
    b3.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x20)
    b3.line.color.rgb = STR_COLOR
    b3.line.width = Pt(1)
    strip_shape_styles_and_shadows(b3)
    tf_b3 = b3.text_frame
    tf_b3.margin_top = Pt(6)
    p_b3 = tf_b3.paragraphs[0]
    p_b3.alignment = PP_ALIGN.CENTER
    r_b3 = p_b3.add_run()
    r_b3.text = "const total = price * hours;\npricePerHourSpan.textContent =\n  price + ' ₽';\ntotalPriceSpan.textContent =\n  total + ' ₽';"
    r_b3.font.name = "Consolas"
    r_b3.font.size = Pt(7.0)
    r_b3.font.color.rgb = CODE_TEXT
    
    bullets = [
        {"x": Pt(36), "header": "Чтение тарифа комнаты:", "items": [
            "Привязка к событию change выпадающего списка",
            "Извлечение цены из атрибута data-price",
            "Преобразование строки в число через Number()",
            "Авто-выбор комнаты из query-параметра URL",
            "Обновление подписи базового почасового тарифа"
        ]},
        {"x": Pt(262), "header": "Обработка времени аренды:", "items": [
            "Событие input обеспечивает мгновенный расчет",
            "Защита от некорректного ввода и пустых значений",
            "Метод Math.max(1, hours) гарантирует минимум 1 ч",
            "Максимальное ограничение до 24 часов",
            "Корректная обработка стрелочек инпута"
        ]},
        {"x": Pt(488), "header": "Синхронизация с экраном:", "items": [
            "Формула расчета: totalPrice = price * hours",
            "Мгновенный вывод суммы без перезагрузки",
            "Подготовка данных для объекта бронирования",
            "Генерация случайного номера заказа при submit",
            "Отправка брони в массив сессии MOCK_BOOKINGS"
        ]}
    ]
    
    bullets_y = box_y + box_h + Pt(10)
    for b in bullets:
        tb = slide.shapes.add_textbox(b["x"], bullets_y, Pt(186), Pt(165))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_head = tf.paragraphs[0]
        r_head = p_head.add_run()
        r_head.text = b["header"] + "\n"
        r_head.font.name = "Montserrat"
        r_head.font.size = Pt(8.5)
        r_head.font.bold = True
        r_head.font.color.rgb = WHITE
        
        for it in b["items"]:
            p_it = tf.add_paragraph()
            p_it.space_after = Pt(2)
            r_dot = p_it.add_run()
            r_dot.text = "• "
            r_dot.font.name = "Inter"
            r_dot.font.size = Pt(8.0)
            r_dot.font.color.rgb = ORANGE_PILL
            
            r_txt = p_it.add_run()
            r_txt.text = it
            r_txt.font.name = "Inter"
            r_txt.font.size = Pt(8.0)
            r_txt.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD8)

# -------------------------------------------------------------
# SLIDE CONTENT DEFINITIONS (21 content slides)
# -------------------------------------------------------------
slides_data = [
    # Slide 3: Module 1
    {
        "type": "code",
        "category": "ФИЛЬТРАЦИЯ ДАННЫХ",
        "question": "Почему метод filter() считается иммутабельным и безопасным для исходных данных?",
        "left_title": "Метод Array.prototype.filter()",
        "code_lines": [
            [("// Фильтрация комнат по максимальной цене", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("budgetRooms = OFFICE_ROOMS.filter(room => {", FUNC_COLOR)],
            [("  return ", KEYWORD_COLOR), ("room.pricePerHour <= ", CODE_TEXT), ("500", NUM_COLOR), (";", CODE_TEXT)],
            [("});", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("console.log(", CODE_TEXT), ("'Найдено комнат:'", STR_COLOR), (", budgetRooms.length);", CODE_TEXT)],
            [("console.log(", CODE_TEXT), ("'Исходный массив:'", STR_COLOR), (", OFFICE_ROOMS.length);", CODE_TEXT)],
            [("// Исходный массив OFFICE_ROOMS не изменился!", COMMENT_COLOR)],
        ],
        "right_title": "Иммутабельный отбор filter()",
        "bullet_items": [
            {"title": "Создание нового массива", "desc": "filter() создает и возвращает новый массив, не изменяя элементы исходного массива.", "height": 38},
            {"title": "Функция-предикат", "desc": "Колбэк возвращает true (элемент включается в выборку) или false (элемент пропускается).", "height": 38},
            {"title": "Сохранение эталона", "desc": "Исходный массив OFFICE_ROOMS остается нетронутым, что позволяет сбрасывать фильтры.", "height": 38},
            {"title": "Основа витрины", "desc": "Метод используется для живого поиска по названию, вместимости и оборудованию.", "height": 38},
        ]
    },

    # Slide 4: Module 1
    {
        "type": "code",
        "category": "ЖИВОЙ ПОИСК",
        "question": "Зачем при текстовом поиске использовать toLowerCase() и trim()?",
        "left_title": "Регистронезависимый поиск",
        "code_lines": [
            [("// Поиск без учета регистра и лишних пробелов", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("filterRoomsBySearch", FUNC_COLOR), ("(rooms, query) {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("q = query.toLowerCase().trim();", CODE_TEXT)],
            [("  return ", KEYWORD_COLOR), ("rooms.filter(room => {", FUNC_COLOR)],
            [("    return ", KEYWORD_COLOR), ("room.title.toLowerCase().includes(q);", CODE_TEXT)],
            [("  });", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Пример: найдет 'Мини-офис Focus'", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("res = filterRoomsBySearch(OFFICE_ROOMS, ", CODE_TEXT), ("'  FOCUS  '", STR_COLOR), (");", CODE_TEXT)],
        ],
        "right_title": "Поиск подстроки через includes()",
        "bullet_items": [
            {"title": "Метод .includes()", "desc": "Проверяет, входит ли искомая строка в название комнаты в любой позиции.", "height": 38},
            {"title": "toLowerCase()", "desc": "Приводит текст и запрос к нижнему регистру: «Фокус» совпадет с «фокус» и «FOCUS».", "height": 38},
            {"title": "Метод .trim()", "desc": "Удаляет случайные пробелы по краям строки, введенные пользователем.", "height": 38},
            {"title": "Реакция на ввод", "desc": "Вызывается мгновенно при наступлении события input в поисковой строке.", "height": 38},
        ]
    },

    # Slide 5: Diagram 1
    {
        "type": "diagram_filter_pipeline",
        "category": "СХЕМА: ФИЛЬТРАЦИЯ И СОРТИРОВКА",
        "question": "Как исходный массив данных трансформируется перед отрисовкой в каталоге?"
    },

    # Slide 6: Module 1
    {
        "type": "code",
        "category": "UX КАТАЛОГА",
        "question": "Почему при разработке поиска критически важно предусматривать Empty State?",
        "left_title": "Обработка пустого результата",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("renderCatalog", FUNC_COLOR), ("(rooms) {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("container = document.getElementById(", CODE_TEXT), ("'catalogContainer'", STR_COLOR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(!rooms.length) {", CODE_TEXT)],
            [("    container.innerHTML = ", CODE_TEXT), ("'<p class=\"empty-msg\">Комнаты не найдены</p>'", STR_COLOR), (";", CODE_TEXT)],
            [("    return", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  container.innerHTML = rooms.map(renderCard).join(", CODE_TEXT), ("''", STR_COLOR), (");", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Состояние «Ничего не найдено»",
        "bullet_items": [
            {"title": "Empty State", "desc": "Специальный блок интерфейса, сообщающий пользователю об отсутствии совпадений.", "height": 38},
            {"title": "Понятная обратная связь", "desc": "Пользователь видит, что поиск успешно отработал, а страница не сломалась.", "height": 38},
            {"title": "Защита от пустоты", "desc": "Исключает появление пустого белого пространства в верстке сетки каталога.", "height": 38},
            {"title": "Сброс запроса", "desc": "При стирании поисковой строки в поле каталог автоматически возвращает все 6 комнат.", "height": 38},
        ]
    },

    # Slide 7: Module 2
    {
        "type": "code",
        "category": "СОРТИРОВКА ДАННЫХ",
        "question": "Почему метод sort() без функции сравнения некорректно сортирует числа?",
        "left_title": "Метод Array.prototype.sort()",
        "code_lines": [
            [("// Массив цен для сортировки", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("prices = [", CODE_TEXT), ("1200", NUM_COLOR), (", ", CODE_TEXT), ("450", NUM_COLOR), (", ", CODE_TEXT), ("900", NUM_COLOR), (", ", CODE_TEXT), ("250", NUM_COLOR), ("];", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Ошибка по умолчанию: лексикографический порядок!", COMMENT_COLOR)],
            [("// prices.sort(); // [1200, 250, 450, 900] -> '1' перед '2'", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("// Правильная числовая сортировка с компаратором:", COMMENT_COLOR)],
            [("prices.sort((a, b) => a - b);", FUNC_COLOR)],
            [("console.log(prices); ", CODE_TEXT), ("// [250, 450, 900, 1200]", COMMENT_COLOR)],
        ],
        "right_title": "Функция-компаратор в sort()",
        "bullet_items": [
            {"title": "Лексикографический дефолт", "desc": "По умолчанию sort() преобразует элементы в строки, поэтому «1200» встает перед «250».", "height": 38},
            {"title": "Правило компаратора", "desc": "Функция (a, b) возвращает отрицательное число, 0 или положительное число.", "height": 38},
            {"title": "Мутирующий метод", "desc": "sort() сортирует элементы прямо на месте, изменяя порядок ячеек исходного массива.", "height": 38},
            {"title": "Безопасное копирование", "desc": "Для защиты эталона сортируют копию массива: [...rooms].sort(...).", "height": 38},
        ]
    },

    # Slide 8: Module 2
    {
        "type": "code",
        "category": "СОРТИРОВКА ЦЕН",
        "question": "Как знак разности (a - b против b - a) меняет порядок элементов в массиве?",
        "left_title": "Сортировка по возрастанию и убыванию",
        "code_lines": [
            [("// 1. По возрастанию (от дешевых к дорогим):", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("sortByPriceAsc", FUNC_COLOR), ("(rooms) {", CODE_TEXT)],
            [("  return ", KEYWORD_COLOR), ("[...rooms].sort((a, b) => {", FUNC_COLOR)],
            [("    return ", KEYWORD_COLOR), ("a.pricePerHour - b.pricePerHour;", CODE_TEXT)],
            [("  });", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 2. По убыванию (от дорогих к дешевым):", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("sortByPriceDesc", FUNC_COLOR), ("(rooms) {", CODE_TEXT)],
            [("  return ", KEYWORD_COLOR), ("[...rooms].sort((a, b) => {", FUNC_COLOR)],
            [("    return ", KEYWORD_COLOR), ("b.pricePerHour - a.pricePerHour;", CODE_TEXT)],
            [("  });", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Направление сортировки цен",
        "bullet_items": [
            {"title": "По возрастанию (Asc)", "desc": "Выражение a.price - b.price ставит меньшую стоимость в начало массива.", "height": 38},
            {"title": "По убыванию (Desc)", "desc": "Выражение b.price - a.price выводит премиальные конференц-залы наверх.", "height": 38},
            {"title": "Спред-оператор [...]", "desc": "Создает поверхностную копию массива перед сортировкой, предотвращая побочные эффекты.", "height": 38},
            {"title": "Интерфейсные кнопки", "desc": "Привязаны к кнопкам «Цена: по возрастанию ↑» и «Цена: по убыванию ↓».", "height": 38},
        ]
    },

    # Slide 9: Module 2
    {
        "type": "code",
        "category": "КОМБИНИРОВАННАЯ ЛОГИКА",
        "question": "В каком порядке логичнее выполнять фильтрацию и сортировку: до или после?",
        "left_title": "Связка поиска и сортировки",
        "code_lines": [
            [("let ", KEYWORD_COLOR), ("displayedRooms = [...OFFICE_ROOMS];", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Фильтрация текущего набора:", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("applyFilter", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("q = searchInput.value.toLowerCase().trim();", CODE_TEXT)],
            [("  displayedRooms = OFFICE_ROOMS.filter(r => ", CODE_TEXT)],
            [("    r.title.toLowerCase().includes(q)", CODE_TEXT)],
            [("  );", CODE_TEXT)],
            [("  render(displayedRooms);", FUNC_COLOR)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Сортировка уже отфильтрованного набора:", COMMENT_COLOR)],
            [("sortAscBtn.addEventListener(", CODE_TEXT), ("'click'", STR_COLOR), (", () => {", CODE_TEXT)],
            [("  displayedRooms.sort((a, b) => a.pricePerHour - b.pricePerHour);", FUNC_COLOR)],
            [("  render(displayedRooms);", FUNC_COLOR)],
            [("});", CODE_TEXT)],
        ],
        "right_title": "Синхронизация фильтра и сортировки",
        "bullet_items": [
            {"title": "Последовательность", "desc": "Сначала фильтруется оригинальный массив, а затем полученный результат сортируется.", "height": 38},
            {"title": "Экономия вычислений", "desc": "Сортировать отфильтрованный маленький массив гораздо быстрее, чем всю базу.", "height": 38},
            {"title": "Сохранение выборки", "desc": "При клике на сортировку найденные поиском карточки не сбрасываются.", "height": 38},
            {"title": "Единая функция render", "desc": "Отображение в DOM делегировано изолированной функции рендера карточек.", "height": 38},
        ]
    },

    # Slide 10: Module 3
    {
        "type": "code",
        "category": "КАЛЬКУЛЯТОР СТОИМОСТИ",
        "question": "Зачем при расчете стоимости использовать приведение типов Number() и Math.max()?",
        "left_title": "Математическая модель расчета",
        "code_lines": [
            [("// Функция расчета стоимости аренды", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("calculateTotal", FUNC_COLOR), ("(pricePerHour, hours) {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("p = Number(pricePerHour) || ", CODE_TEXT), ("0", NUM_COLOR), (";", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("h = Math.max(", CODE_TEXT), ("1", NUM_COLOR), (", Number(hours) || ", CODE_TEXT), ("1", NUM_COLOR), (");", CODE_TEXT)],
            [("  return ", KEYWORD_COLOR), ("p * h;", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Примеры расчета:", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("sum1 = calculateTotal(", CODE_TEXT), ("450", NUM_COLOR), (", ", CODE_TEXT), ("2", NUM_COLOR), ("); ", CODE_TEXT), ("// 900 ₽", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("sum2 = calculateTotal(", CODE_TEXT), ("1200", NUM_COLOR), (", ", CODE_TEXT), ("3", NUM_COLOR), (");", CODE_TEXT), ("// 3600 ₽", COMMENT_COLOR)],
        ],
        "right_title": "Формула стоимости бронирования",
        "bullet_items": [
            {"title": "Базовая формула", "desc": "Итоговая сумма вычисляется простым умножением: total = pricePerHour * hours.", "height": 38},
            {"title": "Приведение Number()", "desc": "Значения из input.value всегда приходят строками ('2' -> 2), требуя преобразования.", "height": 38},
            {"title": "Метод Math.max(1, h)", "desc": "Гарантирует, что количество часов не станет нулевым или отрицательным числом.", "height": 38},
            {"title": "Чистая функция", "desc": "Функция не привязана к разметке и может использоваться в тестах и на сервере.", "height": 38},
        ]
    },

    # Slide 11: Module 3
    {
        "type": "code",
        "category": "СОБЫТИЯ КАЛЬКУЛЯТОРА",
        "question": "В чем разница между событиями input и change при работе с полями форм?",
        "left_title": "Слушатели событий change и input",
        "code_lines": [
            [("const ", KEYWORD_COLOR), ("roomSelect = document.getElementById(", CODE_TEXT), ("'roomSelect'", STR_COLOR), (");", CODE_TEXT)],
            [("const ", KEYWORD_COLOR), ("hoursInput = document.getElementById(", CODE_TEXT), ("'hoursInput'", STR_COLOR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// При смене комнаты в выпадающем списке:", COMMENT_COLOR)],
            [("roomSelect.addEventListener(", CODE_TEXT), ("'change'", STR_COLOR), (", updatePrice);", FUNC_COLOR)],
            [("", CODE_TEXT)],
            [("// При вводе или изменении количества часов:", COMMENT_COLOR)],
            [("hoursInput.addEventListener(", CODE_TEXT), ("'input'", STR_COLOR), (", updatePrice);", FUNC_COLOR)],
            [("", CODE_TEXT)],
            [("// Первичный расчет при загрузке:", COMMENT_COLOR)],
            [("updatePrice();", FUNC_COLOR)],
        ],
        "right_title": "События change vs input",
        "bullet_items": [
            {"title": "Событие change", "desc": "Генерируется выпадающим списком <select> при выборе пользователем другого пункта.", "height": 38},
            {"title": "Событие input", "desc": "Срабатывает мгновенно при каждом нажатии клавиши или клике по стрелкам числа.", "height": 38},
            {"title": "Реальное время (Realtime)", "desc": "Сумма в блоке #totalPrice пересчитывается без малейших задержек и подвисаний.", "height": 38},
            {"title": "Первичный запуск", "desc": "Вызов updatePrice() на старте сразу рассчитывает стоимость для дефолтных значений.", "height": 38},
        ]
    },

    # Slide 12: Diagram 2
    {
        "type": "diagram_calc_realtime",
        "category": "СХЕМА: КАЛЬКУЛЯТОР БРОНИ",
        "question": "Как изменение полей ввода мгновенно пересчитывает итоговую сумму бронирования?"
    },

    # Slide 13: Module 3
    {
        "type": "code",
        "category": "DATA-АТРИБУТЫ",
        "question": "В чем преимущество хранения тарифа комнаты в атрибуте data-price?",
        "left_title": "Хранение цены в data-атрибутах",
        "code_lines": [
            [("<!-- Генерация options с ценой в data-атрибуте -->", COMMENT_COLOR)],
            [("<option value=\"focus-1\" data-price=\"450\">", TAG_COLOR)],
            [("  Мини-офис Focus (450 ₽/час)", CODE_TEXT)],
            [("</option>", TAG_COLOR)],
            [("", CODE_TEXT)],
            [("// Чтение значения в JavaScript через dataset:", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("selectedOpt = roomSelect.options[roomSelect.selectedIndex];", CODE_TEXT)],
            [("const ", KEYWORD_COLOR), ("price = Number(selectedOpt.dataset.price || ", CODE_TEXT), ("0", NUM_COLOR), (");", CODE_TEXT)],
            [("console.log(", CODE_TEXT), ("'Тариф выбранной комнаты:'", STR_COLOR), (", price);", CODE_TEXT)],
        ],
        "right_title": "Метаданные в dataset",
        "bullet_items": [
            {"title": "Атрибуты data-*", "desc": "Стандартный способ HTML5 привязать любые вспомогательные данные к DOM-элементу.", "height": 38},
            {"title": "Объект dataset", "desc": "Браузер автоматически собирает атрибуты data-price в свойство element.dataset.price.", "height": 38},
            {"title": "Мгновенное чтение", "desc": "Не нужно заново сканировать массив и искать объект комнаты по идентификатору.", "height": 38},
            {"title": "Целостность данных", "desc": "Тариф жестко привязан к конкретному элементу выпадающего списка.", "height": 38},
        ]
    },

    # Slide 14: Module 4
    {
        "type": "code",
        "category": "URL ПАРАМЕТРЫ",
        "question": "Как URLSearchParams упрощает извлечение параметров из адресной строки?",
        "left_title": "Парсинг URLSearchParams",
        "code_lines": [
            [("// URL страницы: booking.html?room=focus-1", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("searchString = window.location.search; ", CODE_TEXT), ("// '?room=focus-1'", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("urlParams = new URLSearchParams(searchString);", FUNC_COLOR)],
            [("", CODE_TEXT)],
            [("// Извлечение значения параметра:", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("roomId = urlParams.get(", CODE_TEXT), ("'room'", STR_COLOR), (");", CODE_TEXT)],
            [("console.log(", CODE_TEXT), ("'ID комнаты из URL:'", STR_COLOR), (", roomId); ", CODE_TEXT), ("// 'focus-1'", COMMENT_COLOR)],
        ],
        "right_title": "Интерфейс URLSearchParams",
        "bullet_items": [
            {"title": "Параметры запроса (Query)", "desc": "Данные передаются в строке адреса после символа ? в формате ключ=значение.", "height": 38},
            {"title": "Встроенный класс", "desc": "URLSearchParams избавляет от написания сложных и хрупких регулярных выражений.", "height": 38},
            {"title": "Метод .get()", "desc": "Возвращает строковое значение указанного параметра или null, если параметра нет.", "height": 38},
            {"title": "Связка компонентов", "desc": "Позволяет странице бронирования узнать, какую именно комнату выбрал пользователь.", "height": 38},
        ]
    },

    # Slide 15: Module 4
    {
        "type": "code",
        "category": "АВТОПОДСТАНОВКА URL",
        "question": "Каким образом связываются кнопка бронирования в каталоге и форма оформления заказа?",
        "left_title": "Автовыбор комнаты из ссылки",
        "code_lines": [
            [("// В каталоге ссылка на кнопке формируется так:", COMMENT_COLOR)],
            [("<a href=\"booking.html?room=${room.id}\" class=\"btn btn-primary\">", TAG_COLOR)],
            [("", CODE_TEXT)],
            [("// На странице booking.html считываем параметр:", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("urlParams = new URLSearchParams(window.location.search);", CODE_TEXT)],
            [("const ", KEYWORD_COLOR), ("roomId = urlParams.get(", CODE_TEXT), ("'room'", STR_COLOR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("if ", KEYWORD_COLOR), ("(roomId && roomSelect) {", CODE_TEXT)],
            [("  roomSelect.value = roomId; ", CODE_TEXT), ("// Автовыбор пункта в <select>", COMMENT_COLOR)],
            [("  updatePrice();             ", FUNC_COLOR), ("// Мгновенный расчет тарифа", COMMENT_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Бесшовный переход в бронирование",
        "bullet_items": [
            {"title": "Динамические ссылки", "desc": "Кнопки в каталоге генерируются с параметром конкретной комнаты: ?room=focus-1.", "height": 38},
            {"title": "Автоподстановка в select", "desc": "Свойство roomSelect.value = roomId автоматически выбирает нужный пункт в форме.", "height": 38},
            {"title": "Инициализация цены", "desc": "Вызов updatePrice() сразу показывает тариф 450 ₽/час и сумму за 2 часа (900 ₽).", "height": 38},
            {"title": "Комфорт клиента", "desc": "Пользователю не нужно повторно искать выбранный офис в длинном выпадающем списке.", "height": 38},
        ]
    },

    # Slide 16: Module 4
    {
        "type": "code",
        "category": "ВАЛИДАЦИЯ ФОРМЫ",
        "question": "Зачем дублировать ограничения min и max в JavaScript-коде калькулятора?",
        "left_title": "Контроль ввода количества часов",
        "code_lines": [
            [("<!-- Атрибуты ограничений в HTML -->", COMMENT_COLOR)],
            [("<input type=\"number\" id=\"hoursInput\" min=\"1\" max=\"24\" value=\"2\" required>", TAG_COLOR)],
            [("", CODE_TEXT)],
            [("// Дополнительный контроль в обработчике JS:", COMMENT_COLOR)],
            [("hoursInput.addEventListener(", CODE_TEXT), ("'input'", STR_COLOR), (", () => {", CODE_TEXT)],
            [("  let ", KEYWORD_COLOR), ("val = parseInt(hoursInput.value, ", CODE_TEXT), ("10", NUM_COLOR), (");", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(val < ", CODE_TEXT), ("1", NUM_COLOR), (") hoursInput.value = ", CODE_TEXT), ("1", NUM_COLOR), (";", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(val > ", CODE_TEXT), ("24", NUM_COLOR), (") hoursInput.value = ", CODE_TEXT), ("24", NUM_COLOR), (";", CODE_TEXT)],
            [("  updatePrice();", FUNC_COLOR)],
            [("});", CODE_TEXT)],
        ],
        "right_title": "Защита границ калькулятора",
        "bullet_items": [
            {"title": "HTML-атрибуты min/max", "desc": "Задают диапазон значений для встроенных стрелочек браузера от 1 до 24 часов.", "height": 38},
            {"title": "Ручной ввод с клавиатуры", "desc": "Пользователь может вручную напечатать '-5' или '99', обойдя ограничения стрелок.", "height": 38},
            {"title": "JS-санитайзер", "desc": "Скрипт проверяет введенное число и принудительно возвращает его в диапазон 1–24.", "height": 38},
            {"title": "Корректность бизнес-логики", "desc": "Исключает создание броней с нулевой стоимостью или на сотни часов аренды.", "height": 38},
        ]
    },

    # Slide 17: Module 5
    {
        "type": "code",
        "category": "ПРАКТИКА: КАТАЛОГ",
        "question": "Почему фильтрацию всегда нужно применять к исходному массиву, а не к отфильтрованному?",
        "left_title": "Функция initCatalogFilters()",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("initCatalogFilters", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("container = document.getElementById(", CODE_TEXT), ("'catalogContainer'", STR_COLOR), (");", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("searchInput = document.getElementById(", CODE_TEXT), ("'searchInput'", STR_COLOR), (");", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(!container || typeof OFFICE_ROOMS === 'undefined') return;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  let ", KEYWORD_COLOR), ("displayedRooms = [...OFFICE_ROOMS];", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  function ", KEYWORD_COLOR), ("applyFilter", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("    const ", KEYWORD_COLOR), ("q = (searchInput ? searchInput.value : '').toLowerCase().trim();", CODE_TEXT)],
            [("    displayedRooms = OFFICE_ROOMS.filter(r => r.title.toLowerCase().includes(q));", FUNC_COLOR)],
            [("    render(displayedRooms);", FUNC_COLOR)],
            [("  }", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Управление фильтрацией каталога",
        "bullet_items": [
            {"title": "Инициализация модуля", "desc": "Проверяет наличие контейнера на странице, защищая другие страницы от ошибок.", "height": 38},
            {"title": "Локальное состояние", "desc": "Переменная displayedRooms хранит актуальный срез данных, видимый пользователю.", "height": 38},
            {"title": "Фильтрация от эталона", "desc": "Поиск всегда фильтрует оригинальный OFFICE_ROOMS, возвращая карточки при стирании букв.", "height": 38},
            {"title": "Слушатель input", "desc": "searchInput.addEventListener('input', applyFilter) дает отклик без нажатия Enter.", "height": 38},
        ]
    },

    # Slide 18: Module 5
    {
        "type": "code",
        "category": "ПРАКТИКА: СОРТИРОВКА",
        "question": "Как обеспечить сохранение результатов поиска при переключении сортировки?",
        "left_title": "Обработчики кнопок сортировки",
        "code_lines": [
            [("const ", KEYWORD_COLOR), ("sortAscBtn = document.getElementById(", CODE_TEXT), ("'sortAsc'", STR_COLOR), (");", CODE_TEXT)],
            [("const ", KEYWORD_COLOR), ("sortDescBtn = document.getElementById(", CODE_TEXT), ("'sortDesc'", STR_COLOR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("if ", KEYWORD_COLOR), ("(sortAscBtn) {", CODE_TEXT)],
            [("  sortAscBtn.addEventListener(", CODE_TEXT), ("'click'", STR_COLOR), (", () => {", CODE_TEXT)],
            [("    displayedRooms.sort((a, b) => a.pricePerHour - b.pricePerHour);", FUNC_COLOR)],
            [("    render(displayedRooms);", FUNC_COLOR)],
            [("  });", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("if ", KEYWORD_COLOR), ("(sortDescBtn) {", CODE_TEXT)],
            [("  sortDescBtn.addEventListener(", CODE_TEXT), ("'click'", STR_COLOR), (", () => {", CODE_TEXT)],
            [("    displayedRooms.sort((a, b) => b.pricePerHour - a.pricePerHour);", FUNC_COLOR)],
            [("    render(displayedRooms);", FUNC_COLOR)],
            [("  });", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Интерактивная сортировка карточек",
        "bullet_items": [
            {"title": "Сортировка выборки", "desc": "Сортируется массив displayedRooms: если применен поиск, порядок меняется внутри найденных.", "height": 38},
            {"title": "Числовой порядок", "desc": "От 250 ₽ до 1200 ₽ (по возрастанию) или от 1200 ₽ до 250 ₽ (по убыванию).", "height": 38},
            {"title": "Мгновенный рендер", "desc": "render(displayedRooms) обновляет DOM-дерево за миллисекунды без мигания экрана.", "height": 38},
            {"title": "Готовность карточек", "desc": "Каждая отрисованная карточка сохраняет корректную ссылку с query-параметром.", "height": 38},
        ]
    },

    # Slide 19: Module 5
    {
        "type": "code",
        "category": "ПРАКТИКА: КАЛЬКУЛЯТОР",
        "question": "Какие ключевые этапы включает в себя инициализация калькулятора бронирования?",
        "left_title": "Функция initBookingCalc()",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("initBookingCalc", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  // Проверка авторизации:", COMMENT_COLOR)],
            [("  const ", KEYWORD_COLOR), ("currentUser = localStorage.getItem(", CODE_TEXT), ("'currentUser'", STR_COLOR), (");", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(!currentUser) { window.location.href = ", CODE_TEXT), ("'login.html'", STR_COLOR), ("; return; }", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // Генерация <option> из OFFICE_ROOMS:", COMMENT_COLOR)],
            [("  roomSelect.innerHTML = OFFICE_ROOMS.map(r => `", STR_COLOR)],
            [("    <option value=\"${r.id}\" data-price=\"${r.pricePerHour}\">", STR_COLOR)],
            [("      ${r.title} (${r.pricePerHour} ₽/час)", STR_COLOR)],
            [("    </option>", STR_COLOR)],
            [("  `).join('');", STR_COLOR)],
            [("", CODE_TEXT)],
            [("  // Чтение URL и привязка слушателей:", COMMENT_COLOR)],
            [("  syncURLParam();", FUNC_COLOR)],
            [("  roomSelect.addEventListener('change', updatePrice);", FUNC_COLOR)],
            [("  hoursInput.addEventListener('input', updatePrice);", FUNC_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Сборка клиентского калькулятора",
        "bullet_items": [
            {"title": "Защита авторизацией", "desc": "Если пользователь не залогинен, происходит автоматический редирект на login.html.", "height": 38},
            {"title": "Динамический select", "desc": "Список комнат генерируется из единого массива данных: изменения в data.js сразу видны.", "height": 38},
            {"title": "Синхронизация параметров", "desc": "При переходе по ссылке ?room=focus-1 форма автоматически выбирает нужный офис.", "height": 38},
            {"title": "Связка слушателей", "desc": "Любое действие в форме мгновенно запускает функцию расчета стоимости updatePrice().", "height": 38},
        ]
    },

    # Slide 20: Module 5
    {
        "type": "code",
        "category": "ОФОРМЛЕНИЕ ЗАКАЗА",
        "question": "В чем преимущество метода unshift() перед push() при сохранении новых заказов?",
        "left_title": "Обработка отправки формы бронирования",
        "code_lines": [
            [("form.addEventListener(", CODE_TEXT), ("'submit'", STR_COLOR), (", (e) => {", CODE_TEXT)],
            [("  e.preventDefault();", FUNC_COLOR)],
            [("  const ", KEYWORD_COLOR), ("appNumber = Math.floor(", CODE_TEXT), ("10000", NUM_COLOR), (" + Math.random() * ", CODE_TEXT), ("90000", NUM_COLOR), (");", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("selectedRoom = OFFICE_ROOMS.find(r => r.id === roomSelect.value);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  // Добавление в массив сессии:", COMMENT_COLOR)],
            [("  MOCK_BOOKINGS.unshift({", FUNC_COLOR)],
            [("    id: String(appNumber),", CODE_TEXT)],
            [("    roomTitle: selectedRoom ? selectedRoom.title : 'Офис',", STR_COLOR)],
            [("    date: document.getElementById('bookingDate').value,", CODE_TEXT)],
            [("    hours: Math.max(1, Number(hoursInput.value || 1)),", CODE_TEXT)],
            [("    totalPrice: selectedRoom.pricePerHour * hours", CODE_TEXT)],
            [("  });", CODE_TEXT)],
            [("  showNotification('Бронь создана! Заявка №' + appNumber, 'success');", FUNC_COLOR)],
            [("});", CODE_TEXT)],
        ],
        "right_title": "Создание и регистрация заявки",
        "bullet_items": [
            {"title": "Перехват submit", "desc": "Отменяет перезагрузку страницы браузером и передает управление в JavaScript.", "height": 38},
            {"title": "Номер заявки", "desc": "Генерация 5-значного номера через Math.random() создает ощущение реального бэкенда.", "height": 38},
            {"title": "Метод unshift()", "desc": "Добавляет новую бронь в начало массива, чтобы свежий заказ сразу отобразился сверху.", "height": 38},
            {"title": "Всплывающий Toast", "desc": "Зеленое уведомление подтверждает успех и переводит клиента в личный кабинет.", "height": 38},
        ]
    },

    # Slide 21: Module 5
    {
        "type": "code",
        "category": "НАДЕЖНОСТЬ БИЗНЕС-ЛОГИКИ",
        "question": "Какие краевые случаи необходимо учитывать при создании клиентского калькулятора?",
        "left_title": "Защита от Edge Cases",
        "code_lines": [
            [("// Защита от отсутствующих или поврежденных данных", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("updatePrice", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("selectedOption = roomSelect.options[roomSelect.selectedIndex];", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("price = selectedOption ? Number(selectedOption.dataset.price || ", CODE_TEXT), ("0", NUM_COLOR), (") : ", CODE_TEXT), ("0", NUM_COLOR), (";", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("hours = Math.max(", CODE_TEXT), ("1", NUM_COLOR), (", Number(hoursInput.value || ", CODE_TEXT), ("1", NUM_COLOR), ("));", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("total = price * hours;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(pricePerHourSpan) pricePerHourSpan.textContent = price + ", CODE_TEXT), ("' ₽'", STR_COLOR), (";", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(totalPriceSpan) totalPriceSpan.textContent = total + ", CODE_TEXT), ("' ₽'", STR_COLOR), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Отказоустойчивость расчетов",
        "bullet_items": [
            {"title": "Защита от undefined", "desc": "Проверка selectedOption предотвращает ошибку TypeError при пустом списке комнат.", "height": 38},
            {"title": "Защита от NaN", "desc": "Логический оператор || 0 и || 1 спасает от сложения пустых строк и появления NaN.", "height": 38},
            {"title": "Проверка элементов", "desc": "Условия if (pricePerHourSpan) гарантируют работу кода на любых страницах сайта.", "height": 38},
            {"title": "Стабильность приложения", "desc": "Приложение работает надежно при любых действиях и случайных ошибках пользователя.", "height": 38},
        ]
    },

    # Slide 22: Module 5
    {
        "type": "code",
        "category": "ПРОИЗВОДИТЕЛЬНОСТЬ",
        "question": "Как паттерн Debounce защищает интерфейс от лишних вычислений при быстром наборе текста?",
        "left_title": "Оптимизация поиска через Debounce",
        "code_lines": [
            [("// Задержка вызова функции при быстром наборе", COMMENT_COLOR)],
            [("function ", KEYWORD_COLOR), ("debounce", FUNC_COLOR), ("(callback, delay = ", CODE_TEXT), ("250", NUM_COLOR), (") {", CODE_TEXT)],
            [("  let ", KEYWORD_COLOR), ("timeoutId = ", CODE_TEXT), ("null", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("  return ", KEYWORD_COLOR), ("function ", KEYWORD_COLOR), ("(...args) {", CODE_TEXT)],
            [("    clearTimeout(timeoutId);", FUNC_COLOR)],
            [("    timeoutId = setTimeout(() => {", KEYWORD_COLOR)],
            [("      callback.apply(this, args);", FUNC_COLOR)],
            [("    }, delay);", CODE_TEXT)],
            [("  };", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Применение к поиску в каталоге:", COMMENT_COLOR)],
            [("searchInput.addEventListener(", CODE_TEXT), ("'input'", STR_COLOR), (", debounce(applyFilter, ", CODE_TEXT), ("250", NUM_COLOR), ("));", CODE_TEXT)],
        ],
        "right_title": "Паттерн Debounce для поиска",
        "bullet_items": [
            {"title": "Проблема частых вызовов", "desc": "При быстром наборе 10 букв без debounce фильтрация и перерисовка сработают 10 раз.", "height": 38},
            {"title": "Пауза в 250 мс", "desc": "Debounce ждет паузы в печати пользователя и выполняет поиск ровно 1 раз.", "height": 38},
            {"title": "Экономия ресурсов", "desc": "Снижает нагрузку на процессор браузера и устраняет микрофризы на слабых ноутбуках.", "height": 38},
            {"title": "Промышленный стандарт", "desc": "Обязательный паттерн для всех поисковых строк в коммерческой веб-разработке.", "height": 38},
        ]
    },

    # Slide 23: Module 5
    {
        "type": "code",
        "category": "ИТОГИ ВЕБИНАРА",
        "question": "Какие ключевые паттерны обеспечивают надежность клиентской бизнес-логики?",
        "left_title": "Итоги клиентской бизнес-логики",
        "code_lines": [
            [("// Архитектурный стек вебинара 8:", COMMENT_COLOR)],
            [("// 1. Поиск: .filter() + .toLowerCase() + .includes()", COMMENT_COLOR)],
            [("// 2. Сортировка: .sort((a, b) => a.price - b.price)", COMMENT_COLOR)],
            [("// 3. Калькулятор: price * hours в реальном времени", COMMENT_COLOR)],
            [("// 4. Связка: URLSearchParams(?room=id)", COMMENT_COLOR)],
            [("// 5. Оформление: генерация заявки + MOCK_BOOKINGS", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("console.log(", CODE_TEXT), ("'Бизнес-логика СмартОфис готова!'", STR_COLOR), (");", CODE_TEXT)],
        ],
        "right_title": "Чек-лист качества бизнес-логики",
        "bullet_items": [
            {"title": "Иммутабельность", "desc": "Оригинальные данные в data.js защищены от случайных изменений при сортировке.", "height": 38},
            {"title": "Живой расчет", "desc": "Калькулятор аренды мгновенно реагирует на изменение комнат и часов без перезагрузки.", "height": 38},
            {"title": "Связность сценариев", "desc": "Клиент выбирает комнату в каталоге и бесшовно переходит к ее бронированию.", "height": 38},
            {"title": "Результат достигнут", "desc": "Клиентская часть портала «СмартОфис» получила полноценную интерактивную логику.", "height": 38},
        ]
    }
]

# -------------------------------------------------------------
# PRESENTATION GENERATOR EXECUTION
# -------------------------------------------------------------
print(f"Total content slides to generate: {len(slides_data)}")

prs = Presentation(base_template_path)
blank_layout = prs.slide_layouts[6]

orig_slide_ids = list(prs.slides._sldIdLst)
title_slide_id = orig_slide_ids[0]
plan_slide_id = orig_slide_ids[1]
result_slide_id = orig_slide_ids[25]
goodbye_slide_id = orig_slide_ids[26]

# 0. Update Layouts footer to Вебинар 8
for layout in prs.slide_layouts:
    for s in layout.shapes:
        if "Google Shape;59;p13" in s.name and s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip() == "2":
                        r.text = "8"

# 1. Update Slide 1 (Title slide)
slide1 = prs.slides[0]
for sh in slide1.shapes:
    if sh.name == "Google Shape;59;p13" and sh.has_text_frame:
        sh.text_frame.text = "Вебинар 8 "
    elif sh.shape_type == MSO_SHAPE_TYPE.GROUP and sh.name == "Группа 12":
        for sub in sh.shapes:
            if sub.name == "TextBox 7":
                sub.text_frame.text = "ИНТЕРАКТИВНАЯ ЛОГИКА И РАСЧЕТЫ"
                sub.text_frame.paragraphs[0].runs[0].font.name = "Montserrat"
                sub.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
                sub.text_frame.paragraphs[0].runs[0].font.bold = True
            elif sub.name == "TextBox 9":
                sub.text_frame.text = "Проектирование \nи разработка интерфейсов пользователя"
print("Slide 1 updated successfully.")

# 2. Update Slide 2 (Plan slide)
slide2 = prs.slides[1]
plan_titles = {
    "TextBox 12": "1. Фильтрация массивов",
    "TextBox 15": "2. Сортировка данных",
    "TextBox 18": "3. Динамические расчеты",
    "TextBox 21": "4. Параметры URL",
    "TextBox 24": "5. Практика: Бизнес-логика"
}
plan_subtitles = {
    "TextBox 13": "Метод filter() и поиск подстроки через includes() без регистра",
    "TextBox 16": "Метод sort(), функция сравнения компаратора и порядок цен",
    "TextBox 19": "Калькулятор стоимости, обработка событий input и пересчет sum",
    "TextBox 22": "Чтение query-параметров через URLSearchParams для автовыбора",
    "TextBox 25": "Интерактивный каталог и оформление бронирования в СмартОфис"
}

for sh in slide2.shapes:
    strip_shape_styles_and_shadows(sh)
    if sh.name in plan_titles and sh.has_text_frame:
        p = sh.text_frame.paragraphs[0]
        p.text = plan_titles[sh.name]
        p.runs[0].font.name = "Montserrat"
        p.runs[0].font.size = Pt(11.25)
        p.runs[0].font.bold = True
    elif sh.name in plan_subtitles and sh.has_text_frame:
        sh.width = Pt(450)
        tf = sh.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = plan_subtitles[sh.name]
        p.runs[0].font.name = "Inter"
        p.runs[0].font.size = Pt(9.75)
        p.runs[0].font.bold = False

print("Slide 2 updated successfully.")

# 3. Create 21 Content Slides
content_slide_elements = []
for idx, data in enumerate(slides_data, 1):
    new_slide = prs.slides.add_slide(blank_layout)
    content_slide_elements.append(new_slide._element)
    
    stype = data.get("type")
    if stype == "code":
        create_code_explanation_slide(
            new_slide,
            data["category"],
            data["question"],
            data["left_title"],
            data["code_lines"],
            data["right_title"],
            data["bullet_items"]
        )
    elif stype == "diagram_filter_pipeline":
        create_filter_pipeline_diagram_slide(
            new_slide,
            data["category"],
            data["question"]
        )
    elif stype == "diagram_calc_realtime":
        create_calc_realtime_diagram_slide(
            new_slide,
            data["category"],
            data["question"]
        )
    print(f"  Created Slide {idx + 2}: {data['category']}")

# 4. Update Result Slide (Slide 26 in original template, prs.slides[25])
slide_res = prs.slides[25]
pic_to_remove = None
tb2_res = None
for sh in slide_res.shapes:
    strip_shape_styles_and_shadows(sh)
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        pic_to_remove = sh
    elif sh.name == "TextBox 2":
        tb2_res = sh

if pic_to_remove is not None:
    spTree = slide_res.shapes._spTree
    spTree.remove(pic_to_remove._element)

# Add Webinar 8 result screenshot: aspect ratio 1210:670 (1.806), centered nicely
res_w = Pt(510)
res_h = Pt(282)
res_left = Pt(105)
res_top = Pt(68)
slide_res.shapes.add_picture(result_img_path, res_left, res_top, res_w, res_h)

# Add or update TextBox 2 on Result slide
if tb2_res is None:
    tb2_res = slide_res.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2_res.name = "TextBox 2"

tf_res = tb2_res.text_frame
tf_res.word_wrap = True
tf_res.margin_left = tf_res.margin_top = tf_res.margin_right = tf_res.margin_bottom = 0
p_res = tf_res.paragraphs[0]
p_res.alignment = PP_ALIGN.LEFT
p_res.text = ""
r_res = p_res.add_run()
r_res.text = "Результат: Бизнес-логика клиентской части портала"
r_res.font.name = "Montserrat"
r_res.font.size = Pt(10)
r_res.font.color.rgb = QUESTION_TEXT
print("Result slide updated with webinar 8 catalog screenshot and TextBox 2.")

# 5. Clean shadows on Goodbye slide (Slide 27 in original template, prs.slides[26])
slide_goodbye = prs.slides[26]
for sh in slide_goodbye.shapes:
    strip_shape_styles_and_shadows(sh)
print("Goodbye slide shadows stripped.")

# 6. Reorder slide list in presentation XML:
# [Title, Plan] + new_slide_ids + [Result, Goodbye]
all_slide_ids = list(prs.slides._sldIdLst)
new_slide_ids = all_slide_ids[len(orig_slide_ids):]
new_order = [title_slide_id, plan_slide_id] + new_slide_ids + [result_slide_id, goodbye_slide_id]
prs.slides._sldIdLst[:] = new_order

# 7. Save presentation
os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
prs.save(output_pptx_path)

print(f"\nFinal presentation slide count: {len(prs.slides)}")
print(f"\nУСПЕХ! Презентация Вебинара 8 сохранена: {output_pptx_path}")
