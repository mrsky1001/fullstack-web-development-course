import sys
import os
import shutil
import pptx
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

sys.stdout.reconfigure(encoding='utf-8')

base_pptx = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx"
output_dir = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-03-flexbox-grid-cards"
output_pptx = os.path.join(output_dir, "вебинар 3.pptx")
result_img_path = r"d:\IEEU International East-European University\fullstack-web-development-course\webinar3_result_exact.png"

prs = Presentation(base_pptx)
blank_layout = prs.slide_layouts[12]

# Brand Design Tokens & Colors
ORANGE_PILL   = RGBColor(0xFE, 0x60, 0x02)  # #FE6002
ORANGE_LINE   = RGBColor(0xFF, 0x6F, 0x03)  # #FF6F03
BORDER_GRAY   = RGBColor(0xEE, 0xEE, 0xEE)  # #EEEEEE
QUESTION_TEXT = RGBColor(0xA6, 0xA1, 0xA1)  # #A6A1A1 - TextBox 2
DARK_TEXT     = RGBColor(0x22, 0x22, 0x22)  # #222222
BODY_TEXT     = RGBColor(0x44, 0x44, 0x44)  # #444444
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)  # #FFFFFF
GREEN_BADGE   = RGBColor(0x10, 0xB9, 0x81)  # #10B981
GREEN_TEXT    = RGBColor(0x16, 0xA3, 0x4A)  # #16A34A

# Code Syntax Colors
CODE_BG       = RGBColor(0x18, 0x1A, 0x1F)  # #181A1F
CODE_BORDER   = RGBColor(0x2D, 0x31, 0x39)  # #2D3139
CODE_TEXT     = RGBColor(0xD4, 0xD4, 0xD8)  # #D4D4D8
CODE_KW       = RGBColor(0x38, 0xBD, 0xF8)  # Cyan
CODE_TAG      = RGBColor(0xF4, 0x72, 0xB6)  # Pink
CODE_STR      = RGBColor(0xFC, 0xD3, 0x4D)  # Amber
CODE_COMM     = RGBColor(0x71, 0x71, 0x7A)  # Gray
CODE_PROP     = RGBColor(0xA7, 0x8B, 0xFA)  # Violet
CODE_VAL      = RGBColor(0x4A, 0xDE, 0x80)  # Green

def strip_shape_styles_and_shadows(shape):
    """Ensure zero default PPT theme shadows or effects on created shapes."""
    elem = shape._element
    style = elem.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
    if style is not None:
        elem.remove(style)
    for effect in elem.xpath('.//a:outerShdw'):
        effect.getparent().remove(effect)

def add_header_and_footer(slide, category_text, question_text):
    """Creates top-right orange category pill badge and bottom-left question."""
    pill_width = Pt(max(150, min(330, len(category_text) * 9.8 + 36)))
    pill_left = Pt(700) - pill_width
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_left, Pt(20.4), pill_width, Pt(22.6))
    pill.name = "Google Shape;69;p14"
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = ORANGE_PILL
    pill.line.color.rgb = BORDER_GRAY
    pill.line.width = Pt(0.75)
    strip_shape_styles_and_shadows(pill)

    tf_p = pill.text_frame
    tf_p.word_wrap = False
    tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = 0
    p_p = tf_p.paragraphs[0]
    p_p.alignment = PP_ALIGN.CENTER
    r_p = p_p.add_run()
    r_p.text = category_text.upper()
    r_p.font.name = "Inter"
    r_p.font.size = Pt(12)
    r_p.font.bold = False
    r_p.font.color.rgb = WHITE

    tb_bot = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(500), Pt(20))
    tb_bot.name = "TextBox 2"
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    tf_bot.margin_left = tf_bot.margin_top = tf_bot.margin_right = tf_bot.margin_bottom = 0
    p_bot = tf_bot.paragraphs[0]
    p_bot.alignment = PP_ALIGN.LEFT
    r_bot = p_bot.add_run()
    r_bot.text = question_text
    r_bot.font.name = "Montserrat"
    r_bot.font.size = Pt(10)
    r_bot.font.color.rgb = QUESTION_TEXT

def create_code_explanation_slide(slide, category, question, left_title, code_lines, right_title, bullet_items):
    """Standard 2-column layout: Left code card, Right clean explanation with bullet items."""
    add_header_and_footer(slide, category, question)
    
    # Left Card
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(24), Pt(56), Pt(325), Pt(298))
    left_card.name = "Rounded Rectangle 3"
    left_card.adjustments[0] = 0.04
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CODE_BG
    left_card.line.color.rgb = CODE_BORDER
    left_card.line.width = Pt(1)
    strip_shape_styles_and_shadows(left_card)

    tb_lh = slide.shapes.add_textbox(Pt(38), Pt(64), Pt(295), Pt(18))
    tb_lh.name = "TextBox 4"
    tf_lh = tb_lh.text_frame
    tf_lh.word_wrap = True
    tf_lh.margin_left = tf_lh.margin_top = tf_lh.margin_right = tf_lh.margin_bottom = 0
    p_lh = tf_lh.paragraphs[0]
    r_lh = p_lh.add_run()
    r_lh.text = "📄 " + left_title
    r_lh.font.name = "Consolas"
    r_lh.font.size = Pt(9)
    r_lh.font.bold = True
    r_lh.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)

    line_lh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(36), Pt(84), Pt(301), Pt(1))
    line_lh.name = "Rectangle 5"
    line_lh.fill.solid()
    line_lh.fill.fore_color.rgb = ORANGE_LINE
    line_lh.line.fill.background()
    strip_shape_styles_and_shadows(line_lh)

    tb_code = slide.shapes.add_textbox(Pt(36), Pt(90), Pt(301), Pt(256))
    tb_code.name = "TextBox 6"
    tf_code = tb_code.text_frame
    tf_code.word_wrap = True
    tf_code.margin_left = tf_code.margin_top = tf_code.margin_right = tf_code.margin_bottom = 0
    
    first_p = True
    for line in code_lines:
        if first_p:
            p = tf_code.paragraphs[0]
            first_p = False
        else:
            p = tf_code.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        
        if isinstance(line, str):
            r = p.add_run()
            r.text = line
            r.font.name = "Consolas"
            r.font.size = Pt(8)
            r.font.color.rgb = CODE_TEXT
        else:
            for text_chunk, chunk_color in line:
                r = p.add_run()
                r.text = text_chunk
                r.font.name = "Consolas"
                r.font.size = Pt(8)
                r.font.color.rgb = chunk_color

    # Right side: Dynamic layout depending on title length & clean vertical spacing
    is_multi_line = len(right_title) > 34
    
    right_x = Pt(374)
    right_w = Pt(322)
    
    if is_multi_line:
        header_top = Pt(54)
        header_h = Pt(28)
        header_size = Pt(10)
        line_top = Pt(85)
        bullets_start_y = Pt(94)
    else:
        header_top = Pt(60)
        header_h = Pt(18)
        header_size = Pt(10.5)
        line_top = Pt(82)
        bullets_start_y = Pt(91)

    tb_rh = slide.shapes.add_textbox(right_x, header_top, right_w, header_h)
    tb_rh.name = "TextBox 7"
    tf_rh = tb_rh.text_frame
    tf_rh.word_wrap = True
    tf_rh.margin_left = tf_rh.margin_top = tf_rh.margin_right = tf_rh.margin_bottom = 0
    p_rh = tf_rh.paragraphs[0]
    p_rh.line_spacing = 1.05
    r_rh = p_rh.add_run()
    r_rh.text = right_title
    r_rh.font.name = "Montserrat"
    r_rh.font.size = header_size
    r_rh.font.bold = True
    r_rh.font.color.rgb = ORANGE_LINE

    line_rh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, line_top, right_w, Pt(1))
    line_rh.name = "Rectangle 8"
    line_rh.fill.solid()
    line_rh.fill.fore_color.rgb = ORANGE_LINE
    line_rh.line.fill.background()
    strip_shape_styles_and_shadows(line_rh)

    # Bullets with generous, balanced vertical spacing (step 62pt)
    y_offset = bullets_start_y
    for b_idx, item in enumerate(bullet_items):
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, y_offset + Pt(3.5), Pt(6.5), Pt(6.5))
        dot.name = f"Rounded Rectangle Dot {b_idx}"
        dot.adjustments[0] = 0.2
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)

        tb_b = slide.shapes.add_textbox(right_x + Pt(15), y_offset, right_w - Pt(15), Pt(48))
        tb_b.name = f"TextBox Bullet {b_idx}"
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        p_b1 = tf_b.paragraphs[0]
        p_b1.space_after = Pt(2)
        r_title = p_b1.add_run()
        r_title.text = item["title"]
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = DARK_TEXT

        p_b2 = tf_b.add_paragraph()
        p_b2.line_spacing = 1.15
        p_b2.space_after = Pt(0)
        
        if isinstance(item["desc"], str):
            r_desc = p_b2.add_run()
            r_desc.text = item["desc"]
            r_desc.font.name = "Inter"
            r_desc.font.size = Pt(8.5)
            r_desc.font.color.rgb = BODY_TEXT
            desc_len = len(item["desc"])
        else:
            desc_len = 0
            for d_text, d_bold, d_color in item["desc"]:
                r_desc = p_b2.add_run()
                r_desc.text = d_text
                r_desc.font.name = "Inter"
                r_desc.font.size = Pt(8.5)
                r_desc.font.bold = d_bold
                r_desc.font.color.rgb = d_color
                desc_len += len(d_text)

        # Dynamic height calculation and guaranteed 10pt gap to prevent text overlapping
        computed_h = 54 if desc_len > 95 else (44 if desc_len > 45 else 34)
        item_h = Pt(max(item.get("height", 44), computed_h))
        gap = Pt(10)
        y_offset += item_h + gap

def create_diagram_flex_grid_slide(slide, category, question):
    """Full-width 3-column diagram comparing 1D Flexbox, 2D Grid, and Interface Synergy."""
    add_header_and_footer(slide, category, question)
    
    col_w = Pt(210)
    col_h = Pt(298)
    top_y = Pt(56)
    
    cols_data = [
        {
            "left": Pt(24),
            "accent": RGBColor(0x4A, 0xDE, 0x80), # Green
            "title": "1. FLEXBOX (1D)",
            "sub": "Линейная ось (ряд / колонка)",
            "code_badge": "display: flex; gap: 12px;",
            "visual_type": "flex",
            "summary_title": "Одномерная цепочка (1D):",
            "summary_text": "Идеально выстраивает элементы вдоль одной линии (шапка, меню, ряд кнопок, футер) и делит пустоту."
        },
        {
            "left": Pt(255),
            "accent": RGBColor(0x38, 0xBD, 0xF8), # Cyan
            "title": "2. CSS GRID (2D)",
            "sub": "Строки + Столбцы (Матрица)",
            "code_badge": "grid-template-columns: repeat(3, 1fr);",
            "visual_type": "grid",
            "summary_title": "Двумерная координатная сетка (2D):",
            "summary_text": "Управляет строками и колонками одновременно. Незаменима для витрин каталога, дашбордов и фотогалерей."
        },
        {
            "left": Pt(486),
            "accent": RGBColor(0xA7, 0x8B, 0xFA), # Violet
            "title": "3. СИНЕРГИЯ СЕТОК",
            "sub": "Совместная работа в UI",
            "code_badge": "Grid (каталог) + Flex (карточка)",
            "visual_type": "synergy",
            "summary_title": "Командная работа в реальном коде:",
            "summary_text": "CSS Grid держит общую 3-колоночную сетку каталога, а Flexbox организует кнопки и контент внутри каждой карточки!"
        }
    ]
    
    for c_idx, c in enumerate(cols_data):
        x = c["left"]
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.name = f"Card Col {c_idx+1}"
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = CODE_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)

        tb_th = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_th = tb_th.text_frame
        tf_th.word_wrap = True
        tf_th.margin_left = tf_th.margin_top = tf_th.margin_right = tf_th.margin_bottom = 0
        p_t1 = tf_th.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = c["title"]
        r_t1.font.name = "Montserrat"
        r_t1.font.size = Pt(11)
        r_t1.font.bold = True
        r_t1.font.color.rgb = c["accent"]

        p_t2 = tf_th.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = c["sub"]
        r_t2.font.name = "Inter"
        r_t2.font.size = Pt(8)
        r_t2.font.color.rgb = CODE_COMM

        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(46), col_w - Pt(24), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = c["accent"]
        div.line.fill.background()
        strip_shape_styles_and_shadows(div)

        vbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(12), top_y + Pt(54), col_w - Pt(24), Pt(176))
        vbox.adjustments[0] = 0.04
        vbox.fill.solid()
        vbox.fill.fore_color.rgb = RGBColor(0x13, 0x14, 0x18)
        vbox.line.color.rgb = CODE_BORDER
        vbox.line.width = Pt(0.75)
        strip_shape_styles_and_shadows(vbox)

        tb_vb = slide.shapes.add_textbox(x + Pt(16), top_y + Pt(56), col_w - Pt(32), Pt(14))
        tf_vb = tb_vb.text_frame
        tf_vb.word_wrap = True
        tf_vb.margin_left = tf_vb.margin_top = tf_vb.margin_right = tf_vb.margin_bottom = 0
        p_vb = tf_vb.paragraphs[0]
        r_vb = p_vb.add_run()
        r_vb.text = c["code_badge"]
        r_vb.font.name = "Consolas"
        r_vb.font.size = Pt(7.5)
        r_vb.font.color.rgb = CODE_COMM

        if c["visual_type"] == "flex":
            arrow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(20), top_y + Pt(74), col_w - Pt(40), Pt(18))
            arrow.adjustments[0] = 0.5
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0x1F, 0x24, 0x22)
            arrow.line.color.rgb = c["accent"]
            arrow.line.width = Pt(0.75)
            strip_shape_styles_and_shadows(arrow)
            p_ar = arrow.text_frame.paragraphs[0]
            p_ar.alignment = PP_ALIGN.CENTER
            r_ar = p_ar.add_run()
            r_ar.text = "ГЛАВНАЯ ОСЬ X (row) ──►"
            r_ar.font.name = "Consolas"
            r_ar.font.size = Pt(7.5)
            r_ar.font.color.rgb = c["accent"]

            item_w = Pt(51)
            item_h = Pt(76)
            for it_i, (it_title, it_prop) in enumerate([("Item 1", "auto"), ("Item 2", "flex: 1"), ("Item 3", "auto")]):
                it_x = x + Pt(20) + Pt(it_i * 58)
                it_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, it_x, top_y + Pt(98), item_w, item_h)
                it_shp.adjustments[0] = 0.08
                it_shp.fill.solid()
                it_shp.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
                it_shp.line.color.rgb = c["accent"]
                it_shp.line.width = Pt(1)
                strip_shape_styles_and_shadows(it_shp)
                
                tf_it = it_shp.text_frame
                tf_it.margin_left = tf_it.margin_top = tf_it.margin_right = tf_it.margin_bottom = 0
                p_it1 = tf_it.paragraphs[0]
                p_it1.alignment = PP_ALIGN.CENTER
                r_it1 = p_it1.add_run()
                r_it1.text = it_title
                r_it1.font.name = "Consolas"
                r_it1.font.size = Pt(8)
                r_it1.font.bold = True
                r_it1.font.color.rgb = WHITE
                
                p_it2 = tf_it.add_paragraph()
                p_it2.alignment = PP_ALIGN.CENTER
                r_it2 = p_it2.add_run()
                r_it2.text = it_prop
                r_it2.font.name = "Consolas"
                r_it2.font.size = Pt(7)
                r_it2.font.color.rgb = c["accent"]

            inf = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(20), top_y + Pt(182), col_w - Pt(40), Pt(38))
            inf.fill.solid()
            inf.fill.fore_color.rgb = RGBColor(0x16, 0x1E, 0x24)
            inf.line.fill.background()
            strip_shape_styles_and_shadows(inf)
            p_inf = inf.text_frame.paragraphs[0]
            p_inf.alignment = PP_ALIGN.CENTER
            r_inf = p_inf.add_run()
            r_inf.text = "justify-content: space-between;\nalign-items: center;"
            r_inf.font.name = "Consolas"
            r_inf.font.size = Pt(7)
            r_inf.font.color.rgb = CODE_TEXT

        elif c["visual_type"] == "grid":
            cell_w = Pt(79)
            cell_h = Pt(66)
            grid_cells = [
                (x + Pt(20), top_y + Pt(76), "Ячейка 1", "col 1 / row 1"),
                (x + Pt(108), top_y + Pt(76), "Ячейка 2", "col 2 / row 1"),
                (x + Pt(20), top_y + Pt(148), "Ячейка 3", "col 1 / row 2"),
                (x + Pt(108), top_y + Pt(148), "Ячейка 4", "col 2 / row 2"),
            ]
            for gx, gy, gtitle, gcoord in grid_cells:
                g_shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, gx, gy, cell_w, cell_h)
                g_shp.adjustments[0] = 0.08
                g_shp.fill.solid()
                g_shp.fill.fore_color.rgb = RGBColor(0x1B, 0x24, 0x33)
                g_shp.line.color.rgb = c["accent"]
                g_shp.line.width = Pt(1)
                strip_shape_styles_and_shadows(g_shp)
                
                tf_g = g_shp.text_frame
                tf_g.margin_left = tf_g.margin_top = tf_g.margin_right = tf_g.margin_bottom = 0
                p_g1 = tf_g.paragraphs[0]
                p_g1.alignment = PP_ALIGN.CENTER
                r_g1 = p_g1.add_run()
                r_g1.text = gtitle
                r_g1.font.name = "Consolas"
                r_g1.font.size = Pt(8)
                r_g1.font.bold = True
                r_g1.font.color.rgb = WHITE
                
                p_g2 = tf_g.add_paragraph()
                p_g2.alignment = PP_ALIGN.CENTER
                r_g2 = p_g2.add_run()
                r_g2.text = gcoord
                r_g2.font.name = "Inter"
                r_g2.font.size = Pt(6.5)
                r_g2.font.color.rgb = c["accent"]

        else: # synergy
            fr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(20), top_y + Pt(74), col_w - Pt(40), Pt(146))
            fr.adjustments[0] = 0.06
            fr.fill.solid()
            fr.fill.fore_color.rgb = RGBColor(0x20, 0x1A, 0x2E)
            fr.line.color.rgb = c["accent"]
            fr.line.width = Pt(1)
            strip_shape_styles_and_shadows(fr)
            
            p_fr = fr.text_frame.paragraphs[0]
            p_fr.alignment = PP_ALIGN.LEFT
            r_fr = p_fr.add_run()
            r_fr.text = " ВНЕШНИЙ GRID (3 колонки)"
            r_fr.font.name = "Consolas"
            r_fr.font.size = Pt(7)
            r_fr.font.color.rgb = c["accent"]

            in1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(28), top_y + Pt(94), col_w - Pt(56), Pt(56))
            in1.adjustments[0] = 0.08
            in1.fill.solid()
            in1.fill.fore_color.rgb = RGBColor(0x2D, 0x25, 0x3F)
            in1.line.color.rgb = RGBColor(0x6D, 0x28, 0xD9)
            in1.line.width = Pt(1)
            strip_shape_styles_and_shadows(in1)
            
            tf_in1 = in1.text_frame
            p_in1 = tf_in1.paragraphs[0]
            r_in1 = p_in1.add_run()
            r_in1.text = "🖼️ .card-img (200px)\n"
            r_in1.font.name = "Consolas"
            r_in1.font.size = Pt(7)
            r_in1.font.color.rgb = WHITE
            r_in2 = p_in1.add_run()
            r_in2.text = "📄 .card-content (flex: 1)"
            r_in2.font.name = "Consolas"
            r_in2.font.size = Pt(7)
            r_in2.font.color.rgb = RGBColor(0x38, 0xBD, 0xF8)

            in2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(28), top_y + Pt(156), col_w - Pt(56), Pt(54))
            in2.adjustments[0] = 0.08
            in2.fill.solid()
            in2.fill.fore_color.rgb = RGBColor(0x1B, 0x24, 0x33)
            in2.line.color.rgb = RGBColor(0x10, 0xB9, 0x81)
            in2.line.width = Pt(1)
            strip_shape_styles_and_shadows(in2)
            
            tf_in2 = in2.text_frame
            p_in2 = tf_in2.paragraphs[0]
            r_in2a = p_in2.add_run()
            r_in2a.text = "💰 450 ₽/час  |  [👁️] [Бронь]\n"
            r_in2a.font.name = "Consolas"
            r_in2a.font.size = Pt(7)
            r_in2a.font.color.rgb = WHITE
            r_in2b = p_in2.add_run()
            r_in2b.text = "⚙️ display: flex; space-between;"
            r_in2b.font.name = "Consolas"
            r_in2b.font.size = Pt(6.5)
            r_in2b.font.color.rgb = RGBColor(0x4A, 0xDE, 0x80)

        tb_sm = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(236), col_w - Pt(24), Pt(54))
        tf_sm = tb_sm.text_frame
        tf_sm.word_wrap = True
        tf_sm.margin_left = tf_sm.margin_top = tf_sm.margin_right = tf_sm.margin_bottom = 0
        p_s1 = tf_sm.paragraphs[0]
        r_s1 = p_s1.add_run()
        r_s1.text = c["summary_title"] + "\n"
        r_s1.font.name = "Inter"
        r_s1.font.size = Pt(8)
        r_s1.font.bold = True
        r_s1.font.color.rgb = c["accent"]

        p_s2 = tf_sm.add_paragraph()
        r_s2 = p_s2.add_run()
        r_s2.text = c["summary_text"]
        r_s2.font.name = "Inter"
        r_s2.font.size = Pt(7.5)
        r_s2.font.color.rgb = CODE_TEXT

def create_card_anatomy_slide(slide, category, question):
    """Full-width visual breakdown of the .room-card component architecture."""
    add_header_and_footer(slide, category, question)
    
    # Left: Mockup of the card
    card_mock = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(36), Pt(56), Pt(280), Pt(298))
    card_mock.adjustments[0] = 0.04
    card_mock.fill.solid()
    card_mock.fill.fore_color.rgb = WHITE
    card_mock.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card_mock.line.width = Pt(1.5)
    strip_shape_styles_and_shadows(card_mock)

    img_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(37), Pt(57), Pt(278), Pt(120))
    img_box.adjustments[0] = 0.04
    img_box.fill.solid()
    img_box.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xF6)
    img_box.line.color.rgb = RGBColor(0x00, 0x7B, 0xFF)
    img_box.line.width = Pt(1)
    strip_shape_styles_and_shadows(img_box)
    p_ib = img_box.text_frame.paragraphs[0]
    p_ib.alignment = PP_ALIGN.CENTER
    r_ib = p_ib.add_run()
    r_ib.text = "🖼️ .card-img-wrap (height: 200px)\nobject-fit: cover; overflow: hidden;"
    r_ib.font.name = "Consolas"
    r_ib.font.size = Pt(8)
    r_ib.font.color.rgb = RGBColor(0x00, 0x7B, 0xFF)

    tb_mt = slide.shapes.add_textbox(Pt(50), Pt(184), Pt(250), Pt(20))
    tf_mt = tb_mt.text_frame
    p_mt = tf_mt.paragraphs[0]
    r_mt = p_mt.add_run()
    r_mt.text = "Мини-офис Focus (.card-title)"
    r_mt.font.name = "Montserrat"
    r_mt.font.size = Pt(10)
    r_mt.font.bold = True
    r_mt.font.color.rgb = DARK_TEXT

    tb_me = slide.shapes.add_textbox(Pt(50), Pt(206), Pt(250), Pt(50))
    tf_me = tb_me.text_frame
    for eq_t in ["• Wi-Fi 500 Мбит/с", "• 4K Монитор", "• Эргономичное кресло"]:
        p_e = tf_me.add_paragraph() if tf_me.paragraphs[0].text else tf_me.paragraphs[0]
        r_e = p_e.add_run()
        r_e.text = eq_t
        r_e.font.name = "Inter"
        r_e.font.size = Pt(8)
        r_e.font.color.rgb = BODY_TEXT

    f_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(37), Pt(265), Pt(278), Pt(1))
    f_line.fill.solid()
    f_line.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    f_line.line.fill.background()
    strip_shape_styles_and_shadows(f_line)

    f_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(37), Pt(268), Pt(278), Pt(46))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = WHITE
    f_box.line.fill.background()
    strip_shape_styles_and_shadows(f_box)
    
    tb_pr = slide.shapes.add_textbox(Pt(48), Pt(276), Pt(100), Pt(24))
    p_pr = tb_pr.text_frame.paragraphs[0]
    r_pr1 = p_pr.add_run()
    r_pr1.text = "450 ₽ "
    r_pr1.font.name = "Montserrat"
    r_pr1.font.size = Pt(10)
    r_pr1.font.bold = True
    r_pr1.font.color.rgb = DARK_TEXT
    r_pr2 = p_pr.add_run()
    r_pr2.text = "/ час"
    r_pr2.font.name = "Inter"
    r_pr2.font.size = Pt(8)
    r_pr2.font.color.rgb = BODY_TEXT

    b_ic = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(170), Pt(274), Pt(22), Pt(22))
    b_ic.adjustments[0] = 0.2
    b_ic.fill.solid()
    b_ic.fill.fore_color.rgb = WHITE
    b_ic.line.color.rgb = RGBColor(0x00, 0x7B, 0xFF)
    b_ic.line.width = Pt(1)
    strip_shape_styles_and_shadows(b_ic)
    p_bic = b_ic.text_frame.paragraphs[0]
    p_bic.alignment = PP_ALIGN.CENTER
    r_bic = p_bic.add_run()
    r_bic.text = "👁️"
    r_bic.font.size = Pt(8)

    b_act = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(200), Pt(274), Pt(100), Pt(22))
    b_act.adjustments[0] = 0.2
    b_act.fill.solid()
    b_act.fill.fore_color.rgb = RGBColor(0x00, 0x7B, 0xFF)
    b_act.line.fill.background()
    strip_shape_styles_and_shadows(b_act)
    p_bact = b_act.text_frame.paragraphs[0]
    p_bact.alignment = PP_ALIGN.CENTER
    r_bact = p_bact.add_run()
    r_bact.text = "Забронировать"
    r_bact.font.name = "Inter"
    r_bact.font.size = Pt(8)
    r_bact.font.bold = True
    r_bact.font.color.rgb = WHITE

    # Right: Architectural Callouts
    tb_rh = slide.shapes.add_textbox(Pt(345), Pt(60), Pt(340), Pt(18))
    tb_rh.name = "TextBox 7"
    tf_rh = tb_rh.text_frame
    p_rh = tf_rh.paragraphs[0]
    r_rh = p_rh.add_run()
    r_rh.text = "4 СЛОЯ КОМПОНЕНТА .ROOM-CARD"
    r_rh.font.name = "Montserrat"
    r_rh.font.size = Pt(10.5)
    r_rh.font.bold = True
    r_rh.font.color.rgb = ORANGE_LINE

    line_rh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(345), Pt(82), Pt(340), Pt(1))
    line_rh.name = "Rectangle 8"
    line_rh.fill.solid()
    line_rh.fill.fore_color.rgb = ORANGE_LINE
    line_rh.line.fill.background()
    strip_shape_styles_and_shadows(line_rh)

    layers = [
        ("1. Медиа-контейнер (.card-img-wrap):", "Фиксированная высота 200px с overflow: hidden. Защищает карточку от вылезания картинок разного разрешения и формы."),
        ("2. Контентное тело (.card-content):", "Внутренний padding 15px и колоночный Flexbox (display: flex; flex-direction: column; flex: 1;). Гарантирует равную высоту всех карточек в каталоге."),
        ("3. Распорка списка (.card-equipment):", "Список тегов или оборудования с flex: 1 забирает излишки пустого места и выталкивает блок стоимости и кнопок строго в нижний край."),
        ("4. Подвал действия (.card-footer):", "Flexbox с justify-content: space-between разносит цену налево, а блок кнопок .card-btns (иконка детального просмотра + primary кнопка) — направо.")
    ]

    y_pos = Pt(92)
    for l_idx, (ltitle, ldesc) in enumerate(layers):
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(345), y_pos + Pt(3), Pt(6.4), Pt(6.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)

        tb_l = slide.shapes.add_textbox(Pt(360), y_pos, Pt(325), Pt(48))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
        
        p1 = tf_l.paragraphs[0]
        p1.space_after = Pt(2)
        r1 = p1.add_run()
        r1.text = ltitle
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT

        p2 = tf_l.add_paragraph()
        p2.line_spacing = 1.15
        r2 = p2.add_run()
        r2.text = ldesc
        r2.font.name = "Inter"
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = BODY_TEXT

        y_pos += Pt(58)

def create_grid_cards_slide(slide, category, question, title_text, cards_data):
    """2x2 error analysis cards without gray/transparent background plates."""
    add_header_and_footer(slide, category, question)
    
    coords = [
        (Pt(28), Pt(56), Pt(320), Pt(142)),
        (Pt(372), Pt(56), Pt(320), Pt(142)),
        (Pt(28), Pt(208), Pt(320), Pt(142)),
        (Pt(372), Pt(208), Pt(320), Pt(142)),
    ]
    
    for idx, (x, y, w, h) in enumerate(coords):
        data = cards_data[idx]
        
        tb_h = slide.shapes.add_textbox(x, y + Pt(6), w, Pt(20))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
        p_h = tf_h.paragraphs[0]
        r_num = p_h.add_run()
        r_num.text = f"{idx + 1}. {data['title']}"
        r_num.font.name = "Montserrat"
        r_num.font.size = Pt(10.5)
        r_num.font.bold = True
        r_num.font.color.rgb = ORANGE_LINE
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Pt(28), w, Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE_LINE
        line.line.fill.background()
        strip_shape_styles_and_shadows(line)
        
        tb_body = slide.shapes.add_textbox(x, y + Pt(36), w, Pt(96))
        tf_b = tb_body.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        p1 = tf_b.paragraphs[0]
        p1.space_after = Pt(4)
        r1_lbl = p1.add_run()
        r1_lbl.text = "Ошибка: "
        r1_lbl.font.name = "Inter"
        r1_lbl.font.size = Pt(8.5)
        r1_lbl.font.bold = True
        r1_lbl.font.color.rgb = DARK_TEXT
        r1_txt = p1.add_run()
        r1_txt.text = data['problem']
        r1_txt.font.name = "Inter"
        r1_txt.font.size = Pt(8.5)
        r1_txt.font.color.rgb = BODY_TEXT
        
        p2 = tf_b.add_paragraph()
        p2.space_after = Pt(0)
        p2.line_spacing = 1.15
        r2_lbl = p2.add_run()
        r2_lbl.text = "Решение: "
        r2_lbl.font.name = "Inter"
        r2_lbl.font.size = Pt(8.5)
        r2_lbl.font.bold = True
        r2_lbl.font.color.rgb = GREEN_TEXT
        r2_txt = p2.add_run()
        r2_txt.text = data['solution']
        r2_txt.font.name = "Inter"
        r2_txt.font.size = Pt(8.5)
        r2_txt.font.color.rgb = BODY_TEXT

def create_checklist_slide(slide, category, question, title_text, items):
    """Clean 6-item checklist with checkmarks."""
    add_header_and_footer(slide, category, question)

    tb_h = slide.shapes.add_textbox(Pt(54), Pt(68), Pt(600), Pt(22))
    tf_h = tb_h.text_frame
    tf_h.word_wrap = True
    tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
    p_h = tf_h.paragraphs[0]
    r_h = p_h.add_run()
    r_h.text = title_text
    r_h.font.name = "Montserrat"
    r_h.font.size = Pt(11)
    r_h.font.bold = True
    r_h.font.color.rgb = ORANGE_LINE

    line_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(54), Pt(92), Pt(612), Pt(1))
    line_h.fill.solid()
    line_h.fill.fore_color.rgb = ORANGE_LINE
    line_h.line.fill.background()
    strip_shape_styles_and_shadows(line_h)

    y_pos = Pt(102)
    for idx, item in enumerate(items):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(54), y_pos + Pt(2), Pt(22), Pt(20))
        badge.adjustments[0] = 0.2
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN_BADGE
        badge.line.fill.background()
        strip_shape_styles_and_shadows(badge)
        
        tb_chk = slide.shapes.add_textbox(Pt(54), y_pos + Pt(4), Pt(22), Pt(20))
        tf_chk = tb_chk.text_frame
        p_chk = tf_chk.paragraphs[0]
        p_chk.alignment = PP_ALIGN.CENTER
        r_chk = p_chk.add_run()
        r_chk.text = "✔"
        r_chk.font.name = "Segoe UI Symbol"
        r_chk.font.size = Pt(10)
        r_chk.font.bold = True
        r_chk.font.color.rgb = WHITE
        
        tb_item = slide.shapes.add_textbox(Pt(86), y_pos, Pt(570), Pt(38))
        tf_item = tb_item.text_frame
        tf_item.word_wrap = True
        tf_item.margin_left = tf_item.margin_top = tf_item.margin_right = tf_item.margin_bottom = 0
        
        p1 = tf_item.paragraphs[0]
        p1.space_after = Pt(2)
        r1 = p1.add_run()
        r1.text = item["title"]
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = DARK_TEXT

        p2 = tf_item.add_paragraph()
        r2 = p2.add_run()
        r2.text = item["desc"]
        r2.font.name = "Inter"
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = BODY_TEXT

        y_pos += Pt(42)

slides_data = [
    # 1. Slide 3: Theory: CSS Variables (:root)
    {
        "type": "code",
        "category": "ПЕРЕМЕННЫЕ :ROOT",
        "question": "Как работают глобальные CSS-переменные и зачем они нужны?",
        "left_title": "css/style.css (:root variables)",
        "code_lines": [
            [("/* Глобальная область видимости :root */", CODE_COMM)],
            [(":root ", CODE_TAG), ("{", CODE_TEXT)],
            [("  --primary-color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), (";", CODE_TEXT)],
            [("  --primary-hover", CODE_PROP), (": ", CODE_TEXT), ("#0056b3", CODE_VAL), (";", CODE_TEXT)],
            [("  --text-color", CODE_PROP), (": ", CODE_TEXT), ("#222222", CODE_VAL), (";", CODE_TEXT)],
            [("  --text-muted", CODE_PROP), (": ", CODE_TEXT), ("#666666", CODE_VAL), (";", CODE_TEXT)],
            [("  --border-color", CODE_PROP), (": ", CODE_TEXT), ("#dddddd", CODE_VAL), (";", CODE_TEXT)],
            [("  --radius-sm", CODE_PROP), (": ", CODE_TEXT), ("4px", CODE_VAL), (";", CODE_TEXT)],
            [("  --radius-md", CODE_PROP), (": ", CODE_TEXT), ("6px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Использование через var() */", CODE_COMM)],
            [(".btn-primary ", CODE_TAG), ("{", CODE_TEXT)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("var", CODE_KW), ("(", CODE_TEXT), ("--primary-color", CODE_PROP), (");", CODE_TEXT)],
            [("  border-radius", CODE_PROP), (": ", CODE_TEXT), ("var", CODE_KW), ("(", CODE_TEXT), ("--radius-sm", CODE_PROP), (");", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ГЛОБАЛЬНАЯ ОБЛАСТЬ ВИДИМОСТИ И СИНТАКСИС",
        "bullet_items": [
            {
                "title": "Псевдокласс :root:",
                "desc": "Соответствует корневому элементу HTML (<html>). Переменные, объявленные здесь, наследуются всеми элементами на странице без исключений.",
                "height": 46
            },
            {
                "title": "Синтаксис с двумя дефисами (--):",
                "desc": "Имена пользовательских CSS-свойств обязаны начинаться с '--' (например, --primary-color), чтобы исключить конфликт со встроенными стилями.",
                "height": 46
            },
            {
                "title": "Функция var() для чтения значений:",
                "desc": "Значение переменной вызывается через var(--name). Можно передать резервное значение: var(--name, #007bff) на случай отсутствия токена.",
                "height": 48
            },
            {
                "title": "Мгновенное обновление проекта:",
                "desc": "Изменив один цвет в блоке :root, вы моментально перекрашиваете все кнопки, ссылки, карточки и рамки на всем веб-сайте!",
                "height": 44
            }
        ]
    },

    # 2. Slide 4: Theory: Design Tokens & Palette
    {
        "type": "code",
        "category": "ДИЗАЙН-СИСТЕМА",
        "question": "Как организовать токены стилей для чистого и масштабируемого кода?",
        "left_title": "css/style.css (Дизайн-токены)",
        "code_lines": [
            [(":root ", CODE_TAG), ("{", CODE_TEXT)],
            [("  /* 1. Цветовая палитра */", CODE_COMM)],
            [("  --color-primary", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), (";", CODE_TEXT)],
            [("  --color-success", CODE_PROP), (": ", CODE_TEXT), ("#28a745", CODE_VAL), (";", CODE_TEXT)],
            [("  --color-dark", CODE_PROP), (": ", CODE_TEXT), ("#222222", CODE_VAL), (";", CODE_TEXT)],
            [("  --color-gray", CODE_PROP), (": ", CODE_TEXT), ("#555555", CODE_VAL), (";", CODE_TEXT)],
            [("  --color-bg-light", CODE_PROP), (": ", CODE_TEXT), ("#f8f9fa", CODE_VAL), (";", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  /* 2. Отступы и геометрия */", CODE_COMM)],
            [("  --gap-grid", CODE_PROP), (": ", CODE_TEXT), ("20px", CODE_VAL), (";", CODE_TEXT)],
            [("  --pad-card", CODE_PROP), (": ", CODE_TEXT), ("15px", CODE_VAL), (";", CODE_TEXT)],
            [("  --radius-card", CODE_PROP), (": ", CODE_TEXT), ("6px", CODE_VAL), (";", CODE_TEXT)],
            [("  --transition-smooth", CODE_PROP), (": ", CODE_TEXT), ("0.2s ease", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ДИЗАЙН-ТОКЕНЫ: МОСТ МЕЖДУ FIGMA И КОДОМ",
        "bullet_items": [
            {
                "title": "Что такое дизайн-токены:",
                "desc": "Это неделимые константы дизайн-системы (цвета, радиусы, интервалы, шрифты), которые стандартизируют внешний вид всех интерфейсов.",
                "height": 46
            },
            {
                "title": "Группировка по назначениям:",
                "desc": "Рекомендуется делить токены на группы: семантические цвета (brand, success, danger), нейтральные оттенки текста, отступы и радиусы.",
                "height": 48
            },
            {
                "title": "Принцип DRY (Не повторяйся):",
                "desc": "Вместо многократного копирования случайных HEX-кодов (#007bff) по файлу стилей, разработчик использует понятные осмысленные имена токенов.",
                "height": 46
            },
            {
                "title": "Готовность к темной теме:",
                "desc": "Дизайн-токены позволяют за 5 минут реализовать темную тему или сменить корпоративную гамму под праздничную акцию сервиса.",
                "height": 44
            }
        ]
    },

    # 3. Slide 5: Diagram: 1D vs 2D Layouts (Diagram slide)
    {
        "type": "diagram_flex_grid",
        "category": "СХЕМА: 1D vs 2D СЕТКИ",
        "question": "В чем принципиальная разница в логике работы Flexbox и CSS Grid?"
    },

    # 4. Slide 6: Theory: Flexbox in Catalog
    {
        "type": "code",
        "category": "FLEXBOX В КАТАЛОГЕ",
        "question": "Как работает Flexbox при верстке витрины карточек каталога?",
        "left_title": "css/style.css (.rooms-grid)",
        "code_lines": [
            [("/* Flex-контейнер витрины комнат */", CODE_COMM)],
            [(".rooms-grid ", CODE_TAG), ("{", CODE_TEXT)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  flex-wrap", CODE_PROP), (": ", CODE_TEXT), ("wrap", CODE_VAL), ("; /* Перенос на новую строку */", CODE_COMM)],
            [("  gap", CODE_PROP), (": ", CODE_TEXT), ("20px", CODE_VAL), (";        /* Шаг между карточками */", CODE_COMM)],
            [("  margin-top", CODE_PROP), (": ", CODE_TEXT), ("30px", CODE_VAL), (";", CODE_TEXT)],
            [("  margin-bottom", CODE_PROP), (": ", CODE_TEXT), ("30px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Карточка в 3-колоночной сетке */", CODE_COMM)],
            [(".room-card ", CODE_TAG), ("{", CODE_TEXT)],
            [("  width", CODE_PROP), (": ", CODE_TEXT), ("373px", CODE_VAL), (";     /* (1200 - 2*20) / 3 */", CODE_COMM)],
            [("  flex-shrink", CODE_PROP), (": ", CODE_TEXT), ("0", CODE_VAL), (";   /* Защита от сжатия */", CODE_COMM)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "МЕХАНИКА МНОГОСТРОЧНОГО FLEX-ПОТОКА",
        "bullet_items": [
            {
                "title": "Свойство flex-wrap: wrap:",
                "desc": "По умолчанию флекс сжимает карточки в одну линию. Значение wrap разрешает автоматический перенос карточек на следующую строку.",
                "height": 46
            },
            {
                "title": "Современное свойство gap:",
                "desc": "gap: 20px создает равномерные промежутки между карточками строго внутри контейнера, избавляя от отступов :last-child и хаков.",
                "height": 46
            },
            {
                "title": "Математика колонок (373px):",
                "desc": "Ширина контейнера 1200px минус 2 зазора по 20px (40px) = 1160px. Делим на 3 карточки и получаем ровно 373.33px на элемент.",
                "height": 48
            },
            {
                "title": "Поведение при добавлении карточек:",
                "desc": "Если в каталог добавить 4-ю и 5-ю карточки, они плавно встанут во второй ряд с сохранением идеальной геометрии сетки.",
                "height": 44
            }
        ]
    },

    # 5. Slide 7: Theory: CSS Grid in Catalog
    {
        "type": "code",
        "category": "CSS GRID В КАТАЛОГЕ",
        "question": "Как построить каталог с помощью двумерной координатной сетки Grid?",
        "left_title": "css/style.css (CSS Grid альтернатива)",
        "code_lines": [
            [("/* Альтернативная сетка на CSS Grid */", CODE_COMM)],
            [(".rooms-grid-grid ", CODE_TAG), ("{", CODE_TEXT)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("grid", CODE_VAL), (";", CODE_TEXT)],
            [("  grid-template-columns", CODE_PROP), (": ", CODE_TEXT), ("repeat(3, 1fr)", CODE_VAL), (";", CODE_TEXT)],
            [("  gap", CODE_PROP), (": ", CODE_TEXT), ("20px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Автоматическая адаптивность без медиа */", CODE_COMM)],
            [(".rooms-grid-autofit ", CODE_TAG), ("{", CODE_TEXT)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("grid", CODE_VAL), (";", CODE_TEXT)],
            [("  grid-template-columns", CODE_PROP), (": ", CODE_TEXT)],
            [("    ", CODE_TEXT), ("repeat", CODE_KW), ("(", CODE_TEXT), ("auto-fit", CODE_VAL), (", ", CODE_TEXT), ("minmax", CODE_KW), ("(", CODE_TEXT), ("320px", CODE_VAL), (", ", CODE_TEXT), ("1fr", CODE_VAL), ("));", CODE_TEXT)],
            [("  gap", CODE_PROP), (": ", CODE_TEXT), ("20px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ПРЕИМУЩЕСТВА СЕТОК CSS GRID",
        "bullet_items": [
            {
                "title": "Дисплей grid (Координатная сетка):",
                "desc": "Активирует двумерную матричную модель. Браузер берет расчет ширины колонок и рядов полностью на себя без ручных вычислений.",
                "height": 46
            },
            {
                "title": "Доли пространства (fr - Fraction):",
                "desc": "1fr означает 'одну равную долю от остатка места'. repeat(3, 1fr) гарантирует ровно 3 колонки идеальной ширины.",
                "height": 46
            },
            {
                "title": "Функция repeat() для чистоты кода:",
                "desc": "Позволяет не перечислять размеры вручную (1fr 1fr 1fr), сокращая синтаксис и облегчая изменение количества колонок.",
                "height": 48
            },
            {
                "title": "Связка repeat(auto-fit, minmax):",
                "desc": "Продвинутая техника: если экран сужается, колонки сами перестраиваются из 3 в 2 и в 1 без единого медиа-запроса!",
                "height": 44
            }
        ]
    },

    # 6. Slide 8: Diagram: Card Anatomy (Diagram slide)
    {
        "type": "card_anatomy",
        "category": "СХЕМА: АНАТОМИЯ КАРТОЧКИ",
        "question": "Из каких структурных слоев состоит профессиональная карточка каталога?"
    },

    # 7. Slide 9: Practice HTML: Popular Section Layout
    {
        "type": "code",
        "category": "СЕКЦИЯ КАТАЛОГА: HTML",
        "question": "Как семантически грамотно структурировать секцию предложений?",
        "left_title": "index.html (<section class=\"popular-section\">)",
        "code_lines": [
            [("<!-- Секция популярных комнат -->", CODE_COMM)],
            [("<section ", CODE_TAG), ("class=", CODE_PROP), ("\"popular-section\"", CODE_STR), (">", CODE_TAG)],
            [("  <h2 ", CODE_TAG), ("class=", CODE_PROP), ("\"page-title\"", CODE_STR), (">", CODE_TAG)],
            [("    Популярные офисные комнаты", CODE_TEXT)],
            [("  </h2>", CODE_TAG)],
            [("  <p ", CODE_TAG), ("class=", CODE_PROP), ("\"page-subtitle\"", CODE_STR), (">", CODE_TAG)],
            [("    Наиболее востребованные пространства", CODE_TEXT)],
            [("  </p>", CODE_TAG)],
            [("", CODE_TEXT)],
            [("  <!-- Контейнер сетки -->", CODE_COMM)],
            [("  <div ", CODE_TAG), ("class=", CODE_PROP), ("\"rooms-grid\"", CODE_STR), (">", CODE_TAG)],
            [("    <!-- 3 карточки .room-card -->", CODE_COMM)],
            [("  </div>", CODE_TAG)],
            [("", CODE_TEXT)],
            [("  <!-- Действие внизу секции -->", CODE_COMM)],
            [("  <div ", CODE_TAG), ("class=", CODE_PROP), ("\"center-action\"", CODE_STR), (">", CODE_TAG)],
            [("    <a ", CODE_TAG), ("href=", CODE_PROP), ("\"pages/catalog.html\"", CODE_STR), (" class=", CODE_PROP), ("\"btn btn-outline\"", CODE_STR), (">", CODE_TAG)],
            [("      Больше офисов", CODE_TEXT)],
            [("    </a>", CODE_TAG)],
            [("  </div>", CODE_TAG)],
            [("</section>", CODE_TAG)]
        ],
        "right_title": "СЕМАНТИЧЕСКАЯ АРХИТЕКТУРА РАЗДЕЛА СТРАНИЦЫ",
        "bullet_items": [
            {
                "title": "Тег <section> для независимого блока:",
                "desc": "Выделяет самостоятельный смысловой раздел каталога. Это улучшает доступность для скринридеров и ранжирование страницы в поисковиках.",
                "height": 46
            },
            {
                "title": "Иерархия заголовков (h2):",
                "desc": "Заголовок h1 зарезервирован для главного экрана (Hero). Все внутренние разделы каталога начинаются строго со второго уровня h2.",
                "height": 46
            },
            {
                "title": "Служебный подзаголовок .page-subtitle:",
                "desc": "Кратко раскрывает ценность предложения. Оформляется единым классом дизайн-системы с приглушенным контрастом текста.",
                "height": 48
            },
            {
                "title": "Центрирующий блок .center-action:",
                "desc": "Изолирует кнопку 'Больше офисов' от сетки карточек, обеспечивая отступ сверху и выравнивание ссылки по центру.",
                "height": 44
            }
        ]
    },

    # 8. Slide 10: Practice HTML: Room Card HTML Structure
    {
        "type": "code",
        "category": "КАРТОЧКА: HTML-РАЗМЕТКА",
        "question": "Как правильно вложить элементы внутрь карточки для удобной стилизации?",
        "left_title": "index.html (<div class=\"room-card\">)",
        "code_lines": [
            [("<div ", CODE_TAG), ("class=", CODE_PROP), ("\"room-card\"", CODE_STR), (">", CODE_TAG)],
            [("  <!-- 1. Превью с защитой от сбоя загрузки -->", CODE_COMM)],
            [("  <div ", CODE_TAG), ("class=", CODE_PROP), ("\"card-img-wrap\"", CODE_STR), (">", CODE_TAG)],
            [("    <a ", CODE_TAG), ("href=", CODE_PROP), ("\"pages/catalog.html\"", CODE_STR), (">", CODE_TAG)],
            [("      <img ", CODE_TAG), ("src=", CODE_PROP), ("\"img/room-1.jpg\"", CODE_STR), (" alt=", CODE_PROP), ("\"Мини-офис Focus\"", CODE_STR)],
            [("           class=", CODE_PROP), ("\"card-img\"", CODE_STR), (" onerror=", CODE_PROP), ("\"this.src='img/no-image.svg'\"", CODE_STR), (">", CODE_TAG)],
            [("    </a>", CODE_TAG)],
            [("  </div>", CODE_TAG)],
            [("", CODE_TEXT)],
            [("  <!-- 2. Контентное тело -->", CODE_COMM)],
            [("  <div ", CODE_TAG), ("class=", CODE_PROP), ("\"card-content\"", CODE_STR), (">", CODE_TAG)],
            [("    <h3 ", CODE_TAG), ("class=", CODE_PROP), ("\"card-title\"", CODE_STR), (">", CODE_TAG)],
            [("      <a ", CODE_TAG), ("href=", CODE_PROP), ("\"pages/catalog.html\"", CODE_STR), (">Мини-офис Focus</a>", CODE_TEXT)],
            [("    </h3>", CODE_TAG)],
            [("    <ul ", CODE_TAG), ("class=", CODE_PROP), ("\"card-equipment\"", CODE_STR), (">", CODE_TAG)],
            [("      <li>Wi-Fi 500 Мбит/с</li>", CODE_TEXT)],
            [("      <li>4K Монитор</li>", CODE_TEXT)],
            [("    </ul>", CODE_TAG)],
            [("    <!-- 3. Футер карточки -->", CODE_COMM)],
            [("  </div>", CODE_TAG)],
            [("</div>", CODE_TAG)]
        ],
        "right_title": "СЛОИ РАЗМЕТКИ КАРТОЧКИ ТОВАРА",
        "bullet_items": [
            {
                "title": "Резервный обработчик onerror:",
                "desc": "Если фотография комнаты временно недоступна на сервере, событие onerror мгновенно подставит нейтральную SVG-заглушку no-image.svg.",
                "height": 46
            },
            {
                "title": "Разделение фото и контента:",
                "desc": "Контейнер .card-content позволяет задать внутренний padding для текста, не создавая нежелательных зазоров вокруг фотографии в шапке.",
                "height": 46
            },
            {
                "title": "Кликабельность фото и заголовка:",
                "desc": "Пользователи интуитивно кликают по картинке и тексту. Оборачивание в теги <a> повышает конверсию переходов в каталог.",
                "height": 48
            },
            {
                "title": "Семантический список <ul> для оснащения:",
                "desc": "Характеристики офиса представляют собой перечень фактов, поэтому тег <ul> с элементами <li> является стандартом семантики.",
                "height": 44
            }
        ]
    },

    # 9. Slide 11: Practice CSS: Room Card CSS Skeleton
    {
        "type": "code",
        "category": "КАРТОЧКА: CSS-КАРКАС",
        "question": "Как сформировать белый блок карточки со скруглением и тенью?",
        "left_title": "css/style.css (.room-card)",
        "code_lines": [
            [(".room-card ", CODE_TAG), ("{", CODE_TEXT)],
            [("  width", CODE_PROP), (": ", CODE_TEXT), ("373px", CODE_VAL), (";", CODE_TEXT)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#ffffff", CODE_VAL), (";", CODE_TEXT)],
            [("  border", CODE_PROP), (": ", CODE_TEXT), ("1px solid #dddddd", CODE_VAL), (";", CODE_TEXT)],
            [("  border-radius", CODE_PROP), (": ", CODE_TEXT), ("6px", CODE_VAL), (";", CODE_TEXT)],
            [("  overflow", CODE_PROP), (": ", CODE_TEXT), ("hidden", CODE_VAL), ("; /* Срезать углы фото */", CODE_COMM)],
            [("", CODE_TEXT)],
            [("  /* Вертикальный Flexbox */", CODE_COMM)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  flex-direction", CODE_PROP), (": ", CODE_TEXT), ("column", CODE_VAL), (";", CODE_TEXT)],
            [("  transition", CODE_PROP), (": ", CODE_TEXT), ("transform 0.2s, box-shadow 0.2s", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Эффект парения при наведении */", CODE_COMM)],
            [(".room-card:hover ", CODE_TAG), ("{", CODE_TEXT)],
            [("  transform", CODE_PROP), (": ", CODE_TEXT), ("translateY(-4px)", CODE_VAL), (";", CODE_TEXT)],
            [("  box-shadow", CODE_PROP), (": ", CODE_TEXT), ("0 8px 24px rgba(0,0,0,0.08)", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ГЕОМЕТРИЯ И ИНТЕРАКТИВНОСТЬ КАРТОЧКИ",
        "bullet_items": [
            {
                "title": "overflow: hidden для скругления углов:",
                "desc": "Прямоугольная картинка вверху перекрывает border-radius карточки. Свойство overflow: hidden аккуратно обрезает выступающие пиксели.",
                "height": 46
            },
            {
                "title": "flex-direction: column (Вертикальная ось):",
                "desc": "Превращает карточку во внутренний вертикальный флекс-поток: фото крепится сверху, контент растягивается, футер опускается вниз.",
                "height": 48
            },
            {
                "title": "Микро-анимация при наведении (:hover):",
                "desc": "Свойства transform: translateY(-4px) и мягкая тень создают ощущение легкости и отзывчивости интерфейса при взаимодействии.",
                "height": 46
            },
            {
                "title": "Плавный переход через transition:",
                "desc": "Задержка 0.2s сглаживает анимацию, исключая неприятное для глаз резкое мигание теней при быстром движении курсора мыши.",
                "height": 44
            }
        ]
    },

    # 10. Slide 12: Theory & CSS: Image Processing (object-fit: cover)
    {
        "type": "code",
        "category": "ГРАФИКА В КАРТОЧКАХ",
        "question": "Почему фотографии не искажаются даже при разных пропорциях кадров?",
        "left_title": "css/style.css (.card-img & object-fit)",
        "code_lines": [
            [("/* Контейнер картинки фиксированной высоты */", CODE_COMM)],
            [(".card-img-wrap ", CODE_TAG), ("{", CODE_TEXT)],
            [("  width", CODE_PROP), (": ", CODE_TEXT), ("100%", CODE_VAL), (";", CODE_TEXT)],
            [("  height", CODE_PROP), (": ", CODE_TEXT), ("200px", CODE_VAL), (";", CODE_TEXT)],
            [("  overflow", CODE_PROP), (": ", CODE_TEXT), ("hidden", CODE_VAL), (";", CODE_TEXT)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#f0f4f8", CODE_VAL), ("; /* Нейтральный фон */", CODE_COMM)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Картинка без сплющивания */", CODE_COMM)],
            [(".card-img ", CODE_TAG), ("{", CODE_TEXT)],
            [("  width", CODE_PROP), (": ", CODE_TEXT), ("100%", CODE_VAL), (";", CODE_TEXT)],
            [("  height", CODE_PROP), (": ", CODE_TEXT), ("100%", CODE_VAL), (";", CODE_TEXT)],
            [("  object-fit", CODE_PROP), (": ", CODE_TEXT), ("cover", CODE_VAL), (";        /* Сохранять пропорции! */", CODE_COMM)],
            [("  object-position", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";   /* Центрировать кадр */", CODE_COMM)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("block", CODE_VAL), (";           /* Убрать зазор под img */", CODE_COMM)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "OBJECT-FIT: COVER — СПАСЕНИЕ ВЕРСТКИ КАРТИНОК",
        "bullet_items": [
            {
                "title": "Опасность сплющивания фото:",
                "desc": "Если пользователь загрузит вертикальное фото, а блок карточки горизонтальный, браузер без cover исказит пропорции мебели и людей.",
                "height": 46
            },
            {
                "title": "Принцип работы object-fit: cover:",
                "desc": "Работает аналогично background-size: cover: картинка пропорционально заполняет весь блок 200px, а лишние края деликатно обрезаются.",
                "height": 48
            },
            {
                "title": "Фокусировка через object-position: center:",
                "desc": "Гарантирует, что центр композиции кадра (стол, кресло, переговорка) всегда останется видимым в центре миниатюры карточки.",
                "height": 46
            },
            {
                "title": "Снятие зазора через display: block:",
                "desc": "По умолчанию тег <img> строчный и оставляет под собой невидимый отступ в 3-4px базовой линии. display: block полностью устраняет дефект.",
                "height": 44
            }
        ]
    },

    # 11. Slide 13: Practice CSS: Content Block and Equal Height
    {
        "type": "code",
        "category": "КОНТЕНТ КАРТОЧКИ: CSS",
        "question": "Как выровнять карточки по высоте, если длина текста у них разная?",
        "left_title": "css/style.css (.card-content & flex: 1)",
        "code_lines": [
            [(".card-content ", CODE_TAG), ("{", CODE_TEXT)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("15px", CODE_VAL), (";", CODE_TEXT)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  flex-direction", CODE_PROP), (": ", CODE_TEXT), ("column", CODE_VAL), (";", CODE_TEXT)],
            [("  flex", CODE_PROP), (": ", CODE_TEXT), ("1", CODE_VAL), ("; /* Занять всю оставшуюся высоту */", CODE_COMM)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [(".card-title ", CODE_TAG), ("{", CODE_TEXT)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("18px", CODE_VAL), (";", CODE_TEXT)],
            [("  font-weight", CODE_PROP), (": ", CODE_TEXT), ("700", CODE_VAL), (";", CODE_TEXT)],
            [("  margin-bottom", CODE_PROP), (": ", CODE_TEXT), ("10px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [(".card-equipment ", CODE_TAG), ("{", CODE_TEXT)],
            [("  flex", CODE_PROP), (": ", CODE_TEXT), ("1", CODE_VAL), ("; /* Распорка выталкивает футер */", CODE_COMM)],
            [("  margin-bottom", CODE_PROP), (": ", CODE_TEXT), ("15px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "МАГИЯ ФЛЕКС-РАСПОРКИ (FLEX: 1)",
        "bullet_items": [
            {
                "title": "Проблема разной длины текста:",
                "desc": "У 'Мини-офиса Focus' заголовок в 1 строку, а у конференц-зала — в 2 строки. Из-за этого кнопки бронирования начинают 'плясать'.",
                "height": 46
            },
            {
                "title": "flex: 1 у блока .card-content:",
                "desc": "Заставляет тело карточки жадно занять всю доступную высоту родителя, уравнивая все карточки в строке по самому высокому соседу.",
                "height": 48
            },
            {
                "title": "Список оборудования как распорка:",
                "desc": "Свойство flex: 1 у списка .card-equipment съедает всю пустоту между текстом и футером, принудительно прижимая кнопки к низу!",
                "height": 46
            },
            {
                "title": "Эстетика идеального каталога:",
                "desc": "В результате все кнопки 'Забронировать' и цены во всех карточках каталога выстраиваются по идеальной горизонтальной линейке.",
                "height": 44
            }
        ]
    },

    # 12. Slide 14: Practice CSS: Equipment List Styling
    {
        "type": "code",
        "category": "СПИСОК ОБОРУДОВАНИЯ",
        "question": "Как оформить список характеристик комнаты компактно и аккуратно?",
        "left_title": "css/style.css (.card-equipment)",
        "code_lines": [
            [("/* Список параметров оснащения */", CODE_COMM)],
            [(".card-equipment ", CODE_TAG), ("{", CODE_TEXT)],
            [("  list-style", CODE_PROP), (": ", CODE_TEXT), ("none", CODE_VAL), (";   /* Убираем стандартные точки */", CODE_COMM)],
            [("  margin-bottom", CODE_PROP), (": ", CODE_TEXT), ("15px", CODE_VAL), (";", CODE_TEXT)],
            [("  flex", CODE_PROP), (": ", CODE_TEXT), ("1", CODE_VAL), (";          /* Выталкивает футер вниз */", CODE_COMM)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [(".card-equipment li ", CODE_TAG), ("{", CODE_TEXT)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("13px", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#555555", CODE_VAL), (";       /* Приглушенный серый */", CODE_COMM)],
            [("  margin-bottom", CODE_PROP), (": ", CODE_TEXT), ("4px", CODE_VAL), (";", CODE_TEXT)],
            [("  line-height", CODE_PROP), (": ", CODE_TEXT), ("1.4", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ВТОРИЧНАЯ ТИПОГРАФИКА ХАРАКТЕРИСТИК",
        "bullet_items": [
            {
                "title": "Сброс браузерных маркеров:",
                "desc": "list-style: none удаляет неуправляемые черные буллеты браузера, позволяя полностью контролировать отступы текста.",
                "height": 46
            },
            {
                "title": "Визуальная иерархия (Кегль 13px):",
                "desc": "Технические параметры не должны отвлекать от заголовка и цены. Размер 13px и цвет #555555 делают блок легко сканируемым.",
                "height": 48
            },
            {
                "title": "Компактный шаг margin-bottom: 4px:",
                "desc": "Сжатый межстрочный интервал группирует пункты Wi-Fi, 4K монитор и кресло в монолитный блок преимуществ помещения.",
                "height": 46
            },
            {
                "title": "Готовность к микро-иконкам:",
                "desc": "Структура позволяет легко добавить перед текстом компактные векторные иконки галочек или цветные буллеты через псевдоэлемент ::before.",
                "height": 44
            }
        ]
    },

    # 13. Slide 15: Practice CSS: Card Footer Space-Between
    {
        "type": "code",
        "category": "ФУТЕР КАРТОЧКИ: FLEXBOX",
        "question": "Как разнести стоимость комнаты и кнопки действий по краям блока?",
        "left_title": "css/style.css (.card-footer)",
        "code_lines": [
            [("/* Разделительный подвал карточки */", CODE_COMM)],
            [(".card-footer ", CODE_TAG), ("{", CODE_TEXT)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  justify-content", CODE_PROP), (": ", CODE_TEXT), ("space-between", CODE_VAL), ("; /* Разнос краев */", CODE_COMM)],
            [("  align-items", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";        /* Центр по высоте */", CODE_COMM)],
            [("  border-top", CODE_PROP), (": ", CODE_TEXT), ("1px solid #eeeeee", CODE_VAL), (";  /* Тонкая линия */", CODE_COMM)],
            [("  padding-top", CODE_PROP), (": ", CODE_TEXT), ("10px", CODE_VAL), (";", CODE_TEXT)],
            [("  gap", CODE_PROP), (": ", CODE_TEXT), ("10px", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Связка кнопок действия справа */", CODE_COMM)],
            [(".card-btns ", CODE_TAG), ("{", CODE_TEXT)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("flex", CODE_VAL), (";", CODE_TEXT)],
            [("  align-items", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";", CODE_TEXT)],
            [("  gap", CODE_PROP), (": ", CODE_TEXT), ("8px", CODE_VAL), (";           /* Зазор между кнопками */", CODE_COMM)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "РАЗНОС И ВЫРАВНИВАНИЕ ЭЛЕМЕНТОВ ПРИНЯТИЯ РЕШЕНИЯ",
        "bullet_items": [
            {
                "title": "justify-content: space-between:",
                "desc": "Автоматически прижимает цену (.card-price) к левому краю, а блок кнопок (.card-btns) — к правому краю карточки.",
                "height": 46
            },
            {
                "title": "Тонкая деликатная линия border-top:",
                "desc": "Разделитель цвета #eeeeee визуально отсекает текстовое описание от коммерческого блока целевого действия пользователя.",
                "height": 48
            },
            {
                "title": "Выравнивание по вертикали (align-items):",
                "desc": "Значение center гарантирует, что цифра стоимости и кнопка действия будут строго выровнены по одной горизонтальной оси.",
                "height": 46
            },
            {
                "title": "Вложенный флекс-ряд .card-btns:",
                "desc": "Связывает кнопку-иконку детального просмотра и кнопку 'Забронировать' в компактную аккуратную пару с зазором gap: 8px.",
                "height": 44
            }
        ]
    },

    # 14. Slide 16: Practice CSS: Price Typography (.card-price)
    {
        "type": "code",
        "category": "БЛОК СТОИМОСТИ",
        "question": "Как оформить цену с четким выделением периода тарификации?",
        "left_title": "css/style.css (.card-price)",
        "code_lines": [
            [("/* Базовое начертание суммы */", CODE_COMM)],
            [(".card-price ", CODE_TAG), ("{", CODE_TEXT)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("18px", CODE_VAL), (";", CODE_TEXT)],
            [("  font-weight", CODE_PROP), (": ", CODE_TEXT), ("700", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#222222", CODE_VAL), (";", CODE_TEXT)],
            [("  letter-spacing", CODE_PROP), (": ", CODE_TEXT), ("-0.01em", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Пояснение периода аренды */", CODE_COMM)],
            [(".card-price span ", CODE_TAG), ("{", CODE_TEXT)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("13px", CODE_VAL), (";", CODE_TEXT)],
            [("  font-weight", CODE_PROP), (": ", CODE_TEXT), ("400", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#666666", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "АКЦЕНТИРОВАНИЕ ЦИФР В ЭЛЕКТРОННОЙ КОММЕРЦИИ",
        "bullet_items": [
            {
                "title": "Контраст веса шрифта (700 против 400):",
                "desc": "Крупный жирный шрифт 18px мгновенно считывается взглядом, сразу сообщая арендатору стоимость конкретной локации.",
                "height": 46
            },
            {
                "title": "Изоляция периода через тег <span>:",
                "desc": "Тег <span> позволяет изменить размер и цвет надписи '/ час' прямо внутри общей строки без создания лишних блочных элементов.",
                "height": 48
            },
            {
                "title": "Приглушение тарифа (#666666):",
                "desc": "Серый цвет периода аренды дает необходимый контекст расчета, не отвлекая внимания от ключевой суммы в рублях.",
                "height": 46
            },
            {
                "title": "Отрицательный letter-spacing:",
                "desc": "Значение -0.01em слегка уплотняет межбуквенный интервал цифр, делая ценник более монолитным и премиальным.",
                "height": 44
            }
        ]
    },

    # 15. Slide 17: Practice CSS: Icon Button Component (.btn-icon)
    {
        "type": "code",
        "category": "КНОПКА-ИКОНКА .BTN-ICON",
        "question": "Как сверстать точную квадратную кнопку для SVG-иконки с эффектом наведения?",
        "left_title": "css/style.css (.btn-icon)",
        "code_lines": [
            [(".btn-icon ", CODE_TAG), ("{", CODE_TEXT)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("inline-flex", CODE_VAL), ("; /* Центр SVG */", CODE_COMM)],
            [("  align-items", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";", CODE_TEXT)],
            [("  justify-content", CODE_PROP), (": ", CODE_TEXT), ("center", CODE_VAL), (";", CODE_TEXT)],
            [("  width", CODE_PROP), (": ", CODE_TEXT), ("36px", CODE_VAL), (";", CODE_TEXT)],
            [("  height", CODE_PROP), (": ", CODE_TEXT), ("36px", CODE_VAL), (";", CODE_TEXT)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("0", CODE_VAL), (";", CODE_TEXT)],
            [("  border-radius", CODE_PROP), (": ", CODE_TEXT), ("4px", CODE_VAL), (";", CODE_TEXT)],
            [("  border", CODE_PROP), (": ", CODE_TEXT), ("1px solid #007bff", CODE_VAL), (";", CODE_TEXT)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("transparent", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), (";", CODE_TEXT)],
            [("  cursor", CODE_PROP), (": ", CODE_TEXT), ("pointer", CODE_VAL), (";", CODE_TEXT)],
            [("  text-decoration", CODE_PROP), (": ", CODE_TEXT), ("none", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [(".btn-icon:hover ", CODE_TAG), ("{", CODE_TEXT)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#eaf2ff", CODE_VAL), (";", CODE_TEXT)],
            [("  color", CODE_PROP), (": ", CODE_TEXT), ("#0056b3", CODE_VAL), (";", CODE_TEXT)],
            [("  border-color", CODE_PROP), (": ", CODE_TEXT), ("#0056b3", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "АНАТОМИЯ ИДЕАЛЬНОЙ КНОПКИ-ИКОНКИ",
        "bullet_items": [
            {
                "title": "inline-flex для центрирования SVG:",
                "desc": "Выравнивает векторную иконку глазка ровно по центру квадрата 36x36px без подбора марджинов и паддингов.",
                "height": 46
            },
            {
                "title": "Точные габариты (36x36px):",
                "desc": "Высота 36px в точности совпадает с высотой соседней кнопки 'Забронировать', обеспечивая гармонию в футере.",
                "height": 48
            },
            {
                "title": "Доступность для скринридеров:",
                "desc": "Поскольку в кнопке нет текстовой подписи, обязательны атрибуты aria-label='Подробнее' и всплывающая подсказка title.",
                "height": 46
            },
            {
                "title": "Деликатный отклик при наведении:",
                "desc": "Мягкий голубой фон #eaf2ff и усиление цвета рамки подтверждают пользователю кликабельность элемента.",
                "height": 44
            }
        ]
    },

    # 16. Slide 18: Practice CSS: Button System (.btn-primary & .btn-outline)
    {
        "type": "code",
        "category": "СИСТЕМА КНОПОК",
        "question": "Как разделить геометрию и цветовую схему для кнопок интерфейса?",
        "left_title": "css/style.css (.btn-primary & .btn-outline)",
        "code_lines": [
            [("/* Базовый класс кнопки */", CODE_COMM)],
            [(".btn ", CODE_TAG), ("{", CODE_TEXT)],
            [("  display", CODE_PROP), (": ", CODE_TEXT), ("inline-block", CODE_VAL), (";", CODE_TEXT)],
            [("  padding", CODE_PROP), (": ", CODE_TEXT), ("8px 16px", CODE_VAL), (";", CODE_TEXT)],
            [("  font-size", CODE_PROP), (": ", CODE_TEXT), ("14px", CODE_VAL), (";", CODE_TEXT)],
            [("  font-weight", CODE_PROP), (": ", CODE_TEXT), ("600", CODE_VAL), (";", CODE_TEXT)],
            [("  border-radius", CODE_PROP), (": ", CODE_TEXT), ("4px", CODE_VAL), (";", CODE_TEXT)],
            [("  border", CODE_PROP), (": ", CODE_TEXT), ("1px solid transparent", CODE_VAL), (";", CODE_TEXT)],
            [("  cursor", CODE_PROP), (": ", CODE_TEXT), ("pointer", CODE_VAL), (";", CODE_TEXT)],
            [("  text-decoration", CODE_PROP), (": ", CODE_TEXT), ("none", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("/* Основная акцентная кнопка */", CODE_COMM)],
            [(".btn-primary ", CODE_TAG), ("{", CODE_TEXT)],
            [("  background-color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), ("; color: #fff;", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [(".btn-primary:hover ", CODE_TAG), ("{", CODE_TEXT), (" background-color: #0056b3; }", CODE_TEXT)],
            [("/* Контурная обводная кнопка */", CODE_COMM)],
            [(".btn-outline ", CODE_TAG), ("{", CODE_TEXT)],
            [("  border-color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), ("; color: #007bff;", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [(".btn-outline:hover ", CODE_TAG), ("{", CODE_TEXT), (" background-color: #007bff; color: #fff; }", CODE_TEXT)]
        ],
        "right_title": "ПРИНЦИП КОМПОНЕНТНОГО РАЗДЕЛЕНИЯ КНОПОК",
        "bullet_items": [
            {
                "title": "Базовый класс .btn (Форма):",
                "desc": "Определяет общую геометрию: внутренние отступы, начертание шрифта 600, скругление углов 4px и курсор-руку (cursor: pointer).",
                "height": 46
            },
            {
                "title": "Основное действие (.btn-primary):",
                "desc": "Залита фирменным синим цветом, побуждая пользователя к ключевому конверсионному шагу — бронированию рабочего места.",
                "height": 48
            },
            {
                "title": "Второстепенное действие (.btn-outline):",
                "desc": "Прозрачный фон с цветным контуром. Используется внизу секции для перехода к полному каталогу без перегрузки визуала.",
                "height": 46
            },
            {
                "title": "Универсальность классов:",
                "desc": "Классы .btn применимы как к ссылкам <a>, так и к тегам <button> или <input type='submit'>, сохраняя 100% единый вид.",
                "height": 44
            }
        ]
    },

    # 17. Slide 19: Theory: Media Queries & Responsive Design
    {
        "type": "code",
        "category": "АДАПТИВНАЯ ВЕРСТКА",
        "question": "Как медиа-запросы перестраивают сайт под экраны телефонов и планшетов?",
        "left_title": "css/style.css (@media синтаксис)",
        "code_lines": [
            [("/* 1. Ноутбуки и компактные экраны */", CODE_COMM)],
            [("@media ", CODE_KW), ("(max-width: 992px) ", CODE_TAG), ("{", CODE_TEXT)],
            [("  /* Правила до 992px */", CODE_COMM)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* 2. Планшеты */", CODE_COMM)],
            [("@media ", CODE_KW), ("(max-width: 768px) ", CODE_TAG), ("{", CODE_TEXT)],
            [("  /* Правила до 768px */", CODE_COMM)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* 3. Мобильные телефоны */", CODE_COMM)],
            [("@media ", CODE_KW), ("(max-width: 480px) ", CODE_TAG), ("{", CODE_TEXT)],
            [("  /* Правила до 480px */", CODE_COMM)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Обязательный метатег в HTML head: */", CODE_COMM)],
            [("<!-- <meta name=\"viewport\" ", CODE_COMM)],
            [("         content=\"width=device-width, initial-scale=1.0\"> -->", CODE_COMM)]
        ],
        "right_title": "МЕХАНИКА МЕДИА-ЗАПРОСОВ И BREAKPOINTS",
        "bullet_items": [
            {
                "title": "Директива @media screen:",
                "desc": "Инструмент CSS, который активирует заданные свойства стилей только при выполнении условий ширины окна браузера пользователя.",
                "height": 46
            },
            {
                "title": "Точки останова (Breakpoints):",
                "desc": "Индустриальные стандарты контрольных точек: 1200px (десктоп), 992px (ноутбук), 768px (планшет), 480px (смартфон).",
                "height": 46
            },
            {
                "title": "Критическая важность Viewport:",
                "desc": "Без метатега viewport мобильные браузеры отмасштабируют сайт как микроскопическую картинку, сделав текст нечитаемым!",
                "height": 48
            },
            {
                "title": "Стратегия Desktop-First:",
                "desc": "Мы проектируем базовые стили для десктопа (1200px), а директивами max-width шаг за шагом оптимизируем сетку для меньших устройств.",
                "height": 44
            }
        ]
    },

    # 18. Slide 20: Practice CSS: Responsive Media Queries (768px & 480px)
    {
        "type": "code",
        "category": "АДАПТИВНОСТЬ: МЕДИА-ЗАПРОСЫ",
        "question": "Как перестроить витрину и шапку сайта под экраны планшетов и смартфонов?",
        "left_title": "css/style.css (@media 768px & 480px)",
        "code_lines": [
            [("/* 1. Планшет (max-width: 768px): 2 колонки */", CODE_COMM)],
            [("@media ", CODE_KW), ("(max-width: 768px) ", CODE_TAG), ("{", CODE_TEXT)],
            [("  .container ", CODE_TAG), ("{ width: 100%; padding: 0 16px; }", CODE_TEXT)],
            [("  .room-card ", CODE_TAG), ("{ width: calc(50% - 10px); }", CODE_TEXT)],
            [("  .hero-title ", CODE_TAG), ("{ font-size: 52px; line-height: 1.1; }", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* 2. Смартфон (max-width: 480px): 1 колонка */", CODE_COMM)],
            [("@media ", CODE_KW), ("(max-width: 480px) ", CODE_TAG), ("{", CODE_TEXT)],
            [("  .room-card ", CODE_TAG), ("{ width: 100%; }", CODE_TEXT)],
            [("  .header-container ", CODE_TAG), ("{ flex-direction: column; gap: 12px; }", CODE_TEXT)],
            [("  .btn ", CODE_TAG), ("{ padding: 12px 20px; font-size: 15px; }", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ТАКТИКА АДАПТАЦИИ ПОД ВСЕ ЭКРАНЫ",
        "bullet_items": [
            {
                "title": "Снятие жесткой ширины 1200px:",
                "desc": "Свойство width: 100% устраняет фиксированную ширину и защищает мобильные экраны от появления горизонтального скролла.",
                "height": 46
            },
            {
                "title": "Сетка колонок (2 -> 1):",
                "desc": "На планшетах карточки выстраиваются по две в ряд через calc(50% - 10px), а на смартфонах разворачиваются в 100% ленту.",
                "height": 46
            },
            {
                "title": "Вертикальный стек в шапке:",
                "desc": "Директива flex-direction: column гармонично размещает логотип и ссылки навигации друг под другом на узких дисплеях.",
                "height": 46
            },
            {
                "title": "Увеличение тач-зон (Touch Targets):",
                "desc": "Кнопки и ссылки получают увеличенные отступы (padding), делая интерфейс удобным для быстрого нажатия пальцем.",
                "height": 46
            }
        ]
    },

    # 20. Slide 22: Theory & Practice: Theme Switching via :root
    {
        "type": "code",
        "category": "ПЕРЕИСПОЛЬЗУЕМОСТЬ",
        "question": "Как мгновенно изменить цветовую тему проекта заменой одной строчки в :root?",
        "left_title": "css/style.css (Динамические темы)",
        "code_lines": [
            [("/* Базовая тема (СмартОфис) */", CODE_COMM)],
            [(":root ", CODE_TAG), ("{", CODE_TEXT)],
            [("  --primary-color", CODE_PROP), (": ", CODE_TEXT), ("#007bff", CODE_VAL), (";", CODE_TEXT)],
            [("  --primary-hover", CODE_PROP), (": ", CODE_TEXT), ("#0056b3", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Тема 'Изумрудный Коворкинг' */", CODE_COMM)],
            [("body.theme-emerald ", CODE_TAG), ("{", CODE_TEXT)],
            [("  --primary-color", CODE_PROP), (": ", CODE_TEXT), ("#10b981", CODE_VAL), (";", CODE_TEXT)],
            [("  --primary-hover", CODE_PROP), (": ", CODE_TEXT), ("#059669", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* Тема 'Премиум Фиолетовый' */", CODE_COMM)],
            [("body.theme-purple ", CODE_TAG), ("{", CODE_TEXT)],
            [("  --primary-color", CODE_PROP), (": ", CODE_TEXT), ("#8b5cf6", CODE_VAL), (";", CODE_TEXT)],
            [("  --primary-hover", CODE_PROP), (": ", CODE_TEXT), ("#7c3aed", CODE_VAL), (";", CODE_TEXT)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "КАСКАД ПЕРЕМЕННЫХ И WHITE-LABEL АРХИТЕКТУРА",
        "bullet_items": [
            {
                "title": "Локальное переопределение токенов:",
                "desc": "Селектор body.theme-emerald переопределяет переменные для себя и всех дочерних элементов благодаря каскаду CSS.",
                "height": 46
            },
            {
                "title": "0% переписывания стилей компонентов:",
                "desc": "Кнопки, иконки, ссылки и фокусные рамки карточек читают var(--primary-color) и преображаются автоматически без единой правки.",
                "height": 48
            },
            {
                "title": "Основа White-label разработки:",
                "desc": "Один и тот же движок бронирования офисов можно продавать разным заказчикам, меняя лишь файл переменных брендовой гаммы.",
                "height": 46
            },
            {
                "title": "Переключение в 1 строчку через JS:",
                "desc": "В JavaScript достаточно выполнить document.body.className = 'theme-emerald', чтобы весь интерфейс сменил стиль прямо на лету!",
                "height": 44
            }
        ]
    },

    # 21. Slide 23: Student Mistakes (Grid Cards 2x2)
    {
        "type": "grid_cards",
        "category": "РАЗБОР ОШИБОК",
        "question": "На чем чаще всего спотыкаются студенты при первой верстке сеток и карточек?",
        "title_text": "ТОП-4 ТИПИЧНЫХ ОШИБОК СТУДЕНТОВ ПРИ ВЕРСТКЕ КАРТОЧЕК",
        "cards_data": [
            {
                "title": "Сплющенные и растянутые фото",
                "problem": "Задание жестких width: 100% и height: 200px тегу <img> без свойства object-fit искажает естественные пропорции интерьеров.",
                "solution": "Всегда оборачивать картинку в .card-img-wrap и задавать object-fit: cover с display: block."
            },
            {
                "title": "«Пляшущие» кнопки в футере",
                "problem": "Из-за разного количества строк в описании футеры карточек находятся на разной высоте, ломая нижнюю линию витрины.",
                "solution": "Задать карточке flex-direction: column, а контенту или списку оборудования присвоить распорку flex: 1."
            },
            {
                "title": "Выпадение карточки из строки",
                "problem": "При добавлении border: 1px или padding карточка расширяется на пару пикселей, и третья карточка падает во второй ряд.",
                "solution": "Обязательно подключать box-sizing: border-box в сбросе стилей селектора * в начале файла CSS."
            },
            {
                "title": "Жесткий height у блоков с текстом",
                "problem": "Установка height: 100px для заголовка или списка оборудования приводит к обрезанию или вылезанию текста наружу.",
                "solution": "Задавать только min-height, либо позволять флекс-распоркам автоматически распределять высоту по контенту."
            }
        ]
    },

    # 22. Slide 24: Checklist
    {
        "type": "checklist",
        "category": "ЧЕК-ЛИСТ КАЧЕСТВА",
        "question": "Как убедиться, что секция каталога выполнена на 100% правильно?",
        "title_text": "КРИТЕРИИ КАЧЕСТВА ВЕРСТКИ СЕКЦИИ ПОПУЛЯРНЫХ КОМНАТ",
        "items": [
            {
                "title": "Ровная 3-колоночная сетка витрины:",
                "desc": "Контейнер .rooms-grid использует Flexbox (или Grid) с аккуратными одинаковыми промежутками gap: 20px."
            },
            {
                "title": "Идеальные пропорции всех фотографий:",
                "desc": "Изображения сохраняют естественные пропорции без сплющивания благодаря свойству object-fit: cover."
            },
            {
                "title": "Единая горизонтальная линия кнопок:",
                "desc": "Подвалы всех карточек находятся на одной высоте благодаря колоночному флексу и распорке flex: 1."
            },
            {
                "title": "Интерактивный отклик компонентов (:hover):",
                "desc": "Карточки приподнимаются при наведении, кнопки плавно меняют цвет фона и контура с переходом transition."
            },
            {
                "title": "Централизация палитры через переменные :root:",
                "desc": "Цвета кнопок, ссылок и акцентов подключены через функцию var(), исключая дублирование HEX-кодов."
            },
            {
                "title": "Безупречная адаптивность без полосы скролла:",
                "desc": "При уменьшении экрана карточки аккуратно перестраиваются, а контейнер растягивается без горизонтального скролла."
            }
        ]
    }
]

print(f"Total content slides to generate: {len(slides_data)}")

orig_slide_ids = list(prs.slides._sldIdLst)
title_slide_id = orig_slide_ids[0]
plan_slide_id = orig_slide_ids[1]
result_slide_id = orig_slide_ids[25]
goodbye_slide_id = orig_slide_ids[26]

# 1. Update Slide 1 (Title slide)
slide1 = prs.slides[0]
for sh in slide1.shapes:
    if "Google Shape;59;p13" in sh.name:
        sh.text_frame.text = "Вебинар 3 "
        p = sh.text_frame.paragraphs[0]
        p.runs[0].font.name = "PT Sans Caption"
        p.runs[0].font.size = Pt(10)
    elif sh.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in sh.shapes:
            if sub.name == "TextBox 7":
                sub.text_frame.text = "СОВРЕМЕННЫЕ МЕТОДЫ ВЕРСТКИ"
                sub.text_frame.paragraphs[0].runs[0].font.name = "Montserrat"
                sub.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
                sub.text_frame.paragraphs[0].runs[0].font.bold = True
            elif sub.name == "TextBox 9":
                sub.text_frame.text = "Проектирование \nи разработка интерфейсов пользователя"
print("Slide 1 updated successfully.")

# 2. Update Slide 2 (Plan slide)
slide2 = prs.slides[1]

# Fix layout footers so "Вебинар 2" becomes "Вебинар 3" without resetting fonts
for lyt in prs.slide_layouts:
    for sh in lyt.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                if "Вебинар" in p.text:
                    for r in p.runs:
                        if r.text.strip() == "2":
                            r.text = "3"

plan_titles = {
    "TextBox 12": "1. Сетки Flexbox и CSS Grid",
    "TextBox 15": "2. Адаптивная верстка",
    "TextBox 18": "3. Глобальные переменные",
    "TextBox 21": "4. Практика: Карточки каталога",
    "TextBox 24": "5. Практика: Цветовая схема",
}
plan_subtitles = {
    "TextBox 13": "Технологии Flexbox и CSS Grid",
    "TextBox 16": "Адаптивное отображение через медиа-запросы",
    "TextBox 19": "Работа с глобальными переменными",
    "TextBox 22": "Проектирование гибких карточек для каталога",
    "TextBox 25": "Настройка цветовой схемы через :root",
}

for sh in slide2.shapes:
    if sh.name in plan_titles:
        sh.width = Pt(400)
        sh.text_frame.clear()
        sh.text_frame.word_wrap = False
        p = sh.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = plan_titles[sh.name]
        r.font.name = "Montserrat"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x18, 0x1A, 0x1F)
    elif sh.name in plan_subtitles:
        sh.width = Pt(480)
        sh.text_frame.clear()
        sh.text_frame.word_wrap = False
        p = sh.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = plan_subtitles[sh.name]
        r.font.name = "Inter"
        r.font.size = Pt(9.5)
        r.font.bold = False
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# Strip shadows and theme styles from Slide 2
for sh in slide2.shapes:
    strip_shape_styles_and_shadows(sh)
print("Slide 2 updated successfully.")

# 3. Create all 22 new content slides
for idx, sdata in enumerate(slides_data):
    s = prs.slides.add_slide(blank_layout)
    stype = sdata["type"]
    
    if stype == "code":
        create_code_explanation_slide(
            s, sdata["category"], sdata["question"],
            sdata["left_title"], sdata["code_lines"],
            sdata["right_title"], sdata["bullet_items"]
        )
    elif stype == "diagram_flex_grid":
        create_diagram_flex_grid_slide(s, sdata["category"], sdata["question"])
    elif stype == "card_anatomy":
        create_card_anatomy_slide(s, sdata["category"], sdata["question"])
    elif stype == "grid_cards":
        create_grid_cards_slide(s, sdata["category"], sdata["question"], sdata["title_text"], sdata["cards_data"])
    elif stype == "checklist":
        create_checklist_slide(s, sdata["category"], sdata["question"], sdata["title_text"], sdata["items"])
    
    print(f"  Created Slide {idx+3}: {sdata['category']}")

# 4. Update Result Slide (Slide 26 of original, now to become Slide 25)
slide_res = prs.slides[25]
pic_to_remove = None
tb2_res = None
for sh in slide_res.shapes:
    strip_shape_styles_and_shadows(sh)
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        pic_to_remove = sh
    elif sh.name == "TextBox 2":
        tb2_res = sh

if pic_to_remove is not None:
    spTree = slide_res.shapes._spTree
    spTree.remove(pic_to_remove._element)

# Add Webinar 3 result screenshot: aspect ratio 1210:670 (1.806), centered nicely
res_w = Pt(510)
res_h = Pt(282)
res_left = Pt(105)
res_top = Pt(68)
slide_res.shapes.add_picture(result_img_path, res_left, res_top, res_w, res_h)

# Add or update TextBox 2 on Result slide
if tb2_res is None:
    tb2_res = slide_res.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2_res.name = "TextBox 2"

tf_res = tb2_res.text_frame
tf_res.word_wrap = True
tf_res.margin_left = tf_res.margin_top = tf_res.margin_right = tf_res.margin_bottom = 0
p_res = tf_res.paragraphs[0]
p_res.alignment = PP_ALIGN.LEFT
p_res.text = ""
r_res = p_res.add_run()
r_res.text = "Результат: Адаптивная секция популярных предложений с переиспользуемыми компонентами"
r_res.font.name = "Montserrat"
r_res.font.size = Pt(10)
r_res.font.color.rgb = QUESTION_TEXT
print("Result slide updated with webinar 3 screenshot and TextBox 2.")

# Clean any shadows on Goodbye slide (Slide 26)
slide_goodbye = prs.slides[26]
for sh in slide_goodbye.shapes:
    strip_shape_styles_and_shadows(sh)
print("Goodbye slide shadows stripped.")

# 5. Reorder slide list in presentation XML:
# [Title, Plan] + new_slide_ids + [Result, Goodbye]
all_slide_ids = list(prs.slides._sldIdLst)
new_slide_ids = all_slide_ids[len(orig_slide_ids):]
new_order = [title_slide_id, plan_slide_id] + new_slide_ids + [result_slide_id, goodbye_slide_id]
prs.slides._sldIdLst[:] = new_order

print(f"\nFinal presentation slide count: {len(prs.slides)}")

# 6. Save presentation atomically
tmp_output = output_pptx + ".tmp"
prs.save(tmp_output)
if os.path.exists(output_pptx):
    os.remove(output_pptx)
os.replace(tmp_output, output_pptx)

print(f"\nУСПЕХ! Презентация Вебинара 3 сохранена: {output_pptx}")
