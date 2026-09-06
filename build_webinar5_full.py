import os
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement

# -------------------------------------------------------------
# PATH CONFIGURATION
# -------------------------------------------------------------
base_template_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx"
output_pptx_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-05-dynamic-catalog-data\вебинар 5.pptx"
result_img_path = r"d:\IEEU International East-European University\fullstack-web-development-course\webinar5_catalog_exact.png"

# -------------------------------------------------------------
# PALETTE & STYLES (Strict corporate college theme)
# -------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x18, 0x1A, 0x1F)
CODE_BORDER = RGBColor(0x2D, 0x31, 0x39)
CODE_TEXT = RGBColor(0xD4, 0xD4, 0xD8)
ORANGE_PILL = RGBColor(0xFE, 0x60, 0x02)
ORANGE_LINE = RGBColor(0xFE, 0x60, 0x02)
QUESTION_TEXT = RGBColor(0xA6, 0xA1, 0xA1) # Strictly #A6A1A1
TEXT_PRIMARY = RGBColor(0x18, 0x1A, 0x1F)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
KEYWORD_COLOR = RGBColor(0xC6, 0x78, 0xDD) # Purple
FUNC_COLOR = RGBColor(0x61, 0xAF, 0xEF)    # Blue
STR_COLOR = RGBColor(0x98, 0xC3, 0x79)     # Green
COMMENT_COLOR = RGBColor(0x5C, 0x63, 0x70) # Gray
NUM_COLOR = RGBColor(0xD1, 0x9A, 0x66)     # Orange / number
TAG_COLOR = RGBColor(0xE0, 0x6C, 0x75)     # Red/Coral
ATTR_COLOR = RGBColor(0xD1, 0x9A, 0x66)    # Yellow/Orange

def strip_shape_styles_and_shadows(shape):
    """Zero tolerance to shadows and PowerPoint theme styles."""
    spPr = shape._element.spPr
    for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst'):
        spPr.remove(shd)
    for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw'):
        spPr.remove(shd)
    style = shape._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
    if style is not None:
        shape._element.remove(style)

def add_header_and_footer(slide, category_text, question_text):
    """Standard header pill badge and bottom question prompt."""
    pill_w = Pt(max(140, len(category_text) * 7.8 + 24))
    pill_h = Pt(22.6)
    pill_left = Pt(700) - pill_w
    pill_top = Pt(20.4)
    
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_left, pill_top, pill_w, pill_h)
    badge.name = "Google Shape;69;p14"
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE_PILL
    badge.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    badge.line.width = Pt(1)
    strip_shape_styles_and_shadows(badge)
    
    tf_b = badge.text_frame
    tf_b.word_wrap = False
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = category_text
    r_b.font.name = "Inter"
    r_b.font.size = Pt(12)
    r_b.font.bold = False
    r_b.font.color.rgb = WHITE
    
    # TextBox 2 at bottom
    tb2 = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2.name = "TextBox 2"
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = question_text
    r2.font.name = "Montserrat"
    r2.font.size = Pt(10)
    r2.font.bold = False
    r2.font.color.rgb = QUESTION_TEXT

def create_code_explanation_slide(slide, category, question, left_title, code_lines, right_title, bullet_items):
    """Standard 2-column layout: Left code card, Right clean explanation with bullet items."""
    add_header_and_footer(slide, category, question)
    
    # Left Card
    card_x = Pt(24)
    card_y = Pt(56)
    card_w = Pt(325)
    card_h = Pt(296)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
    card.name = "Rounded Rectangle 3"
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG
    card.line.color.rgb = CODE_BORDER
    card.line.width = Pt(1)
    strip_shape_styles_and_shadows(card)
    
    # Header inside card
    tb_hdr = slide.shapes.add_textbox(card_x + Pt(14), card_y + Pt(10), card_w - Pt(28), Pt(18))
    tb_hdr.name = "TextBox 4"
    tf_h = tb_hdr.text_frame
    tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
    p_h = tf_h.paragraphs[0]
    r_h = p_h.add_run()
    r_h.text = left_title
    r_h.font.name = "Consolas"
    r_h.font.size = Pt(9.5)
    r_h.font.bold = True
    r_h.font.color.rgb = FUNC_COLOR
    
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_x + Pt(14), card_y + Pt(30), card_w - Pt(28), Pt(1))
    div.name = "Rectangle 5"
    div.fill.solid()
    div.fill.fore_color.rgb = ORANGE_LINE
    div.line.fill.background()
    strip_shape_styles_and_shadows(div)
    
    tb_code = slide.shapes.add_textbox(card_x + Pt(14), card_y + Pt(36), card_w - Pt(28), card_h - Pt(44))
    tb_code.name = "TextBox 6"
    tf_c = tb_code.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    
    for i, line_data in enumerate(code_lines):
        p = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
        p.space_after = Pt(2)
        if isinstance(line_data, str):
            r = p.add_run()
            r.text = line_data
            r.font.name = "Consolas"
            r.font.size = Pt(8.5)
            r.font.color.rgb = CODE_TEXT
        elif isinstance(line_data, list):
            for token_text, token_color in line_data:
                r = p.add_run()
                r.text = token_text
                r.font.name = "Consolas"
                r.font.size = Pt(8.5)
                r.font.color.rgb = token_color
                
    # Right Column
    rcol_x = Pt(365)
    rcol_y = Pt(66)
    rcol_w = Pt(330)
    
    tb_rtitle = slide.shapes.add_textbox(rcol_x, rcol_y, rcol_w, Pt(22))
    tb_rtitle.name = "TextBox 7"
    tf_rt = tb_rtitle.text_frame
    tf_rt.word_wrap = True
    tf_rt.margin_left = tf_rt.margin_top = tf_rt.margin_right = tf_rt.margin_bottom = 0
    p_rt = tf_rt.paragraphs[0]
    r_rt = p_rt.add_run()
    r_rt.text = right_title
    r_rt.font.name = "Montserrat"
    r_rt.font.size = Pt(12)
    r_rt.font.bold = True
    r_rt.font.color.rgb = ORANGE_LINE
    
    rdiv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rcol_x, rcol_y + Pt(25), rcol_w, Pt(1))
    rdiv.name = "Rectangle 8"
    rdiv.fill.solid()
    rdiv.fill.fore_color.rgb = ORANGE_LINE
    rdiv.line.fill.background()
    strip_shape_styles_and_shadows(rdiv)
    
    cur_y = rcol_y + Pt(34)
    for idx, item in enumerate(bullet_items):
        item_h = Pt(item.get("height", 42))
        
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rcol_x, cur_y + Pt(3), Pt(7), Pt(7))
        dot.name = f"Rounded Rectangle Dot {idx}"
        dot.adjustments[0] = 0.2
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)
        
        tb_item = slide.shapes.add_textbox(rcol_x + Pt(14), cur_y, rcol_w - Pt(14), item_h)
        tb_item.name = f"TextBox Bullet {idx}"
        tf_i = tb_item.text_frame
        tf_i.word_wrap = True
        tf_i.margin_left = tf_i.margin_top = tf_i.margin_right = tf_i.margin_bottom = 0
        
        p = tf_i.paragraphs[0]
        p.space_after = Pt(0)
        
        r_title = p.add_run()
        r_title.text = item["title"] + ":\n"
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = TEXT_PRIMARY
        
        r_desc = p.add_run()
        r_desc.text = item["desc"]
        r_desc.font.name = "Inter"
        r_desc.font.size = Pt(8.5)
        r_desc.font.bold = False
        r_desc.font.color.rgb = TEXT_MUTED
        
        cur_y += item_h + Pt(6)

def create_data_flow_diagram_slide(slide, category, question):
    """Custom Diagram 1: Data Flow pipeline from data.js to DOM."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. ИСТОЧНИК ДАННЫХ", "subtitle": "Файл js/data.js", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. ОБРАБОТКА И РЕНДЕР", "subtitle": "Функция renderCatalog", "accent": NUM_COLOR},
        {"x": Pt(476), "title": "3. ИНЪЕКЦИЯ В DOM", "subtitle": "Сетка #catalogContainer", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r1 = p_t1.add_run()
        r1.text = c["title"]
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r2 = p_t2.add_run()
        r2.text = c["subtitle"]
        r2.font.name = "Inter"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CODE_TEXT
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = c["accent"]
        line.line.fill.background()
        strip_shape_styles_and_shadows(line)
        
        if c["title"].startswith("1."):
            b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b1.adjustments[0] = 0.1
            b1.fill.solid()
            b1.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b1.line.color.rgb = c["accent"]
            b1.line.width = Pt(1)
            strip_shape_styles_and_shadows(b1)
            p_b1 = b1.text_frame.paragraphs[0]
            p_b1.alignment = PP_ALIGN.CENTER
            r_b1 = p_b1.add_run()
            r_b1.text = "OFFICE_ROOMS = [\n  { id: 'focus-1', ... },\n  { id: 'alpha-2', ... }\n]"
            r_b1.font.name = "Consolas"
            r_b1.font.size = Pt(7)
            r_b1.font.color.rgb = WHITE
            
            tb_desc = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d = tb_desc.text_frame
            tf_d.word_wrap = True
            tf_d.margin_left = tf_d.margin_top = tf_d.margin_right = tf_d.margin_bottom = 0
            p_d = tf_d.paragraphs[0]
            r_d = p_d.add_run()
            r_d.text = "Единый источник правды:\n\n• Массив из 6 объектов\n• Цены за час аренды\n• Вместимость и площадь\n• Список оснащения (Wi-Fi, 4K)\n• Пути к фото в папке img/\n• Флаг isPopular: true/false"
            r_d.font.name = "Inter"
            r_d.font.size = Pt(8)
            r_d.font.color.rgb = CODE_TEXT
            
        elif c["title"].startswith("2."):
            b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b2.adjustments[0] = 0.1
            b2.fill.solid()
            b2.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b2.line.color.rgb = c["accent"]
            b2.line.width = Pt(1)
            strip_shape_styles_and_shadows(b2)
            p_b2 = b2.text_frame.paragraphs[0]
            p_b2.alignment = PP_ALIGN.CENTER
            r_b2 = p_b2.add_run()
            r_b2.text = "OFFICE_ROOMS.map(room => `\n  <div class=\"room-card\">\n    ...${room.title}...\n  </div>\n`).join('')"
            r_b2.font.name = "Consolas"
            r_b2.font.size = Pt(6.5)
            r_b2.font.color.rgb = NUM_COLOR
            
            tb_desc2 = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d2 = tb_desc2.text_frame
            tf_d2.word_wrap = True
            tf_d2.margin_left = tf_d2.margin_top = tf_d2.margin_right = tf_d2.margin_bottom = 0
            p_d2 = tf_d2.paragraphs[0]
            r_d2 = p_d2.add_run()
            r_d2.text = "Конвейер трансформации:\n\n• Обход каждого объекта в массиве\n• Интерполяция ${room.id} в ссылки\n• Вложенный map по удобствам\n• Склейка через .join('')\n• Исключение запятых между блоками\n• Подготовка единой строки HTML"
            r_d2.font.name = "Inter"
            r_d2.font.size = Pt(8)
            r_d2.font.color.rgb = CODE_TEXT
            
        elif c["title"].startswith("3."):
            b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b3.adjustments[0] = 0.1
            b3.fill.solid()
            b3.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b3.line.color.rgb = c["accent"]
            b3.line.width = Pt(1)
            strip_shape_styles_and_shadows(b3)
            p_b3 = b3.text_frame.paragraphs[0]
            p_b3.alignment = PP_ALIGN.CENTER
            r_b3 = p_b3.add_run()
            r_b3.text = "container = #catalogContainer;\ncontainer.innerHTML =\n  htmlCardsMarkup;"
            r_b3.font.name = "Consolas"
            r_b3.font.size = Pt(7)
            r_b3.font.color.rgb = STR_COLOR
            
            tb_desc3 = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d3 = tb_desc3.text_frame
            tf_d3.word_wrap = True
            tf_d3.margin_left = tf_d3.margin_top = tf_d3.margin_right = tf_d3.margin_bottom = 0
            p_d3 = tf_d3.paragraphs[0]
            r_d3 = p_d3.add_run()
            r_d3.text = "Отрисовка в интерфейсе:\n\n• Браузер парсит HTML в узлы DOM\n• Раскладка в адаптивный CSS Grid\n• Активация ссылок на room-details\n• Подстановка onerror для фото\n• Кнопки быстрого бронирования\n• Автоматический пересчет сетки"
            r_d3.font.name = "Inter"
            r_d3.font.size = Pt(8)
            r_d3.font.color.rgb = CODE_TEXT

def create_routing_diagram_slide(slide, category, question):
    """Custom Diagram 2: Dynamic URL routing and params extraction."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. КЛИК В КАТАЛОГЕ", "subtitle": "Переход по ссылке", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. АНАЛИЗ АДРЕСА", "subtitle": "Чтение URLSearchParams", "accent": KEYWORD_COLOR},
        {"x": Pt(476), "title": "3. ДЕТАЛИ ИЛИ 404", "subtitle": "Отрисовка результата", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r1 = p_t1.add_run()
        r1.text = c["title"]
        r1.font.name = "Montserrat"
        r1.font.size = Pt(9.5)
        r1.font.bold = True
        r1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r2 = p_t2.add_run()
        r2.text = c["subtitle"]
        r2.font.name = "Inter"
        r2.font.size = Pt(7.5)
        r2.font.color.rgb = CODE_TEXT
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = c["accent"]
        line.line.fill.background()
        strip_shape_styles_and_shadows(line)
        
        if c["title"].startswith("1."):
            b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b1.adjustments[0] = 0.1
            b1.fill.solid()
            b1.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b1.line.color.rgb = c["accent"]
            b1.line.width = Pt(1)
            strip_shape_styles_and_shadows(b1)
            p_b1 = b1.text_frame.paragraphs[0]
            p_b1.alignment = PP_ALIGN.CENTER
            r_b1 = p_b1.add_run()
            r_b1.text = "<a href=\"room-details.html\n?id=focus-1\">\n  Подробнее\n</a>"
            r_b1.font.name = "Consolas"
            r_b1.font.size = Pt(7)
            r_b1.font.color.rgb = WHITE
            
            tb_desc = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d = tb_desc.text_frame
            tf_d.word_wrap = True
            tf_d.margin_left = tf_d.margin_top = tf_d.margin_right = tf_d.margin_bottom = 0
            p_d = tf_d.paragraphs[0]
            r_d = p_d.add_run()
            r_d.text = "Инициация запроса:\n\n• Пользователь кликает по карточке\n• Браузер формирует GET-запрос\n• Знак '?' отделяет query-строку\n• Передается параметр id=focus-1\n• Открывается общий шаблон страницы\n• Адрес доступен для шеринга"
            r_d.font.name = "Inter"
            r_d.font.size = Pt(8)
            r_d.font.color.rgb = CODE_TEXT
            
        elif c["title"].startswith("2."):
            b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b2.adjustments[0] = 0.1
            b2.fill.solid()
            b2.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b2.line.color.rgb = c["accent"]
            b2.line.width = Pt(1)
            strip_shape_styles_and_shadows(b2)
            p_b2 = b2.text_frame.paragraphs[0]
            p_b2.alignment = PP_ALIGN.CENTER
            r_b2 = p_b2.add_run()
            r_b2.text = "params = new URLSearchParams(\n  window.location.search\n);\nroomId = params.get('id');\nroom = OFFICE_ROOMS.find(...);"
            r_b2.font.name = "Consolas"
            r_b2.font.size = Pt(6.5)
            r_b2.font.color.rgb = KEYWORD_COLOR
            
            tb_desc2 = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d2 = tb_desc2.text_frame
            tf_d2.word_wrap = True
            tf_d2.margin_left = tf_d2.margin_top = tf_d2.margin_right = tf_d2.margin_bottom = 0
            p_d2 = tf_d2.paragraphs[0]
            r_d2 = p_d2.add_run()
            r_d2.text = "Парсинг и поиск сущности:\n\n• Считывание параметров адресной строки\n• Извлечение ключа: params.get('id')\n• Метод .find() сканирует массив\n• Поиск совпадения: r.id === roomId\n• Возврат найденного объекта\n• Если ID неверен: возвращает undefined"
            r_d2.font.name = "Inter"
            r_d2.font.size = Pt(8)
            r_d2.font.color.rgb = CODE_TEXT
            
        elif c["title"].startswith("3."):
            b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Pt(14), top_y + Pt(54), col_w - Pt(28), Pt(54))
            b3.adjustments[0] = 0.1
            b3.fill.solid()
            b3.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            b3.line.color.rgb = c["accent"]
            b3.line.width = Pt(1)
            strip_shape_styles_and_shadows(b3)
            p_b3 = b3.text_frame.paragraphs[0]
            p_b3.alignment = PP_ALIGN.CENTER
            r_b3 = p_b3.add_run()
            r_b3.text = "if (room) {\n  renderDetailsCard(room);\n} else {\n  renderEmptyNotFound();\n}"
            r_b3.font.name = "Consolas"
            r_b3.font.size = Pt(7)
            r_b3.font.color.rgb = STR_COLOR
            
            tb_desc3 = slide.shapes.add_textbox(x + Pt(14), top_y + Pt(120), col_w - Pt(28), Pt(140))
            tf_d3 = tb_desc3.text_frame
            tf_d3.word_wrap = True
            tf_d3.margin_left = tf_d3.margin_top = tf_d3.margin_right = tf_d3.margin_bottom = 0
            p_d3 = tf_d3.paragraphs[0]
            r_d3 = p_d3.add_run()
            r_d3.text = "Безопасное ветвление UX:\n\n• УСПЕХ: отображение фото, бейджей,\n  описания, списка оборудования и\n  кнопки бронирования с ID\n• ОШИБКА: дружелюбный fallback-блок\n  'Комната не найдена' с кнопкой\n  возврата в каталог"
            r_d3.font.name = "Inter"
            r_d3.font.size = Pt(8)
            r_d3.font.color.rgb = CODE_TEXT

def create_grid_cards_slide(slide, category, question, title_text, cards_data):
    """2x2 error analysis cards."""
    add_header_and_footer(slide, category, question)
    
    coords = [
        (Pt(24), Pt(56), Pt(325), Pt(142)),
        (Pt(365), Pt(56), Pt(330), Pt(142)),
        (Pt(24), Pt(210), Pt(325), Pt(142)),
        (Pt(365), Pt(210), Pt(330), Pt(142)),
    ]
    
    for idx, (x, y, w, h) in enumerate(coords):
        data = cards_data[idx]
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.adjustments[0] = 0.04
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        box.line.width = Pt(1)
        strip_shape_styles_and_shadows(box)

        tb_h = slide.shapes.add_textbox(x + Pt(14), y + Pt(10), w - Pt(28), Pt(20))
        tf_h = tb_h.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
        p_h = tf_h.paragraphs[0]
        r_num = p_h.add_run()
        r_num.text = f"{idx + 1}. {data['title']}"
        r_num.font.name = "Montserrat"
        r_num.font.size = Pt(10)
        r_num.font.bold = True
        r_num.font.color.rgb = ORANGE_LINE
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(14), y + Pt(32), w - Pt(28), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE_LINE
        line.line.fill.background()
        strip_shape_styles_and_shadows(line)
        
        tb_body = slide.shapes.add_textbox(x + Pt(14), y + Pt(38), w - Pt(28), Pt(94))
        tf_b = tb_body.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        p1 = tf_b.paragraphs[0]
        p1.space_after = Pt(3)
        r1_lbl = p1.add_run()
        r1_lbl.text = "Ошибка: "
        r1_lbl.font.name = "Inter"
        r1_lbl.font.size = Pt(8.5)
        r1_lbl.font.bold = True
        r1_lbl.font.color.rgb = TEXT_PRIMARY
        
        r1_txt = p1.add_run()
        r1_txt.text = data["problem"]
        r1_txt.font.name = "Inter"
        r1_txt.font.size = Pt(8)
        r1_txt.font.color.rgb = TEXT_MUTED
        
        p2 = tf_b.add_paragraph()
        r2_lbl = p2.add_run()
        r2_lbl.text = "Решение: "
        r2_lbl.font.name = "Inter"
        r2_lbl.font.size = Pt(8.5)
        r2_lbl.font.bold = True
        r2_lbl.font.color.rgb = RGBColor(0x0E, 0x9F, 0x6E) # Green
        
        r2_txt = p2.add_run()
        r2_txt.text = data["solution"]
        r2_txt.font.name = "Inter"
        r2_txt.font.size = Pt(8)
        r2_txt.font.color.rgb = TEXT_MUTED

def create_checklist_slide(slide, category, question, title_text, items):
    """Quality checklist with clean checkmark cards."""
    add_header_and_footer(slide, category, question)
    
    tb_t = slide.shapes.add_textbox(Pt(24), Pt(56), Pt(670), Pt(22))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    r_t = p_t.add_run()
    r_t.text = title_text
    r_t.font.name = "Montserrat"
    r_t.font.size = Pt(11)
    r_t.font.bold = True
    r_t.font.color.rgb = ORANGE_LINE
    
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(24), Pt(80), Pt(670), Pt(1))
    div.fill.solid()
    div.fill.fore_color.rgb = ORANGE_LINE
    div.line.fill.background()
    strip_shape_styles_and_shadows(div)
    
    y_start = Pt(92)
    row_h = Pt(40)
    for idx, item in enumerate(items):
        cur_y = y_start + idx * (row_h + Pt(4))
        
        # Check badge
        chk = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(24), cur_y + Pt(2), Pt(24), Pt(24))
        chk.adjustments[0] = 0.2
        chk.fill.solid()
        chk.fill.fore_color.rgb = RGBColor(0x0E, 0x9F, 0x6E) # Green
        chk.line.fill.background()
        strip_shape_styles_and_shadows(chk)
        
        p_c = chk.text_frame.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        r_c = p_c.add_run()
        r_c.text = "✔"
        r_c.font.name = "Segoe UI"
        r_c.font.size = Pt(10)
        r_c.font.color.rgb = WHITE
        
        # Text
        tb_item = slide.shapes.add_textbox(Pt(58), cur_y, Pt(636), row_h)
        tf_i = tb_item.text_frame
        tf_i.word_wrap = True
        tf_i.margin_left = tf_i.margin_top = tf_i.margin_right = tf_i.margin_bottom = 0
        
        p = tf_i.paragraphs[0]
        r_lbl = p.add_run()
        r_lbl.text = item["title"] + ": "
        r_lbl.font.name = "Montserrat"
        r_lbl.font.size = Pt(9.5)
        r_lbl.font.bold = True
        r_lbl.font.color.rgb = TEXT_PRIMARY
        
        p2 = tf_i.add_paragraph()
        r_desc = p2.add_run()
        r_desc.text = item["desc"]
        r_desc.font.name = "Inter"
        r_desc.font.size = Pt(8.5)
        r_desc.font.color.rgb = TEXT_MUTED

# -------------------------------------------------------------
# PRESENTATION DATA (21 content slides)
# -------------------------------------------------------------
slides_data = [
    # 1. Slide 3: Code/Theory: Object Model & Structure
    {
        "type": "code",
        "category": "СТРУКТУРА ДАННЫХ",
        "question": "Как правильно организовать хранение каталога комнат на стороне клиента?",
        "left_title": "📄 js/data.js (Модель сущности комнаты)",
        "code_lines": [
            [("// Модель объекта офисной комнаты", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("room = {", CODE_TEXT)],
            [("  id: ", FUNC_COLOR), ("'focus-1'", STR_COLOR), (",", CODE_TEXT)],
            [("  title: ", FUNC_COLOR), ("'Мини-офис Focus'", STR_COLOR), (",", CODE_TEXT)],
            [("  pricePerHour: ", FUNC_COLOR), ("450", NUM_COLOR), (",", CODE_TEXT)],
            [("  capacity: ", FUNC_COLOR), ("'1-2 человека'", STR_COLOR), (",", CODE_TEXT)],
            [("  area: ", FUNC_COLOR), ("'12 м²'", STR_COLOR), (",", CODE_TEXT)],
            [("  description: ", FUNC_COLOR), ("'Тихое пространство...'", STR_COLOR), (",", CODE_TEXT)],
            [("  equipment: [", FUNC_COLOR)],
            [("    'Wi-Fi 500 Мбит/с'", STR_COLOR), (",", CODE_TEXT)],
            [("    '4K Монитор'", STR_COLOR), (",", CODE_TEXT)],
            [("    'Эргономичное кресло'", STR_COLOR)],
            [("  ],", CODE_TEXT)],
            [("  image: ", FUNC_COLOR), ("'../img/room-1.jpg'", STR_COLOR), (",", CODE_TEXT)],
            [("  isPopular: ", FUNC_COLOR), ("true", KEYWORD_COLOR)],
            [("};", CODE_TEXT)]
        ],
        "right_title": "МОДЕЛИРОВАНИЕ СУЩНОСТИ: КАРТОЧКА ОБЪЕКТА",
        "bullet_items": [
            {
                "title": "Уникальный ID (Идентификатор)",
                "desc": "Первичный ключ id: 'focus-1' однозначно адресует комнату при формировании ссылок и поиске в массиве.",
                "height": 40
            },
            {
                "title": "Типизированные поля параметров",
                "desc": "Числовая цена pricePerHour используется в расчетах калькулятора, а строки описывают вместимость и площадь.",
                "height": 40
            },
            {
                "title": "Вложенные списки (equipment)",
                "desc": "Массив строк equipment хранит список преимуществ и техники, которые динамически выводятся в теги <li>.",
                "height": 40
            },
            {
                "title": "Флаги состояния (isPopular)",
                "desc": "Булево значение true / false управляет показом акцентного бейджа «Популярное» на карточке.",
                "height": 40
            }
        ]
    },

    # 2. Slide 4: Code/Practice: Database Module (data.js)
    {
        "type": "code",
        "category": "БАЗА ДАННЫХ JS",
        "question": "Почему данные выносятся в отдельный файл data.js перед основным скриптом?",
        "left_title": "📄 js/data.js (Массив всех комнат)",
        "code_lines": [
            [("// Глобальный массив доступных помещений", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("OFFICE_ROOMS = [", CODE_TEXT)],
            [("  {", CODE_TEXT)],
            [("    id: ", FUNC_COLOR), ("'focus-1'", STR_COLOR), (",", CODE_TEXT)],
            [("    title: ", FUNC_COLOR), ("'Мини-офис Focus'", STR_COLOR), (",", CODE_TEXT)],
            [("    pricePerHour: ", FUNC_COLOR), ("450", NUM_COLOR), (",", CODE_TEXT)],
            [("    // ...", COMMENT_COLOR)],
            [("  },", CODE_TEXT)],
            [("  {", CODE_TEXT)],
            [("    id: ", FUNC_COLOR), ("'alpha-2'", STR_COLOR), (",", CODE_TEXT)],
            [("    title: ", FUNC_COLOR), ("'Конференц-зал Alpha'", STR_COLOR), (",", CODE_TEXT)],
            [("    pricePerHour: ", FUNC_COLOR), ("1200", NUM_COLOR), (",", CODE_TEXT)],
            [("    // ...", COMMENT_COLOR)],
            [("  }", CODE_TEXT)],
            [("  // Всего 6 офисных помещений", COMMENT_COLOR)],
            [("];", CODE_TEXT)]
        ],
        "right_title": "РАЗДЕЛЕНИЕ ДАННЫХ И ЛОГИКИ: МОДУЛЬ DATA.JS",
        "bullet_items": [
            {
                "title": "Single Source of Truth",
                "desc": "Каталог, главная страница и страница описания используют единый источник данных, исключая рассинхрон.",
                "height": 40
            },
            {
                "title": "Отказ от ручной верстки",
                "desc": "Вместо создания 6 отдельных HTML-файлов для каждой комнаты мы генерируем интерфейс из одного массива.",
                "height": 40
            },
            {
                "title": "Порядок подключения скриптов",
                "desc": "Файл data.js обязательно подключается перед main.js, чтобы переменная OFFICE_ROOMS была доступна.",
                "height": 40
            },
            {
                "title": "Подготовка к работе с API",
                "desc": "В будущем этот массив будет приходить от настоящего бэкенд-сервера через fetch без изменения логики.",
                "height": 40
            }
        ]
    },

    # 3. Slide 5: Diagram: Data Flow Pipeline
    {
        "type": "data_flow",
        "category": "СХЕМА: DATA FLOW",
        "question": "Каков путь данных от исходного массива в data.js до отображения на экране?"
    },

    # 4. Slide 6: Code/Theory: Array.map
    {
        "type": "code",
        "category": "МЕТОД MAP",
        "question": "Чем метод map() превосходит классический цикл for при генерации интерфейса?",
        "left_title": "📄 js/main.js (Трансформация через .map)",
        "code_lines": [
            [("// Преобразование массива объектов в разметку", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("htmlCards = OFFICE_ROOMS.", CODE_TEXT), ("map", FUNC_COLOR), ("(room => {", CODE_TEXT)],
            [("  return `", STR_COLOR)],
            [("    <div class=\"room-card\">", STR_COLOR)],
            [("      <img src=\"${room.image}\" alt=\"${room.title}\">", STR_COLOR)],
            [("      <h3>${room.title}</h3>", STR_COLOR)],
            [("      <div class=\"price\">${room.pricePerHour} ₽/ч</div>", STR_COLOR)],
            [("    </div>", STR_COLOR)],
            [("  `;", STR_COLOR)],
            [("});", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Результат: массив строк ['<div>...</div>', ...]", COMMENT_COLOR)],
            [("console.", FUNC_COLOR), ("log", FUNC_COLOR), ("(htmlCards.length); ", CODE_TEXT), ("// 6 элементов", COMMENT_COLOR)]
        ],
        "right_title": "ДЕКЛАРАТИВНАЯ ГЕНЕРАЦИЯ: ARRAY.MAP()",
        "bullet_items": [
            {
                "title": "Принцип неизменяемости (Immutability)",
                "desc": "Метод .map() не меняет исходный массив комнат, а создает новый массив с результатами преобразования.",
                "height": 40
            },
            {
                "title": "Лаконичный стрелочный синтаксис",
                "desc": "Запись room => `...` избавляет от громоздких циклов с ручным отслеживанием индексов i++.",
                "height": 40
            },
            {
                "title": "Интерполяция выражений",
                "desc": "Внутри шаблонной строки можно подставлять любые переменные, тернарные операторы и формулы.",
                "height": 40
            },
            {
                "title": "Промышленный стандарт веб-разработки",
                "desc": "Трансформация данных через .map() лежит в основе генерации списков во всех фреймворках (React, Vue).",
                "height": 40
            }
        ]
    },

    # 5. Slide 7: Code/Theory: Array.join('')
    {
        "type": "code",
        "category": "МЕТОД JOIN",
        "question": "Почему без вызова .join('') на веб-странице появляются лишние запятые?",
        "left_title": "📄 js/main.js (Опасность метода toString)",
        "code_lines": [
            [("// ❌ ОШИБКА: Массив вставляется напрямую:", COMMENT_COLOR)],
            [("container.", CODE_TEXT), ("innerHTML", FUNC_COLOR), (" = OFFICE_ROOMS.", CODE_TEXT), ("map", FUNC_COLOR), ("(...);", CODE_TEXT)],
            [("// JS вызовет toString() -> на экране появятся запятые!", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("//  ПРАВИЛЬНО: Склейка без разделителя:", COMMENT_COLOR)],
            [("container.", CODE_TEXT), ("innerHTML", FUNC_COLOR), (" = OFFICE_ROOMS", CODE_TEXT)],
            [("  .", CODE_TEXT), ("map", FUNC_COLOR), ("(room => `", CODE_TEXT)],
            [("    <div class=\"room-card\">", STR_COLOR)],
            [("      <!-- разметка карточки -->", COMMENT_COLOR)],
            [("    </div>", STR_COLOR)],
            [("  `)", CODE_TEXT)],
            [("  .", CODE_TEXT), ("join", FUNC_COLOR), ("(''); ", KEYWORD_COLOR), ("// Склейка пустой строкой", COMMENT_COLOR)]
        ],
        "right_title": "СКЛЕЙКА СТРОК: РОЛЬ МЕТОДА JOIN('')",
        "bullet_items": [
            {
                "title": "Поведение innerHTML с массивами",
                "desc": "Свойство innerHTML ожидает единую строку. Если передать массив, браузер неявно преобразует его в строку.",
                "height": 40
            },
            {
                "title": "Нежелательные символы запятых",
                "desc": "Стандартный toString() разделяет элементы массива запятыми, которые отображаются как мусор в интерфейсе.",
                "height": 40
            },
            {
                "title": "Склейка пустой строкой join('')",
                "desc": "Метод объединяет массив строк в единую монолитную HTML-разметку без каких-либо промежуточных знаков.",
                "height": 40
            },
            {
                "title": "Высокая производительность",
                "desc": "Операция .map().join('') работает в десятки раз быстрее последовательного прибавления innerHTML += в цикле.",
                "height": 40
            }
        ]
    },

    # 6. Slide 8: Code/Practice: Catalog Renderer
    {
        "type": "code",
        "category": "РЕНДЕРИНГ КАТАЛОГА",
        "question": "Как безопасно наполнить контейнер каталога сгенерированными карточками?",
        "left_title": "📄 js/main.js (Функция renderCatalog)",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("renderCatalog", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("container = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('catalogContainer');", STR_COLOR)],
            [("  if ", KEYWORD_COLOR), ("(!container || ", CODE_TEXT), ("typeof ", KEYWORD_COLOR), ("OFFICE_ROOMS === 'undefined') ", STR_COLOR), ("return;", KEYWORD_COLOR)],
            [("", CODE_TEXT)],
            [("  container.", CODE_TEXT), ("innerHTML", FUNC_COLOR), (" = OFFICE_ROOMS.", CODE_TEXT), ("map", FUNC_COLOR), ("(room => `", CODE_TEXT)],
            [("    <div class=\"room-card\">", STR_COLOR)],
            [("      <div class=\"card-img-wrap\">", STR_COLOR)],
            [("        <a href=\"room-details.html?id=${room.id}\">", STR_COLOR)],
            [("          <img src=\"${room.image}\" alt=\"${room.title}\"", STR_COLOR)],
            [("               onerror=\"this.src='../img/no-image.svg'\">", STR_COLOR)],
            [("        </a>", STR_COLOR)],
            [("      </div>", STR_COLOR)],
            [("      <!-- тело и подвал карточки -->", COMMENT_COLOR)],
            [("    </div>", STR_COLOR)],
            [("  `).", CODE_TEXT), ("join", FUNC_COLOR), ("('');", STR_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ФУНКЦИЯ RENDERCATALOG: МОНТИРОВАНИЕ В DOM",
        "bullet_items": [
            {
                "title": "Защитные проверки (Guard Clauses)",
                "desc": "Проверка if (!container) гарантирует, что функция не упадет с ошибкой на страницах, где каталога нет.",
                "height": 40
            },
            {
                "title": "Атомарное обновление DOM",
                "desc": "Присваивание container.innerHTML происходит за одну операцию, исключая мигание интерфейса.",
                "height": 40
            },
            {
                "title": "Динамические ссылки с ID",
                "desc": "Каждая ссылка формируется с уникальным GET-параметром ?id=${room.id} для точного перехода.",
                "height": 40
            },
            {
                "title": "Трехуровневая структура карточки",
                "desc": "Карточка логически разделена на обертку фото (card-img-wrap), контент (card-content) и подвал с кнопками.",
                "height": 40
            }
        ]
    },

    # 7. Slide 9: Code/Practice: Nested List Rendering (equipment)
    {
        "type": "code",
        "category": "ВЛОЖЕННЫЕ СПИСКИ",
        "question": "Как сгенерировать теги <li> для массива характеристик внутри карточки?",
        "left_title": "📄 js/main.js (Вложенный маппинг)",
        "code_lines": [
            [("// Вывод списка оснащения внутри карточки", COMMENT_COLOR)],
            [("<ul class=\"card-equipment\">", STR_COLOR)],
            [("  ${room.equipment.", STR_COLOR), ("map", FUNC_COLOR), ("(item => `", STR_COLOR)],
            [("    <li>${item}</li>", STR_COLOR)],
            [("  `).", STR_COLOR), ("join", FUNC_COLOR), ("('')}", STR_COLOR)],
            [("</ul>", STR_COLOR)],
            [("", CODE_TEXT)],
            [("// Пример исходных данных:", COMMENT_COLOR)],
            [("// equipment: ['Wi-Fi 500 Мбит/с', '4K Монитор']", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("// Сгенерированный HTML в DOM:", COMMENT_COLOR)],
            [("// <li>Wi-Fi 500 Мбит/с</li>", COMMENT_COLOR)],
            [("// <li>4K Монитор</li>", COMMENT_COLOR)]
        ],
        "right_title": "ВЛОЖЕННЫЕ ТРАНСФОРМАЦИИ: СПИСОК В СПИСКЕ",
        "bullet_items": [
            {
                "title": "Вложенный вызов .map()",
                "desc": "Внутри интерполяции внешней карточки запускается внутренний цикл по массиву удобств конкретной комнаты.",
                "height": 40
            },
            {
                "title": "Обязательный внутренний .join('')",
                "desc": "Вложенный маппинг также требует вызова .join(''), иначе между тегами <li> появятся запятые.",
                "height": 40
            },
            {
                "title": "Адаптивность к числу элементов",
                "desc": "Если в комнате 3 опции, отрисуется 3 пункта; если 5 — отрисуется 5, верстка подстраивается автоматически.",
                "height": 40
            },
            {
                "title": "Семантическая разметка <ul> и <li>",
                "desc": "Использование стандартных тегов списков обеспечивает доступность сайта для программ чтения экрана.",
                "height": 40
            }
        ]
    },

    # 8. Slide 10: Code/Practice: Media Fallbacks (onerror)
    {
        "type": "code",
        "category": "ОБРАБОТКА ФОТО",
        "question": "Что делать, если фотография комнаты не загрузилась или путь указан с опечаткой?",
        "left_title": "📄 HTML / JS (Атрибут onerror)",
        "code_lines": [
            [("// Защита от битых изображений (Broken Images)", COMMENT_COLOR)],
            [("<img", STR_COLOR)],
            [("  src=\"${room.image}\"", STR_COLOR)],
            [("  alt=\"${room.title}\"", STR_COLOR)],
            [("  class=\"card-img\"", STR_COLOR)],
            [("  onerror=\"this.src='../img/no-image.svg'\"", KEYWORD_COLOR)],
            [(">", STR_COLOR)],
            [("", CODE_TEXT)],
            [("// Механика работы:", COMMENT_COLOR)],
            [("// 1. Браузер пытается загрузить картинку по src", COMMENT_COLOR)],
            [("// 2. При ошибке 404 срабатывает событие 'error'", COMMENT_COLOR)],
            [("// 3. Код подменяет src на векторную заглушку", COMMENT_COLOR)]
        ],
        "right_title": "БЕЗОПАСНАЯ ЗАГРУЗКА: АТРИБУТ ONERROR",
        "bullet_items": [
            {
                "title": "Проблема битых изображений",
                "desc": "При ошибке в пути или сетевом сбое браузер показывает сломанную иконку, разрушающую дизайн карточки.",
                "height": 40
            },
            {
                "title": "Событие ошибки загрузки (error)",
                "desc": "Тег <img> генерирует событие ошибки, если ресурс недоступен или сервер вернул статус 404/500.",
                "height": 40
            },
            {
                "title": "Подмена источника через this.src",
                "desc": "Ключевое слово this ссылается на сам элемент <img>, мгновенно переключая источник на резервный SVG.",
                "height": 40
            },
            {
                "title": "Премиальный пользовательский опыт",
                "desc": "Пользователь всегда видит аккуратный нейтральный плейсхолдер с логотипом сервиса вместо пустоты.",
                "height": 40
            }
        ]
    },

    # 9. Slide 11: Code/Theory: GET Parameters & Query String
    {
        "type": "code",
        "category": "GET-ПАРАМЕТРЫ URL",
        "question": "Как связать карточку каталога со страницей детального описания?",
        "left_title": "📄 Анатомия ссылки с параметром",
        "code_lines": [
            [("// Ссылка перехода в карточке каталога", COMMENT_COLOR)],
            [("<a href=\"room-details.html?id=${room.id}\" class=\"btn-icon\">", STR_COLOR)],
            [("  <svg>...</svg>", STR_COLOR)],
            [("</a>", STR_COLOR)],
            [("", CODE_TEXT)],
            [("// Анатомия итогового адреса в браузере:", COMMENT_COLOR)],
            [("// /pages/room-details.html?id=focus-1", FUNC_COLOR)],
            [("// ──────────┬─────────── ─┬─ ───┬───", COMMENT_COLOR)],
            [("//          │             │     └── Значение: focus-1", COMMENT_COLOR)],
            [("//          │             └── Разделитель параметров", COMMENT_COLOR)],
            [("//          └── Путь к файлу шаблона", COMMENT_COLOR)]
        ],
        "right_title": "QUERY STRING: ПЕРЕДАЧА СОСТОЯНИЯ В URL",
        "bullet_items": [
            {
                "title": "Разделитель параметров (?)",
                "desc": "Символ вопросительного знака отделяет путь к файлу страницы от блока передаваемых параметров.",
                "height": 40
            },
            {
                "title": "Формат ключ-значение (key=value)",
                "desc": "Параметры передаются парами. Знак равенства связывает имя переменной (id) с ее значением (focus-1).",
                "height": 40
            },
            {
                "title": "Сохраняемость ссылок (Deep Linking)",
                "desc": "Ссылку на конкретную комнату можно отправить в мессенджере или сохранить в закладки браузера.",
                "height": 40
            },
            {
                "title": "Динамический шаблон",
                "desc": "Браузер загружает один HTML-файл, а JavaScript считывает ID и наполняет страницу нужным контентом.",
                "height": 40
            }
        ]
    },

    # 10. Slide 12: Diagram: Routing Flow
    {
        "type": "routing",
        "category": "СХЕМА: МАРШРУТИЗАЦИЯ",
        "question": "Какова последовательность перехода из каталога на страницу деталей комнаты?"
    },

    # 11. Slide 13: Code/Theory: URLSearchParams API
    {
        "type": "code",
        "category": "URLSEARCHPARAMS",
        "question": "Как браузер извлекает значение конкретного параметра из строки search?",
        "left_title": "📄 js/main.js (Парсинг URLSearchParams)",
        "code_lines": [
            [("// 1. Получаем поисковую строку запроса: '?id=focus-1'", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("searchStr = window.location.search;", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 2. Создаем объект-парсер URLSearchParams", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("urlParams = ", CODE_TEXT), ("new ", KEYWORD_COLOR), ("URLSearchParams", FUNC_COLOR), ("(searchStr);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// 3. Извлекаем значение нужного параметра по ключу", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("roomId = urlParams.", CODE_TEXT), ("get", FUNC_COLOR), ("('id'); ", STR_COLOR), ("// 'focus-1'", COMMENT_COLOR)],
            [("", CODE_TEXT)],
            [("// Страховка от альтернативных имен ключа:", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("safeId = urlParams.", CODE_TEXT), ("get", FUNC_COLOR), ("('id') || urlParams.", STR_COLOR), ("get", FUNC_COLOR), ("('room');", STR_COLOR)]
        ],
        "right_title": "ВСТРОЕННЫЙ ИНТЕРФЕЙС URLSEARCHPARAMS",
        "bullet_items": [
            {
                "title": "Современный Web API",
                "desc": "Интерфейс URLSearchParams избавляет от написания громоздких регулярных выражений для разбора URL.",
                "height": 40
            },
            {
                "title": "Метод urlParams.get('key')",
                "desc": "Возвращает декодированное значение параметра либо null, если такой параметр не был передан в строке.",
                "height": 40
            },
            {
                "title": "Автоматическое декодирование",
                "desc": "Класс автоматически переводит закодированные символы (пробелы %20, кириллицу) в читаемый текст.",
                "height": 40
            },
            {
                "title": "Мульти-ключевой захват",
                "desc": "Конструкция get('id') || get('room') делает скрипт устойчивым к разным вариантам ссылок.",
                "height": 40
            }
        ]
    },

    # 12. Slide 14: Code/Theory: Array.find
    {
        "type": "code",
        "category": "ПОИСК ЧЕРЕЗ FIND",
        "question": "В чем разница между точечным поиском через find() и фильтрацией filter()?",
        "left_title": "📄 js/main.js (Метод Array.find)",
        "code_lines": [
            [("// Поиск единственной комнаты по совпадению ID", COMMENT_COLOR)],
            [("const ", KEYWORD_COLOR), ("room = OFFICE_ROOMS.", CODE_TEXT), ("find", FUNC_COLOR), ("(r => r.id === roomId);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("if ", KEYWORD_COLOR), ("(room) {", CODE_TEXT)],
            [("  console.", FUNC_COLOR), ("log", FUNC_COLOR), ("('Найдено:', room.title);", STR_COLOR)],
            [("  // room — это конкретный объект {}", COMMENT_COLOR)],
            [("} ", CODE_TEXT), ("else ", KEYWORD_COLOR), ("{", CODE_TEXT)],
            [("  console.", FUNC_COLOR), ("warn", FUNC_COLOR), ("('Комната не найдена в базе!');", STR_COLOR)],
            [("  // room имеет значение undefined", COMMENT_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "ТОЧНЫЙ ПОИСК В МАССИВЕ: ARRAY.FIND()",
        "bullet_items": [
            {
                "title": "Остановка на первом совпадении",
                "desc": "Метод .find() немедленно прерывает перебор, как только предикат вернул true, работая максимально быстро.",
                "height": 40
            },
            {
                "title": "Возврат объекта вместо массива",
                "desc": "В отличие от .filter(), который всегда возвращает массив, .find() возвращает сам целевой объект.",
                "height": 40
            },
            {
                "title": "Строгое сравнение идентификаторов",
                "desc": "Условие r.id === roomId сравнивает уникальные строки без неявного приведения типов.",
                "height": 40
            },
            {
                "title": "Результат при отсутствии элемента",
                "desc": "Если элемента с таким ID нет, метод возвращает undefined, что позволяет легко обработать ошибку 404.",
                "height": 40
            }
        ]
    },

    # 13. Slide 15: Code/Practice: Room Details HTML Scaffold
    {
        "type": "code",
        "category": "СТРАНИЦА ДЕТАЛЕЙ",
        "question": "Какую базовую структуру должен иметь HTML-шаблон для динамической страницы?",
        "left_title": "📄 pages/room-details.html (Каркас)",
        "code_lines": [
            [("<!DOCTYPE html>", COMMENT_COLOR)],
            [("<html lang=\"ru\">", TAG_COLOR)],
            [("<head>", TAG_COLOR)],
            [("  <link rel=\"stylesheet\" href=\"../css/style.css\">", TAG_COLOR)],
            [("  <script src=\"../js/data.js\" defer></script>", TAG_COLOR)],
            [("  <script src=\"../js/main.js\" defer></script>", TAG_COLOR)],
            [("</head>", TAG_COLOR)],
            [("<body>", TAG_COLOR)],
            [("  <!-- Шапка сайта -->", COMMENT_COLOR)],
            [("  <main class=\"main\">", TAG_COLOR)],
            [("    <div class=\"container\" id=\"roomDetailsContainer\">", TAG_COLOR)],
            [("      <!-- Контент монтирует JavaScript -->", COMMENT_COLOR)],
            [("    </div>", TAG_COLOR)],
            [("  </main>", TAG_COLOR)],
            [("  <!-- Подвал сайта -->", COMMENT_COLOR)],
            [("</body>", TAG_COLOR)]
        ],
        "right_title": "ШАБЛОН СТРАНИЦЫ: ROOM-DETAILS.HTML",
        "bullet_items": [
            {
                "title": "Один шаблон для всех объектов",
                "desc": "Вместо десятков статических страниц поддерживается один универсальный каркас с контейнером.",
                "height": 40
            },
            {
                "title": "Относительные пути ../",
                "desc": "Поскольку страница лежит в папке pages/, пути к стилям и скриптам поднимаются на уровень выше (../).",
                "height": 40
            },
            {
                "title": "Строгий порядок defer-скриптов",
                "desc": "Скрипт data.js подключается перед main.js, гарантируя готовность массива комнат к моменту запуска.",
                "height": 40
            },
            {
                "title": "Целевой контейнер инъекции",
                "desc": "Блок #roomDetailsContainer служит точкой входа, куда JavaScript помещает разметку карточки.",
                "height": 40
            }
        ]
    },

    # 14. Slide 16: Code/Practice: initRoomDetails logic
    {
        "type": "code",
        "category": "ЛОГИКА ДЕТАЛЕЙ",
        "question": "Как связать получение параметров URL и вывод детальной карточки?",
        "left_title": "📄 js/main.js (Функция initRoomDetails)",
        "code_lines": [
            [("function ", KEYWORD_COLOR), ("initRoomDetails", FUNC_COLOR), ("() {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("container = document.", CODE_TEXT), ("getElementById", FUNC_COLOR), ("('roomDetailsContainer');", STR_COLOR)],
            [("  if ", KEYWORD_COLOR), ("(!container || ", CODE_TEXT), ("typeof ", KEYWORD_COLOR), ("OFFICE_ROOMS === 'undefined') ", STR_COLOR), ("return;", KEYWORD_COLOR)],
            [("", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("urlParams = ", CODE_TEXT), ("new ", KEYWORD_COLOR), ("URLSearchParams", FUNC_COLOR), ("(window.location.search);", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("roomId = urlParams.", CODE_TEXT), ("get", FUNC_COLOR), ("('id') || urlParams.", STR_COLOR), ("get", FUNC_COLOR), ("('room');", STR_COLOR)],
            [("  const ", KEYWORD_COLOR), ("room = OFFICE_ROOMS.", CODE_TEXT), ("find", FUNC_COLOR), ("(r => r.id === roomId);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  if ", KEYWORD_COLOR), ("(!room) {", CODE_TEXT)],
            [("    // Вывод сообщения об ошибке (Empty State)", COMMENT_COLOR)],
            [("    return;", KEYWORD_COLOR)],
            [("  }", CODE_TEXT)],
            [("  // Рендеринг подробного описания комнаты", COMMENT_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "АРХИТЕКТУРА ИНИЦИАЛИЗАЦИИ СТРАНИЦЫ",
        "bullet_items": [
            {
                "title": "Изоляция выполнения скрипта",
                "desc": "Проверка if (!container) останавливает выполнение функции на всех остальных страницах проекта.",
                "height": 40
            },
            {
                "title": "Извлечение идентификатора",
                "desc": "Скрипт читает параметр id из адресной строки и передает его в поиск по базе данных комнат.",
                "height": 40
            },
            {
                "title": "Разделение сценариев (Branching)",
                "desc": "Четкая граница: успешный показ характеристик комнаты либо переход в сценарий обработки ошибки 404.",
                "height": 40
            },
            {
                "title": "Интеграция в DOMContentLoaded",
                "desc": "Функция запускается автоматически при готовности DOM-дерева в центральном обработчике событий.",
                "height": 40
            }
        ]
    },

    # 15. Slide 17: Code/Practice: Fallback / 404 Empty State
    {
        "type": "code",
        "category": "ОБРАБОТКА 404",
        "question": "Что увидит пользователь, если перейдет по несуществующему адресу ?id=xxx?",
        "left_title": "📄 js/main.js (Обработка ошибки 404)",
        "code_lines": [
            [("if ", KEYWORD_COLOR), ("(!room) {", CODE_TEXT)],
            [("  container.", CODE_TEXT), ("innerHTML", FUNC_COLOR), (" = `", CODE_TEXT)],
            [("    <div class=\"empty-message\">", STR_COLOR)],
            [("      <h2>Комната не найдена</h2>", STR_COLOR)],
            [("      <p style=\"margin: 10px 0 20px 0;\">", STR_COLOR)],
            [("        Возможно, ссылка устарела или комната была удалена.", STR_COLOR)],
            [("      </p>", STR_COLOR)],
            [("      <a href=\"catalog.html\" class=\"btn btn-primary\">", STR_COLOR)],
            [("        Вернуться в каталог", STR_COLOR)],
            [("      </a>", STR_COLOR)],
            [("    </div>", STR_COLOR)],
            [("  `;", STR_COLOR)],
            [("  return;", KEYWORD_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "EMPTY STATE: ДРУЖЕЛЮБНЫЙ FALLBACK",
        "bullet_items": [
            {
                "title": "Защита от падения скрипта",
                "desc": "Без проверки обращение к room.title вызовет TypeError, оставив пользователя перед пустым белым экраном.",
                "height": 40
            },
            {
                "title": "Понятное объяснение проблемы",
                "desc": "Интерфейс вежливо сообщает студенту или посетителю причину отсутствия карточки на русском языке.",
                "height": 40
            },
            {
                "title": "Кнопка целевого возврата (CTA)",
                "desc": "Ссылка «Вернуться в каталог» помогает посетителю продолжить работу, не уходя с сайта.",
                "height": 40
            },
            {
                "title": "Ранний возврат (Early Return)",
                "desc": "Инструкция return прерывает функцию, предотвращая попытку генерации несуществующих данных.",
                "height": 40
            }
        ]
    },

    # 16. Slide 18: Code/Practice: Render Details Card & Badges
    {
        "type": "code",
        "category": "ДЕТАЛЬНАЯ КАРТОЧКА",
        "question": "Как отобразить бейджи характеристик и статус популярности через тернарный оператор?",
        "left_title": "📄 js/main.js (Рендеринг деталей)",
        "code_lines": [
            [("container.", CODE_TEXT), ("innerHTML", FUNC_COLOR), (" = `", CODE_TEXT)],
            [("  <div class=\"room-details-card\">", STR_COLOR)],
            [("    <div class=\"room-details-gallery\">", STR_COLOR)],
            [("      <img src=\"${room.image}\" alt=\"${room.title}\">", STR_COLOR)],
            [("    </div>", STR_COLOR)],
            [("    <div class=\"room-details-info\">", STR_COLOR)],
            [("      <h1>${room.title}</h1>", STR_COLOR)],
            [("      <div class=\"price\">${room.pricePerHour} ₽/ч</div>", STR_COLOR)],
            [("      <div class=\"room-badges\">", STR_COLOR)],
            [("        <span class=\"badge\">${room.capacity}</span>", STR_COLOR)],
            [("        <span class=\"badge\">${room.area}</span>", STR_COLOR)],
            [("        ${room.isPopular ? ", STR_COLOR), ("'<span class=\"popular\">Популярное</span>'", KEYWORD_COLOR), (" : ''}", STR_COLOR)],
            [("      </div>", STR_COLOR)],
            [("      <p class=\"desc\">${room.description}</p>", STR_COLOR)],
            [("    </div>", STR_COLOR)],
            [("  </div>", STR_COLOR)],
            [("`;", STR_COLOR)]
        ],
        "right_title": "УСЛОВНЫЙ РЕНДЕРИНГ В ШАБЛОНЕ",
        "bullet_items": [
            {
                "title": "Тернарный оператор для бейджей",
                "desc": "Выражение room.isPopular ? '<span...>' : '' добавляет плашку «Популярное» только нужным помещениям.",
                "height": 40
            },
            {
                "title": "Двухколоночная верстка страницы",
                "desc": "Галерея комнаты слева и текстовые характеристики справа создают баланс восприятия контента.",
                "height": 40
            },
            {
                "title": "Прямая подстановка свойств",
                "desc": "Заголовок, стоимость аренды, вместимость и площадь подставляются без ручной манипуляции каждым элементом.",
                "height": 40
            },
            {
                "title": "Единый шаблонный блок",
                "desc": "Карточка собирается в одну строку, сводя количество обращений к DOM к минимальному числу.",
                "height": 40
            }
        ]
    },

    # 17. Slide 19: Code/Practice: Booking Link Integration
    {
        "type": "code",
        "category": "БРОНИРОВАНИЕ",
        "question": "Как передать выбранную комнату дальше в форму бронирования?",
        "left_title": "📄 js/main.js (Кнопки действий)",
        "code_lines": [
            [("<div class=\"room-details-actions\">", STR_COLOR)],
            [("  <!-- Ссылка с предвыбором комнаты в форме -->", COMMENT_COLOR)],
            [("  <a href=\"booking.html?room=${room.id}\"", STR_COLOR)],
            [("     class=\"btn btn-primary\">", STR_COLOR)],
            [("    Забронировать эту комнату", STR_COLOR)],
            [("  </a>", STR_COLOR)],
            [("", CODE_TEXT)],
            [("  <!-- Кнопка быстрого возврата в каталог -->", COMMENT_COLOR)],
            [("  <a href=\"catalog.html\"", STR_COLOR)],
            [("     class=\"btn btn-outline\">", STR_COLOR)],
            [("    ← Назад в каталог", STR_COLOR)],
            [("  </a>", STR_COLOR)],
            [("</div>", STR_COLOR)]
        ],
        "right_title": "СКВОЗНАЯ ПЕРЕДАЧА ID В БРОНИРОВАНИЕ",
        "bullet_items": [
            {
                "title": "Непрерывная цепочка данных",
                "desc": "ID комнаты передается из каталога в детали, а затем в форму создания брони (booking.html?room=focus-1).",
                "height": 40
            },
            {
                "title": "Автовыбор комнаты в форме",
                "desc": "На следующем занятии страница бронирования автоматически выберет эту комнату по значению из URL.",
                "height": 40
            },
            {
                "title": "Контурная кнопка btn-outline",
                "desc": "Второстепенное действие возврата оформлено менее контрастной кнопкой, не отвлекая от бронирования.",
                "height": 40
            },
            {
                "title": "Архитектура воронки оформления",
                "desc": "Плавный и логичный путь пользователя: поиск в каталоге → изучение деталей → оформление заявки.",
                "height": 40
            }
        ]
    },

    # 18. Slide 20: Code/Practice: CSS Styling for Details Card
    {
        "type": "code",
        "category": "СТИЛИЗАЦИЯ ДЕТАЛЕЙ",
        "question": "Как сделать верстку страницы описания удобной и читаемой на любых экранах?",
        "left_title": "📄 css/style.css (Стили детальной карточки)",
        "code_lines": [
            [(".room-details-card {", FUNC_COLOR)],
            [("  display: ", KEYWORD_COLOR), ("grid;", CODE_TEXT)],
            [("  grid-template-columns: ", KEYWORD_COLOR), ("1fr 1fr;", NUM_COLOR)],
            [("  gap: ", KEYWORD_COLOR), ("32px;", NUM_COLOR)],
            [("  background: ", KEYWORD_COLOR), ("#FFFFFF;", CODE_TEXT)],
            [("  border: ", KEYWORD_COLOR), ("1px solid #EEEEEE;", CODE_TEXT)],
            [("  border-radius: ", KEYWORD_COLOR), ("12px;", NUM_COLOR)],
            [("  padding: ", KEYWORD_COLOR), ("32px;", NUM_COLOR)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [(".room-details-img {", FUNC_COLOR)],
            [("  width: ", KEYWORD_COLOR), ("100%; ", NUM_COLOR), ("height: ", KEYWORD_COLOR), ("380px;", NUM_COLOR)],
            [("  object-fit: ", KEYWORD_COLOR), ("cover; ", KEYWORD_COLOR), ("border-radius: ", KEYWORD_COLOR), ("8px;", NUM_COLOR)],
            [("}", CODE_TEXT)]
        ],
        "right_title": "СОВРЕМЕННЫЙ ДИЗАЙН: GRID ДЕТАЛЬНОЙ КАРТОЧКИ",
        "bullet_items": [
            {
                "title": "Сетка 1fr 1fr в CSS Grid",
                "desc": "Пропорциональное разделение блока 50/50 между фотографией помещения и колонкой с описанием.",
                "height": 40
            },
            {
                "title": "Свойство object-fit: cover",
                "desc": "Защищает фотографии комнат от искажений и сплющивания, аккуратно кадрируя их под единый размер.",
                "height": 40
            },
            {
                "title": "Единый визуальный язык",
                "desc": "Белый фон, тонкая граница #EEEEEE и радиус 12px полностью согласованы с карточками каталога.",
                "height": 40
            },
            {
                "title": "Адаптивность для смартфонов",
                "desc": "Через медиа-запрос @media (max-width: 768px) сетка автоматически перестраивается в одну колонку.",
                "height": 40
            }
        ]
    },

    # 19. Slide 21: Code/Practice: Active Nav Sync for Details Page
    {
        "type": "code",
        "category": "УМНАЯ НАВИГАЦИЯ",
        "question": "Какой пункт меню должен быть подсвечен, когда пользователь изучает комнату?",
        "left_title": "📄 js/main.js (Подсветка раздела Каталог)",
        "code_lines": [
            [("// Подсветка родительского раздела в initNavigation", COMMENT_COLOR)],
            [("links.", CODE_TEXT), ("forEach", FUNC_COLOR), ("(link => {", CODE_TEXT)],
            [("  const ", KEYWORD_COLOR), ("href = link.", CODE_TEXT), ("getAttribute", FUNC_COLOR), ("('href');", STR_COLOR)],
            [("", CODE_TEXT)],
            [("  if (", KEYWORD_COLOR)],
            [("    (href.", CODE_TEXT), ("includes", FUNC_COLOR), ("('catalog.html') || href.", STR_COLOR), ("includes", FUNC_COLOR), ("('room-details.html')) &&", STR_COLOR)],
            [("    (current.", CODE_TEXT), ("includes", FUNC_COLOR), ("('catalog.html') || current.", STR_COLOR), ("includes", FUNC_COLOR), ("('room-details.html'))", STR_COLOR)],
            [("  ) {", CODE_TEXT)],
            [("    link.classList.", CODE_TEXT), ("add", FUNC_COLOR), ("('active');", STR_COLOR)],
            [("  }", CODE_TEXT)],
            [("});", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Пункт 'Каталог' подсвечен даже на room-details!", COMMENT_COLOR)]
        ],
        "right_title": "СИНХРОНИЗАЦИЯ НАВИГАЦИИ С ДОЧЕРНИМИ СТРАНИЦАМИ",
        "bullet_items": [
            {
                "title": "Проблема страниц без пунктов меню",
                "desc": "Страница room-details отсутствует в шапке как отдельная ссылка, поэтому меню теряло бы подсветку.",
                "height": 40
            },
            {
                "title": "Контекстная принадлежность",
                "desc": "Детальный просмотр является дочерним экраном каталога и должен визуально ассоциироваться с ним.",
                "height": 40
            },
            {
                "title": "Проверка обоих адресов через includes()",
                "desc": "Скрипт проверяет наличие room-details в адресе и активирует пункт «Каталог» в шапке сайта.",
                "height": 40
            },
            {
                "title": "Интуитивное понимание контекста",
                "desc": "Студенты и пользователи в любой момент понимают, в каком разделе веб-сервиса они находятся.",
                "height": 40
            }
        ]
    },

    # 20. Slide 22: Student Mistakes (Grid Cards 2x2)
    {
        "type": "grid_cards",
        "category": "РАЗБОР ОШИБОК",
        "question": "На чем чаще всего спотыкаются студенты при выводе каталога и страниц деталей?",
        "title_text": "ТОП-4 ТИПИЧНЫХ ОШИБОК СТУДЕНТОВ ПРИ РАБОТЕ С ДАННЫМИ И URL",
        "cards_data": [
            {
                "title": "Забытый вызов .join('') после метода .map()",
                "problem": "На веб-странице между карточками комнат появляются лишние запятые, превращая дизайн в неопрятный список.",
                "solution": "Всегда дописывать .join('') в конце цепочки маппинга при передаче результата в свойство .innerHTML."
            },
            {
                "title": "Ошибка порядка: data.js подключен после main.js",
                "problem": "Консоль браузера выдает фатальную ошибку Uncaught ReferenceError: OFFICE_ROOMS is not defined.",
                "solution": "Тег <script src='../js/data.js' defer> обязательно должен стоять в HTML перед скриптом main.js."
            },
            {
                "title": "Опечатка в имени GET-параметра: urlParams.get('id')",
                "problem": "Переменная roomId оказывается null, потому что в ссылке было написано ?room_id= вместо ?id=.",
                "solution": "Использовать страховку urlParams.get('id') || urlParams.get('room') и проверять совпадение ключей."
            },
            {
                "title": "Чтение свойств у undefined при ненайденной комнате",
                "problem": "При опечатке в адресе код падает с ошибкой Cannot read properties of undefined (reading 'title').",
                "solution": "Всегда добавлять ранний выход if (!room) с выводом дружелюбного сообщения «Комната не найдена»."
            }
        ]
    },

    # 21. Slide 23: Checklist
    {
        "type": "checklist",
        "category": "ЧЕК-ЛИСТ КАЧЕСТВА",
        "question": "Как проверить надежность работы каталога и страницы описания?",
        "title_text": "КРИТЕРИИ КАЧЕСТВА ДИНАМИЧЕСКОГО КАТАЛОГА И ДЕТАЛЕЙ",
        "items": [
            {
                "title": "Все 6 комнат динамически рендерятся из data.js",
                "desc": "В каталоге выводятся реальные карточки без ручной верстки дублей в HTML."
            },
            {
                "title": "Отсутствуют лишние запятые между карточками",
                "desc": "Метод map() завершается чистым объединением строк через пустой разделитель .join('')."
            },
            {
                "title": "Клик по карточке открывает страницу с GET-параметром",
                "desc": "Ссылка формирует корректный адрес вида room-details.html?id=focus-1."
            },
            {
                "title": "Страница деталей выводит полные характеристики",
                "desc": "Отображаются фото, цена, площадь, вместимость, бейджи и маркированный список удобств."
            },
            {
                "title": "Надежная обработка некорректного ID (404 Fallback)",
                "desc": "При вводе несуществующего ID выводится блок «Комната не найдена» и ссылка возврата в каталог."
            },
            {
                "title": "Кнопка бронирования передает ID в форму",
                "desc": "Ссылка ведет на booking.html?room=... для следующего шага оформления бронирования."
            }
        ]
    }
]

print(f"Total content slides to generate: {len(slides_data)}")

# -------------------------------------------------------------
# PRESENTATION BUILDER EXECUTION
# -------------------------------------------------------------
prs = Presentation(base_template_path)
blank_layout = prs.slide_layouts[6]

orig_slide_ids = list(prs.slides._sldIdLst)
title_slide_id = orig_slide_ids[0]
plan_slide_id = orig_slide_ids[1]
result_slide_id = orig_slide_ids[25] # Slide 26 in original template
goodbye_slide_id = orig_slide_ids[26] # Slide 27 in original template

# 1. Update Slide 1 (Title slide)
slide1 = prs.slides[0]
for sh in slide1.shapes:
    if sh.name == "Google Shape;59;p13" and sh.has_text_frame:
        sh.text_frame.text = "Вебинар 5 "
    elif sh.shape_type == MSO_SHAPE_TYPE.GROUP and sh.name == "Группа 12":
        for sub in sh.shapes:
            if sub.name == "TextBox 7":
                sub.text_frame.text = "РАБОТА С ДИНАМИЧЕСКИМИ ДАННЫМИ"
                sub.text_frame.paragraphs[0].runs[0].font.name = "Montserrat"
                sub.text_frame.paragraphs[0].runs[0].font.size = Pt(15)
                sub.text_frame.paragraphs[0].runs[0].font.bold = True
            elif sub.name == "TextBox 9":
                sub.text_frame.text = "Проектирование \nи разработка интерфейсов пользователя"
print("Slide 1 updated successfully.")

# 2. Update Slide 2 (Plan slide)
slide2 = prs.slides[1]
plan_titles = {
    "TextBox 12": "1. Структуры данных",
    "TextBox 15": "2. Генерация контента",
    "TextBox 18": "3. Подготовка mock-данных",
    "TextBox 21": "4. Практика: Отрисовка каталога",
    "TextBox 24": "5. Результат и итоги",
}
plan_subtitles = {
    "TextBox 13": "Объекты и массивы в JavaScript: моделирование сущностей",
    "TextBox 16": "Циклы, условия и метод map() для создания интерфейса",
    "TextBox 19": "Создание файла js/data.js с массивом офисных комнат",
    "TextBox 22": "Автоматическая отрисовка каталога на основе массива объектов",
    "TextBox 25": "Функциональная страница каталога, разбор ошибок и чек-лист",
}

for sh in slide2.shapes:
    if sh.name in plan_titles:
        sh.text_frame.clear()
        p = sh.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = plan_titles[sh.name]
        r.font.name = "Montserrat SemiBold"
        r.font.size = Pt(11.25)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x18, 0x1A, 0x1F)
    elif sh.name in plan_subtitles:
        sh.width = Pt(450)
        sh.text_frame.clear()
        sh.text_frame.word_wrap = False
        p = sh.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = plan_subtitles[sh.name]
        r.font.name = "Inter"
        r.font.size = Pt(9.75)
        r.font.bold = False
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
print("Slide 2 updated successfully.")

# 3. Create all 21 new content slides
for idx, sdata in enumerate(slides_data):
    s = prs.slides.add_slide(blank_layout)
    stype = sdata["type"]
    
    if stype == "code":
        create_code_explanation_slide(
            s, sdata["category"], sdata["question"],
            sdata["left_title"], sdata["code_lines"],
            sdata["right_title"], sdata["bullet_items"]
        )
    elif stype == "data_flow":
        create_data_flow_diagram_slide(s, sdata["category"], sdata["question"])
    elif stype == "routing":
        create_routing_diagram_slide(s, sdata["category"], sdata["question"])
    elif stype == "grid_cards":
        create_grid_cards_slide(s, sdata["category"], sdata["question"], sdata["title_text"], sdata["cards_data"])
    elif stype == "checklist":
        create_checklist_slide(s, sdata["category"], sdata["question"], sdata["title_text"], sdata["items"])
    
    print(f"  Created Slide {idx+3}: {sdata['category']}")

# 4. Update Result Slide (Slide 24)
slide_res = prs.slides[25]
pic_to_remove = None
tb2_res = None
for sh in slide_res.shapes:
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        pic_to_remove = sh
    elif sh.name == "TextBox 2":
        tb2_res = sh

if pic_to_remove is not None:
    spTree = slide_res.shapes._spTree
    spTree.remove(pic_to_remove._element)

# Add Webinar 5 result screenshot: aspect ratio 1210:670 (1.806), centered nicely
res_w = Pt(510)
res_h = Pt(282)
res_left = Pt(105)
res_top = Pt(68)
slide_res.shapes.add_picture(result_img_path, res_left, res_top, res_w, res_h)

# Add or update TextBox 2 on Result slide
if tb2_res is None:
    tb2_res = slide_res.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(550), Pt(20))
    tb2_res.name = "TextBox 2"

tf_res = tb2_res.text_frame
tf_res.word_wrap = True
tf_res.margin_left = tf_res.margin_top = tf_res.margin_right = tf_res.margin_bottom = 0
p_res = tf_res.paragraphs[0]
p_res.alignment = PP_ALIGN.LEFT
p_res.text = ""
r_res = p_res.add_run()
r_res.text = "Результат: Функциональная страница каталога с динамическим рендерингом"
r_res.font.name = "Montserrat"
r_res.font.size = Pt(10)
r_res.font.color.rgb = QUESTION_TEXT
print("Result slide updated with webinar 5 screenshot and TextBox 2.")

# Clean shadows on Goodbye slide (Slide 25)
slide_goodbye = prs.slides[26]
for sh in slide_goodbye.shapes:
    strip_shape_styles_and_shadows(sh)
print("Goodbye slide shadows stripped.")

# 5. Reorder slide list in presentation XML:
# [Title, Plan] + new_slide_ids + [Result, Goodbye]
all_slide_ids = list(prs.slides._sldIdLst)
new_slide_ids = all_slide_ids[len(orig_slide_ids):]
new_order = [title_slide_id, plan_slide_id] + new_slide_ids + [result_slide_id, goodbye_slide_id]
prs.slides._sldIdLst[:] = new_order

# 6. Save presentation
os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
prs.save(output_pptx_path)

print(f"\nFinal presentation slide count: {len(prs.slides)}")
print(f"\nУСПЕХ! Презентация Вебинара 5 сохранена: {output_pptx_path}")
