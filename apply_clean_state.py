import sys
import os
from pptx import Presentation
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

targets = [
    r"05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx",
    r"05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - дополненный.pptx"
]

for path in targets:
    if not os.path.exists(path):
        continue
    prs = Presentation(path)
    print(f"\nProcessing {os.path.basename(path)}...")
    
    # Ensure Slide 3 points to rId54 (Theory slide)
    # Ensure Slide 22 points to rId22 (Goodbye slide)
    sldIdLst = prs.slides._sldIdLst
    rid_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    
    sldIdLst[2].attrib[rid_attr] = 'rId54'
    sldIdLst[21].attrib[rid_attr] = 'rId22'
    
    # Check if rId55 is in rels and pop it to avoid duplicate part reference
    if 'rId55' in prs.part.rels:
        prs.part.rels.pop('rId55')
        print("  Removed unused rId55.")

    # Now inspect and polish Slide 3
    slide3 = prs.slides[2]
    for s in slide3.shapes:
        if s.name == "TextBox 6":
            # Fix first line to HTML comment
            p0 = s.text_frame.paragraphs[0]
            if len(p0.runs) > 0:
                p0.runs[0].text = "<!-- === 1. HTML: АНАТОМИЯ ЭЛЕМЕНТА === -->"
                p0.runs[0].font.color.rgb = RGBColor(0x71, 0x71, 0x7A)
        elif s.name == "TextBox 2":
            # Ensure #A6A1A1 color
            for p in s.text_frame.paragraphs:
                p.font.color.rgb = RGBColor(0xA6, 0xA1, 0xA1)
                for r in p.runs:
                    r.font.color.rgb = RGBColor(0xA6, 0xA1, 0xA1)

    prs.save(path)
    print(f"  Successfully saved {os.path.basename(path)}.")

print("\nFinished apply_clean_state.")
