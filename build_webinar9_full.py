import os
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

# -------------------------------------------------------------
# PATH CONFIGURATION
# -------------------------------------------------------------
base_template_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx"
output_pptx_path = r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-09-final-assembly-my-bookings\вебинар 9.pptx"
result_img_path = r"d:\IEEU International East-European University\fullstack-web-development-course\webinar9_result_exact.png"

# -------------------------------------------------------------
# PALETTE & STYLES (Strict corporate college theme)
# -------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x18, 0x1A, 0x1F)
CODE_BORDER = RGBColor(0x2D, 0x31, 0x39)
CODE_TEXT = RGBColor(0xD4, 0xD4, 0xD8)
ORANGE_PILL = RGBColor(0xFE, 0x60, 0x02)
ORANGE_LINE = RGBColor(0xFE, 0x60, 0x02)
QUESTION_TEXT = RGBColor(0xA6, 0xA1, 0xA1) # Strictly #A6A1A1
TEXT_PRIMARY = RGBColor(0x18, 0x1A, 0x1F)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
KEYWORD_COLOR = RGBColor(0xC6, 0x78, 0xDD) # Purple
FUNC_COLOR = RGBColor(0x61, 0xAF, 0xEF)    # Blue
STR_COLOR = RGBColor(0x98, 0xC3, 0x79)     # Green
COMMENT_COLOR = RGBColor(0x5C, 0x63, 0x70) # Gray
NUM_COLOR = RGBColor(0xD1, 0x9A, 0x66)     # Orange / number
TAG_COLOR = RGBColor(0xE0, 0x6C, 0x75)     # Red/Coral
DANGER_COLOR = RGBColor(0xDC, 0x26, 0x26)  # Red for invalid state

def strip_shape_styles_and_shadows(shape):
    """Zero tolerance to shadows and PowerPoint theme styles."""
    spPr = getattr(shape._element, 'spPr', None)
    if spPr is not None:
        for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst'):
            spPr.remove(shd)
        for shd in spPr.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw'):
            spPr.remove(shd)
    style = shape._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}style')
    if style is not None:
        shape._element.remove(style)

def add_header_and_footer(slide, category_text, question_text):
    """Standard header pill badge with dynamic width and bottom question prompt."""
    pill_w = Pt(max(140, len(category_text) * 7.8 + 24))
    pill_h = Pt(22.6)
    pill_left = Pt(700) - pill_w
    pill_top = Pt(20.4)
    
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_left, pill_top, pill_w, pill_h)
    badge.name = "Google Shape;69;p14"
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE_PILL
    badge.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    badge.line.width = Pt(1)
    strip_shape_styles_and_shadows(badge)
    
    tf_b = badge.text_frame
    tf_b.word_wrap = False
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = category_text
    r_b.font.name = "Inter"
    r_b.font.size = Pt(12)
    r_b.font.bold = False
    r_b.font.color.rgb = WHITE
    
    # TextBox 2 at bottom
    tb2 = slide.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2.name = "TextBox 2"
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = question_text
    r2.font.name = "Montserrat"
    r2.font.size = Pt(10)
    r2.font.bold = False
    r2.font.color.rgb = QUESTION_TEXT

def create_code_explanation_slide(slide, category, question, left_title, code_lines, right_title, bullet_items):
    """Standard 2-column layout: Left code card, Right clean explanation with bullet items."""
    add_header_and_footer(slide, category, question)
    
    # Left Card
    card_x = Pt(24)
    card_y = Pt(56)
    card_w = Pt(325)
    card_h = Pt(296)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
    card.name = "Rounded Rectangle 1"
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG
    card.line.color.rgb = CODE_BORDER
    card.line.width = Pt(1)
    strip_shape_styles_and_shadows(card)
    
    # Left Card Title
    tb_lt = slide.shapes.add_textbox(card_x + Pt(14), card_y + Pt(10), card_w - Pt(28), Pt(20))
    tb_lt.name = "TextBox Left Title"
    tf_lt = tb_lt.text_frame
    tf_lt.margin_left = tf_lt.margin_top = tf_lt.margin_right = tf_lt.margin_bottom = 0
    p_lt = tf_lt.paragraphs[0]
    r_lt = p_lt.add_run()
    r_lt.text = left_title
    r_lt.font.name = "Montserrat"
    r_lt.font.size = Pt(11)
    r_lt.font.bold = True
    r_lt.font.color.rgb = WHITE
    
    # Inner Code Box
    box_x = card_x + Pt(14)
    box_y = card_y + Pt(32)
    box_w = card_w - Pt(28)
    box_h = card_h - Pt(42)
    
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, box_y, box_w, box_h)
    code_box.name = "Rounded Rectangle Code Box"
    code_box.adjustments[0] = 0.03
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x2A)
    code_box.line.color.rgb = RGBColor(0x36, 0x3B, 0x44)
    code_box.line.width = Pt(1)
    strip_shape_styles_and_shadows(code_box)
    
    # Code Text
    tb_code = slide.shapes.add_textbox(box_x + Pt(10), box_y + Pt(8), box_w - Pt(20), box_h - Pt(16))
    tb_code.name = "TextBox Code Text"
    tf_c = tb_code.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
    
    for i, line_runs in enumerate(code_lines):
        p = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
        p.space_after = Pt(2)
        for token_text, token_color in line_runs:
            r = p.add_run()
            r.text = token_text
            r.font.name = "Consolas"
            r.font.size = Pt(8.5)
            r.font.bold = False
            r.font.color.rgb = token_color
            
    # Right Column
    rcol_x = Pt(365)
    rcol_y = Pt(56)
    rcol_w = Pt(335)
    
    tb_rt = slide.shapes.add_textbox(rcol_x, rcol_y, rcol_w, Pt(22))
    tb_rt.name = "TextBox Right Title"
    tf_rt = tb_rt.text_frame
    tf_rt.margin_left = tf_rt.margin_top = tf_rt.margin_right = tf_rt.margin_bottom = 0
    p_rt = tf_rt.paragraphs[0]
    r_rt = p_rt.add_run()
    r_rt.text = right_title
    r_rt.font.name = "Montserrat"
    r_rt.font.size = Pt(13)
    r_rt.font.bold = True
    r_rt.font.color.rgb = TEXT_PRIMARY
    
    # Orange divider
    rdiv = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rcol_x, rcol_y + Pt(25), rcol_w, Pt(1))
    rdiv.name = "Rectangle 8"
    rdiv.fill.solid()
    rdiv.fill.fore_color.rgb = ORANGE_LINE
    rdiv.line.fill.background()
    strip_shape_styles_and_shadows(rdiv)
    
    cur_y = rcol_y + Pt(34)
    for idx, item in enumerate(bullet_items):
        item_h = Pt(item.get("height", 42))
        
        dot = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rcol_x, cur_y + Pt(3), Pt(7), Pt(7))
        dot.name = f"Rounded Rectangle Dot {idx}"
        dot.adjustments[0] = 0.2
        dot.fill.solid()
        dot.fill.fore_color.rgb = ORANGE_PILL
        dot.line.fill.background()
        strip_shape_styles_and_shadows(dot)
        
        tb_item = slide.shapes.add_textbox(rcol_x + Pt(14), cur_y, rcol_w - Pt(14), item_h)
        tb_item.name = f"TextBox Bullet {idx}"
        tf_i = tb_item.text_frame
        tf_i.word_wrap = True
        tf_i.margin_left = tf_i.margin_top = tf_i.margin_right = tf_i.margin_bottom = 0
        
        p = tf_i.paragraphs[0]
        p.space_after = Pt(0)
        
        r_title = p.add_run()
        r_title.text = item["title"] + ":\n"
        r_title.font.name = "Montserrat"
        r_title.font.size = Pt(9.5)
        r_title.font.bold = True
        r_title.font.color.rgb = TEXT_PRIMARY
        
        r_desc = p.add_run()
        r_desc.text = item["desc"]
        r_desc.font.name = "Inter"
        r_desc.font.size = Pt(8.5)
        r_desc.font.bold = False
        r_desc.font.color.rgb = TEXT_MUTED
        
        cur_y += item_h + Pt(6)

def create_architecture_diagram_slide(slide, category, question):
    """Custom Diagram 1: Full App Architecture & Inter-Module Connections."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. СЛОЙ ДАННЫХ", "subtitle": "data.js & localStorage", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. ЯДРО ПРИЛОЖЕНИЯ", "subtitle": "main.js & Обработчики", "accent": NUM_COLOR},
        {"x": Pt(476), "title": "3. СТРАНИЦЫ ИНТЕРФЕЙСА", "subtitle": "7 связанных HTML-файлов", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = c["title"]
        r_t1.font.name = "Montserrat"
        r_t1.font.size = Pt(10)
        r_t1.font.bold = True
        r_t1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = c["subtitle"]
        r_t2.font.name = "Inter"
        r_t2.font.size = Pt(8.5)
        r_t2.font.color.rgb = RGBColor(0x9E, 0xA2, 0xAC)
        
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = c["accent"]
        div.line.fill.background()
        strip_shape_styles_and_shadows(div)
        
    # Column 1 Box & Bullets
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(36), top_y + Pt(52), col_w - Pt(24), Pt(56))
    box1.adjustments[0] = 0.05
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x2A)
    box1.line.color.rgb = FUNC_COLOR
    strip_shape_styles_and_shadows(box1)
    
    tf_b1 = box1.text_frame
    p_b1 = tf_b1.paragraphs[0]
    r_b1 = p_b1.add_run()
    r_b1.text = "const OFFICE_ROOMS = [...];\nconst MOCK_BOOKINGS = [...];\nlocalStorage.getItem('user');"
    r_b1.font.name = "Consolas"
    r_b1.font.size = Pt(7.5)
    r_b1.font.color.rgb = CODE_TEXT
    
    tb_c1 = slide.shapes.add_textbox(Pt(36), top_y + Pt(114), col_w - Pt(24), Pt(170))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    c1_items = [
        "Эталонный массив каталога комнат",
        "Хранение активных заявок аренды",
        "Изоляция данных от бизнес-логики",
        "Сохранение сессии в браузере",
        "Единый источник истины (SSOT)"
    ]
    for i, itm in enumerate(c1_items):
        p = tf_c1.paragraphs[0] if i == 0 else tf_c1.add_paragraph()
        p.space_after = Pt(4)
        r_dot = p.add_run()
        r_dot.text = "• "
        r_dot.font.color.rgb = FUNC_COLOR
        r_txt = p.add_run()
        r_txt.text = itm
        r_txt.font.name = "Inter"
        r_txt.font.size = Pt(8.5)
        r_txt.font.color.rgb = WHITE
        
    # Column 2 Box & Bullets
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(262), top_y + Pt(52), col_w - Pt(24), Pt(56))
    box2.adjustments[0] = 0.05
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x2A)
    box2.line.color.rgb = NUM_COLOR
    strip_shape_styles_and_shadows(box2)
    
    tf_b2 = box2.text_frame
    p_b2 = tf_b2.paragraphs[0]
    r_b2 = p_b2.add_run()
    r_b2.text = "initSlider();\ninitCatalogFilters();\ninitBookingCalc();\ninitMyBookings();"
    r_b2.font.name = "Consolas"
    r_b2.font.size = Pt(7.5)
    r_b2.font.color.rgb = CODE_TEXT
    
    tb_c2 = slide.shapes.add_textbox(Pt(262), top_y + Pt(114), col_w - Pt(24), Pt(170))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    c2_items = [
        "Безопасная инициализация модулей",
        "Защита от ошибок null при смене страниц",
        "Динамические фильтры и сортировка",
        "События input и живой расчет цены",
        "Всплывающие нотификации (Toasts)"
    ]
    for i, itm in enumerate(c2_items):
        p = tf_c2.paragraphs[0] if i == 0 else tf_c2.add_paragraph()
        p.space_after = Pt(4)
        r_dot = p.add_run()
        r_dot.text = "• "
        r_dot.font.color.rgb = NUM_COLOR
        r_txt = p.add_run()
        r_txt.text = itm
        r_txt.font.name = "Inter"
        r_txt.font.size = Pt(8.5)
        r_txt.font.color.rgb = WHITE
        
    # Column 3 Box & Bullets
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(488), top_y + Pt(52), col_w - Pt(24), Pt(56))
    box3.adjustments[0] = 0.05
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x2A)
    box3.line.color.rgb = STR_COLOR
    strip_shape_styles_and_shadows(box3)
    
    tf_b3 = box3.text_frame
    p_b3 = tf_b3.paragraphs[0]
    r_b3 = p_b3.add_run()
    r_b3.text = "index.html -> catalog.html\n-> room-details.html?id=1\n-> booking.html?room=1\n-> my-bookings.html"
    r_b3.font.name = "Consolas"
    r_b3.font.size = Pt(7.5)
    r_b3.font.color.rgb = CODE_TEXT
    
    tb_c3 = slide.shapes.add_textbox(Pt(488), top_y + Pt(114), col_w - Pt(24), Pt(170))
    tf_c3 = tb_c3.text_frame
    tf_c3.word_wrap = True
    c3_items = [
        "7 семантических страниц HTML5",
        "Бесшовная передача ID через URL",
        "Единая сквозная шапка и футер",
        "Авторизация и страница кабинета",
        "Полный цикл аренды от клика до чека"
    ]
    for i, itm in enumerate(c3_items):
        p = tf_c3.paragraphs[0] if i == 0 else tf_c3.add_paragraph()
        p.space_after = Pt(4)
        r_dot = p.add_run()
        r_dot.text = "• "
        r_dot.font.color.rgb = STR_COLOR
        r_txt = p.add_run()
        r_txt.text = itm
        r_txt.font.name = "Inter"
        r_txt.font.size = Pt(8.5)
        r_txt.font.color.rgb = WHITE

def create_checklist_diagram_slide(slide, category, question):
    """Custom Diagram 2: Examination Checklist (100 Points)."""
    add_header_and_footer(slide, category, question)
    
    top_y = Pt(56)
    col_w = Pt(210)
    col_h = Pt(296)
    
    cols = [
        {"x": Pt(24), "title": "1. СЕМАНТИКА И ВЕРСТКА", "subtitle": "Критерии 1-3 (30 баллов)", "accent": FUNC_COLOR},
        {"x": Pt(250), "title": "2. ИНТЕРАКТИВНОСТЬ", "subtitle": "Критерии 4-6 (40 баллов)", "accent": NUM_COLOR},
        {"x": Pt(476), "title": "3. БИЗНЕС-ЛОГИКА И UX", "subtitle": "Критерии 7-10 (30 баллов)", "accent": STR_COLOR},
    ]
    
    for c in cols:
        x = c["x"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, col_w, col_h)
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = CODE_BORDER
        card.line.width = Pt(1)
        strip_shape_styles_and_shadows(card)
        
        tb_t = slide.shapes.add_textbox(x + Pt(12), top_y + Pt(10), col_w - Pt(24), Pt(32))
        tf_t = tb_t.text_frame
        p_t1 = tf_t.paragraphs[0]
        r_t1 = p_t1.add_run()
        r_t1.text = c["title"]
        r_t1.font.name = "Montserrat"
        r_t1.font.size = Pt(10)
        r_t1.font.bold = True
        r_t1.font.color.rgb = c["accent"]
        
        p_t2 = tf_t.add_paragraph()
        r_t2 = p_t2.add_run()
        r_t2.text = c["subtitle"]
        r_t2.font.name = "Inter"
        r_t2.font.size = Pt(8.5)
        r_t2.font.color.rgb = RGBColor(0x9E, 0xA2, 0xAC)
        
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Pt(12), top_y + Pt(44), col_w - Pt(24), Pt(1))
        div.fill.solid()
        div.fill.fore_color.rgb = c["accent"]
        div.line.fill.background()
        strip_shape_styles_and_shadows(div)
        
    # Column 1 Box & Bullets
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(36), top_y + Pt(52), col_w - Pt(24), Pt(46))
    box1.adjustments[0] = 0.05
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x2A)
    box1.line.color.rgb = FUNC_COLOR
    strip_shape_styles_and_shadows(box1)
    
    tf_b1 = box1.text_frame
    p_b1 = tf_b1.paragraphs[0]
    r_b1 = p_b1.add_run()
    r_b1.text = "ИТОГО ЗА БЛОК 1: 30 БАЛЛОВ\nКачество верстки и архитектура"
    r_b1.font.name = "Consolas"
    r_b1.font.size = Pt(8)
    r_b1.font.color.rgb = CODE_TEXT
    
    tb_c1 = slide.shapes.add_textbox(Pt(36), top_y + Pt(104), col_w - Pt(24), Pt(180))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    c1_items = [
        "1. Структура проекта (10 б.): 7 страниц, папки css, js, img, относительные пути",
        "2. Семантика и стили (10 б.): теги header, nav, main, footer, шрифт Inter, палитра",
        "3. Главная страница (10 б.): Hero-баннер, 3 карточки популярных комнат, object-fit"
    ]
    for i, itm in enumerate(c1_items):
        p = tf_c1.paragraphs[0] if i == 0 else tf_c1.add_paragraph()
        p.space_after = Pt(6)
        r_dot = p.add_run()
        r_dot.text = "✔ "
        r_dot.font.color.rgb = FUNC_COLOR
        r_txt = p.add_run()
        r_txt.text = itm
        r_txt.font.name = "Inter"
        r_txt.font.size = Pt(8.5)
        r_txt.font.color.rgb = WHITE
        
    # Column 2 Box & Bullets
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(262), top_y + Pt(52), col_w - Pt(24), Pt(46))
    box2.adjustments[0] = 0.05
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x2A)
    box2.line.color.rgb = NUM_COLOR
    strip_shape_styles_and_shadows(box2)
    
    tf_b2 = box2.text_frame
    p_b2 = tf_b2.paragraphs[0]
    r_b2 = p_b2.add_run()
    r_b2.text = "ИТОГО ЗА БЛОК 2: 40 БАЛЛОВ\nИнтерактивные интерфейсы"
    r_b2.font.name = "Consolas"
    r_b2.font.size = Pt(8)
    r_b2.font.color.rgb = CODE_TEXT
    
    tb_c2 = slide.shapes.add_textbox(Pt(262), top_y + Pt(104), col_w - Pt(24), Pt(180))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    c2_items = [
        "4. Слайдер на главной (10 б.): автопрокрутка 3с, стрелки, индикаторы, сброс таймера",
        "5. Динамический каталог (15 б.): рендер из OFFICE_ROOMS через map(), карточки",
        "6. Страница комнаты (15 б.): динамическая загрузка по ?id=..., галерея, характеристики"
    ]
    for i, itm in enumerate(c2_items):
        p = tf_c2.paragraphs[0] if i == 0 else tf_c2.add_paragraph()
        p.space_after = Pt(6)
        r_dot = p.add_run()
        r_dot.text = "✔ "
        r_dot.font.color.rgb = NUM_COLOR
        r_txt = p.add_run()
        r_txt.text = itm
        r_txt.font.name = "Inter"
        r_txt.font.size = Pt(8.5)
        r_txt.font.color.rgb = WHITE
        
    # Column 3 Box & Bullets
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(488), top_y + Pt(52), col_w - Pt(24), Pt(46))
    box3.adjustments[0] = 0.05
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x2A)
    box3.line.color.rgb = STR_COLOR
    strip_shape_styles_and_shadows(box3)
    
    tf_b3 = box3.text_frame
    p_b3 = tf_b3.paragraphs[0]
    r_b3 = p_b3.add_run()
    r_b3.text = "ИТОГО ЗА БЛОК 3: 30 БАЛЛОВ\nБизнес-логика, UX и чистота"
    r_b3.font.name = "Consolas"
    r_b3.font.size = Pt(8)
    r_b3.font.color.rgb = CODE_TEXT
    
    tb_c3 = slide.shapes.add_textbox(Pt(488), top_y + Pt(104), col_w - Pt(24), Pt(180))
    tf_c3 = tb_c3.text_frame
    tf_c3.word_wrap = True
    c3_items = [
        "7. Поиск и сортировка (10 б.): живой фильтр по имени и сортировка цен (asc/desc)",
        "8. Валидация форм (10 б.): индикация .is-invalid, проверка пароля, тосты",
        "9. Калькулятор брони (5 б.): живой пересчет price * hours, автовыбор комнаты",
        "10. Кабинет и чистота (5 б.): список броней, empty state, 0 ошибок в консоли F12"
    ]
    for i, itm in enumerate(c3_items):
        p = tf_c3.paragraphs[0] if i == 0 else tf_c3.add_paragraph()
        p.space_after = Pt(4)
        r_dot = p.add_run()
        r_dot.text = "✔ "
        r_dot.font.color.rgb = STR_COLOR
        r_txt = p.add_run()
        r_txt.text = itm
        r_txt.font.name = "Inter"
        r_txt.font.size = Pt(8.5)
        r_txt.font.color.rgb = WHITE

# -------------------------------------------------------------
# SLIDES DATA (21 CONTENT SLIDES: Slides 3 to 23)
# -------------------------------------------------------------
slides_data = [
    # Slide 3: Code slide
    {
        "type": "code",
        "category": "АНАЛИЗ ОШИБОК",
        "question": "Почему проверка на существование элемента в DOM обязательна при многостраничной структуре?",
        "left_title": "Защита от ошибки TypeError: null",
        "code_lines": [
            [("// Ошибка: скрипт запущен на странице без слайдера", COMMENT_COLOR)],
            [("const", KEYWORD_COLOR), (" slider = document.querySelector(", CODE_TEXT), ("'.slider'", STR_COLOR), (");", CODE_TEXT)],
            [("// ПАДЕНИЕ: Uncaught TypeError: Cannot read properties of null", DANGER_COLOR)],
            [("slider.addEventListener(", CODE_TEXT), ("'click'", STR_COLOR), (", handler);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// НАДЕЖНО: обязательный Guard Clause", COMMENT_COLOR)],
            [("function", KEYWORD_COLOR), (" initSlider() {", FUNC_COLOR)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" slides = document.querySelectorAll(", CODE_TEXT), ("'.slide'", STR_COLOR), (");", CODE_TEXT)],
            [("  ", CODE_TEXT), ("if", KEYWORD_COLOR), (" (!slides.length) ", CODE_TEXT), ("return", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("  ", CODE_TEXT), ("// Логика инициализируется только при наличии элементов", COMMENT_COLOR)],
            [("  startAutoPlay();", FUNC_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Типичная ошибка: null pointer в DOM",
        "bullet_items": [
            {"title": "Общий файл main.js", "desc": "Единый скрипт подключается ко всем 7 страницам, но элементы слайдера есть только на главной.", "height": 42},
            {"title": "Остановка скрипта", "desc": "Необработанная ошибка обращения к null блокирует выполнение всех последующих функций.", "height": 42},
            {"title": "Паттерн Guard Clause", "desc": "Каждая функция модуля обязана начинаться с проверки наличия контейнера в DOM через if (!el) return.", "height": 42},
            {"title": "Чистота консоли", "desc": "В консоли разработчика (F12) на любой странице должно быть ровно ноль ошибок.", "height": 42},
        ]
    },

    # Slide 4: Code slide
    {
        "type": "code",
        "category": "АНАЛИЗ ОШИБОК",
        "question": "К каким последствиям приводит прямая мутация исходных массивов данных в веб-приложениях?",
        "left_title": "Иммутабельность vs Мутация данных",
        "code_lines": [
            [("// ОШИБКА: метод .sort() мутирует оригинальный массив!", DANGER_COLOR)],
            [("function", KEYWORD_COLOR), (" sortRoomsUnsafe() {", FUNC_COLOR)],
            [("  OFFICE_ROOMS.sort((a, b) => a.price - b.price);", CODE_TEXT)],
            [("  // Теперь исходный порядок комнат безвозвратно утерян!", COMMENT_COLOR)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// ПРАВИЛЬНО: поверхностное клонирование через спред", COMMENT_COLOR)],
            [("function", KEYWORD_COLOR), (" sortRoomsSafe() {", FUNC_COLOR)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" copy = [...OFFICE_ROOMS];", CODE_TEXT)],
            [("  copy.sort((a, b) => a.price - b.price);", CODE_TEXT)],
            [("  renderCatalog(copy); ", FUNC_COLOR), ("// data.js защищен!", COMMENT_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Мутация данных и побочные эффекты",
        "bullet_items": [
            {"title": "Потеря эталона", "desc": "Если мутировать массив в data.js, пользователь не сможет сбросить фильтры к исходному состоянию.", "height": 42},
            {"title": "Спред-оператор [...]", "desc": "Создает копию массива верхнего уровня, сохраняя первоначальный список комнат нетронутым.", "height": 42},
            {"title": "Предсказуемость логики", "desc": "Функции сортировки и фильтрации должны быть чистыми: получать данные и возвращать новую выборку.", "height": 42},
            {"title": "Безопасность состояния", "desc": "Глобальные константы приложения остаются стабильными на протяжении всей сессии работы.", "height": 42},
        ]
    },

    # Slide 5: Custom Diagram 1
    {
        "type": "diagram_arch",
        "category": "СХЕМА: АРХИТЕКТУРА И СВЯЗИ ПРОЕКТА",
        "question": "Как модульное разделение JavaScript и данных обеспечивает масштабируемость проекта?",
    },

    # Slide 6: Code slide
    {
        "type": "code",
        "category": "АНАЛИЗ ОШИБОК",
        "question": "В чем главное преимущество делегирования событий при частой перерисовке списков в DOM?",
        "left_title": "Делегирование событий в каталоге",
        "code_lines": [
            [("// Неэффективно: вешать слушатель на каждую новую кнопку", COMMENT_COLOR)],
            [("document.querySelectorAll(", CODE_TEXT), ("'.card-btn'", STR_COLOR), (").forEach(...);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// ЭФФЕКТИВНО: Делегирование на родительском контейнере", COMMENT_COLOR)],
            [("catalogGrid.addEventListener(", CODE_TEXT), ("'click'", STR_COLOR), (", (e) => {", FUNC_COLOR)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" btn = e.target.closest(", CODE_TEXT), ("'.btn-book'", STR_COLOR), (");", CODE_TEXT)],
            [("  ", CODE_TEXT), ("if", KEYWORD_COLOR), (" (!btn) ", CODE_TEXT), ("return", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" roomId = btn.dataset.id;", CODE_TEXT)],
            [("  window.location.href = `booking.html?room=${roomId}`;", CODE_TEXT)],
            [("});", CODE_TEXT)],
        ],
        "right_title": "Управление обработчиками событий",
        "bullet_items": [
            {"title": "Проблема перерисовки", "desc": "При каждом поиске или сортировке innerHTML перезаписывается, и прямые слушатели уничтожаются.", "height": 42},
            {"title": "Всплытие (Event Bubbling)", "desc": "Событие click поднимается от кнопки к контейнеру, где перехватывается единственным обработчиком.", "height": 42},
            {"title": "Метод e.target.closest()", "desc": "Точно определяет целевую кнопку внутри карточки, даже если клик пришелся по вложенной иконке.", "height": 42},
            {"title": "Экономия ресурсов", "desc": "Один слушатель на сетку вместо десятков обработчиков предотвращает утечки оперативной памяти.", "height": 42},
        ]
    },

    # Slide 7: Code slide
    {
        "type": "code",
        "category": "ЗАКРЕПЛЕНИЕ МАТЕРИАЛА",
        "question": "Какую роль семантическая разметка играет в реальных проектах и техническом аудите?",
        "left_title": "Семантическое ядро интерфейса",
        "code_lines": [
            [("<!-- Стандарт семантики портала СмартОфис -->", COMMENT_COLOR)],
            [("<header", TAG_COLOR), (" class=\"header\">", CODE_TEXT), ("<!-- Шапка и навигация -->", COMMENT_COLOR), ("</header>", TAG_COLOR)],
            [("<main", TAG_COLOR), (" class=\"main\">", CODE_TEXT)],
            [("  <div", TAG_COLOR), (" class=\"container\">", CODE_TEXT)],
            [("    <h1", TAG_COLOR), (" class=\"page-title\">Мои бронирования</h1>", CODE_TEXT)],
            [("    <section", TAG_COLOR), (" class=\"bookings-section\">", CODE_TEXT)],
            [("      <div", TAG_COLOR), (" class=\"bookings-list\" id=\"myBookingsList\">", CODE_TEXT)],
            [("        <!-- Динамические карточки -->", COMMENT_COLOR)],
            [("      </div>", TAG_COLOR)],
            [("    </section>", TAG_COLOR)],
            [("  </div>", TAG_COLOR)],
            [("</main>", TAG_COLOR)],
            [("<footer", TAG_COLOR), (" class=\"footer\">", CODE_TEXT), ("<!-- Контакты и копирайт -->", COMMENT_COLOR), ("</footer>", TAG_COLOR)],
        ],
        "right_title": "HTML5 разметка коммерческого уровня",
        "bullet_items": [
            {"title": "Семантические ориентиры", "desc": "Теги header, nav, main, section и footer однозначно структурируют документ для браузера.", "height": 42},
            {"title": "Иерархия заголовков", "desc": "Строго один тег <h1> на страницу, четкая вложенность <h2> и <h3> внутри карточек и блоков.", "height": 42},
            {"title": "Доступность (a11y)", "desc": "Скринридеры безошибочно озвучивают интерфейс незрячим пользователям без лишних костылей.", "height": 42},
            {"title": "SEO-оптимизация", "desc": "Поисковые роботы корректно индексируют контент страниц, повышая позицию сайта в поисковой выдаче.", "height": 42},
        ]
    },

    # Slide 8: Code slide
    {
        "type": "code",
        "category": "ЗАКРЕПЛЕНИЕ МАТЕРИАЛА",
        "question": "Когда в верстке интерфейса эффективнее использовать CSS Grid, а когда Flexbox?",
        "left_title": "Комбинация Flexbox и Grid сеток",
        "code_lines": [
            [("/* 1. CSS Grid для двумерной витрины каталога */", COMMENT_COLOR)],
            [(".catalog-grid {", CODE_TEXT)],
            [("  display: ", KEYWORD_COLOR), ("grid;", CODE_TEXT)],
            [("  grid-template-columns: ", KEYWORD_COLOR), ("repeat(auto-fit, minmax(280px, 1fr));", STR_COLOR)],
            [("  gap: ", KEYWORD_COLOR), ("24px;", NUM_COLOR)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("/* 2. Flexbox для одномерного тулбара и карточки */", COMMENT_COLOR)],
            [(".booking-item {", CODE_TEXT)],
            [("  display: ", KEYWORD_COLOR), ("flex;", CODE_TEXT)],
            [("  justify-content: ", KEYWORD_COLOR), ("space-between;", STR_COLOR)],
            [("  align-items: ", KEYWORD_COLOR), ("center;", STR_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Современная верстка без костылей",
        "bullet_items": [
            {"title": "CSS Grid для каталога", "desc": "Автоматическая раскладка repeat(auto-fit, minmax) адаптирует колонки под любую ширину экрана.", "height": 42},
            {"title": "Flexbox для компонентов", "desc": "Идеально распределяет пространство в шапке, тулбаре поиска, карточках и строках бронирования.", "height": 42},
            {"title": "Отказ от флоатов", "desc": "Полное отсутствие устаревших float и clearfix: верстка остается гибкой и предсказуемой.", "height": 42},
            {"title": "CSS-переменные :root", "desc": "Централизованное управление цветовой палитрой (#007bff, #28a745) и радиусами скругления.", "height": 42},
        ]
    },

    # Slide 9: Code slide
    {
        "type": "code",
        "category": "ЗАКРЕПЛЕНИЕ МАТЕРИАЛА",
        "question": "Как синтаксические конструкции ES6+ повышают читаемость и поддерживаемость кодовой базы?",
        "left_title": "Современный синтаксис ES6+",
        "code_lines": [
            [("// Деструктуризация параметров объекта", COMMENT_COLOR)],
            [("function", KEYWORD_COLOR), (" renderBookingCard({ id, roomTitle, date, hours, totalPrice }) {", FUNC_COLOR)],
            [("  ", CODE_TEXT), ("return", KEYWORD_COLOR), (" `", STR_COLOR)],
            [("    <div class=\"booking-item\">", STR_COLOR)],
            [("      <div>", STR_COLOR)],
            [("        <h3>${roomTitle}</h3>", STR_COLOR)],
            [("        <p>Дата: ${date} | Время: ${hours} ч. | Заказ №${id}</p>", STR_COLOR)],
            [("      </div>", STR_COLOR)],
            [("      <div class=\"price-block\">", STR_COLOR)],
            [("        <span class=\"price\">${totalPrice} ₽</span>", STR_COLOR)],
            [("        <span class=\"badge-success\">Подтверждено</span>", STR_COLOR)],
            [("      </div>", STR_COLOR)],
            [("    </div>", STR_COLOR)],
            [("  `;", STR_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Чистый JavaScript без библиотек",
        "bullet_items": [
            {"title": "Шаблонные строки", "desc": "Интерполяция ${...} позволяет создавать читаемую HTML-разметку прямо в коде функций рендеринга.", "height": 42},
            {"title": "Деструктуризация", "desc": "Извлечение нужных свойств объекта прямо в аргументах функции делает сигнатуру прозрачной.", "height": 42},
            {"title": "Стрелочные функции", "desc": "Компактный синтаксис для методов массивов map(), filter(), sort() и обработчиков событий.", "height": 42},
            {"title": "Стандарты ES6+", "desc": "Полный отказ от устаревшего var в пользу блочных и неизменяемых let / const.", "height": 42},
        ]
    },

    # Slide 10: Code slide
    {
        "type": "code",
        "category": "ФАЙЛОВАЯ АРХИТЕКТУРА",
        "question": "Почему изоляция статических страниц в папку pages/ упрощает масштабирование портала?",
        "left_title": "Эталонная структура проекта",
        "code_lines": [
            [("smartoffice-portal/", COMMENT_COLOR)],
            [("├── index.html              ", STR_COLOR), ("# Главная (Hero, слайдер, топ)", COMMENT_COLOR)],
            [("├── pages/                  ", FUNC_COLOR), ("# Дополнительные страницы", COMMENT_COLOR)],
            [("│   ├── catalog.html        ", STR_COLOR), ("# Каталог комнат и фильтры", COMMENT_COLOR)],
            [("│   ├── room-details.html   ", STR_COLOR), ("# Страница конкретного офиса", COMMENT_COLOR)],
            [("│   ├── booking.html        ", STR_COLOR), ("# Калькулятор бронирования", COMMENT_COLOR)],
            [("│   ├── my-bookings.html    ", STR_COLOR), ("# Личный кабинет (заказы)", COMMENT_COLOR)],
            [("│   ├── login.html          ", STR_COLOR), ("# Форма входа в систему", COMMENT_COLOR)],
            [("│   └── register.html       ", STR_COLOR), ("# Форма регистрации", COMMENT_COLOR)],
            [("├── css/style.css           ", NUM_COLOR), ("# Единый файл стилей проекта", COMMENT_COLOR)],
            [("├── js/                     ", KEYWORD_COLOR), ("# Скрипты и структуры данных", COMMENT_COLOR)],
            [("│   ├── data.js             ", KEYWORD_COLOR), ("# Мок-данные и база комнат", COMMENT_COLOR)],
            [("│   └── main.js             ", KEYWORD_COLOR), ("# Бизнес-логика приложения", COMMENT_COLOR)],
            [("└── img/                    ", TAG_COLOR), ("# Иконки SVG и фото комнат", COMMENT_COLOR)],
        ],
        "right_title": "Организация файлов веб-приложения",
        "bullet_items": [
            {"title": "Единая точка входа", "desc": "Файл index.html располагается в корневом каталоге и открывается веб-сервером по умолчанию.", "height": 42},
            {"title": "Папка pages/", "desc": "Все вложенные страницы изолированы в отдельной директории, устраняя хаос в корне проекта.", "height": 42},
            {"title": "Разделение логики", "desc": "Файл data.js отвечает исключительно за структуры данных, а main.js — за интерактивность и DOM.", "height": 42},
            {"title": "Медиа-ресурсы img/", "desc": "Все векторные и растровые изображения собраны в одном месте с унифицированными именами.", "height": 42},
        ]
    },

    # Slide 11: Code slide
    {
        "type": "code",
        "category": "ФАЙЛОВАЯ АРХИТЕКТУРА",
        "question": "Как корректно организовать относительные пути в стилях и скриптах для вложенных страниц?",
        "left_title": "Правила относительной адресации",
        "code_lines": [
            [("<!-- 1. Из корневого index.html: -->", COMMENT_COLOR)],
            [("<link rel=\"stylesheet\" href=\"css/style.css\">", CODE_TEXT)],
            [("<script src=\"js/data.js\" defer></script>", CODE_TEXT)],
            [("<a href=\"pages/catalog.html\">Каталог</a>", STR_COLOR)],
            [("", CODE_TEXT)],
            [("<!-- 2. Из вложенного файла pages/my-bookings.html: -->", COMMENT_COLOR)],
            [("<link rel=\"stylesheet\" href=\"../css/style.css\">", CODE_TEXT)],
            [("<script src=\"../js/data.js\" defer></script>", CODE_TEXT)],
            [("<img src=\"../img/logo.svg\" alt=\"Логотип\">", NUM_COLOR)],
            [("<a href=\"../index.html\">Главная</a>", STR_COLOR)],
            [("<a href=\"catalog.html\">Каталог (соседний файл)</a>", FUNC_COLOR)],
        ],
        "right_title": "Навигация и относительные пути",
        "bullet_items": [
            {"title": "Символ выхода ../", "desc": "Две точки со слешем поднимают браузер на один уровень вверх — из папки pages/ в корень проекта.", "height": 42},
            {"title": "Соседние файлы", "desc": "Ссылки между страницами внутри pages/ оформляются без префикса (например: href=\"catalog.html\").", "height": 42},
            {"title": "Опасность абсолютных путей", "desc": "Пути от корня (/css/style.css) ломаются при открытии проекта локально или на поддоменах.", "height": 42},
            {"title": "Атрибут defer", "desc": "Гарантирует загрузку скриптов после полного построения DOM-дерева страницы в браузере.", "height": 42},
        ]
    },

    # Slide 12: Custom Diagram 2
    {
        "type": "diagram_checklist",
        "category": "СХЕМА: ЭКЗАМЕНАЦИОННЫЙ ЧЕК-ЛИСТ",
        "question": "Какие критерии являются ключевыми при оценке качества и стабильности веб-приложения?",
    },

    # Slide 13: Code slide
    {
        "type": "code",
        "category": "ПОДГОТОВКА К ЗАЩИТЕ",
        "question": "Как структурировать выступление на защите, чтобы показать все сильные стороны проекта?",
        "left_title": "Сквозной пользовательский сценарий",
        "code_lines": [
            [("// Сценарий демонстрации на защите проекта:", COMMENT_COLOR)],
            [("1. Главная страница:", FUNC_COLOR)],
            [("   - Демонстрация Hero-баннера и автослайдера (3 сек)", CODE_TEXT)],
            [("   - Ручное переключение слайдов стрелками и точками", CODE_TEXT)],
            [("2. Каталог комнат:", NUM_COLOR)],
            [("   - Живой поиск: ввод 'Focus' -> мгновенная фильтрация", CODE_TEXT)],
            [("   - Сортировка по цене: возрастание / убывание", CODE_TEXT)],
            [("3. Страница комнаты (room-details.html?id=...):", STR_COLOR)],
            [("   - Переход по ссылке, чтение параметров URL", CODE_TEXT)],
            [("4. Калькулятор бронирования:", KEYWORD_COLOR)],
            [("   - Автовыбор комнаты, ввод часов -> пересчет цены", CODE_TEXT)],
            [("5. Кабинет 'Мои бронирования':", TAG_COLOR)],
            [("   - Проверка созданной заявки в истории заказов", CODE_TEXT)],
        ],
        "right_title": "План успешной демонстрации проекта",
        "bullet_items": [
            {"title": "Пользовательская история", "desc": "Показывайте проект не как разрозненный код, а как законченный коммерческий сценарий клиента.", "height": 42},
            {"title": "Интерактивный отклик", "desc": "Обратите внимание комиссии на скорость работы интерфейса: фильтрация и расчеты без перезагрузок.", "height": 42},
            {"title": "Обработка ошибок", "desc": "Продемонстрируйте валидацию пустых полей формы и отображение блока Empty State в поиске.", "height": 42},
            {"title": "Уверенность в коде", "desc": "Будьте готовы открыть любой скрипт и объяснить, какая функция отвечает за конкретное действие.", "height": 42},
        ]
    },

    # Slide 14: Code slide
    {
        "type": "code",
        "category": "ПОДГОТОВКА К ЗАЩИТЕ",
        "question": "К каким техническим вопросам преподавателя нужно быть готовым в процессе защиты?",
        "left_title": "Типовые вопросы на защите",
        "code_lines": [
            [("// Вопрос 1: Как работает поиск комнат?", COMMENT_COLOR)],
            [("const", KEYWORD_COLOR), (" q = input.value.toLowerCase().trim();", CODE_TEXT)],
            [("const", KEYWORD_COLOR), (" res = data.filter(r => r.title.toLowerCase().includes(q));", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Вопрос 2: Как сохраняется сессия пользователя?", COMMENT_COLOR)],
            [("localStorage.setItem(", CODE_TEXT), ("'currentUser'", STR_COLOR), (", email);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Вопрос 3: Зачем нужен URLSearchParams?", COMMENT_COLOR)],
            [("const", KEYWORD_COLOR), (" id = new URLSearchParams(location.search).get('id');", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("// Вопрос 4: Как очистить интервал слайдера?", COMMENT_COLOR)],
            [("clearInterval(sliderTimer);", FUNC_COLOR)],
        ],
        "right_title": "Технические аргументы для комиссии",
        "bullet_items": [
            {"title": "Чистый JavaScript", "desc": "Обоснуйте, почему проект собран на Vanilla JS: максимальная производительность и понимание базы.", "height": 42},
            {"title": "Событийная модель", "desc": "Поясните разницу между событиями change (выбор из списка) и input (живой ввод с клавиатуры).", "height": 42},
            {"title": "Хранилище localStorage", "desc": "Объясните, почему данные сессии не сбрасываются при закрытии вкладки или перезагрузке страницы.", "height": 42},
            {"title": "Защита маршрутов", "desc": "Расскажите, как скрипт блокирует доступ к оформлению брони для неавторизованных гостей.", "height": 42},
        ]
    },

    # Slide 15: Code slide
    {
        "type": "code",
        "category": "ЛИЧНЫЙ КАБИНЕТ",
        "question": "Зачем для истории заказов выделять отдельный семантический контейнер с уникальным id?",
        "left_title": "Разметка pages/my-bookings.html",
        "code_lines": [
            [("<!-- Контейнер личного кабинета пользователя -->", COMMENT_COLOR)],
            [("<main", TAG_COLOR), (" class=\"main\">", CODE_TEXT)],
            [("  <div", TAG_COLOR), (" class=\"container\">", CODE_TEXT)],
            [("    <h1", TAG_COLOR), (" class=\"page-title\">Мои бронирования</h1>", CODE_TEXT)],
            [("    <p", TAG_COLOR), (" class=\"page-subtitle\">", CODE_TEXT)],
            [("      История ваших заявок на аренду рабочих пространств", STR_COLOR)],
            [("    </p>", TAG_COLOR)],
            [("", CODE_TEXT)],
            [("    <!-- Точка монтирования списка броней -->", COMMENT_COLOR)],
            [("    <div", TAG_COLOR), (" class=\"bookings-list\" id=\"myBookingsList\"></div>", NUM_COLOR)],
            [("  </div>", TAG_COLOR)],
            [("</main>", TAG_COLOR)],
        ],
        "right_title": "Страница истории бронирований",
        "bullet_items": [
            {"title": "Целевой контейнер", "desc": "Элемент #myBookingsList служит точкой динамического рендеринга карточек из массива MOCK_BOOKINGS.", "height": 42},
            {"title": "Заголовок и подзаголовок", "desc": "Четкая смысловая иерархия страницы информирует пользователя о назначении раздела кабинета.", "height": 42},
            {"title": "Стилистическое единство", "desc": "Классы container, main, page-title используют общие правила оформления из style.css.", "height": 42},
            {"title": "Изолированный компонент", "desc": "Верстка подготовлена как для отображения заполненного списка, так и для Empty State.", "height": 42},
        ]
    },

    # Slide 16: Code slide
    {
        "type": "code",
        "category": "ЛИЧНЫЙ КАБИНЕТ",
        "question": "Как реализовать базовую защиту приватных страниц от неавторизованного доступа на клиенте?",
        "left_title": "Проверка авторизации в кабинете",
        "code_lines": [
            [("function", KEYWORD_COLOR), (" initMyBookings() {", FUNC_COLOR)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" container = document.getElementById(", CODE_TEXT), ("'myBookingsList'", STR_COLOR), (");", CODE_TEXT)],
            [("  ", CODE_TEXT), ("if", KEYWORD_COLOR), (" (!container || ", CODE_TEXT), ("typeof", KEYWORD_COLOR), (" MOCK_BOOKINGS === ", CODE_TEXT), ("'undefined'", STR_COLOR), (") ", CODE_TEXT), ("return", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  ", CODE_TEXT), ("// Защита приватного раздела: проверка сессии", COMMENT_COLOR)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" currentUser = localStorage.getItem(", CODE_TEXT), ("'currentUser'", STR_COLOR), (");", CODE_TEXT)],
            [("  ", CODE_TEXT), ("if", KEYWORD_COLOR), (" (!currentUser) {", CODE_TEXT)],
            [("    ", CODE_TEXT), ("// Редирект неавторизованного пользователя на вход", COMMENT_COLOR)],
            [("    window.location.href = ", CODE_TEXT), ("'login.html'", STR_COLOR), (";", CODE_TEXT)],
            [("    ", CODE_TEXT), ("return", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("  // Пользователь авторизован — запускаем рендеринг", COMMENT_COLOR)],
            [("  renderBookings(container);", FUNC_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Клиентский контроль доступа",
        "bullet_items": [
            {"title": "Чтение localStorage", "desc": "Ключ currentUser хранит признак активной сессии (email или статус авторизации пользователя).", "height": 42},
            {"title": "Автоматический редирект", "desc": "Если гость пытается зайти по прямой ссылке в кабинет, браузер мгновенно направляет его на login.html.", "height": 42},
            {"title": "Защита от сбоя MOCK_BOOKINGS", "desc": "Проверка typeof гарантирует отсутствие фатальных ошибок, если файл data.js не загрузился.", "height": 42},
            {"title": "Приватный пользовательский опыт", "desc": "Личная информация и история бронирований доступны только после успешной авторизации.", "height": 42},
        ]
    },

    # Slide 17: Code slide
    {
        "type": "code",
        "category": "ПРАКТИКА: СБОРКА",
        "question": "Как метод map() позволяет лаконично преобразовать массив объектов броней в готовый DOM?",
        "left_title": "Рендеринг истории бронирований",
        "code_lines": [
            [("function", KEYWORD_COLOR), (" renderBookings(container) {", FUNC_COLOR)],
            [("  ", CODE_TEXT), ("if", KEYWORD_COLOR), (" (!MOCK_BOOKINGS || MOCK_BOOKINGS.length === 0) {", CODE_TEXT)],
            [("    container.innerHTML = ", CODE_TEXT), ("'<div class=\"empty-message\">У вас пока нет бронирований</div>'", STR_COLOR), (";", CODE_TEXT)],
            [("    ", CODE_TEXT), ("return", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  container.innerHTML = MOCK_BOOKINGS.map(item => `", CODE_TEXT)],
            [("    <div class=\"booking-item\">", STR_COLOR)],
            [("      <div>", STR_COLOR)],
            [("        <h3>${item.roomTitle}</h3>", STR_COLOR)],
            [("        <div class=\"meta\">Дата: <strong>${item.date}</strong> | ${item.hours} ч. | №${item.id}</div>", STR_COLOR)],
            [("      </div>", STR_COLOR)],
            [("      <div class=\"summary\">", STR_COLOR)],
            [("        <div class=\"price\">${item.totalPrice} ₽</div>", STR_COLOR)],
            [("        <span class=\"badge-confirmed\">Подтверждено</span>", STR_COLOR)],
            [("      </div>", STR_COLOR)],
            [("    </div>", STR_COLOR)],
            [("  `).join('');", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Генерация карточек в кабинете",
        "bullet_items": [
            {"title": "Декларативный подход", "desc": "Метод map() трансформирует каждый объект заявки в готовую HTML-строку карточки брони.", "height": 42},
            {"title": "Склейка методом join('')", "desc": "Объединяет элементы массива в единый HTML-фрагмент без разделительных запятых.", "height": 42},
            {"title": "Быстродействие", "desc": "Единоразовое присвоение innerHTML обновляет DOM за один цикл перерисовки браузера.", "height": 42},
            {"title": "Связка с датой и временем", "desc": "Пользователь видит точную дату, длительность и присвоенный 5-значный номер заказа.", "height": 42},
        ]
    },

    # Slide 18: Code slide
    {
        "type": "code",
        "category": "ПРАКТИКА: СБОРКА",
        "question": "Почему оформление статуса заявки цветовым бейджем повышает информативность интерфейса?",
        "left_title": "Стилизация карточки бронирования",
        "code_lines": [
            [("/* Стили элемента бронирования в style.css */", COMMENT_COLOR)],
            [(".booking-item {", CODE_TEXT)],
            [("  background: ", KEYWORD_COLOR), ("#ffffff;", CODE_TEXT)],
            [("  border: ", KEYWORD_COLOR), ("1px solid #e2e8f0;", STR_COLOR)],
            [("  border-radius: ", KEYWORD_COLOR), ("8px;", NUM_COLOR)],
            [("  padding: ", KEYWORD_COLOR), ("18px 24px;", NUM_COLOR)],
            [("  margin-bottom: ", KEYWORD_COLOR), ("16px;", NUM_COLOR)],
            [("  display: ", KEYWORD_COLOR), ("flex;", CODE_TEXT)],
            [("  justify-content: ", KEYWORD_COLOR), ("space-between;", STR_COLOR)],
            [("  align-items: ", KEYWORD_COLOR), ("center;", STR_COLOR)],
            [("}", CODE_TEXT)],
            [("", CODE_TEXT)],
            [(".badge-confirmed {", CODE_TEXT)],
            [("  color: ", KEYWORD_COLOR), ("#28a745; ", CODE_TEXT), ("/* Зеленый акцент успеха */", COMMENT_COLOR)],
            [("  font-size: ", KEYWORD_COLOR), ("12px;", NUM_COLOR)],
            [("  font-weight: ", KEYWORD_COLOR), ("600;", NUM_COLOR)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Визуальное оформление кабинета",
        "bullet_items": [
            {"title": "Карточная компоновка", "desc": "Белый фон, тонкая граница #e2e8f0 и скругление углов создают легкий современный стиль.", "height": 42},
            {"title": "Двустороннее выравнивание", "desc": "Слева располагаются детали офиса и дата, справа — крупная цена и статус подтверждения.", "height": 42},
            {"title": "Цветовая дифференциация", "desc": "Статус «Подтверждено» выделен зеленым цветом (#28a745), давая мгновенный позитивный отклик.", "height": 42},
            {"title": "Типографика Inter", "desc": "Четкий контраст между полужирными заголовками комнат и аккуратными мета-данными серым цветом.", "height": 42},
        ]
    },

    # Slide 19: Code slide
    {
        "type": "code",
        "category": "ПРАКТИКА: СБОРКА",
        "question": "Почему отсутствие Empty State считается критической ошибкой в пользовательских интерфейсах?",
        "left_title": "Обработка пустого состояния (Empty State)",
        "code_lines": [
            [("// Защита от пустого белого экрана", COMMENT_COLOR)],
            [("if", KEYWORD_COLOR), (" (!MOCK_BOOKINGS || MOCK_BOOKINGS.length === 0) {", CODE_TEXT)],
            [("  container.innerHTML = `", CODE_TEXT)],
            [("    <div class=\"empty-message\">", STR_COLOR)],
            [("      <div class=\"empty-icon\">📅</div>", STR_COLOR)],
            [("      <h2>У вас пока нет активных бронирований</h2>", STR_COLOR)],
            [("      <p>Выберите подходящее рабочее пространство в каталоге</p>", STR_COLOR)],
            [("      <a href=\"catalog.html\" class=\"btn btn-primary\">", STR_COLOR)],
            [("        Перейти в каталог комнат", STR_COLOR)],
            [("      </a>", STR_COLOR)],
            [("    </div>", STR_COLOR)],
            [("  `;", STR_COLOR)],
            [("  ", CODE_TEXT), ("return", KEYWORD_COLOR), (";", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "UX пустого состояния кабинета",
        "bullet_items": [
            {"title": "Понятная обратная связь", "desc": "Вместо пустого белого пространства пользователь видит дружелюбное сообщение с иконкой.", "height": 42},
            {"title": "Призыв к действию (CTA)", "desc": "Кнопка «Перейти в каталог» помогает посетителю продолжить сценарий, не теряясь на странице.", "height": 42},
            {"title": "Критерий чек-листа", "desc": "Экзаменационная комиссия специально проверяет реакцию интерфейса на пустой массив данных.", "height": 42},
            {"title": "Профессиональный UX", "desc": "Интерфейс направляет пользователя к следующему шагу при любых граничных состояниях.", "height": 42},
        ]
    },

    # Slide 20: Code slide
    {
        "type": "code",
        "category": "ПРАКТИКА: СБОРКА",
        "question": "Как динамическое управление классами и стилями навигации отражает статус пользователя?",
        "left_title": "Сквозная навигация портала",
        "code_lines": [
            [("function", KEYWORD_COLOR), (" updateAuthNav() {", FUNC_COLOR)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" user = localStorage.getItem(", CODE_TEXT), ("'currentUser'", STR_COLOR), (");", CODE_TEXT)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" navItem = document.getElementById(", CODE_TEXT), ("'myBookingsNavItem'", STR_COLOR), (");", CODE_TEXT)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" btn = document.getElementById(", CODE_TEXT), ("'authNavBtn'", STR_COLOR), (");", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  ", CODE_TEXT), ("if", KEYWORD_COLOR), (" (user) {", CODE_TEXT)],
            [("    navItem.style.display = ", CODE_TEXT), ("'block'", STR_COLOR), ("; ", COMMENT_COLOR), ("// Показываем кабинет", COMMENT_COLOR)],
            [("    btn.textContent = ", CODE_TEXT), ("'Выйти'", STR_COLOR), (";", CODE_TEXT)],
            [("    btn.onclick = (e) => {", FUNC_COLOR)],
            [("      e.preventDefault();", CODE_TEXT)],
            [("      localStorage.removeItem(", CODE_TEXT), ("'currentUser'", STR_COLOR), (");", CODE_TEXT)],
            [("      window.location.reload();", FUNC_COLOR)],
            [("    };", CODE_TEXT)],
            [("  } ", CODE_TEXT), ("else", KEYWORD_COLOR), (" {", CODE_TEXT)],
            [("    navItem.style.display = ", CODE_TEXT), ("'none'", STR_COLOR), (";", CODE_TEXT)],
            [("    btn.textContent = ", CODE_TEXT), ("'Войти'", STR_COLOR), (";", CODE_TEXT)],
            [("  }", CODE_TEXT)],
            [("}", CODE_TEXT)],
        ],
        "right_title": "Синхронизация шапки сайта",
        "bullet_items": [
            {"title": "Управление состоянием", "desc": "Шапка сайта на всех 7 страницах автоматически переключается между гостевым и авторским режимами.", "height": 42},
            {"title": "Пункт «Мои бронирования»", "desc": "Отображается только в том случае, если пользователь вошел в систему под своей учетной записью.", "height": 42},
            {"title": "Кнопка «Выйти»", "desc": "Стирает сессию из localStorage, показывает уведомление и перезагружает интерфейс в гостевой режим.", "height": 42},
            {"title": "Подсветка активной ссылки", "desc": "Функция initNavigation() сопоставляет location.pathname и добавляет класс active текущей вкладке.", "height": 42},
        ]
    },

    # Slide 21: Code slide
    {
        "type": "code",
        "category": "ПРАКТИКА: СБОРКА",
        "question": "Почему метод unshift() предпочтительнее push() при добавлении свежих пользовательских заказов?",
        "left_title": "Бесшовное добавление заказа в историю",
        "code_lines": [
            [("// Обработчик отправки формы booking.html", COMMENT_COLOR)],
            [("form.addEventListener(", CODE_TEXT), ("'submit'", STR_COLOR), (", (e) => {", FUNC_COLOR)],
            [("  e.preventDefault();", CODE_TEXT)],
            [("  ", CODE_TEXT), ("const", KEYWORD_COLOR), (" orderNumber = Math.floor(10000 + Math.random() * 90000);", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  ", CODE_TEXT), ("// unshift() добавляет новую бронь в начало списка", COMMENT_COLOR)],
            [("  MOCK_BOOKINGS.unshift({", FUNC_COLOR)],
            [("    id: String(orderNumber),", STR_COLOR)],
            [("    roomTitle: selectedRoom.title,", STR_COLOR)],
            [("    date: bookingDateInput.value,", STR_COLOR)],
            [("    hours: bookingHours,", NUM_COLOR)],
            [("    totalPrice: finalCalculatedPrice", NUM_COLOR)],
            [("  });", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("  showNotification(`Заказ №${orderNumber} успешно оформлен!`, 'success');", FUNC_COLOR)],
            [("  setTimeout(() => window.location.href = 'my-bookings.html', 1000);", CODE_TEXT)],
            [("});", CODE_TEXT)],
        ],
        "right_title": "Полный сквозной цикл бронирования",
        "bullet_items": [
            {"title": "Хронологический порядок", "desc": "Метод unshift() ставит свежую бронь на первое место, чтобы клиент сразу увидел ее в кабинете.", "height": 42},
            {"title": "Генерация ID заказа", "desc": "Случайный 5-значный номер заказа симулирует реальную серверную систему управления заявками.", "height": 42},
            {"title": "Компонент всплывающих тостов", "desc": "Всплывающее зеленое уведомление подтверждает успех бронирования перед автоматическим переходом.", "height": 42},
            {"title": "Автоматический переход", "desc": "Задержка setTimeout плавно переводит пользователя на страницу истории его бронирований.", "height": 42},
        ]
    },

    # Slide 22: Code slide
    {
        "type": "code",
        "category": "ПРАКТИКА: СБОРКА",
        "question": "Как методическая самопроверка по чек-листу исключает потерю баллов на итоговой защите?",
        "left_title": "Предзащитный аудит проекта",
        "code_lines": [
            [("// Контрольный чек-лист самопроверки (100 баллов):", COMMENT_COLOR)],
            [("[OK] 1. Файловая структура: 7 страниц, папки css/js/img  (10 б.)", STR_COLOR)],
            [("[OK] 2. Семантика и стили: теги HTML5, шрифт Inter, палитра (10 б.)", STR_COLOR)],
            [("[OK] 3. Главная страница: Hero + популярные комнаты        (10 б.)", STR_COLOR)],
            [("[OK] 4. Слайдер на главной: автоскролл 3с + кнопки        (10 б.)", STR_COLOR)],
            [("[OK] 5. Каталог комнат: динамический рендер из данных    (15 б.)", STR_COLOR)],
            [("[OK] 6. Страница комнаты: загрузка по query-параметру ?id (15 б.)", STR_COLOR)],
            [("[OK] 7. Поиск и сортировка: живой фильтр + кнопки цен    (10 б.)", STR_COLOR)],
            [("[OK] 8. Валидация форм: регулярки, .is-invalid, тосты     (10 б.)", STR_COLOR)],
            [("[OK] 9. Калькулятор: мгновенный расчет часов и тарифа      (5 б.)", STR_COLOR)],
            [("[OK] 10. Мои бронирования: кабинет + 0 ошибок в F12        (5 б.)", STR_COLOR)],
        ],
        "right_title": "Финальная приемка и контроль качества",
        "bullet_items": [
            {"title": "100 из 100 баллов", "desc": "Все 10 экзаменационных требований реализованы и протестированы во всех популярных браузерах.", "height": 42},
            {"title": "Тестирование граничных условий", "desc": "Проверены пустой ввод, несуществующий ID комнаты в URL (?id=999) и неверный пароль.", "height": 42},
            {"title": "Адаптивная верстка", "desc": "Интерфейс корректно отображается как на экранах ноутбуков (1440px), так и на смартфонах.", "height": 42},
            {"title": "Готовность к защите", "desc": "Проект полностью упакован, содержит чистый код и готов к публичной защите перед экспертами.", "height": 42},
        ]
    },

    # Slide 23: Code slide
    {
        "type": "code",
        "category": "ИТОГИ КУРСА",
        "question": "Какие навыки разработки интерфейсов станут фундаментом для изучения фреймворков (React, Vue)?",
        "left_title": "Стек компетенций фронтенд-разработчика",
        "code_lines": [
            [("// Портфолио студента: стек технологий фронтенда", COMMENT_COLOR)],
            [("const", KEYWORD_COLOR), (" FRONTEND_STACK = {", CODE_TEXT)],
            [("  markup: ", CODE_TEXT), ("'HTML5 Semantic, Forms, Validation, SEO, a11y'", STR_COLOR), (",", CODE_TEXT)],
            [("  styling: ", CODE_TEXT), ("'CSS3, Flexbox, Grid, Custom Properties, Responsive'", STR_COLOR), (",", CODE_TEXT)],
            [("  scripting: ", CODE_TEXT), ("'JavaScript ES6+, DOM API, Event Delegation'", STR_COLOR), (",", CODE_TEXT)],
            [("  dataHandling: ", CODE_TEXT), ("'Arrays, Objects, map, filter, sort, immutability'", STR_COLOR), (",", CODE_TEXT)],
            [("  architecture: ", CODE_TEXT), ("'Multi-page SPA elements, URL params, localStorage'", STR_COLOR), (",", CODE_TEXT)],
            [("  project: ", CODE_TEXT), ("'Портал СмартОфис (7 полноценных страниц)'", NUM_COLOR)],
            [("};", CODE_TEXT)],
            [("", CODE_TEXT)],
            [("console.log(", CODE_TEXT), ("'Финальная сборка веб-приложения завершена!'", STR_COLOR), (");", CODE_TEXT)],
        ],
        "right_title": "Чему мы научились за курс",
        "bullet_items": [
            {"title": "Путь от нуля до портала", "desc": "От первых тегов <div> и стилей текста до сложного коммерческого многостраничного приложения.", "height": 42},
            {"title": "Глубокое понимание базы", "desc": "Уверенная работа с чистым DOM и событиями без зависимости от тяжелых сторонних фреймворков.", "height": 42},
            {"title": "Архитектурное мышление", "desc": "Умение проектировать структуры данных, организовывать файловую систему и разделять модули.", "height": 42},
            {"title": "Готовое портфолио", "desc": "Полноценный проект «СмартОфис» готов для размещения на GitHub и демонстрации работодателям.", "height": 42},
        ]
    }
]

# -------------------------------------------------------------
# PRESENTATION GENERATOR EXECUTION
# -------------------------------------------------------------
print(f"Total content slides to generate: {len(slides_data)}")

prs = Presentation(base_template_path)
blank_layout = prs.slide_layouts[6]

orig_slide_ids = list(prs.slides._sldIdLst)
title_slide_id = orig_slide_ids[0]
plan_slide_id = orig_slide_ids[1]
result_slide_id = orig_slide_ids[25]
goodbye_slide_id = orig_slide_ids[26]

# 0. Update Layouts footer to Вебинар 9
for layout in prs.slide_layouts:
    for s in layout.shapes:
        if "Google Shape;59;p13" in s.name and s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip() == "2":
                        r.text = "9"

# 1. Update Slide 1 (Title slide)
slide1 = prs.slides[0]
for sh in slide1.shapes:
    if sh.name == "Google Shape;59;p13" and sh.has_text_frame:
        sh.text_frame.text = "Вебинар 9 "
    elif sh.shape_type == MSO_SHAPE_TYPE.GROUP and sh.name == "Группа 12":
        for sub in sh.shapes:
            if sub.name == "TextBox 7":
                sub.text_frame.text = "ИТОГОВАЯ СБОРКА И РЕВЬЮ"
                sub.text_frame.paragraphs[0].runs[0].font.name = "Montserrat"
                sub.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
                sub.text_frame.paragraphs[0].runs[0].font.bold = True
            elif sub.name == "TextBox 9":
                sub.text_frame.text = "Проектирование \nи разработка интерфейсов пользователя"
print("Slide 1 updated successfully.")

# 2. Update Slide 2 (Plan slide)
slide2 = prs.slides[1]
plan_titles = {
    "TextBox 12": "1. Анализ типичных ошибок",
    "TextBox 15": "2. Закрепление материала",
    "TextBox 18": "3. Файловая архитектура",
    "TextBox 21": "4. Подготовка к защите",
    "TextBox 24": "5. Практика: Итоговый проект"
}
plan_subtitles = {
    "TextBox 13": "Ошибки в DOM, null pointer, обработчиках событий и утечках памяти",
    "TextBox 16": "HTML5 семантика, адаптивный CSS и компонентная архитектура JS",
    "TextBox 19": "Структура каталогов, относительные пути и разделение модулей",
    "TextBox 22": "Демонстрация сценариев, чек-лист 100 баллов и защита кода",
    "TextBox 25": "Сборка страницы «Мои бронирования» и финальное ревью приложения"
}

for sh in slide2.shapes:
    strip_shape_styles_and_shadows(sh)
    if sh.name in plan_titles and sh.has_text_frame:
        p = sh.text_frame.paragraphs[0]
        p.text = plan_titles[sh.name]
        p.runs[0].font.name = "Montserrat"
        p.runs[0].font.size = Pt(11.25)
        p.runs[0].font.bold = True
    elif sh.name in plan_subtitles and sh.has_text_frame:
        sh.width = Pt(450)
        tf = sh.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = plan_subtitles[sh.name]
        p.runs[0].font.name = "Inter"
        p.runs[0].font.size = Pt(9.75)
        p.runs[0].font.bold = False

print("Slide 2 updated successfully.")

# 3. Create 21 Content Slides
content_slide_elements = []
for idx, data in enumerate(slides_data, 1):
    new_slide = prs.slides.add_slide(blank_layout)
    content_slide_elements.append(new_slide._element)
    
    stype = data.get("type")
    if stype == "code":
        create_code_explanation_slide(
            new_slide,
            data["category"],
            data["question"],
            data["left_title"],
            data["code_lines"],
            data["right_title"],
            data["bullet_items"]
        )
    elif stype == "diagram_arch":
        create_architecture_diagram_slide(
            new_slide,
            data["category"],
            data["question"]
        )
    elif stype == "diagram_checklist":
        create_checklist_diagram_slide(
            new_slide,
            data["category"],
            data["question"]
        )
    print(f"  Created Slide {idx + 2}: {data['category']}")

# 4. Update Result Slide (Slide 26 in original template, prs.slides[25])
slide_res = prs.slides[25]
pic_to_remove = None
tb2_res = None
for sh in slide_res.shapes:
    strip_shape_styles_and_shadows(sh)
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        pic_to_remove = sh
    elif sh.name == "TextBox 2":
        tb2_res = sh

if pic_to_remove is not None:
    spTree = slide_res.shapes._spTree
    spTree.remove(pic_to_remove._element)

# Add Webinar 9 result screenshot: aspect ratio 1210:670 (1.806), centered nicely
res_w = Pt(510)
res_h = Pt(282)
res_left = Pt(105)
res_top = Pt(68)
slide_res.shapes.add_picture(result_img_path, res_left, res_top, res_w, res_h)

# Add or update TextBox 2 on Result slide
if tb2_res is None:
    tb2_res = slide_res.shapes.add_textbox(Pt(24.1), Pt(368.1), Pt(650), Pt(20))
    tb2_res.name = "TextBox 2"

tf_res = tb2_res.text_frame
tf_res.word_wrap = True
tf_res.margin_left = tf_res.margin_top = tf_res.margin_right = tf_res.margin_bottom = 0
p_res = tf_res.paragraphs[0]
p_res.alignment = PP_ALIGN.LEFT
p_res.text = ""
r_res = p_res.add_run()
r_res.text = "Результат: Финальная сборка веб-приложения"
r_res.font.name = "Montserrat"
r_res.font.size = Pt(10)
r_res.font.color.rgb = QUESTION_TEXT
print("Result slide updated with webinar 9 my-bookings screenshot and TextBox 2.")

# 5. Clean shadows on Goodbye slide (Slide 27 in original template, prs.slides[26])
slide_goodbye = prs.slides[26]
for sh in slide_goodbye.shapes:
    strip_shape_styles_and_shadows(sh)
print("Goodbye slide shadows stripped.")

# 6. Reorder slide list in presentation XML:
# [Title, Plan] + new_slide_ids + [Result, Goodbye]
all_slide_ids = list(prs.slides._sldIdLst)
new_slide_ids = all_slide_ids[len(orig_slide_ids):]
new_order = [title_slide_id, plan_slide_id] + new_slide_ids + [result_slide_id, goodbye_slide_id]
prs.slides._sldIdLst[:] = new_order

# 7. Save presentation
os.makedirs(os.path.dirname(output_pptx_path), exist_ok=True)
prs.save(output_pptx_path)

print(f"\nFinal presentation slide count: {len(prs.slides)}")
print(f"\nУСПЕХ! Презентация Вебинара 9 сохранена: {output_pptx_path}")
