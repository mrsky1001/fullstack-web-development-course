import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation(r"d:\IEEU International East-European University\fullstack-web-development-course\05. step by step 2026 & 2027\01. frontend\webinar-02-markup-and-styles\вебинар 2 - обновленный.pptx")
slide3 = prs.slides[2]

for shape in slide3.shapes:
    if shape.name == "Google Shape;69;p14":
        print(shape._element.xml)
